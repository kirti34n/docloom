"""Source ingestion: parse -> sanitize -> chunk -> embed.

Parsers stay lightweight (pdfplumber/pypdf, python-docx, trafilatura) -- no
torch. Text is sanitized at the boundary (control/bidi/zero-width chars) and
treated as data, never instructions."""

from __future__ import annotations

import asyncio
import ipaddress
import json
import re
import socket
from pathlib import Path
from urllib.parse import urljoin, urlparse

import httpx

from .db import execute, owner_of_source, query_one
from .settings import data_dir, get_setting

# control chars (minus tab/newline/CR), NEL and line/paragraph separators,
# zero-width chars, bidi overrides, BOM
_UNSAFE = re.compile(
    "[\x00-\x08\x0b\x0c\x0e-\x1f\x7f\x85\u2028\u2029"
    "​‎-‏‪-‮⁦-⁩﻿]"
)

CHUNK_CHARS = 1000
CHUNK_OVERLAP = 150


def sanitize(text: str) -> str:
    return _UNSAFE.sub("", text).replace("\r\n", "\n").replace("\r", "\n")


def _source_dir(source_id: str) -> Path:
    d = data_dir() / "sources" / source_id
    d.mkdir(parents=True, exist_ok=True)
    return d


# ------------------------------------------------------------------ parsers


def parse_pdf(path: Path) -> list[tuple[int, str]]:
    """Return (page_number, text) pairs."""
    pages: list[tuple[int, str]] = []
    try:
        import pdfplumber

        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                pages.append((i, page.extract_text() or ""))
        if any(t.strip() for _, t in pages):
            return pages
    except Exception:
        pass
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    return [(i, (p.extract_text() or "")) for i, p in enumerate(reader.pages, start=1)]


MAX_ZIP_UNCOMPRESSED_BYTES = 200 * 1024 * 1024   # 200 MB total inflated
MAX_ZIP_ENTRIES = 10_000
MAX_ZIP_RATIO = 100                              # inflated/compressed ceiling
ZIP_RATIO_FLOOR_BYTES = 32 * 1024 * 1024         # only ratio-check above this


def _zip_guard(path: Path) -> None:
    """docx/pptx/xlsx/epub are all ZIP containers; a small upload can still
    decompress to gigabytes of XML (a zip bomb) and exhaust server memory.
    zf.infolist() only reads the central directory -- no decompression -- so
    this check is cheap and runs before the real parser touches the archive."""
    import zipfile

    with zipfile.ZipFile(path) as zf:
        infos = zf.infolist()
        if len(infos) > MAX_ZIP_ENTRIES:
            raise ValueError("archive has too many entries")
        total_unc = sum(zi.file_size for zi in infos)
        total_comp = sum(zi.compress_size for zi in infos) or 1
        if total_unc > MAX_ZIP_UNCOMPRESSED_BYTES:
            raise ValueError("archive decompresses to too much data")
        if total_unc > ZIP_RATIO_FLOOR_BYTES and total_unc / total_comp > MAX_ZIP_RATIO:
            raise ValueError("suspicious archive compression ratio")


def parse_docx(path: Path) -> str:
    """Extract paragraph text plus any table content (tables are not part of
    doc.paragraphs in python-docx and would otherwise be silently dropped)."""
    import docx

    _zip_guard(path)
    doc = docx.Document(str(path))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def parse_pptx(path: Path) -> str:
    """Extract text from every shape (and speaker notes) of a .pptx deck."""
    from pptx import Presentation

    _zip_guard(path)
    out: list[str] = []
    for i, slide in enumerate(Presentation(str(path)).slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                parts.append(f"[notes] {notes}")
        if parts:
            out.append(f"Slide {i}\n" + "\n".join(parts))
    return "\n\n".join(out)


def parse_xlsx(path: Path) -> str:
    """Flatten every sheet of a workbook to `col | col | col` rows."""
    from openpyxl import load_workbook

    _zip_guard(path)
    wb = load_workbook(str(path), read_only=True, data_only=True)
    try:
        out: list[str] = []
        for ws in wb.worksheets:
            rows = [
                " | ".join("" if v is None else str(v) for v in row).rstrip(" |")
                for row in ws.iter_rows(values_only=True)
                if any(v is not None and str(v).strip() for v in row)
            ]
            if rows:
                out.append(f"Sheet: {ws.title}\n" + "\n".join(rows))
        return "\n\n".join(out)
    finally:
        wb.close()


def parse_epub(path: Path) -> str:
    """Extract reading-order text from an EPUB (ebooklib + a tiny HTML strip)."""
    from ebooklib import ITEM_DOCUMENT, epub

    _zip_guard(path)
    book = epub.read_epub(str(path))
    tag = re.compile(r"<[^>]+>")
    out: list[str] = []
    for item in book.get_items_of_type(ITEM_DOCUMENT):
        html = item.get_content().decode("utf-8", "ignore")
        # drop scripts/styles, then strip tags; entities are left as-is (rare)
        html = re.sub(r"(?is)<(script|style)[^>]*>.*?</\1>", " ", html)
        text = re.sub(r"(?is)<br\s*/?>|</p>|</div>|</h[1-6]>", "\n", html)
        text = tag.sub("", text)
        text = re.sub(r"[ \t]+", " ", text)
        text = re.sub(r"\n{3,}", "\n\n", text).strip()
        if text:
            out.append(text)
    return "\n\n".join(out)


def parse_csv(text: str) -> str:
    import csv
    import io

    text = text.replace("\r\n", "\n").replace("\r", "\n")
    sample = text[:4096]
    try:
        dialect = csv.Sniffer().sniff(sample) if sample.strip() else csv.excel
    except csv.Error:
        dialect = csv.excel
    prev = csv.field_size_limit()
    try:
        # csv's default field_size_limit (131072) rejects a valid CSV with one
        # large cell (embedded JSON/base64/long free-text). The whole text is
        # already in memory, so raising the limit adds no memory risk; cap it
        # Windows-safe (< 2**31, the C long on LLP64) so field_size_limit does
        # not raise OverflowError.
        csv.field_size_limit(min(2**31 - 1, len(text) + 1))
        return "\n".join(
            " | ".join(cell.strip() for cell in row)
            for row in csv.reader(io.StringIO(text), dialect)
            if any(cell.strip() for cell in row)
        )
    finally:
        csv.field_size_limit(prev)


def parse_html(text: str) -> str:
    import trafilatura

    return trafilatura.extract(text, include_comments=False, include_tables=True) or text


def read_text_smart(path: Path) -> str:
    """Read a text file without mojibake: UTF-8 first (the common case), then
    sniff the encoding, then a lossless latin-1 fallback (never raises, never
    inserts U+FFFD replacement chars for western text)."""
    raw = path.read_bytes()
    try:
        return raw.decode("utf-8")
    except UnicodeDecodeError:
        pass
    try:
        from charset_normalizer import from_bytes

        best = from_bytes(raw).best()
        if best is not None:
            return str(best)
    except Exception:
        pass
    return raw.decode("latin-1")


_YT = re.compile(
    r"(?:youtube\.com/(?:watch\?(?:.*&)?v=|shorts/|embed/)|youtu\.be/)"
    r"([A-Za-z0-9_-]{11})"
)
_YT_HOSTS = {"youtube.com", "www.youtube.com", "m.youtube.com", "youtu.be"}


def youtube_id(url: str) -> str | None:
    # Only treat this as a YouTube link when the URL's actual HOST is YouTube.
    # Matching the pattern anywhere in the string would let a URL like
    # http://169.254.169.254/x/youtube.com/watch?v=ID take the (unguarded)
    # transcript path against an internal host: an SSRF bypass.
    host = (urlparse(url).hostname or "").lower()
    if host not in _YT_HOSTS:
        return None
    m = _YT.search(url)
    return m.group(1) if m else None


def fetch_youtube(url: str, video_id: str) -> tuple[str, str]:
    """(title, transcript) for a YouTube link. Needs youtube-transcript-api.

    youtube-transcript-api 1.0 replaced the static
    YouTubeTranscriptApi.get_transcript(video_id) with an instance API:
    YouTubeTranscriptApi().fetch(video_id) returns a FetchedTranscript of
    FetchedTranscriptSnippet objects (attribute access, not dict indexing)."""
    from youtube_transcript_api import YouTubeTranscriptApi

    snippets = YouTubeTranscriptApi().fetch(video_id)
    text = " ".join(s.text.strip() for s in snippets if s.text.strip())
    if not text:
        raise ValueError("no transcript available for this video")
    # cheap title: the page <title>, falling back to the id
    title = f"YouTube {video_id}"
    try:
        with httpx.Client(timeout=10, follow_redirects=True,
                          headers={"User-Agent": "docloom-studio/0.1"}) as client:
            html = client.get(url).text
        m = re.search(r"<title>(.*?)</title>", html, re.S)
        if m:
            title = re.sub(r"\s*-\s*YouTube\s*$", "", m.group(1)).strip() or title
    except Exception:
        pass
    return title, text


_ALLOWED_SCHEMES = {"http", "https"}
_MAX_REDIRECTS = 5
_MAX_FETCH_BYTES = 25 * 1024 * 1024   # 25 MB cap on a fetched page body


def _is_public_host(host: str) -> bool:
    """True only if every address `host` resolves to is a public, routable
    address. Rejects loopback (127.0.0.1), private ranges (10/8, 172.16/12,
    192.168/16), link-local (169.254.0.0/16, the AWS/GCP/Azure metadata
    endpoint lives at 169.254.169.254), and other reserved/multicast ranges.
    Fails closed: an unresolvable host is not public."""
    try:
        infos = socket.getaddrinfo(host, None)
    except socket.gaierror:
        return False
    if not infos:
        return False
    for info in infos:
        raw = info[4][0].split("%", 1)[0]  # strip an IPv6 zone id if present
        try:
            ip = ipaddress.ip_address(raw)
        except ValueError:
            return False
        if (ip.is_private or ip.is_loopback or ip.is_link_local
                or ip.is_reserved or ip.is_multicast or ip.is_unspecified):
            return False
    return True


def _guard_url(url: str) -> None:
    """SSRF guard for a URL this server is about to fetch on a user's behalf:
    reject non-http(s) schemes and any host that doesn't resolve exclusively
    to public addresses. Callers must re-run this on every redirect hop too:
    a URL that looks public can still 302 to localhost or a cloud metadata
    address."""
    parsed = urlparse(url)
    if parsed.scheme not in _ALLOWED_SCHEMES:
        raise ValueError(f"refusing to fetch URL scheme {parsed.scheme!r}")
    host = parsed.hostname
    if not host or not _is_public_host(host):
        raise ValueError(f"refusing to fetch a non-public address ({host!r})")


def _reject_non_public_peer(resp: httpx.Response) -> None:
    """Re-validate the address httpx actually connected to. _guard_url resolves
    the host in a getaddrinfo separate from httpx's connect-time resolution, so
    a low-TTL attacker domain can pass the guard and then rebind to an internal
    address (DNS-rebinding TOCTOU). Checking the real peer address closes that
    gap; we abort before reading the body, so nothing internal reaches the user.
    Falls closed only when an address is present -- if the transport exposes no
    peer (e.g. a test double), there is nothing to re-validate."""
    stream = resp.extensions.get("network_stream")
    if stream is None:
        return
    addr = stream.get_extra_info("server_addr")
    if not addr:
        return
    try:
        ip = ipaddress.ip_address(addr[0].split("%", 1)[0])
    except ValueError:
        raise ValueError("connection resolved to an unparseable address")
    if (ip.is_private or ip.is_loopback or ip.is_link_local
            or ip.is_reserved or ip.is_multicast or ip.is_unspecified):
        raise ValueError(f"connection rebound to a non-public address ({ip})")


def fetch_url(url: str) -> tuple[str, str]:
    """Return (title, main-text) for a web page, or a transcript for YouTube.

    Redirects are followed by hand (not httpx's follow_redirects=True) so the
    SSRF guard re-validates every hop: a URL that looks public but redirects
    to localhost, an RFC1918 range, or the cloud metadata address must be
    refused exactly as if it had been given directly."""
    vid = youtube_id(url)
    if vid:
        return fetch_youtube(url, vid)

    import trafilatura

    html = ""
    with httpx.Client(timeout=20, follow_redirects=False,
                      headers={"User-Agent": "docloom-studio/0.1"}) as client:
        for _ in range(_MAX_REDIRECTS + 1):
            _guard_url(url)
            with client.stream("GET", url) as resp:
                _reject_non_public_peer(resp)
                if resp.has_redirect_location:
                    url = urljoin(url, resp.headers["location"])
                    continue
                resp.raise_for_status()  # 404/403/500 must fail ingestion, not feed error HTML
                encoding = resp.encoding or "utf-8"
                total = 0
                blocks: list[bytes] = []
                for block in resp.iter_bytes():
                    total += len(block)
                    if total > _MAX_FETCH_BYTES:
                        raise ValueError(
                            f"page exceeds {_MAX_FETCH_BYTES // (1024 * 1024)} MB limit")
                    blocks.append(block)
            html = b"".join(blocks).decode(encoding, errors="replace")
            break
        else:
            raise ValueError("too many redirects")
    text = trafilatura.extract(html, include_comments=False,
                               include_tables=True) or ""
    meta = trafilatura.extract_metadata(html)
    title = (getattr(meta, "title", None) if meta else None) or url
    return title, text


# ------------------------------------------------------------------ chunking


def chunk_text(text: str, source_id: str, page: int | None = None) -> list[dict]:
    """Split into ~CHUNK_CHARS windows with overlap, preferring paragraph
    boundaries. Carries source/page/section metadata."""
    text = text.strip()
    if not text:
        return []
    paras = re.split(r"\n\s*\n", text)
    chunks: list[dict] = []
    buf = ""
    section = ""
    for para in paras:
        para = para.strip()
        if not para:
            continue
        # a short title-case line is treated as a section heading
        if len(para) < 80 and "\n" not in para and para[:1].isupper():
            section = para
        # a paragraph longer than a whole chunk is windowed on its own so it
        # is never emitted as one oversized chunk
        if len(para) > CHUNK_CHARS:
            if buf:
                chunks.append({"text": buf.strip(), "section": section, "page": page})
                buf = ""
            start = 0
            while start < len(para):
                piece = para[start : start + CHUNK_CHARS].strip()
                if piece:  # an all-whitespace window must not become an empty chunk
                    chunks.append({"text": piece, "section": section, "page": page})
                if start + CHUNK_CHARS >= len(para):
                    break  # this window reached the end; no redundant tail window
                start += CHUNK_CHARS - CHUNK_OVERLAP
            continue
        if len(buf) + len(para) + 2 > CHUNK_CHARS and buf:
            chunks.append({"text": buf.strip(), "section": section, "page": page})
            buf = buf[-CHUNK_OVERLAP:] + "\n\n" + para
        else:
            buf = (buf + "\n\n" + para) if buf else para
    if buf.strip():
        chunks.append({"text": buf.strip(), "section": section, "page": page})
    return chunks


# ------------------------------------------------------------------ pipeline

INSIGHT_SYSTEM = (
    "Summarize the following source in 3-5 sentences: the key facts, figures, "
    "and claims someone would need to answer questions about it. Plain prose, "
    "no preamble, no markdown."
)


async def _summarize_source(source_id: str, chunks: list[dict]) -> str:
    """A short standing summary of a source, used when its context_mode is
    'insights' (feed the gist instead of every chunk at retrieval time).
    Best-effort: never raises, returns '' on any failure so a source without
    a provider configured (or a flaky one) still ingests normally and just
    falls back to 'full' chunk retrieval."""
    from .providers import ProviderConfig, complete

    text = "\n\n".join(c["text"] for c in chunks[:12])[:8000]
    if not text.strip():
        return ""
    try:
        cfg = ProviderConfig(**get_setting("provider.generation", owner_of_source(source_id)))
        summary = await complete(
            cfg,
            [{"role": "system", "content": INSIGHT_SYSTEM},
             {"role": "user", "content": text}],
            temperature=0.2, max_tokens=400,
        )
        return summary.strip()
    except Exception:
        return ""


async def ingest_source(source_id: str, ctx=None) -> None:
    """Parse the source, chunk it, embed it. Updates source status."""
    row = query_one("SELECT * FROM sources WHERE id = ?", (source_id,))
    if row is None:
        return
    kind, path, url = row["kind"], row["path"], row["url"]
    meta = json.loads(row["meta_json"])

    try:
        chunks: list[dict] = []
        if kind == "file":
            p = Path(path)
            ext = p.suffix.lower()
            if ext == ".pdf":
                for page_no, page_text in parse_pdf(p):
                    chunks += chunk_text(sanitize(page_text), source_id, page=page_no)
            elif ext == ".docx":
                chunks += chunk_text(sanitize(parse_docx(p)), source_id)
            elif ext == ".pptx":
                chunks += chunk_text(sanitize(parse_pptx(p)), source_id)
            elif ext in (".xlsx", ".xlsm"):
                chunks += chunk_text(sanitize(parse_xlsx(p)), source_id)
            elif ext == ".csv":
                chunks += chunk_text(sanitize(parse_csv(read_text_smart(p))), source_id)
            elif ext in (".html", ".htm"):
                chunks += chunk_text(sanitize(parse_html(read_text_smart(p))), source_id)
            elif ext == ".epub":
                chunks += chunk_text(sanitize(parse_epub(p)), source_id)
            else:  # txt, md, and other text — encoding-sniffed
                chunks += chunk_text(sanitize(read_text_smart(p)), source_id)
        elif kind == "url":
            title, text = fetch_url(url)
            meta["fetched_title"] = title
            chunks += chunk_text(sanitize(text), source_id)
            if not row["title"] or row["title"] == url:
                execute("UPDATE sources SET title = ? WHERE id = ?",
                        (title[:200], source_id))
        elif kind in ("text", "research"):
            txt = meta.get("text", "")
            if txt:
                chunks += chunk_text(sanitize(txt), source_id)
            else:
                # text was dropped after the first ingest (meta.pop below);
                # chunks.jsonl is the durable copy -- reload the already-parsed
                # chunks and re-embed them. Do NOT re-run chunk_text on these
                # (they are already chunked; re-chunking would double the
                # overlap window).
                chunks = load_chunks(source_id)

        if not chunks:
            raise ValueError("no extractable text")

        for i, c in enumerate(chunks):
            c["source_id"] = source_id
            c["chunk_ix"] = i
        (_source_dir(source_id) / "chunks.jsonl").write_text(
            "\n".join(json.dumps(c, ensure_ascii=False) for c in chunks),
            encoding="utf-8",
        )
        if ctx:
            ctx.emit("chunk", "done", detail=f"{len(chunks)} chunks")

        from .embeddings import embed_source

        await embed_source(source_id, [c["text"] for c in chunks])
        # the chunks just written to chunks.jsonl are now the durable copy of
        # this text; stop also carrying it in meta_json, which is re-parsed on
        # every sources-list poll and would otherwise bloat that forever
        meta.pop("text", None)
        try:
            summary = await _summarize_source(source_id, chunks)
            if summary:
                meta["insight_summary"] = summary
                await embed_source(source_id, [summary], name="summary")
        except Exception:
            pass  # best-effort enrichment: 'insights' mode falls back to 'full'
        execute("UPDATE sources SET status = 'ready', meta_json = ? WHERE id = ?",
                (json.dumps(meta), source_id))
        if ctx:
            ctx.emit("embed", "done")
    except asyncio.CancelledError:
        # job cancel or server shutdown: don't leave the source stuck in 'pending'
        meta["error"] = "ingestion cancelled"
        execute("UPDATE sources SET status = 'failed', meta_json = ? WHERE id = ?",
                (json.dumps(meta), source_id))
        raise  # propagate cancellation
    except Exception as e:
        meta["error"] = str(e)[:300]
        execute("UPDATE sources SET status = 'failed', meta_json = ? WHERE id = ?",
                (json.dumps(meta), source_id))
        if ctx:
            ctx.emit("ingest", "failed", detail=str(e)[:200])


def load_chunks(source_id: str) -> list[dict]:
    path = _source_dir(source_id) / "chunks.jsonl"
    if not path.is_file():
        return []
    # split only on real record separators: str.splitlines() would also break
    # on NEL/LINE SEPARATOR/PARAGRAPH SEPARATOR that json.dumps emits literally
    # inside a record, corrupting the JSON. read_text() has already normalized
    # \r\n and \r to \n, and json.dumps escapes any real newline inside strings.
    return [json.loads(line) for line in path.read_text(encoding="utf-8").split("\n") if line]
