"""Per-user asset resolver + brand kit feed generation and export."""

import asyncio
import json
import os
import tempfile

os.environ.setdefault("DOCLOOM_STUDIO_HOME", tempfile.mkdtemp(prefix="dstudio-assets-"))

import pytest  # noqa: E402

from docloom import Slide  # noqa: E402
from docloom_studio import generate as gen  # noqa: E402
from docloom_studio.assets import (  # noqa: E402
    apply_brand, brand_logo_image, resolve_image)
from docloom_studio.db import execute, init_db, new_id, now, query_one  # noqa: E402
from docloom_studio.irx import bake, load_document  # noqa: E402
from docloom_studio.settings import data_dir, set_setting  # noqa: E402


@pytest.fixture(autouse=True)
def _db():
    init_db()
    for t in ("artifact_versions", "artifacts", "sources", "notebooks",
              "assets", "user_settings", "auth_sessions", "workspaces", "users"):
        execute(f"DELETE FROM {t}")


def _user() -> str:
    uid = new_id()
    execute("INSERT INTO users (id, email, password_hash, created) VALUES (?, ?, ?, ?)",
            (uid, f"{uid}@t.local", "x", now()))
    return uid


def _notebook(user_id: str) -> str:
    wid = new_id()
    execute("INSERT INTO workspaces (id, user_id, name, created) VALUES (?, ?, ?, ?)",
            (wid, user_id, "w", now()))
    nb = new_id()
    execute("INSERT INTO notebooks (id, name, workspace_id, created, updated) "
            "VALUES (?, ?, ?, ?, ?)", (nb, "nb", wid, now(), now()))
    return nb


def _image_asset(tags: str, user_id: str) -> str:
    aid = new_id()
    adir = data_dir() / "assets" / aid
    adir.mkdir(parents=True, exist_ok=True)
    (adir / "pic.png").write_bytes(b"\x89PNG\r\n\x1a\n")
    execute("INSERT INTO assets (id, type, filename, tags, user_id, created) "
            "VALUES (?, 'image', 'pic.png', ?, ?, ?)", (aid, tags, user_id, now()))
    return aid


def _font_asset(user_id: str, name="brand.woff2") -> str:
    aid = new_id()
    adir = data_dir() / "assets" / aid
    adir.mkdir(parents=True, exist_ok=True)
    (adir / name).write_bytes(b"wOFF" + b"\x00" * 40)
    execute("INSERT INTO assets (id, type, filename, tags, user_id, created) "
            "VALUES (?, 'font', ?, '', ?, ?)", (aid, name, user_id, now()))
    return aid


def test_resolver_matches_on_tags():
    u = _user()
    _image_asset("volcano, lava", u)
    team = _image_asset("remote, team, collaboration", u)
    assert resolve_image("a remote team on a call", u) == team
    assert resolve_image("nothing relevant xyzzy", u) is None


def test_resolver_is_per_user():
    u1, u2 = _user(), _user()
    a1 = _image_asset("remote, team", u1)
    assert resolve_image("remote team", u1) == a1
    assert resolve_image("remote team", u2) is None  # u2 can't see u1's assets


def test_generation_fills_image_slot_and_bake_resolves(monkeypatch):
    u = _user()
    aid = _image_asset("remote, team", u)
    outline = gen.Outline(deck_title="Async", slides=[
        gen.OutlineItem(title="The team", layout="image_left", intent="who")])

    async def fake_gv(cfg, messages, schema, parse, lint_fn=None, max_rounds=3,
                      on_round=None):
        if "Draft slide" not in messages[-1]["content"]:
            return outline
        return Slide(layout="image_left", title="The team",
                     image={"query": "remote team"},
                     blocks=[{"type": "paragraph", "text": "We are distributed."}])

    monkeypatch.setattr(gen, "generate_validated", fake_gv)

    nb = _notebook(u)
    art = gen.create_artifact(nb, "deck")

    class Ctx:
        def emit(self, *a, **k):
            pass

    asyncio.run(gen.run_deck_pipeline(Ctx(), nb, art, "async work"))
    payload = json.loads(query_one(
        "SELECT payload_json FROM artifacts WHERE id = ?", (art,))["payload_json"])

    img = payload["ir"]["slides"][1]["image"]
    assert img["path"] == f"asset://{aid}"
    baked = bake(load_document(payload))
    assert baked.slides[1].image.path.endswith("pic.png")
    assert "asset://" not in baked.slides[1].image.path


def test_brand_logo_lands_on_title_slide(monkeypatch):
    u = _user()
    logo = _image_asset("logo", u)
    set_setting("brand.active", {"accent": "#ff0066", "logo_asset_id": logo}, u)
    outline = gen.Outline(deck_title="Branded", slides=[
        gen.OutlineItem(title="One", layout="content", intent="x")])

    async def fake_gv(cfg, messages, schema, parse, lint_fn=None, max_rounds=3,
                      on_round=None):
        if "Draft slide" not in messages[-1]["content"]:
            return outline
        return Slide(layout="content", title="One",
                     blocks=[{"type": "paragraph", "text": "hi"}])

    monkeypatch.setattr(gen, "generate_validated", fake_gv)
    nb = _notebook(u)
    art = gen.create_artifact(nb, "deck")

    class Ctx:
        def emit(self, *a, **k):
            pass

    asyncio.run(gen.run_deck_pipeline(Ctx(), nb, art, "branded deck"))
    payload = json.loads(query_one(
        "SELECT payload_json FROM artifacts WHERE id = ?", (art,))["payload_json"])
    assert payload["ir"]["slides"][0]["image"]["path"] == f"asset://{logo}"
    # brand accent overrides the theme at export, for this user
    assert apply_brand({"primary": "#000"}, u)["accent"] == "#ff0066"
    # another user's export is unbranded
    assert "accent" not in apply_brand({"primary": "#000"}, _user()) or \
        apply_brand({"primary": "#000"}, _user()).get("accent") != "#ff0066"


def test_brand_fonts_thread_into_theme():
    u = _user()
    hf, bf = _font_asset(u, "head.otf"), _font_asset(u, "body.woff2")
    set_setting("brand.active", {
        "heading_family": "BrandSerif", "heading_asset_id": hf,
        "body_family": "BrandSans", "body_asset_id": bf,
    }, u)
    themed = apply_brand({"primary": "#000"}, u)
    assert themed["font_heading"] == "BrandSerif"
    assert themed["font_body"] == "BrandSans"
    assert themed["font_heading_src"].endswith("head.otf")
    assert themed["font_body_src"].endswith("body.woff2")
    # a user without the font assets gets no src (can't see another user's files)
    assert "font_body_src" not in apply_brand({"primary": "#000"}, _user())


def test_brand_logo_image_is_per_user():
    u1, u2 = _user(), _user()
    logo = _image_asset("logo", u1)
    set_setting("brand.active", {"logo_asset_id": logo}, u1)
    img = brand_logo_image(u1)
    assert img and img["path"].endswith("pic.png")
    # u2 has no brand set → no logo; and even pointing at u1's asset id fails
    assert brand_logo_image(u2) is None


def test_export_stamps_brand_logo_on_document(monkeypatch):
    """End-to-end: export sets doc.logo from the active brand and renders."""
    from docloom import Document, Slide
    from docloom_studio.artifacts import export_artifact, ExportRequest
    from docloom_studio.generate import save_artifact, create_artifact

    u = _user()
    logo = _image_asset("logo", u)
    # python-pptx needs a genuinely-decodable image to embed
    from PIL import Image as PILImage
    logo_row = query_one("SELECT filename FROM assets WHERE id = ?", (logo,))
    PILImage.new("RGB", (64, 32), (10, 120, 200)).save(
        data_dir() / "assets" / logo / logo_row["filename"])
    set_setting("brand.active", {"logo_asset_id": logo}, u)
    nb = _notebook(u)
    aid = create_artifact(nb, "deck")
    doc = Document(title="Deck", slides=[
        Slide(layout="title", title="Deck"),
        Slide(layout="content", title="Body",
              blocks=[{"type": "bullets", "items": [{"text": "one"}]}]),
    ])
    save_artifact(aid, "Deck", {"ir": doc.model_dump(exclude_none=True),
                                "theme_name": "slate"})
    res = asyncio.run(export_artifact(aid, ExportRequest(format="pptx"),
                                      user={"id": u}))
    out = data_dir() / "exports" / res["filename"]
    assert out.is_file()
    # the content slide carries the stamped logo picture
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    prs = Presentation(str(out))
    pics = [s for s in prs.slides[1].shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1
