"""Universal ingestion: PPTX/XLSX/CSV/HTML parse correctly and non-UTF-8 text
is decoded by sniffing rather than mojibaking."""

import os
import tempfile

os.environ.setdefault("DOCLOOM_STUDIO_HOME", tempfile.mkdtemp(prefix="ds-ingest-"))

from pathlib import Path  # noqa: E402

from docloom_studio.ingest import (  # noqa: E402
    parse_csv, parse_html, parse_pptx, parse_xlsx, read_text_smart,
)


def test_parse_pptx(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Roadmap Q3"
    s.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1)).text_frame.text = (
        "Ship the podcast feature"
    )
    s.notes_slide.notes_text_frame.text = "kokoro tts"
    f = tmp_path / "deck.pptx"
    prs.save(str(f))
    t = parse_pptx(f)
    assert "Roadmap Q3" in t
    assert "podcast" in t
    assert "[notes] kokoro tts" in t


def test_parse_xlsx(tmp_path):
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Metrics"
    for row in (["Region", "Revenue"], ["EU", 1200], ["US", 3400]):
        ws.append(row)
    f = tmp_path / "data.xlsx"
    wb.save(str(f))
    t = parse_xlsx(f)
    assert "Sheet: Metrics" in t
    assert "Region | Revenue" in t
    assert "US | 3400" in t


def test_parse_csv():
    t = parse_csv("name,score\nalice,10\nbob,7\n")
    assert "name | score" in t
    assert "alice | 10" in t


def test_parse_html():
    t = parse_html(
        "<html><body><article><h1>T</h1>"
        "<p>Real body text kept by trafilatura.</p></article></body></html>"
    )
    assert "Real body text" in t


def test_read_text_smart_utf8_exact(tmp_path):
    f = tmp_path / "u.txt"
    f.write_text("café résumé señor", encoding="utf-8")
    assert read_text_smart(f) == "café résumé señor"  # exact, the common case


def test_read_text_smart_non_utf8_no_mojibake(tmp_path):
    f = tmp_path / "latin.txt"
    f.write_bytes("café résumé señor".encode("latin-1"))  # invalid UTF-8
    t = read_text_smart(f)
    assert "�" not in t  # no replacement-char mojibake
    assert "caf" in t and "sum" in t  # readable western text survived
