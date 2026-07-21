"""Native, editable PPTX shapes for the Diagram IR block (P2:
docs/diagram-plan.md section 4b -- "the feature the owner asked for: a very
good architecture diagram that you can edit too").

One solved layout (diagram_svg.solve()), many emitters. This module never
lays anything out itself: it reads a SolvedDiagram (already positioned, px
canvas space) and turns it into add_shape rounded rectangles / flowchart
shapes with text frames, add_connector ELBOW connectors GLUED to their
endpoint shapes with begin_connect/end_connect (so a node dragged in
PowerPoint drags its edges with it -- the one thing that makes this a real
editable diagram instead of a picture), and a whole-diagram add_group_shape
stamped with the Tier 1/Tier 2 content-hash contract (docs/diagram-plan.md
section 1): `docloom:diagram:{id or 'anon'}:{diagram_hash(d)}`.

Entry point (matches the plan's signature exactly):

    add_diagram(slide, d, solved, theme, x_in, y_in, w_in, max_h_in,
                *, mode="attached") -> float   # inches consumed

`solved` is the caller's own first solve() (at detail="full"); this module
climbs the font-floor degradation ladder ("full" -> "label+sub" -> "label",
re-solving as needed) and, if even the sparsest detail level cannot clear an
8pt node label, falls back to a rasterized picture (raster.svg_to_png) and
then, if the optional [diagrams] extra is absent, a visible placeholder box.
Never raises: every failure mode still leaves something legible on the
slide (P5 audit defect 1's "no silent drops" rule applies here too).
"""
from __future__ import annotations

import io
import warnings

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from ..ir import Diagram, diagram_hash
from ..theme import contrast_ratio
from . import raster
from .diagram_svg import (
    BAR,
    EDGE_STYLE,
    MARGIN,
    TITLE_H,
    SolvedDiagram,
    kind_palette,
    measure,
    paint_svg,
    solve,
    solve_ir,
)

EMU_IN = 914400
PX_PER_IN = 96.0          # solved geometry lives in the painter's canvas
                           # units, treated as CSS px at 96dpi (matches the
                           # plan's own k = w_in / (solved.width/96) formula)
PX_TO_EMU = int(EMU_IN / PX_PER_IN)  # 9525, the standard EMU-per-pixel constant

# font sizes (SVG px units) baked into diagram_svg.node_box()/paint_svg(),
# duplicated here because the native emitter re-derives its own pt sizes
# from the fit scale k rather than reading them back out of solved geometry
# (SolvedNode carries no font metadata, only text and box size).
LABEL_PX, SUB_PX, TAG_PX, ELAB_PX, TITLE_PX = 14.5, 10.5, 9.2, 10.5, 21.0
MIN_LABEL_PT = 8.0          # docs/diagram-plan.md section 4b font floor
FONT_FLOOR_PT = 5.0         # absolute safety net for sub/tag/edge-label text
DETAIL_LADDER = ("full", "label+sub", "label")
# Grid packing (diagram_svg.py's ROW_LIMIT/BAND_GAP, 2026-07-16) does NOT
# make this ladder -- or _raster_fallback below -- unnecessary. Measured: it
# turns a purpose-built hub-fanout fixture from 6.0pt/RASTER to 8.5pt/NATIVE
# (a rank whose real nodes all converge on the same one or two downstream
# neighbors, so banding actually compresses the cross stack instead of just
# relabeling it -- see diagram_svg.py's ROW_LIMIT docstring), but it is a
# genuinely conditional win, not a universal one: none of the 5 bake-off
# specs benefit at all (their widest ranks are either one indivisible group,
# which grid packing can never split, or dominated by long-edge dummy
# congestion elsewhere in the diagram, which is a different problem this
# does not touch), and a pure N-deep chain (one node per rank, the
# _dense_diagram fixture below) has nothing to band in the first place.
# solve()'s own banded-vs-unbanded comparison (_fit_score in diagram_svg.py)
# guarantees grid packing is never a net loss, but "never a loss" is not
# "always clears the floor" -- this ladder, and the raster/placeholder
# fallback past it, are still the only thing standing between a genuinely
# dense diagram and either an illegible native render or (with no raster
# extra installed) nothing at all. Keep it.

DIAGRAM_RASTER_PX = 1600    # matches docx.py's DIAGRAM_RASTER_PX

# caption strip: a 0.22in textbox with a 0.04in gap above it, reserved
# BEFORE any fit scale is computed (docs/diagram-status.md finding 5 --
# see _reserve_caption()).
CAPTION_GAP_IN = 0.04
CAPTION_BOX_H_IN = 0.22
CAPTION_RESERVE_IN = CAPTION_GAP_IN + CAPTION_BOX_H_IN

# node kind -> flowchart preset; everything not listed is ROUNDED_RECTANGLE.
# FLOWCHART_DATABASE does not exist in python-pptx's 182-member enum, so
# store uses FLOWCHART_MAGNETIC_DISK (a cylinder) instead, per the plan.
_SHAPE_PRESET = {
    "store": MSO_SHAPE.FLOWCHART_MAGNETIC_DISK,
    "external": MSO_SHAPE.FLOWCHART_TERMINATOR,
    "queue": MSO_SHAPE.FLOWCHART_MULTIDOCUMENT,
}
_DASHED_KINDS = {"external"}

# connection-site index for ROUNDED_RECTANGLE and every other preset this
# module uses EXCEPT the cylinder: 0=top, 1=left, 2=bottom, 3=right (proven
# 21/21 real stCxn/endCxn glue across ROUNDED_RECTANGLE, FLOWCHART_TERMINATOR
# and FLOWCHART_MULTIDOCUMENT on the spec3 payments diagram in the research
# PoC).
_CXN_BEGIN_END = {"LR": (3, 1), "TB": (2, 0)}

# FLOWCHART_MAGNETIC_DISK (store/cylinder) does NOT follow that convention
# and has no right-mid connection site at all -- verified empirically by
# rendering every index 0-3 through LibreOffice (docs/diagram-status.md
# finding 10): idx 0 and 1 both land near the top of the cap, idx 2 is the
# clean left-mid entry point on the body, idx 3 is the clean bottom-mid exit
# point. There is no index that lands on the right edge, so an LR diagram
# with a store as the SOURCE has to leave from the bottom (numerically the
# same index python-pptx calls "right" for a rectangle, but a different
# physical site on this preset); a store as the TARGET must receive on
# index 2, not the rectangle's index 1, or the connector rests on top of
# the cap pointing sideways instead of entering the body -- the exact bug
# the finding described ("write txn" terminating on the Ledger's cap).
_STORE_CXN_BEGIN_END = {"LR": (3, 2), "TB": (3, 0)}

# SVG dasharray strings (diagram_svg's EDGE_STYLE third element, plus the
# group/node border patterns paint_svg uses) mapped to the nearest built-in
# OOXML dash style, so the native path honors the same dash language as the
# raster path (docs/diagram-status.md finding 11: a "secure" edge used to
# lose its dash pattern entirely because the native emitter unpacked it into
# a throwaway `_dash` and never used it).
_DASH_MAP = {
    "": None,
    "6 3": MSO_LINE_DASH_STYLE.DASH,           # external node border
    "6 4": MSO_LINE_DASH_STYLE.DASH,           # "dashed" edge style
    "7 4": MSO_LINE_DASH_STYLE.DASH,           # security-group border
    "9 3 2 3": MSO_LINE_DASH_STYLE.DASH_DOT,   # "secure" edge style
}


def _rgb(hexcolor: str) -> RGBColor:
    h = hexcolor.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _readable_fg(theme, fill_hex: str) -> str:
    """theme.text or theme.background, whichever contrasts more against
    `fill_hex` -- mirrors render/pptx.py's own _label_fg (chart data
    labels), reimplemented locally because pptx.py imports THIS module (a
    reverse import would be circular).

    Node/sublabel text used to hardcode theme.text (label) / theme.muted
    (sublabel) regardless of what is actually painted behind them.
    kind_palette's node fills (diagram_svg.py) are derived from hue alone at
    a fixed ~0.955 lightness -- always near-white -- independent of
    theme.background/surface, so on an inverted band (pptx.py's
    _band_theme, where "theme.text" becomes the band's own light
    foreground) the hardcoded choice put equally-light text on that
    already-near-white fill: the measured "Go 1.23"/"PostgreSQL 16"
    invisible-sublabel bug. Resolving against the fill actually behind the
    text, instead of trusting the caller's "text" token to always mean
    dark, fixes this independent of whatever band the diagram happens to be
    drawn on."""
    if contrast_ratio(theme.background, fill_hex) >= contrast_ratio(theme.text, fill_hex):
        return theme.background
    return theme.text


def theme_dict(theme) -> dict:
    """Adapt a docloom Theme model to the plain 6-key dict overlay
    diagram_svg's solve()/paint_svg() expect (docs/diagram-plan.md section
    3: "the docloom Theme model is adapted by callers"). Same shape as the
    identically-named helper every other diagram-consuming renderer builds
    (docx.py's _diagram_theme, html.py, markdown.py, typst.py); exposed here
    as a public function so pptx.py's dispatch hook can stay a one-line
    call instead of duplicating the 6-key literal."""
    return {
        "primary": theme.primary,
        "accent": theme.accent,
        "surface": theme.surface,
        "text": theme.text,
        "muted": theme.muted,
        "background": theme.background,
    }


def _set_line(shape, color_hex: str, *, dash: str = "", weight: float = 1.25,
             arrow: bool = False) -> None:
    """Set a shape's or connector's line color/width/dash, and optionally a
    triangle tail arrowhead. python-pptx has no high-level arrowhead API, so
    the arrow is written directly into the a:ln XML (proven in the research
    PoC's set_arrow()). `dash` is one of diagram_svg's own SVG dasharray
    strings, mapped through _DASH_MAP to the nearest built-in OOXML dash
    style; "" (or any pattern not worth distinguishing) stays solid."""
    ln = shape.line._get_or_add_ln()
    if arrow:
        for e in ln.findall(qn("a:tailEnd")):
            ln.remove(e)
        tail = ln.makeelement(qn("a:tailEnd"), {})
        tail.set("type", "triangle")
        tail.set("w", "med")
        tail.set("len", "med")
        ln.append(tail)
    shape.line.color.rgb = _rgb(color_hex)
    shape.line.width = Pt(weight)
    style = _DASH_MAP.get(dash, MSO_LINE_DASH_STYLE.DASH if dash else None)
    if style is not None:
        shape.line.dash_style = style


def _straight_line(slide, x1, y1, x2, y2):
    """A plain line-shaped autoshape from (x1,y1) to (x2,y2), built via
    build_freeform rather than add_connector so it emits a normal `<p:sp>`
    (custom geometry) instead of a `<p:cxnSp>` connector -- used for the
    legend's edge-style key so it never inflates the diagram's own
    stCxn/endCxn/cxnSp counts (mirrors the existing freeform edge mode
    below)."""
    fb = slide.shapes.build_freeform(x1, y1)
    fb.add_line_segments([(x2, y2)], close=False)
    shp = fb.convert_to_shape()
    shp.fill.background()
    return shp


def _fit(s: SolvedDiagram, w_in: float, max_h_in: float) -> tuple[float, float, float]:
    """Fit scale k, and the two canvas dimensions in inches, for solved
    geometry `s` inside a (w_in, max_h_in) box (docs/diagram-plan.md section
    4b: "k = min(w_in / (solved.width/96), max_h_in / (solved.height/96))").
    Returns (k, canvas_w_in, canvas_h_in); k is 0.0 for a degenerate
    (zero-extent) canvas so callers never divide by zero downstream."""
    canvas_w_in = s.width / PX_PER_IN
    canvas_h_in = s.height / PX_PER_IN
    if canvas_w_in <= 0 or canvas_h_in <= 0 or w_in <= 0 or max_h_in <= 0:
        return 0.0, canvas_w_in, canvas_h_in
    k = min(w_in / canvas_w_in, max_h_in / canvas_h_in)
    return k, canvas_w_in, canvas_h_in


def _label_pt(k: float) -> float:
    return LABEL_PX * k * 72.0 / 96.0


def _reserve_caption(max_h_in: float, has_caption: bool) -> tuple[float, bool]:
    """Reserve room for the caption strip BEFORE any fit scale is computed
    (docs/diagram-status.md finding 5). The previous code fit the diagram
    into the FULL max_h_in and only afterward checked
    "if h + 0.26 <= max_h_in", but the fit scale k is solved so that
    whenever height binds, canvas_h_in * k == max_h_in exactly -- h ==
    max_h_in always in that case, so the guard was structurally always
    false. Height binds for any solved aspect below ~2.63, and the painter
    targets 2.0-2.2, so the caption was dead code for essentially every
    diagram. Reserving space up front instead means the diagram itself is
    fit into a strictly smaller box, leaving a slot the caption always fits
    into.

    Returns (height available to the diagram content itself, whether a
    caption will actually be drawn -- false only when max_h_in is too small
    to hold both a sliver of diagram and the caption strip, in which case
    the caption is dropped rather than starving the diagram to nothing)."""
    if not has_caption or max_h_in <= CAPTION_RESERVE_IN:
        return max_h_in, False
    return max_h_in - CAPTION_RESERVE_IN, True


def _draw_caption(slide, d: Diagram, theme, x_in: float, y_in: float,
                  w_in: float, content_h_in: float) -> float:
    """Draw the caption strip immediately below already-placed diagram
    content and return the additional height it consumes
    (CAPTION_RESERVE_IN). Shared by the native, raster-fallback, and
    placeholder paths so a caption looks identical no matter which path
    produced the diagram above it (docs/diagram-status.md finding 5).
    Callers are responsible for having already reserved this much room via
    _reserve_caption()."""
    cap = slide.shapes.add_textbox(
        Inches(x_in), Inches(y_in + content_h_in + CAPTION_GAP_IN),
        Inches(w_in), Inches(CAPTION_BOX_H_IN),
    )
    ctf = cap.text_frame
    ctf.word_wrap = True
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    cr.text = d.caption
    cr.font.size = Pt(11)
    cr.font.italic = True
    cr.font.name = theme.font_body
    cr.font.color.rgb = _rgb(theme.muted)
    return CAPTION_RESERVE_IN


def _warn_illegible(d: Diagram, fitted_pt: float | None) -> None:
    """docs/diagram-status.md finding 3: the native 8pt floor correctly
    rejected an over-dense diagram from the vector path, but the raster
    fallback it handed off to had no floor of its own AND no warning either
    -- render(doc, 'pptx') emitted zero warnings even though node labels
    came out at 3-4pt, unreadable at presentation distance. This is the
    warning that was missing; stacklevel=3 points at add_diagram's own
    caller (warnings.warn in here -> _warn_illegible -> add_diagram ->
    caller), matching this module's other user-facing warnings."""
    if fitted_pt is None:
        msg = (f"pptx: diagram {d.id!r} could not be laid out at any detail "
               "level within its box; rendering as a raster image so the "
               "content is not silently dropped")
    else:
        msg = (f"pptx: diagram {d.id!r} does not clear the "
               f"{MIN_LABEL_PT:.0f}pt node-label floor even at the sparsest "
               f"detail level (best {fitted_pt:.1f}pt); rendering as a "
               "raster image at that same sparsest layout so the content is "
               "not silently dropped, but labels may be hard to read at "
               "presentation distance -- consider splitting the diagram or "
               "shortening labels")
    warnings.warn(msg, stacklevel=3)


def placeholder(slide, d: Diagram, theme, x_in: float, y_in: float,
                w_in: float, max_h_in: float) -> float:
    """A visible stand-in for a diagram that could not be emitted in any
    form (solve() itself raised, or rasterization is unavailable AND the
    degradation ladder never cleared the font floor). Mirrors pptx.py's own
    _placeholder_block styling so a diagram failure reads the same as any
    other unresolved block, never a silent gap (P5 audit defect 1's "no
    silent drops" rule extends to diagrams)."""
    draw_h, want_caption = _reserve_caption(max_h_in, bool(d.caption))
    h = min(draw_h, 1.6)
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_in), Inches(y_in), Inches(w_in), Inches(h)
    )
    box.name = f"docloom:diagram:{d.id or 'anon'}:{diagram_hash(d)}"
    box.fill.solid()
    box.fill.fore_color.rgb = _rgb(theme.surface)
    box.line.color.rgb = _rgb(theme.muted)
    box.line.width = Pt(1)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = d.alt or d.title or "Diagram unavailable"
    run.font.size = Pt(12)
    run.font.italic = True
    run.font.name = theme.font_body
    run.font.color.rgb = _rgb(theme.muted)
    if want_caption:
        h += _draw_caption(slide, d, theme, x_in, y_in, w_in, h)
    return h


def _raster_fallback(slide, d: Diagram, theme, td: dict, x_in: float, y_in: float,
                     w_in: float, max_h_in: float, *, solved: SolvedDiagram | None = None,
                     target_aspect: float | None = None, legend: bool = True) -> float | None:
    """Full-fidelity PNG fallback for a diagram whose degradation ladder
    never cleared the node-label font floor. Returns None (never raises) so
    the caller can fall through to placeholder() when the optional
    [diagrams] extra is absent or rasterization fails for any reason.

    `solved` is the caller's own best (sparsest-detail) layout attempt from
    the native degradation ladder, when one exists: rasterizing that same
    geometry -- rather than blowing back up to "full" detail -- is the
    "respect a minimum legible scale" half of finding 3, since a sparser
    layout has a smaller canvas and therefore a larger fitted scale for the
    same box. Only re-solves (at the sparsest "label" detail) when the
    caller has no usable geometry at all, e.g. solve() itself raised on
    every rung of the ladder -- `legend` (propagated from the caller's own
    original solve(), same seam-preservation reasoning as add_diagram's own
    ladder re-solve) governs that re-solve exactly as it would the caller's."""
    draw_h, want_caption = _reserve_caption(max_h_in, bool(d.caption))
    s = solved
    if s is None:
        ta = target_aspect if target_aspect else (w_in / draw_h if draw_h > 0 else 2.0)
        try:
            s = solve_ir(d, td, target_aspect=ta, detail="label", legend=legend)
        except Exception:
            return None
    try:
        svg = paint_svg(s, td)
    except Exception:
        return None
    png = raster.svg_to_png(
        svg, width=DIAGRAM_RASTER_PX, font_files=raster.theme_font_files(theme)
    )
    if not png:
        return None
    try:
        pic = slide.shapes.add_picture(io.BytesIO(png), Inches(x_in), Inches(y_in))
    except Exception:
        return None
    scale = min(Inches(w_in) / pic.width, Inches(draw_h) / pic.height, 1.0)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = Inches(x_in) + (Inches(w_in) - pic.width) // 2
    pic.top = Inches(y_in)
    pic.name = f"docloom:diagram:{d.id or 'anon'}:{diagram_hash(d)}"
    # docs/diagram-status.md finding 4: the multi-sentence Diagram.alt was
    # discarded here, leaving python-pptx's own filename-derived default
    # ("image.png") as the only thing a screen reader would announce.
    # python-pptx has no high-level alt-text setter, so the descr attribute
    # is written directly (proven: add_picture always creates
    # p:nvPicPr/p:cNvPr, so this element always exists).
    pic._element.nvPicPr.cNvPr.set("descr", d.alt or d.title or "Diagram")
    h = pic.height / EMU_IN
    if want_caption:
        h += _draw_caption(slide, d, theme, x_in, y_in, w_in, h)
    return h


def add_diagram(slide, d: Diagram, solved: SolvedDiagram, theme,
                x_in: float, y_in: float, w_in: float, max_h_in: float,
                *, mode: str = "attached") -> float:
    """Emit `solved` (the caller's own solve(d, ..., detail="full")) as
    native, editable PowerPoint shapes fit into the (w_in, max_h_in) box at
    (x_in, y_in). Returns the height actually consumed, in inches. Never
    raises: every failure mode (below-floor legibility even at the sparsest
    detail, a raster extra that is not installed, a solve() that raises on
    re-attempt) still degrades to something visible on the slide.

    `mode` is "attached" (default: add_connector ELBOW glued to each node
    via begin_connect/end_connect, so PowerPoint re-routes edges when a user
    drags a node -- the actual editability promise) or "freeform" (follows
    the painter's own routed polyline via build_freeform, visually cleaner
    at density but NOT glued to its endpoints)."""
    if not d.nodes:
        return 0.0
    td = theme_dict(theme)
    draw_max_h, want_caption = _reserve_caption(max_h_in, bool(d.caption))
    target_aspect = w_in / draw_max_h if draw_max_h > 0 else 2.0
    # Preserve the CALLER's own legend choice (read off the `solved` they
    # already produced) across every re-solve this function does below: the
    # degradation ladder and the raster fallback both call solve() again,
    # and a re-solve that silently dropped back to legend=True (solve()'s
    # default) would hand paint_svg/the raster path a legend band the
    # caller never asked to reserve -- the same legend_h seam this module
    # otherwise fixes at the drawing end (see part 6 of _emit_native), just
    # at the solving end instead.
    legend = solved.legend_h > 0

    s = solved
    detail_idx = 0
    k, canvas_w_in, canvas_h_in = _fit(s, w_in, draw_max_h)
    lab_pt = _label_pt(k)
    while lab_pt < MIN_LABEL_PT and detail_idx < len(DETAIL_LADDER) - 1:
        detail_idx += 1
        try:
            s = solve_ir(d, td, target_aspect=target_aspect,
                        detail=DETAIL_LADDER[detail_idx], legend=legend)
        except Exception:
            s = None
            break
        k, canvas_w_in, canvas_h_in = _fit(s, w_in, draw_max_h)
        lab_pt = _label_pt(k)

    if s is None or lab_pt < MIN_LABEL_PT:
        fitted_pt = lab_pt if s is not None else None
        h = _raster_fallback(slide, d, theme, td, x_in, y_in, w_in, max_h_in,
                             solved=s, target_aspect=target_aspect, legend=legend)
        if h is not None:
            _warn_illegible(d, fitted_pt)
            return h
        warnings.warn(
            f"pptx: diagram {d.id!r} could not reach a legible fitted font "
            "size and the [diagrams] raster extra is unavailable; "
            "placeholder shown",
            stacklevel=2,
        )
        return placeholder(slide, d, theme, x_in, y_in, w_in, max_h_in)

    return _emit_native(slide, d, s, theme, td, x_in, y_in, w_in, draw_max_h,
                        k, canvas_w_in, canvas_h_in, mode, want_caption)


def _bar_span(n) -> tuple[float, float] | None:
    """Vertical extent (canvas px, absolute) of the kind-accent bar drawn on
    a node's left edge, mirroring diagram_svg.node_shape()'s own barpath()
    calls exactly. None for "external" (the painter draws no bar on it
    either). docs/diagram-status.md finding 11: the native path drew no
    kind bars at all, so the same Diagram block looked like two different
    products depending on whether it went native or raster."""
    if n.type == "external":
        return None
    if n.type == "store":
        ry = 9.0
        return n.y + ry + 2, n.y + n.h - ry - 2
    if n.type == "queue":
        return n.y + 8, n.y + n.h - 8
    if n.type == "security":
        return n.y + 9, n.y + n.h - 9
    return n.y + 8, n.y + n.h - 8


def _site_point(n, idx: int) -> tuple[float, float]:
    """Canvas-space (x, y) of connection-site `idx` on node `n`, matching
    the SAME index convention _CXN_BEGIN_END / _STORE_CXN_BEGIN_END encode
    (0=top,1=left,2=bottom,3=right for a plain rectangle-like preset; a
    store/cylinder has no right-mid site at all, so its own table only ever
    passes 0, 2, or 3 -- see _STORE_CXN_BEGIN_END's docstring). Used only to
    PREDICT where PowerPoint's own default elbow routing would draw a
    connector (see _default_elbow_crosses_obstacle below); it never
    influences the actual glued connector, which still uses begin_connect/
    end_connect's own site indices exactly as before."""
    if n.type == "store":
        return {0: (n.x + n.w / 2, n.y), 2: (n.x, n.y + n.h / 2),
                3: (n.x + n.w / 2, n.y + n.h)}.get(idx, (n.x + n.w / 2, n.y + n.h / 2))
    return {0: (n.x + n.w / 2, n.y), 1: (n.x, n.y + n.h / 2),
            2: (n.x + n.w / 2, n.y + n.h), 3: (n.x + n.w, n.y + n.h / 2)}[idx]


def _default_elbow_crosses_obstacle(begin, end, obstacles) -> bool:
    """True if an unrelated node's rect (`obstacles`: every OTHER node's
    rect -- never the edge's own source or target, which the caller has
    already excluded) overlaps the axis-aligned bounding box of `begin` and
    `end`.

    Why a bounding-box test, not a predicted exact path: an early version
    of this function predicted PowerPoint's default elbow bend (a single
    bend at the literal midpoint of the first axis -- the OOXML
    bentConnector's documented adj1=50% default) and only flagged an edge
    whose PREDICTED path crossed an obstacle. Verified by rendering that
    this under-caught: LibreOffice's actual rendered bend for spec3's Card
    Vault -> Card Network edge did not land at the predicted 50% midpoint,
    so the predicted 2-segment path missed the Postgres Replica cylinder by
    ~6px while the ACTUAL rendered connector cut straight through it. The
    bounding-box test does not depend on knowing which exact bend rule any
    given renderer uses: an orthogonal (horizontal/vertical-only) elbow
    between two points can never leave the axis-aligned rectangle those two
    points define, REGARDLESS of where its bend(s) land, so "does an
    unrelated node overlap that rectangle at all" is a sound, renderer-
    agnostic upper bound for "could this connector visually cross it" --
    zero false negatives (never misses a real crossing), at the cost of
    occasionally flagging an edge that a specific renderer's actual bend
    choice would have dodged (losing that one edge's glue is a strictly
    better failure mode than shipping a broken-looking connector).

    Why this exists (docs/diagram-status.md re-audit, native-emitter
    routing finding): attached-mode connectors are GLUED (begin_connect/
    end_connect) but NOT given custom bend geometry, so PowerPoint/
    LibreOffice draw their own default elbow between the two connection
    points with zero awareness of any other shape on the slide -- unlike
    the shared solve()'s own route(), which threads every edge through
    dummy-node lanes specifically to avoid this. solve()'s own routed
    polyline (SolvedEdge.pts) already avoids every node by construction; an
    edge flagged here falls back to that polyline via build_freeform
    instead of add_connector (see _emit_native part 4) -- unglued for just
    that one edge, but visually correct, rather than glued and broken."""
    bx0, bx1 = sorted((begin[0], end[0]))
    by0, by1 = sorted((begin[1], end[1]))
    box = (bx0, by0, bx1 - bx0, by1 - by0)
    return any(_rects_overlap(box, rect) for rect in obstacles)


def _rects_overlap(r1, r2) -> bool:
    x1, y1, w1, h1 = r1
    x2, y2, w2, h2 = r2
    return not (x1 + w1 <= x2 or x2 + w2 <= x1 or y1 + h1 <= y2 or y2 + h2 <= y1)


def _emit_native(slide, d: Diagram, s: SolvedDiagram, theme, td: dict,
                 x_in: float, y_in: float, w_in: float, max_h_in: float,
                 k: float, canvas_w_in: float, canvas_h_in: float,
                 mode: str, want_caption: bool) -> float:
    pal = kind_palette(td)
    offx_in = x_in + max(0.0, (w_in - canvas_w_in * k) / 2)
    offy_in = y_in + max(0.0, (max_h_in - canvas_h_in * k) / 2)

    def X(v: float):
        return Emu(int(round(offx_in * EMU_IN + v * k * PX_TO_EMU)))

    def Y(v: float):
        return Emu(int(round(offy_in * EMU_IN + v * k * PX_TO_EMU)))

    def D(v: float):
        return Emu(max(1, int(round(v * k * PX_TO_EMU))))

    def PT(px: float, floor: float = FONT_FLOOR_PT) -> float:
        return max(floor, px * k * 72.0 / 96.0)

    all_shapes = []

    # ---- 1. title (if any): drawn inside the TITLE_H band solve() already
    # reserved at the top of the canvas, mirroring paint_svg's own position. ----
    if s.title:
        tb = slide.shapes.add_textbox(
            X(MARGIN), Y(2), D(max(10.0, s.width - 2 * MARGIN)), D(TITLE_H - 6)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = s.title
        r.font.size = Pt(PT(TITLE_PX, floor=10.0))
        r.font.bold = True
        r.font.name = theme.font_heading
        r.font.color.rgb = _rgb(theme.text)
        all_shapes.append(tb)
        rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, X(MARGIN), Y(46), D(48), D(3))
        rule.fill.solid()
        rule.fill.fore_color.rgb = _rgb(theme.primary)
        rule.line.fill.background()
        rule.shadow.inherit = False
        all_shapes.append(rule)

    # ---- 2. group containers first, so they are furthest back ----
    for g in s.groups:
        secure = g.kind == "security-group"
        col = theme.accent if secure else theme.primary
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, X(g.x), Y(g.y), D(g.w), D(g.h))
        try:
            box.adjustments[0] = 0.06
        except (IndexError, ValueError):
            pass
        box.fill.solid()
        box.fill.fore_color.rgb = _rgb(theme.surface)
        _set_line(box, col, dash=("7 4" if secure else ""), weight=1.1)
        box.shadow.inherit = False
        tf = box.text_frame
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Inches(0.1)
        tf.margin_top = Inches(0.03)
        tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = g.label
        r.font.size = Pt(PT(SUB_PX))
        r.font.bold = True
        r.font.color.rgb = _rgb(col)
        all_shapes.append(box)

    # ---- 3. nodes, with a kind-accent bar on the left edge (matches
    # diagram_svg.node_shape()'s barpath(); finding 11) ----
    shp_by_id: dict[str, object] = {}
    type_by_id: dict[str, str] = {}
    for n in s.nodes:
        preset = _SHAPE_PRESET.get(n.type, MSO_SHAPE.ROUNDED_RECTANGLE)
        p = pal.get(n.type, pal["service"])
        shape = slide.shapes.add_shape(preset, X(n.x), Y(n.y), D(n.w), D(n.h))
        shape.name = f"docloom:node:{n.id}"
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(p["fill"])
        _set_line(shape, p["line"], dash=("6 3" if n.type in _DASHED_KINDS else ""),
                 weight=1.25)
        shape.shadow.inherit = False
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        m = Emu(max(1, int(0.04 * EMU_IN)))
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = m
        # Resolve label/sublabel text color against the node's OWN fill,
        # not a hardcoded theme.text/theme.muted: kind_palette's fills stay
        # near-white regardless of the caller's theme (see _readable_fg's
        # docstring), so on an inverted band the old hardcoded choice
        # produced near-invisible text -- the measured "Go 1.23"/
        # "PostgreSQL 16" bug.
        label_fg = _readable_fg(theme, p["fill"])
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        r0 = p0.add_run()
        r0.text = n.label
        r0.font.size = Pt(PT(LABEL_PX, floor=MIN_LABEL_PT))
        r0.font.bold = True
        r0.font.color.rgb = _rgb(label_fg)
        if n.sublabel:
            p1 = tf.add_paragraph()
            p1.alignment = PP_ALIGN.CENTER
            r1 = p1.add_run()
            r1.text = n.sublabel
            r1.font.size = Pt(PT(SUB_PX))
            r1.font.color.rgb = _rgb(label_fg)
        if n.tag:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = n.tag
            r2.font.size = Pt(PT(TAG_PX))
            r2.font.bold = True
            r2.font.color.rgb = _rgb(label_fg)
        shp_by_id[n.id] = shape
        type_by_id[n.id] = n.type
        all_shapes.append(shape)

        bar_span = _bar_span(n)
        if bar_span is not None:
            btop, bbot = bar_span
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, X(n.x + 0.8), Y(btop), D(BAR - 1), D(max(1.0, bbot - btop))
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = _rgb(p["bar"])
            bar.line.fill.background()
            bar.shadow.inherit = False
            all_shapes.append(bar)

    # ---- 4. connectors. Connection-site index depends on the endpoint's
    # node kind: the cylinder (store) preset has no right-mid site, so it
    # gets its own table (docs/diagram-status.md finding 10). ----
    node_by_id = {n.id: n for n in s.nodes}
    for e in s.edges:
        a, b = shp_by_id.get(e.source), shp_by_id.get(e.target)
        if a is None or b is None:
            continue  # a dangling edge in solved geometry shouldn't happen
                      # (lint's diagram/dangling-edge gates it upstream) but
                      # this must never raise on one
        beg_default, end_default = _CXN_BEGIN_END.get(s.direction, _CXN_BEGIN_END["LR"])
        beg_store, end_store = _STORE_CXN_BEGIN_END.get(s.direction, _STORE_CXN_BEGIN_END["LR"])
        beg = beg_store if type_by_id.get(e.source) == "store" else beg_default
        end = end_store if type_by_id.get(e.target) == "store" else end_default
        ck, sw, dash = EDGE_STYLE.get(e.style, EDGE_STYLE["solid"])
        color = {"muted": theme.muted, "primary": theme.primary,
                 "accent": theme.accent}.get(ck, theme.muted)
        weight = 2.3 if e.style == "emphasis" else (1.9 if e.style == "secure" else 1.5)
        use_freeform = mode == "freeform"
        if mode == "attached":
            # Only take the glued connector if an unrelated node's rect does
            # NOT overlap the bounding box between the two connection
            # points -- otherwise fall through to solve()'s own obstacle-
            # avoiding polyline instead (see _default_elbow_crosses_
            # obstacle's docstring for the verified repro and why a
            # bounding-box test, not a predicted exact bend). This trades
            # glue for correctness on the rare edge where they conflict;
            # every other edge is unaffected and stays fully glued.
            src_node, tgt_node = node_by_id.get(e.source), node_by_id.get(e.target)
            if src_node is not None and tgt_node is not None:
                obstacles = [(n.x, n.y, n.w, n.h) for n in s.nodes
                            if n.id not in (e.source, e.target)]
                begin_pt = _site_point(src_node, beg)
                end_pt = _site_point(tgt_node, end)
                if _default_elbow_crosses_obstacle(begin_pt, end_pt, obstacles):
                    use_freeform = True
        if use_freeform:
            pts = e.pts
            if len(pts) < 2:
                continue
            fb = slide.shapes.build_freeform(X(pts[0][0]), Y(pts[0][1]))
            fb.add_line_segments([(X(px), Y(py)) for px, py in pts[1:]], close=False)
            conn = fb.convert_to_shape()
            conn.fill.background()
        else:
            conn = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, X(0), Y(0), X(10), Y(10))
            conn.begin_connect(a, beg)
            conn.end_connect(b, end)
        _set_line(conn, color, dash=dash, weight=weight, arrow=True)
        conn.shadow.inherit = False
        all_shapes.append(conn)

    # ---- 5. edge labels, with an opaque halo so a crossing line never
    # strikes the text through (the same fix paint_svg applies) ----
    for e in s.edges:
        if not e.label_box or not (e.label or "").strip():
            continue
        lx, ly, lw, lh = e.label_box
        tb = slide.shapes.add_textbox(X(lx), Y(ly), D(lw), D(lh))
        tb.fill.solid()
        tb.fill.fore_color.rgb = _rgb(theme.background)
        tb.line.fill.background()
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = e.label
        r.font.size = Pt(PT(ELAB_PX))
        r.font.color.rgb = _rgb(theme.muted)
        all_shapes.append(tb)

    # ---- 6. legend: kind swatches + edge-style key, drawn natively inside
    # the s.legend_h canvas band solve() reserves at the bottom of the
    # diagram, NOT the module constant (docs/diagram-status.md findings 11
    # and 16, plus the legend_h seam residual found in the 2026-07-16
    # re-audit): gating on `s.legend` (the always-populated kind list) and
    # drawing at a position derived from the module's LEGEND_H constant
    # made the drawn band and the reserved band two independent things that
    # happened to agree only because every caller in this file always
    # solve()s with legend=True (the default) -- so s.legend_h always equals
    # LEGEND_H today. A caller that ever solves with legend=False would hit
    # `s.legend` still True (it is unconditionally populated) but
    # `s.legend_h == 0.0`, so this used to draw a legend into a 0-height
    # band positioned with the WRONG (nonzero) offset, landing on top of the
    # diagram's own lowest node row. Gating on `s.legend_h > 0` and
    # positioning from `s.legend_h` itself makes "is a band reserved" and
    # "where is it" the same fact solve() already computed, by construction,
    # so the two can never disagree again. Uses the painter's own measure()
    # so spacing matches paint_svg's layout exactly. Legend key lines are
    # built via _straight_line (freeform), not add_connector, so they never
    # inflate the diagram's own connector/glue counts. ----
    if s.legend_h > 0:
        ly = s.height - s.legend_h + 22
        rule = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, X(MARGIN), Y(ly - 16), D(max(1.0, s.width - 2 * MARGIN)), D(1)
        )
        rule.fill.solid()
        rule.fill.fore_color.rgb = _rgb(theme.muted)
        rule.line.fill.background()
        rule.shadow.inherit = False
        all_shapes.append(rule)

        lx = float(MARGIN)
        for kind in s.legend:
            p = pal.get(kind, pal["service"])
            chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, X(lx), Y(ly - 2), D(12), D(12))
            chip.fill.solid()
            chip.fill.fore_color.rgb = _rgb(p["fill"])
            _set_line(chip, p["line"], weight=0.75)
            chip.shadow.inherit = False
            all_shapes.append(chip)
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, X(lx + 0.5), Y(ly - 2), D(3), D(12))
            bar.fill.solid()
            bar.fill.fore_color.rgb = _rgb(p["bar"])
            bar.line.fill.background()
            bar.shadow.inherit = False
            all_shapes.append(bar)
            lab = slide.shapes.add_textbox(X(lx + 17), Y(ly - 8), D(measure(kind, 10) + 6), D(16))
            ltf = lab.text_frame
            ltf.word_wrap = False
            ltf.margin_left = ltf.margin_right = ltf.margin_top = ltf.margin_bottom = 0
            lr = ltf.paragraphs[0].add_run()
            lr.text = kind
            lr.font.size = Pt(PT(10))
            lr.font.color.rgb = _rgb(theme.muted)
            all_shapes.append(lab)
            lx += 17 + measure(kind, 10) + 20

        lx += 10
        for st, name in (("solid", "flow"), ("dashed", "async / return"),
                         ("emphasis", "primary path"), ("secure", "secure")):
            ck, sw, dash = EDGE_STYLE[st]
            color = {"muted": theme.muted, "primary": theme.primary,
                     "accent": theme.accent}.get(ck, theme.muted)
            key = _straight_line(slide, X(lx), Y(ly + 4), X(lx + 24), Y(ly + 4))
            _set_line(key, color, dash=dash, weight=sw)
            key.shadow.inherit = False
            all_shapes.append(key)
            lab = slide.shapes.add_textbox(X(lx + 30), Y(ly - 2), D(measure(name, 10) + 6), D(16))
            ltf = lab.text_frame
            ltf.word_wrap = False
            ltf.margin_left = ltf.margin_right = ltf.margin_top = ltf.margin_bottom = 0
            lr = ltf.paragraphs[0].add_run()
            lr.text = name
            lr.font.size = Pt(PT(10))
            lr.font.color.rgb = _rgb(theme.muted)
            all_shapes.append(lab)
            lx += 30 + measure(name, 10) + 20

    # ---- 7. whole-diagram group + Tier 1/Tier 2 hash stamp ----
    hash_name = f"docloom:diagram:{d.id or 'anon'}:{diagram_hash(d)}"
    try:
        grp = slide.shapes.add_group_shape(all_shapes)
        grp.name = hash_name
        # d.alt reaches HTML (aria-label on the inlined SVG) and the raster
        # fallback below (pic descr), but the native path never carried it
        # anywhere at all -- an accessibility gap for exactly the same
        # reason Image.alt was (see pptx.py's _set_alt), just for the one
        # emitter that groups its own shapes instead of adding one picture.
        if d.alt:
            grp._element.nvGrpSpPr.cNvPr.set("descr", d.alt)
    except Exception:
        # Grouping is a stretch task (docs/diagram-plan.md section 4b), not
        # a dependency: if it ever fails, stamp the hash on a node shape
        # instead so the Tier 1/Tier 2 contract still holds.
        if shp_by_id:
            next(iter(shp_by_id.values())).name = hash_name

    # h is measured from y_in to the diagram's TRUE rendered bottom edge,
    # not y_in + canvas_h_in * k: whenever width binds (canvas_h_in * k <
    # max_h_in), offy_in already centered the diagram vertically inside the
    # box, so its visual bottom sits offy_in - y_in lower than a naive
    # y_in + canvas_h_in * k would suggest. Getting this wrong meant the
    # caption used to render on top of the diagram's own lower half instead
    # of below it (caught visually converting a 5-node LR diagram: the
    # caption text landed across the Risk Engine node, not under the
    # legend) -- exactly the kind of thing finding 5's dead code hid until
    # the caption actually started rendering.
    h = min(max_h_in, (offy_in - y_in) + canvas_h_in * k)
    if want_caption:
        h += _draw_caption(slide, d, theme, x_in, y_in, w_in, h)
    return h
