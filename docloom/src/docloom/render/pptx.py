"""PPTX renderer: editable-native 16:9 decks built shape-by-shape on the
blank layout with python-pptx, so theme colors and fonts fully apply."""

from __future__ import annotations

import colorsys
import io
import warnings
from pathlib import Path
from urllib.parse import urlsplit

from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Inches, Pt

from ..ir import (
    Artifact,
    Block,
    BulletList,
    Callout,
    Chart,
    Code,
    Diagram,
    Divider,
    Document,
    Heading,
    Image,
    NumberedList,
    Paragraph,
    Quote,
    RichText,
    Slide,
    StatRow,
    Table,
    cited_ids,
    normalize_table,
    plain,
    source_numbers,
    spans,
)
from ..theme import Theme, contrast_ratio, hex_to_rgb
from . import RenderError, chart_svg, diagram_pptx, diagram_svg, raster, textfit

# target pixel width for rasterized SVG (2x a 640pt-wide chart), so the picture
# stays crisp on a projector and in print without bloating the deck
RASTER_PX = chart_svg.DEFAULT_WIDTH * 2

# Geometry/typography constants, exported as plain data for the web app.
LAYOUT = {
    "slide_w_in": 13.333,
    "slide_h_in": 7.5,
    "margin_in": 0.6,
    "gap_in": 0.16,
    # 26->30: a content-title band at 26pt read no bigger than a grown body
    # paragraph beside it (see MAX_GROWN_PT), a "generic PowerPoint" tell --
    # PINNED CONTRACT / audit item 3.
    "title_pt": 30,
    "body_pt": 14,
    "hero_title_pt": 36,
    "image_pane_ratio": 0.45,
    # This is the cap for a chart that SHARES its slide with other blocks; a
    # chart alone on its slide instead fills the full remaining body height
    # (see _chart_block's `solo` mode / _body's solo_chart handling), which
    # is the real fix for the old 4.5in cap leaving a permanent ~1.4in dead
    # void below a solo chart (P5 audit defect 6). This shared-slide cap
    # only got a modest bump, not all the way to the audit's suggested
    # 5.4: a slide with a chart PLUS a trailing block (e.g. an artifact
    # placeholder) needs that block to still get drawable room -- 5.4 left
    # a chart-and-artifact slide (the audit's own defect-1 example) with
    # < 0 remaining inches, silently dropping the artifact again via the
    # unrelated overflow-drop in _body's layout loop. NOTE: lint.py:117
    # (CHART_H_IN) mirrors this constant and was NOT updated here (out of
    # this file's ownership) -- flagged as a handoff.
    "chart_max_h_in": 4.8,
    "stat_card_h_in": 1.4,
    "stat_gap_in": 0.25,
    "stat_max_cards": 5,
}

SLIDE_W = LAYOUT["slide_w_in"]
SLIDE_H = LAYOUT["slide_h_in"]
MARGIN = LAYOUT["margin_in"]
GAP = LAYOUT["gap_in"]
MONO = "Consolas"
BODY_PT = LAYOUT["body_pt"]
HEAD_PT = {1: 22, 2: 19, 3: 17, 4: 15}
LIST_ITEM_GAP_PT = 6  # space_after below each bullet/numbered item
# P5 audit defect 4: 1.7 grew body prose to ~92% of the title's own size
# (14pt * 1.7 = 23.8pt) next to a fixed-size block on the same slide (12pt
# code, a 12pt table), destroying the hierarchy. 1.25 keeps the "fill sparse
# space" behavior while MAX_GROWN_PT below puts a hard, size-relative floor
# under how large that can ever get.
GROW_CAP = 1.25
# Modest cap used INSTEAD of GROW_CAP when the grow pass's blocks also
# include a fixed-size block: some growth keeps a sparse slide from reading
# empty, but a full GROW_CAP-style grow next to a table/chart/diagram is
# exactly what produced the mismatched hierarchy in the P5 audit (defect 4).
# See _grow_scale.
FIXED_NEIGHBOR_GROW_CAP = 1.15
MAX_GROWN_PT = round(LAYOUT["title_pt"] * 0.78)  # grown body text must never rival the title
# Block types whose font size the underfull-slide grow pass in _body scales;
# everything else (quotes, tables, charts, ...) keeps its natural size.
_GROWABLE_BLOCKS = (Heading, Paragraph, BulletList, NumberedList, Callout)
# BulletList.display=="grid" and NumberedList.display=="timeline" are BOTH
# nominally BulletList/NumberedList (so _GROWABLE_BLOCKS' isinstance check
# would treat them as ordinary growable prose), but their Gamma-signature
# treatments (_bullet_grid_block/_numbered_timeline_block, item 4) are
# content-sized mini-cards and a fixed-band timeline, not text that should
# expand to fill a sparse slide's slack: growing them the way a plain bullet
# list grows was the second route (besides the raw height-fill this file's
# item 1/4 fixes) to the "cards stretched too tall" defect. _is_growable/
# _is_fixed_size below are what the grow-scale machinery actually consults.


def _is_growable(b: Block) -> bool:
    if isinstance(b, BulletList) and b.display == "grid" and 3 <= len(b.items) <= 6:
        return False
    if isinstance(b, NumberedList) and b.display == "timeline" and 2 <= len(b.items) <= 6:
        return False
    return isinstance(b, _GROWABLE_BLOCKS)


def _is_fixed_size(b: Block) -> bool:
    if isinstance(b, BulletList) and b.display == "grid" and 3 <= len(b.items) <= 6:
        return True
    if isinstance(b, NumberedList) and b.display == "timeline" and 2 <= len(b.items) <= 6:
        return True
    return isinstance(b, _FIXED_SIZE_BLOCKS)
# Presence of any of these on a slide/column no longer suppresses the grow
# pass entirely (that produced the P5 audit's defect-4 mismatch the OTHER
# way -- prose stuck at its unscaled size next to a fixed block, reading
# smaller than the hierarchy intends); _grow_scale instead caps growth at
# FIXED_NEIGHBOR_GROW_CAP when any of these share the blocks list, so the
# fixed block still dictates the slide's own scale of "big" but the prose
# is not stranded at 1.0x on an otherwise generously-grown slide.
_FIXED_SIZE_BLOCKS = (Table, Code, Chart, StatRow, Image, Artifact, Diagram)
TABLE_PT = 12  # preferred row font size, shrunk toward TABLE_MIN_PT to fit
TABLE_MIN_PT = 9
# Post-layout text-fit floor and trigger. MIN_FIT_PT mirrors TABLE_MIN_PT's
# rationale (9pt is the smallest size that reads from a conference-room
# screen); FIT_SLACK_PT is hysteresis: a frame must measurably overflow by
# more than this before its runs are rewritten, so every frame the 1.3x
# _est_lines allocator already sized generously stays byte-identical.
# Not mirrored into lint.py: lint's height model is unaffected -- this pass
# shrinks font sizes INSIDE unchanged box heights, it never touches geometry.
MIN_FIT_PT = 9.0
FIT_SLACK_PT = 6.0
# Kill switch: flipping this to False makes _fit_text_frames a no-op, so the
# PPTX output is byte-identical to the pre-autofit-fix build. See render().
TEXTFIT_ENABLED = True
TABLE_VPAD = 0.06  # matches the cell.margin_top + cell.margin_bottom set below
CHART_TYPE = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "area": XL_CHART_TYPE.AREA,
    "pie": XL_CHART_TYPE.PIE,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
}
# warning/danger used to map straight to theme.muted/theme.text (gray and
# near-black) -- the two styles that most need to read as urgent signaled
# nothing at all, indistinguishable from a plain divider or body-text color.
# _callout_edge_color below derives real amber/red tones instead, by
# hue-shifting theme.accent's own saturation/lightness rather than injecting
# an unrelated stock stoplight hex, so the result still reads as part of
# this theme instead of fighting the brand.
_WARNING_HUE = 38.0  # amber
_DANGER_HUE = 4.0  # red
_SAFE_SCHEMES = {"http", "https", "mailto"}  # matches html.py
# Shared brand-logo target: 0.5in tall on every layout, matching the docx
# (Inches(0.5)), typst (1.27cm), and html (3rem @ 96dpi) renderers so the
# mark reads as one consistent size across every exported format.
LOGO_MAX_H = 0.5
LOGO_MAX_W = SLIDE_W * 0.28  # width guard: caps a very wide (e.g. horizontal wordmark) logo
# Mirrors lint.py:55's DIAGRAM_H_IN so the physical-height estimate _natural_h
# gives an overflow-checking Diagram block matches what lint already assumes
# when it decides a slide is overfull (docs/diagram-plan.md section 6).
DIAGRAM_H_IN = 4.6

# Named so lint.py's own mirrored copies (plain literals -- lint.py must stay
# import-light and layout-agnostic, see the comment above its SLIDE_BODY_H_IN)
# can be pinned against a real name via test_reaudit_lint.py instead of a
# bare, easy-to-drift number sprinkled across both files. These four used to
# be unnamed literals scattered through this module (0.26 in three different
# functions, 0.3 in a fourth, 0.12 in a fifth) and lint.py's geometry model
# mirrored NONE of them -- the direct cause of the audit's silent-drop
# repro: a subtitle (SUBTITLE_PAD_IN + the estimated line) shrinks the real
# body height, and a chart's own caption (CAPTION_H_IN) adds to its real
# footprint, neither of which lint's SLIDE_BODY_H_IN/CHART_H_IN accounted
# for, so a slide that silently dropped its trailing block still scored as
# safe.
CAPTION_H_IN = 0.26       # table/chart/diagram/placeholder caption strip
IMAGE_CAPTION_H_IN = 0.3  # _place_picture's own (slightly taller) image caption
QUOTE_ATTR_H_IN = 0.28    # quote attribution line
SUBTITLE_PAD_IN = 0.12    # _subtitle_line's fixed pad below its estimated lines
# The floor below which a block isn't worth giving its own shape -- and, as
# of the fix for the "ponytail" class (see _body), the minimum every block is
# now RESERVED ahead of time so an earlier block can never silently eat a
# later authored block's entire share of the slide.
MIN_BLOCK_RESERVE_IN = 0.3
# "Shrink, never drop" (the MIN_BLOCK_RESERVE_IN fix above) stops content
# from vanishing with zero trace, but reserving only 0.3in for a VISUAL block
# (an image, artifact, diagram, chart, or stat row) does not give it a
# legible size -- it gives it a speck. Measured live: a diagram squeezed to
# its 0.3in reserve rendered a 2208x1894 raster picture into a 0.47x0.40in
# box (~4735 effective dpi), a fix wearing a silent failure's clothes: lint
# DOES warn (deck/overflow), but the renderer still emitted something no
# viewer could ever read. Below this floor, a visual block is DROPPED with a
# warning naming it instead -- this is the one place in this file where
# "drop" is the correct answer over "shrink to fit", because a visual block
# has no legible degraded form the way text does (text can still be read at
# a smaller size; a diagram at 0.4in tall cannot).
MIN_VISUAL_BLOCK_H_IN = 1.2
_VISUAL_BLOCKS = (Image, Artifact, Diagram, Chart, StatRow)


def _safe_href(url: str) -> str | None:
    try:
        scheme = urlsplit(url).scheme.lower()
    except ValueError:
        return None
    return url if scheme in _SAFE_SCHEMES else None


def _rgb(color: str) -> RGBColor:
    return RGBColor(*hex_to_rgb(color))


def _tint(color: str, f: float) -> str:
    """Lighten `color` toward white by fraction f (0..1)."""
    r, g, b = hex_to_rgb(color)
    return "#%02X%02X%02X" % tuple(round(c + (255 - c) * f) for c in (r, g, b))


def _shade(color: str, f: float) -> str:
    """Darken `color` toward black by fraction f (0..1)."""
    r, g, b = hex_to_rgb(color)
    return "#%02X%02X%02X" % tuple(round(c * (1 - f)) for c in (r, g, b))


def _mix(color: str, toward: str, f: float) -> str:
    """Blend `color` toward an arbitrary target color by fraction f (0..1).
    Unlike _tint/_shade (which always move toward white/black), this blends
    toward whatever color is actually behind the mix -- used for the callout
    fill wash below, which needs to land near theme.background whether that
    background is light or dark, not near a hardcoded white."""
    r1, g1, b1 = hex_to_rgb(color)
    r2, g2, b2 = hex_to_rgb(toward)
    return "#%02X%02X%02X" % tuple(
        round(a + (c - a) * f) for a, c in ((r1, r2), (g1, g2), (b1, b2))
    )


def _hue_shift(color: str, hue_deg: float) -> str:
    """Recolor `color` to hue `hue_deg` (0-360) while preserving its own
    saturation and lightness, so a derived tone still reads as part of this
    theme's palette instead of an unrelated stock hex."""
    r, g, b = hex_to_rgb(color)
    h, l, s = colorsys.rgb_to_hls(r / 255, g / 255, b / 255)
    r2, g2, b2 = colorsys.hls_to_rgb(hue_deg / 360, l, s)
    return "#%02X%02X%02X" % (round(r2 * 255), round(g2 * 255), round(b2 * 255))


def _callout_edge_color(style: str, theme: Theme) -> str:
    """Per-style callout edge color. info/success reuse the brand's own
    primary/accent; warning/danger hue-shift theme.accent toward amber/red
    (see _WARNING_HUE/_DANGER_HUE) so they actually signal urgency instead of
    the old muted-gray/near-black mapping that made the two most urgent
    styles indistinguishable from ordinary chrome."""
    if style == "success":
        return theme.accent
    if style == "warning":
        return _hue_shift(theme.accent, _WARNING_HUE)
    if style == "danger":
        return _hue_shift(theme.accent, _DANGER_HUE)
    return theme.primary  # info, and any unrecognized style


# Finding D: all four callout fills were the flat theme.surface gray, so only
# the 4px edge bar actually carried the style's color -- across a room, a
# danger callout read exactly as loud as an info one. A heavy wash of the
# edge color itself (not a fixed light/dark stock hex) lets the whole card
# signal urgency while staying inside this theme's own palette.
_CALLOUT_FILL_WASH = 0.86  # fraction of the way from the edge color to theme.background


def _callout_fill_color(style: str, theme: Theme) -> str:
    """A tinted wash of `style`'s own edge color, blended toward
    theme.background (not a fixed white) so it reads correctly on a dark
    theme too, unlike _tint which always moves toward literal white."""
    return _mix(_callout_edge_color(style, theme), theme.background, _CALLOUT_FILL_WASH)


def _band_theme(theme: Theme, fill: str) -> Theme:
    """Theme tokens recolored for blocks drawn on top of an inverted band --
    a solid rectangle painted `fill` that covers part or all of the slide
    (hero's dark title strip, fill=theme.text; section's full-bleed cover,
    fill=theme.primary) -- so any block placed on that band resolves its OWN
    fill/foreground choices against what is ACTUALLY painted behind it,
    instead of against the document's still-light theme.background.

    This is the hero contrast regression's actual root cause: the old swap
    (just {"text": theme.background, "muted": theme.surface}) recolored the
    FOREGROUND tokens a block reads for its text, but left "background"
    pointed at the document's own light color. A callout mixes its wash via
    _callout_fill_color(style, theme) -> _mix(edge_color, theme.background,
    0.86): with only text/muted swapped, that wash still lands near-white
    (mixed toward the document's real light background) while the callout's
    own text, now theme.background too, ALSO renders near-white -- white on
    white, measured 1.1-1.3:1, not just low contrast but genuinely
    invisible. Swapping "background" (and "surface", its lighter sibling) to
    the band's own paint fixes this for ANY block that mixes toward
    "background" -- not just callouts: a diagram edge-label pill (drawn
    filled with theme.background in diagram_pptx.py) and a code block's own
    surface rect are exactly the same seam, fixed once here instead of
    patched per block type.

    "primary" only remaps to "accent" when the band's own fill genuinely IS
    theme.primary (section's full-bleed cover): a bullet marker or group
    label drawn in "primary" would otherwise be theme.primary on a
    theme.primary background -- literally invisible, not just low contrast
    (the original, narrower fix this generalizes). Doing that swap
    unconditionally (for hero's fill=theme.text too) was tried and rejected:
    it collapsed _callout_edge_color's info (theme.primary) and success
    (theme.accent) onto the same accent hex on every hero band, trading one
    contrast bug for a color-distinctness one."""
    update = {
        "background": fill,
        "surface": _mix(fill, theme.background, 0.12),
        "text": theme.background,
        "muted": theme.surface,
    }
    if fill == theme.primary:
        update["primary"] = theme.accent
    return theme.model_copy(update=update)


def _divider_color(theme: Theme) -> str:
    """A rule color that actually reads as a line. theme.surface (the
    default) sits at only a ~1.1 contrast ratio against theme.background on
    the default theme -- a near-invisible ghost, measured on a real deck
    (P5 audit defect 9). Step theme.muted progressively lighter until the
    line clears a real, if still subtle, contrast ratio."""
    for f in (0.6, 0.4, 0.2, 0.0):
        c = _tint(theme.muted, f)
        if contrast_ratio(c, theme.background) >= 1.5:
            return c
    return theme.muted


def _series_palette(theme: Theme) -> list[str]:
    """On-brand categorical colors derived from the theme's primary + accent,
    so native charts match the deck instead of Office's default blue/orange."""
    return [
        theme.primary,
        theme.accent,
        _tint(theme.primary, 0.42),
        _shade(theme.accent, 0.28),
        _tint(theme.accent, 0.5),
        _shade(theme.primary, 0.28),
    ]


def _label_fg(fill: str, theme: Theme) -> str:
    """theme.background or theme.text, whichever contrasts more with `fill`,
    so a data label stays legible on light and dark slice colors alike."""
    if contrast_ratio(theme.background, fill) >= contrast_ratio(theme.text, fill):
        return theme.background
    return theme.text


def _hide_chart_border(chart) -> None:
    """No python-pptx API exposes the chart-area outline; drop to the
    underlying chartSpace XML, the same technique the citation superscript in
    _runs uses. A styling hiccup here must never skip the rest of the theme
    styling, so failures are swallowed locally."""
    try:
        chartSpace = chart._chartSpace
        if chartSpace.find(qn("c:spPr")) is not None:
            return
        spPr = parse_xml(
            "<c:spPr %s><a:ln><a:noFill/></a:ln></c:spPr>" % nsdecls("c", "a")
        )
        chartSpace.find(qn("c:chart")).addnext(spPr)
    except Exception:
        pass


def _style_axes(chart, theme: Theme) -> None:
    """Recolor axis lines, gridlines, and tick labels on-brand instead of
    Office's default gray. Callers only reach this for chart types that have
    axes (pie does not)."""
    for axis in (chart.category_axis, chart.value_axis):
        axis.format.line.fill.background()  # no axis line, just tick labels
        axis.tick_labels.font.size = Pt(10)
        # font_heading, not font_body: theme.font_body defaults to Georgia,
        # whose old-style (text) figures sit below the cap line with
        # descenders -- fine for prose, but an axis tick like "6 5 4 3 2 1 0"
        # renders looking broken (P5 audit defect 13). Numeric chrome reads
        # cleanly in the heading font's lining figures instead.
        axis.tick_labels.font.name = theme.font_heading
        axis.tick_labels.font.color.rgb = _rgb(theme.muted)
    if chart.value_axis.has_major_gridlines:
        gl = chart.value_axis.major_gridlines.format.line
        gl.color.rgb = _rgb(_tint(theme.muted, 0.65))
        gl.width = Pt(0.75)


def _fix_dlbl_show_percent(dLbl_elm) -> None:
    """Force an already-materialized per-point c:dLbl to show a percentage,
    not the raw value. See the call site in _style_chart for why this is
    needed: python-pptx has no high-level API for a data label's show
    flags, so this reaches into the oxml element directly."""
    show_val = dLbl_elm.find(qn("c:showVal"))
    if show_val is not None:
        show_val.set("val", "0")
    show_pct = dLbl_elm.find(qn("c:showPercent"))
    if show_pct is not None:
        show_pct.set("val", "1")


def _style_chart(chart, b: Chart, theme: Theme) -> None:
    """Recolor a native chart from the theme palette and add tidy data
    labels, on-brand gridlines/legend, and no default Office chrome.
    Best-effort: callers wrap this so a styling hiccup never drops the chart."""
    palette = _series_palette(theme)
    plot = chart.plots[0]
    _hide_chart_border(chart)
    if b.chart == "pie":
        ser = plot.series[0]
        fills = [palette[i % len(palette)] for i in range(len(ser.points))]
        for i, pt in enumerate(ser.points):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = _rgb(fills[i])
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.show_percentage = True
        dl.show_value = False
        dl.number_format = "0%"
        dl.number_format_is_linked = False
        dl.font.size = Pt(10)
        dl.font.bold = True
        dl.font.name = theme.font_heading  # numeric chrome; see _style_axes
        for i, pt in enumerate(ser.points):
            lf = pt.data_label.font
            lf.size = Pt(10)
            lf.bold = True
            lf.name = theme.font_heading
            lf.color.rgb = _rgb(_label_fg(fills[i], theme))
            # Materializing pt.data_label.font (above) creates a per-point
            # c:dLbl element that python-pptx emits with its OWN default
            # show flags (showVal=1, showPercent=0), silently overriding the
            # plot-level dl.show_percentage=True set two lines up and turning
            # every slice label back into a raw value like "5.4" instead of
            # "54%" (P5 audit defect 14; confirmed by inspecting chart1.xml,
            # not a LibreOffice rendering artifact). Force this point's own
            # flags to match.
            _fix_dlbl_show_percent(pt.data_label._dLbl)
    else:
        for i, ser in enumerate(plot.series):
            color = _rgb(palette[i % len(palette)])
            if b.chart == "line":
                ser.format.line.color.rgb = color
                ser.format.line.width = Pt(2.25)
            elif b.chart == "scatter":
                ser.marker.format.fill.solid()
                ser.marker.format.fill.fore_color.rgb = color
            else:  # column, bar, area
                ser.format.fill.solid()
                ser.format.fill.fore_color.rgb = color
        # value labels only when uncluttered (<=2 series of bars/columns)
        if b.chart in ("column", "bar") and len(b.series) <= 2:
            plot.has_data_labels = True
            dl = plot.data_labels
            dl.number_format = "General"  # 58 not "58." (a trailing-dot glitch)
            dl.number_format_is_linked = False
            dl.font.size = Pt(11)
            dl.font.bold = True
            dl.font.name = theme.font_heading  # numeric chrome; see _style_axes
            dl.font.color.rgb = _rgb(theme.text)
            try:
                dl.position = XL_LABEL_POSITION.OUTSIDE_END
            except (ValueError, NotImplementedError):
                pass
        _style_axes(chart, theme)
    if chart.has_legend:
        leg = chart.legend
        leg.position = XL_LEGEND_POSITION.BOTTOM
        leg.font.size = Pt(11)
        leg.font.name = theme.font_body
        leg.font.color.rgb = _rgb(theme.text)


def _box(slide, x: float, y: float, w: float, h: float):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def _rect(slide, x: float, y: float, w: float, h: float, fill: str):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _set_fill_alpha(shape, alpha_pct: float) -> None:
    """Make a shape's solid fill semi-transparent via a raw <a:alpha> child
    of its <a:srgbClr> -- python-pptx exposes no high-level fill-
    transparency setter (the same gap _hide_chart_border/
    _fix_dlbl_show_percent reach into the oxml for elsewhere in this file),
    so this writes it by hand; transparency IS writable this way, the same
    technique already used for a chart's own spPr. `alpha_pct` is 0-100
    (100 = fully opaque). Best-effort: a styling hiccup here must never fail
    the render."""
    try:
        srgbClr = shape.fill.fore_color._xFill.find(qn("a:srgbClr"))
        if srgbClr is None:
            return
        alpha = parse_xml('<a:alpha %s val="%d"/>' % (nsdecls("a"), round(alpha_pct * 1000)))
        srgbClr.append(alpha)
    except Exception:
        pass


def _runs(
    p,
    rt: RichText,
    theme: Theme,
    numbers: dict[str, int],
    size: float,
    color: str,
    font: str,
    bold: bool = False,
    italic: bool = False,
) -> None:
    for sp in spans(rt):
        run = p.add_run()
        run.text = sp.text
        f = run.font
        f.name = MONO if sp.code else font
        f.size = Pt(size)
        f.bold = sp.bold or bold
        f.italic = sp.italic or italic
        f.color.rgb = _rgb(color)
        if sp.link and _safe_href(sp.link):
            run.hyperlink.address = sp.link
        if sp.cite and sp.cite in numbers:
            sup = p.add_run()
            sup.text = str(numbers[sp.cite])
            sup.font.name = font
            sup.font.size = Pt(max(8, round(size * 0.65)))
            sup.font.color.rgb = _rgb(theme.muted)
            sup._r.get_or_add_rPr().set("baseline", "30000")


def _est_lines(text: str, size: float, width: float) -> int:
    per = max(8, int(width * 144 / size))
    return sum(max(1, (len(ln) + per - 1) // per) for ln in text.split("\n"))


def _line_h(size: float) -> float:
    return size * 1.3 / 72


# ------------------------------------------------------------- body blocks


def _text_block(
    slide, rt, theme, numbers, x, y, w, max_h,
    size=BODY_PT, color=None, font=None, bold=False, italic=False, indent=0.0,
    scale=1.0,
) -> float:
    color = color or theme.text
    font = font or theme.font_body
    size = size * scale
    h = min(max_h, _est_lines(plain(rt), size, w - indent) * _line_h(size))
    tf = _box(slide, x + indent, y, w - indent, h)
    _runs(tf.paragraphs[0], rt, theme, numbers, size, color, font, bold, italic)
    return h


def _list_block(slide, b, theme, numbers, x, y, w, max_h, ordered: bool,
                scale: float = 1.0) -> float:
    bp = BODY_PT * scale
    item_gap = LIST_ITEM_GAP_PT / 72 * scale
    est = sum(
        _est_lines(plain(it.text), bp, w - 0.35 * (it.level + 1))
        * _line_h(bp)
        + item_gap
        for it in b.items
    )
    h = min(max_h, est)
    tf = _box(slide, x, y, w, h)
    counters: dict[int, int] = {}
    for i, it in enumerate(b.items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(LIST_ITEM_GAP_PT)
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(Inches(0.35 * (it.level + 1))))
        pPr.set("indent", str(-Inches(0.25)))
        if ordered:
            counters[it.level] = counters.get(it.level, 0) + 1
            for deeper in [lv for lv in counters if lv > it.level]:
                del counters[deeper]
            marker = f"{counters[it.level]}. "
        else:
            marker = "\u2022  "
        m = p.add_run()
        m.text = marker
        m.font.name = theme.font_body
        m.font.size = Pt(bp)
        m.font.color.rgb = _rgb(theme.primary)
        _runs(p, it.text, theme, numbers, bp, theme.text, theme.font_body)
    return h


def _grid_card_geometry(items, w: float, scale: float):
    """Shared geometry for the bullet-grid mini-cards: `_bullet_grid_block`
    draws from it, `_natural_h` estimates from it, so the two never disagree
    about how tall the grid actually is (item 4b/item 1). Card height is
    driven by its own longest item's wrapped text, NOT stretched to fill
    whatever `max_h` the slide happens to have -- a 2x2 of short items reads
    as a compact card cluster, not four half-empty tiles."""
    n = len(items)
    cols = 2 if n <= 4 else 3
    rows = -(-n // cols)  # ceil
    gap = GAP
    card_w = (w - gap * (cols - 1)) / cols
    pad = 0.16
    chip_d = 0.18
    chip_gap = 0.14  # space between the chip's bottom edge and the text start
    bp = BODY_PT * scale * 0.95
    text_top = pad + chip_d + chip_gap
    text_h = max(
        (_est_lines(plain(it.text), bp, card_w - 2 * pad) * _line_h(bp) for it in items),
        default=0.0,
    )
    card_h = text_top + text_h + pad
    return cols, rows, card_w, card_h, gap, pad, chip_d, text_top, bp


def _bullet_grid_block(slide, b, theme, numbers, x, y, w, max_h,
                       scale: float = 1.0) -> float:
    """PINNED CONTRACT item 4: BulletList.display=='grid' with 3-6 short
    items lays out as 2x2/3x2 mini-cards instead of a plain vertical list --
    a Gamma-signature treatment reached through the existing bullets block,
    no new Slide.layout."""
    items = b.items[:6]
    n = len(items)
    if n == 0:
        return 0.0
    cols, rows, card_w, card_h, gap, pad, chip_d, text_top, bp = _grid_card_geometry(
        items, w, scale
    )
    # Shrink (never stretch) only if the content-sized grid genuinely does
    # not fit the room it was given -- a rare, very cramped slide.
    if rows and max_h > 0:
        content_h = card_h * rows + gap * (rows - 1)
        if content_h > max_h:
            card_h = max(0.4, (max_h - gap * (rows - 1)) / rows)
    h = card_h * rows + gap * (rows - 1) if rows else 0.0
    # A near-white card (item 4a: ~6% tint, not the old 12% saturated fill)
    # with a 1px accent-tinted border reads as a clean mini-card, not a
    # heavy color block competing with the slide's actual content.
    fill = _mix(theme.background, theme.primary, 0.06)
    border = _tint(theme.primary, 0.35)
    for i, it in enumerate(items):
        r, c = divmod(i, cols)
        cx = x + c * (card_w + gap)
        cy = y + r * (card_h + gap)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(cy), Inches(card_w), Inches(card_h)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb(fill)
        card.line.color.rgb = _rgb(border)
        card.line.width = Pt(0.75)
        card.shadow.inherit = False
        try:
            card.adjustments[0] = 0.1  # ~12px-scale rounded corner, not the stock default
        except (IndexError, ValueError):
            pass
        # The chip is the deck's actual accent color (theme.primary, the
        # brand hue driving the card's own border tint) -- not theme.accent,
        # a secondary hue (often green) that clashed against a pink/blue
        # brand theme (item 4c).
        _rect(slide, cx + pad, cy + pad, chip_d, chip_d, theme.primary)
        tf = _box(
            slide, cx + pad, cy + text_top, card_w - 2 * pad,
            max(0.2, card_h - text_top - pad),
        )
        _runs(tf.paragraphs[0], it.text, theme, numbers, bp, theme.text, theme.font_body, bold=True)
    return h


def _numbered_timeline_block(slide, b, theme, numbers, x, y, w, max_h,
                             scale: float = 1.0) -> float:
    """PINNED CONTRACT item 4: NumberedList.display=='timeline' lays items
    out as horizontal numbered nodes (accent '01/02/03', a connector line,
    dots) instead of a plain vertical numbered list -- a Gamma-signature
    treatment reached through the existing numbered block, no new
    Slide.layout."""
    items = b.items[:6]
    n = len(items)
    if n == 0:
        return 0.0
    band_h = min(max_h, 2.0)
    dot_d = 0.26
    dot_y = y + band_h * 0.3
    seg_w = w / n
    line_y = dot_y + dot_d / 2 - 0.01
    if n > 1:
        # Connector drawn BEFORE the dots so the dots layer on top of it.
        _rect(slide, x + seg_w / 2, line_y, w - seg_w, 0.02, _tint(theme.primary, 0.35))
    for i, it in enumerate(items):
        cx = x + i * seg_w + seg_w / 2 - dot_d / 2
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(cx), Inches(dot_y), Inches(dot_d), Inches(dot_d)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = _rgb(theme.accent)
        dot.line.fill.background()
        dot.shadow.inherit = False
        num_tf = _box(slide, x + i * seg_w, dot_y - 0.34, seg_w, 0.3)
        num_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        _runs(
            num_tf.paragraphs[0], f"{i + 1:02d}", theme, {},
            14 * scale, theme.accent, theme.font_heading, bold=True,
        )
        label_top = dot_y + dot_d + 0.12
        label_tf = _box(
            slide, x + i * seg_w + 0.06, label_top, max(0.2, seg_w - 0.12),
            max(0.2, y + band_h - label_top),
        )
        label_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        _runs(
            label_tf.paragraphs[0], it.text, theme, numbers,
            BODY_PT * scale * 0.9, theme.text, theme.font_body,
        )
    return band_h


def _quote_block(slide, b, theme, numbers, x, y, w, max_h) -> float:
    # Reserve the attribution's own strip BEFORE laying out the quote text
    # (mirroring _place_picture's already-correct caption-reserved-upfront
    # pattern), so an authored attribution is never silently dropped just
    # because the quote text itself used up the whole box -- the same class
    # of silent loss as the trailing-block "ponytail" bug in _body, one
    # field down.
    attr_h = QUOTE_ATTR_H_IN if b.attribution else 0.0
    h = _text_block(
        slide, b.text, theme, numbers, x, y, w, max(0.1, max_h - attr_h),
        size=15, italic=True, indent=0.4,
    )
    if b.attribution:
        tf = _box(slide, x + 0.4, y + h + 0.04, w - 0.4, 0.24)
        _runs(
            tf.paragraphs[0], "\u2014 " + b.attribution, theme, numbers,
            12, theme.muted, theme.font_body,
        )
        h += QUOTE_ATTR_H_IN
    return h


def _code_block(slide, b, theme, x, y, w, max_h) -> float:
    size, pad = 12, 0.12
    lines = b.code.split("\n")
    h = min(max_h, len(lines) * _line_h(size) + 2 * pad)
    _rect(slide, x, y, w, h, theme.surface)
    tf = _box(slide, x + pad, y + pad, w - 2 * pad, h - 2 * pad)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = ln
        run.font.name = MONO
        run.font.size = Pt(size)
        run.font.color.rgb = _rgb(theme.text)
    return h


def _table_row_h(size: float, vpad: float = TABLE_VPAD) -> float:
    """The real minimum row height PowerPoint will honor for a run at `size`
    with the cell top/bottom margins set below: text height plus padding."""
    return _line_h(size) + vpad


def _table_col_widths(header, rows, cols: int, tw: float, solo: bool = False) -> list[float]:
    """Column widths weighted by each column's longest plain-text content,
    instead of an equal tw/cols split. An equal split gives "Vendor" and
    "Time to value" the same track as "$18", so the long label crowds its
    cell while the short value floats in a mostly-empty one (P5 audit
    defect 8). Clamp to [0.9in, a per-column cap] so no column collapses
    unreadably narrow or swallows the whole table, then rescale to land
    exactly on tw.

    The cap used to be a PERCENTAGE of tw (0.4 * tw), which shrinks right
    along with a narrow two_column pane's own tw; raised to an absolute
    ~4in ceiling instead, so a long-label column can claim more of its fair
    share before the rescale below normalizes the total back to tw exactly
    (item 2). A `solo` table (the dominant/only block on its slide) drops
    the cap altogether -- nothing should artificially cap a column's share
    of a table that already owns the whole slide width.

    `cols` is the caller's own column count (it may exceed len(header) for a
    genuinely empty table, which normalize_table leaves at width 0 -- the
    caller still renders a 1-column table for it, so this must match)."""
    if cols <= 0:
        return []
    lengths = [
        max(
            [1]
            + ([len(plain(header[c]))] if c < len(header) else [])
            + [len(plain(r[c])) for r in rows if c < len(r)]
        )
        for c in range(cols)
    ]
    total = sum(lengths)
    min_w = 0.9
    max_w = tw if solo else min(4.0, tw)
    raw = [max(min_w, min(max_w, tw * n / total)) for n in lengths]
    scale = tw / sum(raw)
    return [rw * scale for rw in raw]


def _table_block(slide, b, theme, numbers, x, y, w, max_h, solo: bool = False) -> float:
    header, rows = normalize_table(b.header, b.rows)
    cols = max(len(header), 1)
    # A solo/dominant table is the "fixed block owns a sparse slide" case
    # (item 2): widen it toward the full column width instead of floating
    # narrow at the shared-slide cols*2.5in cap.
    tw = w if solo else min(w, max(3.0, cols * 2.5))
    tx = x + (w - tw) / 2 if tw < w else x  # center a narrower-than-column table
    cap_h = CAPTION_H_IN if b.caption else 0.0
    budget = max(_table_row_h(TABLE_MIN_PT), max_h - cap_h)

    # A solo table also grows its row padding modestly (a taller vpad from
    # the start, so the shrink loop below already accounts for it rather
    # than fighting the budget after the fact).
    vpad = TABLE_VPAD * 1.6 if solo else TABLE_VPAD

    # shrink the font (down to a floor) until every row fits at that size,
    # since row height is tied to what the font actually needs, not an
    # independent division PowerPoint won't honor
    size = TABLE_PT
    while size > TABLE_MIN_PT and _table_row_h(size, vpad) * (len(rows) + 1) > budget:
        size -= 1
    row_h = _table_row_h(size, vpad)

    # Item 2, inverted: a solo/dominant table on a sparse slide GROWS its
    # font (up to a 16pt ceiling) toward ~70-75% fill instead of floating
    # small at the top with a dead band below -- the shrink loop above only
    # ever brings size DOWN to fit; this grows it back UP when there is
    # room to spare.
    if solo and rows:
        target_h = budget * 0.72
        while size < 16 and _table_row_h(size + 1, vpad) * (len(rows) + 1) <= target_h:
            size += 1
        row_h = _table_row_h(size, vpad)
        # A table with few rows can't reach the fill target through font
        # growth alone (the 16pt ceiling above is reached long before 5-ish
        # rows add up to 65-75% of a typical body height) -- the same
        # fundamental limit prose growth hits on a very sparse slide. Grow
        # the row height itself (padding) the rest of the way instead: a
        # spacious, airy row reads as a deliberate "this table owns the
        # slide" card instead of a tight 12-16pt table stranded small at
        # the top with a dead band below it. Capped so a 1-2 row table
        # doesn't blow up into absurdly tall rows.
        n_rows_total = len(rows) + 1
        row_h = max(row_h, min(target_h / n_rows_total, 1.3))

    # still too many rows at the floor size: cap the count and add a visible
    # "+N more rows" notice instead of silently overflowing the slide
    room = max(1, int(budget / row_h))
    more = 0
    if len(rows) + 1 > room:
        keep = max(0, room - 2)  # header row + the notice row both cost a slot
        more = len(rows) - keep
        # The "+N more rows" notice IS visible (never a blank gap), but the
        # dropped rows' own cell text carries no trace anywhere else in the
        # deck -- this used to be a genuinely silent loss with zero warning,
        # caught only by rendering-and-reading the actual XML, not by any
        # structural test (none asserted on truncated row CONTENT, only
        # counts).
        ident = f" (id={b.id!r})" if b.id else ""
        warnings.warn(
            f"pptx: table{ident} truncated to {keep} of {len(rows)} rows to "
            f"fit; the dropped rows are summarized by a \"+{more} more "
            "row(s)\" notice but their own cell text is not rendered -- "
            "move the full table to the report or a sheet",
            stacklevel=2,
        )
        rows = rows[:keep]

    n_rows = len(rows) + 1 + (1 if more else 0)
    th = n_rows * row_h
    frame = slide.shapes.add_table(
        n_rows, cols, Inches(tx), Inches(y), Inches(tw), Inches(th)
    )
    tbl = frame.table
    # Table has no authored alt field (unlike Image/Artifact); the caption is
    # the closest authored description, falling back to the column headers
    # so a screen reader still gets SOMETHING beyond "table, 3 by 4" (see
    # _set_alt's docstring for why this reaches into the oxml by hand).
    cols_desc = ", ".join(t for t in (plain(c) for c in header) if t)
    _set_alt(frame, b.caption or (f"Table with columns: {cols_desc}" if cols_desc else "Table"),
             title="Table")
    tbl.first_row = False
    tbl.horz_banding = False
    col_w = _table_col_widths(header, rows, cols, tw, solo=solo)
    for c in range(cols):
        tbl.columns[c].width = Inches(col_w[c])
    for r in range(n_rows):
        tbl.rows[r].height = Inches(row_h)
    for r in range(len(rows) + 1):
        cells = header if r == 0 else rows[r - 1]
        fill = theme.primary if r == 0 else (theme.background if r % 2 else theme.surface)
        color = theme.background if r == 0 else theme.text
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(fill)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(vpad / 2)
            if c < len(cells):
                # font_heading, not font_body: a table is data-dense chrome
                # (prices, dates, short labels), and Georgia's old-style
                # figures make its numbers look like typos (P5 audit
                # defect 13), same reasoning as _style_axes.
                _runs(
                    cell.text_frame.paragraphs[0], cells[c], theme, numbers,
                    size, color, theme.font_heading, bold=(r == 0),
                )
    if more:
        r = n_rows - 1
        note = tbl.cell(r, 0)
        if cols > 1:
            note.merge(tbl.cell(r, cols - 1))
        note.fill.solid()
        note.fill.fore_color.rgb = _rgb(theme.surface)
        note.vertical_anchor = MSO_ANCHOR.MIDDLE
        note.margin_left = note.margin_right = Inches(0.08)
        note.margin_top = note.margin_bottom = Inches(0.03)
        _runs(
            note.text_frame.paragraphs[0],
            f"+ {more} more row{'s' if more != 1 else ''}",
            theme, numbers, size, theme.muted, theme.font_body, italic=True,
        )
    h = th
    # cap_h was already reserved above (before the row/font-shrink budget
    # was even computed), so the caption always fits and is always drawn --
    # not gated on "if there happens to be room left over", which is exactly
    # the pattern that let a caption vanish silently when a slide ran tight.
    if b.caption:
        tf = _box(slide, tx, y + h + 0.04, tw, 0.22)
        _runs(
            tf.paragraphs[0], b.caption, theme, numbers,
            11, theme.muted, theme.font_body, italic=True,
        )
        h += CAPTION_H_IN
    return h


def _callout_block(slide, b, theme, numbers, x, y, w, max_h,
                   scale: float = 1.0) -> float:
    # Item 5: a callout used to render at 13pt -- SMALLER than the 14pt
    # BODY_PT text sitting right under it, the opposite of "emphasis". It
    # now renders at BODY_PT+1 (scaling with the same grow factor as its
    # neighbors), bold in the heading font as a semibold stand-in (no true
    # semibold weight is available), with a heavier pad and edge bar, and
    # capped to ~85% of the column so it reads as a distinct card rather
    # than a full-width strip weaker than the bullets under it.
    size = (BODY_PT + 1) * scale
    pad, edge_w = 0.24, 0.1
    card_w = min(w, w * 0.85)
    est = _est_lines(plain(b.text), size, card_w - edge_w - 2 * pad) * _line_h(size) + 2 * pad
    h = min(max_h, est)
    _rect(slide, x, y, card_w, h, _callout_fill_color(b.style, theme))
    _rect(slide, x, y, edge_w, h, _callout_edge_color(b.style, theme))
    tf = _box(slide, x + edge_w + pad, y + pad, card_w - edge_w - 2 * pad, h - 2 * pad)
    _runs(
        tf.paragraphs[0], b.text, theme, numbers, size, theme.text, theme.font_heading,
        bold=True,
    )
    return h


def _placeholder_block(slide, x, y, w, max_h, theme, alt: str, caption: str | None) -> float:
    """A visible stand-in for a block that could not be embedded (an
    unresolved reference, a missing file, or a format python-pptx cannot
    decode), so content never silently disappears from the deck (P5 audit
    defect 1). Mirrors the DOCX placeholder paragraph's intent, as a native
    PPTX shape with its alt text and caption legible on the slide."""
    # Reserve the caption's own strip BEFORE sizing the placeholder box
    # (same upfront-reservation fix as _quote_block/_table_block/
    # _chart_block above), so a caption on a placeholder is never silently
    # dropped just because max_h was tight.
    cap_h = CAPTION_H_IN if caption else 0.0
    h = min(max(0.1, max_h - cap_h), 1.6)
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = _rgb(theme.surface)
    box.line.color.rgb = _rgb(_tint(theme.muted, 0.35))
    box.line.width = Pt(1)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = alt or "Content unavailable"
    run.font.size = Pt(12)
    run.font.italic = True
    run.font.name = theme.font_body
    run.font.color.rgb = _rgb(theme.muted)
    if caption:
        cap_tf = _box(slide, x, y + h + 0.04, w, 0.22)
        _runs(
            cap_tf.paragraphs[0], caption, theme, {}, 11, theme.muted, theme.font_body,
            italic=True,
        )
        h += CAPTION_H_IN
    return h


def _set_alt(shape, alt: str, title: str | None = None) -> None:
    """Set a shape's accessible name (docPr title) and description (docPr
    descr) from an authored alt/caption string. python-pptx exposes no
    high-level alt-text setter for ANY shape type (the same reason
    diagram_pptx.py's raster fallback reaches into the oxml element
    directly), so this writes the docPr attributes by hand. Originally
    pictures only (`pic._element.nvPicPr.cNvPr`): an Image/Artifact's alt
    text reached HTML (alt=) and Markdown (![alt]) but never PPTX, a genuine
    accessibility gap invisible to every existing structural test because
    alt was never expected to render as slide TEXT in the first place.
    Generalized via `_element._nvXxPr` -- python-pptx names the non-visual
    properties container differently per shape type (nvPicPr for pictures,
    nvGraphicFramePr for charts/tables, nvSpPr for plain autoshapes,
    nvGrpSpPr for groups), but every one of those XML classes exposes the
    same `_nvXxPr` alias to its own container, so one call site now covers
    all of them. Best-effort: a metadata hiccup here must never fail the
    embed itself."""
    if not alt:
        return
    try:
        shape._element._nvXxPr.cNvPr.set("descr", alt)
        shape._element._nvXxPr.cNvPr.set("title", title or alt)
    except Exception:
        pass


def _add_png(slide, png: bytes | None, x, y):
    """Add PNG bytes as a picture shape; None if python-pptx refuses them.
    Best-effort: callers that must never render silently should check for
    None and fall back to _placeholder_block themselves (see
    _image_block_or_placeholder)."""
    if not png:
        return None
    try:
        return slide.shapes.add_picture(io.BytesIO(png), Inches(x), Inches(y))
    except Exception:
        return None


def _add_picture(slide, path: Path, theme, x, y):
    """Add an image file as a picture shape, rasterizing an SVG first:
    python-pptx has no SVG decoder, so without the optional rasterizer extra
    (docloom[diagrams]) an SVG cannot be embedded and None is returned, which
    keeps the caller's existing fallback (silent for _chart_block's own
    tiered fallback chain, a placeholder for _image_block_or_placeholder)."""
    if raster.is_svg(path):
        return _add_png(
            slide,
            raster.svg_file_to_png(
                path, width=RASTER_PX, font_files=raster.theme_font_files(theme)
            ),
            x, y,
        )
    try:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y))
    except Exception:
        return None


def _place_picture(slide, pic, caption, theme, x, y, w, max_h) -> float:
    """Fit an already-added picture into the (w, max_h) slot, center it, and
    draw its caption. Returns the height consumed, in inches."""
    cap_h = IMAGE_CAPTION_H_IN if caption else 0.0
    scale = min(
        Inches(w) / pic.width, Inches(max(0.4, max_h - cap_h)) / pic.height, 1.0
    )
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = Inches(x) + (Inches(w) - pic.width) // 2
    h = pic.height / 914400
    if caption:
        tf = _box(slide, x, y + h + 0.05, w, 0.22)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        _runs(
            tf.paragraphs[0], caption, theme, {}, 11, theme.muted, theme.font_body,
            italic=True,
        )
        h += cap_h
    return h


def _image_block(slide, b, theme, x, y, w, max_h) -> float:
    """Embed b.path as a picture; 0.0 (no shape at all) if there is nothing
    to embed or embedding fails. Used directly by _chart_block's own
    tiered fallback chain (which has a further fallback of its own -- a
    painted chart, then a data table -- so a placeholder here would
    pre-empt a strictly better outcome). Callers that want a visible
    placeholder instead of silence on failure should use
    _image_block_or_placeholder."""
    # slots carrying only query/asset_id render as nothing in v0.2 pptx: a
    # deliberate empty slot, never a failure, so always silent regardless
    # of which wrapper the caller uses.
    if not b.path:
        return 0.0
    path = Path(b.path)
    if not path.is_file():
        return 0.0
    pic = _add_picture(slide, path, theme, x, y)
    if pic is None:
        return 0.0
    _set_alt(pic, b.alt)
    return _place_picture(slide, pic, b.caption, theme, x, y, w, max_h)


def _image_block_or_placeholder(slide, b, theme, x, y, w, max_h) -> float:
    """Like _image_block, but a reference that looked resolvable and then
    failed to embed (a path given but the file is missing, or the file
    exists but python-pptx cannot decode it, e.g. an SVG without the raster
    extra) draws a visible placeholder and warns instead of the block just
    vanishing (P5 audit defect 1). A block with no path at all is left to
    _image_block's silent no-op: that is a deliberate empty slot, not a
    failure, matching the docx renderer's convention."""
    h = _image_block(slide, b, theme, x, y, w, max_h)
    if h > 0.0 or not b.path:
        return h
    warnings.warn(
        f"pptx: image could not be embedded ({b.path!r}); placeholder shown",
        stacklevel=2,
    )
    return _placeholder_block(slide, x, y, w, max_h, theme, b.alt, b.caption)


def _chart_data(b: Chart):
    """Chart block -> python-pptx chart data. Raises on anything the native
    chart path cannot represent; callers fall back to image/table."""
    if b.chart == "scatter":
        xs = [float(lb) for lb in b.labels]  # non-numeric labels -> ValueError
        if not xs:
            raise ValueError("scatter chart has no labels")
        data = XyChartData()
        for s in b.series:
            ser = data.add_series(s.name or "")
            for xv, yv in zip(xs, s.values):
                if yv is not None:
                    ser.add_data_point(xv, yv)
        return data
    n = max(len(b.labels), max((len(s.values) for s in b.series), default=0))
    if n == 0 or not b.series:
        raise ValueError("empty chart data")
    if b.chart == "pie" and len(b.series) > 1:
        # a native pie can only carry one series; raise so the caller falls
        # back to a data table instead of silently dropping the rest
        raise ValueError("pie chart supports only a single series")
    data = CategoryChartData()
    data.categories = (list(b.labels) + [""] * n)[:n]
    for s in b.series:
        # pad ragged series with None (blank points); pptx accepts None values
        data.add_series(s.name or "", (list(s.values) + [None] * n)[:n])
    return data


def _chart_table(b: Chart) -> Table:
    header: list[RichText] = [b.title or ""] + [
        s.name or f"Series {i + 1}" for i, s in enumerate(b.series)
    ]
    rows: list[list[RichText]] = []
    n = max([len(b.labels)] + [len(s.values) for s in b.series], default=0)
    for i in range(n):
        label = b.labels[i] if i < len(b.labels) else ""
        row: list[RichText] = [label]
        for s in b.series:
            v = s.values[i] if i < len(s.values) else None
            row.append("" if v is None else f"{v:g}")
        rows.append(row)
    return Table(header=header, rows=rows, caption=b.caption)


def _chart_block(slide, b: Chart, theme, numbers, x, y, w, max_h,
                 solo: bool = False) -> float:
    cap_h = CAPTION_H_IN if b.caption else 0.0
    try:
        data = _chart_data(b)
        room = max_h - cap_h
        # a chart alone on its slide fills the whole body instead of
        # stopping at the general cap: capping it there left a permanent
        # dead void below the chart (P5 audit defect 6).
        h = room if solo else min(room, LAYOUT["chart_max_h_in"])
        if h < 1.0:
            raise ValueError("not enough room for a native chart")
        frame = slide.shapes.add_chart(
            CHART_TYPE[b.chart], Inches(x), Inches(y), Inches(w), Inches(h), data
        )
        # Chart has no authored alt field (unlike Image/Artifact); the title
        # is the closest authored accessible name and the caption the closest
        # description, same reasoning as the table's descr above.
        _set_alt(frame, b.caption or b.title or f"{b.chart.capitalize()} chart",
                 title=b.title or "Chart")
        chart = frame.chart
        chart.font.name = theme.font_body
        chart.font.size = Pt(11)
        chart.font.color.rgb = _rgb(theme.text)
        if b.title:
            chart.has_title = True
            chart.chart_title.text_frame.text = b.title
        else:
            chart.has_title = False
        if b.chart == "pie" or len(b.series) > 1:
            chart.has_legend = True
            chart.legend.include_in_layout = False
        else:
            chart.has_legend = False
        if b.chart == "bar":
            # Office's native default draws horizontal-bar categories
            # first-at-bottom; every other docloom painter (chart_svg, used
            # by HTML/DOCX/Typst and this file's own fallback) draws
            # first-at-top. Flip only the category axis so the same Chart IR
            # reads the same way across every format.
            chart.category_axis.reverse_order = True
        try:
            _style_chart(chart, b, theme)
        except Exception:
            pass  # on-brand recolor is best-effort; keep the native chart
    except Exception:
        # fallback chain: pre-rendered image if present and embeddable, else a data table
        if b.path and Path(b.path).is_file():
            try:
                h = _image_block(
                    slide, Image(path=b.path, alt=b.title or "", caption=b.caption),
                    theme, x, y, w, max_h,
                )
            except Exception:
                h = 0.0  # unreadable: fall through to the table
            if h > 0.0:
                return h
            # unembeddable (corrupt, or SVG with no rasterizer installed):
            # fall through to the painted chart, then to the table
        pic = _add_png(
            slide,
            raster.svg_to_png(
                chart_svg.render_svg(b, theme),
                width=RASTER_PX,
                font_files=raster.theme_font_files(theme),
            ),
            x, y,
        )
        if pic is not None:  # docloom's own chart painter, rasterized
            return _place_picture(slide, pic, b.caption, theme, x, y, w, max_h)
        return _table_block(slide, _chart_table(b), theme, numbers, x, y, w, max_h)
    # cap_h was already reserved above (room = max_h - cap_h), so this always
    # fits and is always drawn -- see the same fix on _table_block/
    # _quote_block/_placeholder_block above.
    if b.caption:
        tf = _box(slide, x, y + h + 0.04, w, 0.22)
        _runs(
            tf.paragraphs[0], b.caption, theme, numbers,
            11, theme.muted, theme.font_body, italic=True,
        )
        h += CAPTION_H_IN
    return h


def _big_number_pt(value: str, w: float, max_h: float) -> int:
    """Pick the numeral's point size for the "big number" hero-stat
    treatment: target the oversized 96-150pt range and shrink only as far as
    the box actually forces -- a value string wide enough to wrap at w, or a
    max_h too tight to fit the numeral plus a label/delta strip beneath it.
    Deriving the size straight off `max_h` (the old `int(max_h * 34)`, capped
    120-150) meant the numeral's own reported natural height depended on
    whatever leftover room `_body`'s reserve math happened to hand this
    block, not on the value itself -- undersizing it whenever that room fell
    short of the ~4.4in the old formula assumed (defect item 3)."""
    label_room = 0.45  # rough space a label/delta strip needs beneath the numeral
    for pt in (150, 132, 116, 96):
        if _est_lines(value, pt, w) <= 1 and _line_h(pt) <= max(0.3, max_h - label_room):
            return pt
    return 96


def _big_number_block(slide, st, theme, x, y, w, max_h) -> float:
    """PINNED CONTRACT item 4: a StatRow with exactly one stat renders as ONE
    oversized numeral (a Gamma-style "big number" slide) instead of a single
    small stat card floating alone in a mostly-empty column -- unconditionally
    (item 3): a lone stat is always the hero treatment, whether or not it
    shares its slide with other blocks (the old `dominant`/solo gate is what
    let a StatRow sharing a slide with a caption paragraph fall through to
    the small compact card instead)."""
    value_pt = _big_number_pt(st.value, w, max_h)
    label_pt, delta_pt = 18, 14
    line_h = _line_h(value_pt)
    label_h = (label_pt * 1.3 / 72 + 0.08) if st.label else 0.0
    delta_h = (delta_pt * 1.3 / 72 + 0.06) if st.delta else 0.0
    total_h = line_h + label_h + delta_h
    h = min(max_h, total_h)
    # Top-aligned, not self-centered within `max_h`: the group's OWN natural
    # height (this same total_h, via _natural_h) is what _body's group-level
    # offset centers in the slide's available room, whether this block is
    # solo or shares the slide with others. Self-centering here too (as the
    # single-stat-solo case used to) double-counts that centering -- and,
    # worse, silently under-reports the height this block actually draws
    # into (only `total_h`, not the extra shift), letting a trailing sibling
    # block start high enough to overlap the shifted-down label/delta text.
    yy = y
    tf = _box(slide, x, yy, w, line_h)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _runs(
        tf.paragraphs[0], st.value, theme, {},
        value_pt, theme.primary, theme.font_heading, bold=True,
    )
    yy += line_h
    if st.label:
        tf = _box(slide, x, yy, w, label_h)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        _runs(tf.paragraphs[0], st.label, theme, {}, label_pt, theme.muted, theme.font_body)
        yy += label_h
    if st.delta:
        tf = _box(slide, x, yy, w, delta_h)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        delta_color = theme.muted if st.delta.strip().startswith("-") else theme.accent
        _runs(tf.paragraphs[0], st.delta, theme, {}, delta_pt, delta_color, theme.font_body)
    return h


def _stat_card_row(slide, items, theme, x, y, w, max_h, *, upgraded: bool) -> float:
    """A row of stat cards. `upgraded` (PINNED CONTRACT item 4: 2-4 stats
    sharing a slide as its dominant block) renders a larger numeral on a
    tinted accent card that scales with the room available (item 2's
    "fixed blocks own sparse slides"); otherwise this is the original
    compact card used for 5+ stats or a stat row sharing its slide with
    other content."""
    gap = LAYOUT["stat_gap_in"]
    if upgraded:
        h = max(LAYOUT["stat_card_h_in"], min(2.4, max_h))
        value_pt = 64 if len(items) == 2 else (56 if len(items) == 3 else 44)
        pad = 0.26
        fill = _mix(theme.background, theme.accent, 0.10)
    else:
        h = min(max_h, LAYOUT["stat_card_h_in"])
        value_pt = 24
        pad = 0.18
        fill = theme.surface
    h = min(h, max_h) if max_h > 0 else h
    cw = (w - gap * (len(items) - 1)) / len(items)
    for i, st in enumerate(items):
        cx = x + i * (cw + gap)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(y), Inches(cw), Inches(h)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb(fill)
        card.line.fill.background()
        card.shadow.inherit = False
        if upgraded:
            try:
                card.adjustments[0] = 0.08  # a subtler corner than the stock default
            except (IndexError, ValueError):
                pass
        value_h = _line_h(value_pt) + 0.1
        ty = y + pad
        tf = _box(slide, cx + pad, ty, cw - 2 * pad, value_h)
        _runs(
            tf.paragraphs[0], st.value, theme, {},
            value_pt, theme.primary, theme.font_heading, bold=True,
        )
        ty += value_h
        # label/delta raised from 11/10pt: footnote-scale on a 13.3in slide,
        # unreadable from the back of a room (P5 audit defect 12). The card
        # has the room: measured content bottoms out well inside its height.
        if ty + 0.26 <= y + h:
            tf = _box(slide, cx + pad, ty, cw - 2 * pad, 0.26)
            _runs(tf.paragraphs[0], st.label, theme, {}, 13, theme.muted, theme.font_body)
            ty += 0.32
        if st.delta and ty + 0.22 <= y + h:
            tf = _box(slide, cx + pad, ty, cw - 2 * pad, 0.22)
            delta_color = theme.muted if st.delta.strip().startswith("-") else theme.accent
            _runs(tf.paragraphs[0], st.delta, theme, {}, 12, delta_color, theme.font_body)
    return h


def _stats_block(slide, b: StatRow, theme, x, y, w, max_h, solo: bool = False) -> float:
    items = b.items[: LAYOUT["stat_max_cards"]]  # extras dropped: more cards than this don't fit legibly on one row
    if not items:
        return 0.0
    # PINNED CONTRACT item 4 -- Gamma-signature treatments via existing
    # blocks, no new Slide.layout: exactly one stat always becomes one
    # oversized numeral (item 3 -- a single stat reads as a hero statistic
    # regardless of whether it shares its slide with other blocks); 2-4
    # stats become an upgraded card row, but only when the stat row is truly
    # the slide's dominant block (solo -- the ONLY block on its slide/
    # column): the reserve-based layout in _body always leaves an earlier
    # block a generous max_h regardless of what shares the slide with it, so
    # max_h alone is not a reliable "dominant" signal there. A 5-up row, or a
    # 2-4 row sharing its slide with other blocks, keeps the original compact
    # card so it does not overpower its neighbors.
    dominant = solo
    if len(items) == 1:
        return _big_number_block(slide, items[0], theme, x, y, w, max_h)
    if 2 <= len(items) <= 4 and dominant:
        return _stat_card_row(slide, items, theme, x, y, w, max_h, upgraded=True)
    return _stat_card_row(slide, items, theme, x, y, w, max_h, upgraded=False)


def _block(slide, b: Block, theme, numbers, x, y, w, max_h,
          scale: float = 1.0, solo: bool = False) -> float:
    if isinstance(b, Heading):
        return _text_block(
            slide, b.text, theme, numbers, x, y, w, max_h,
            size=HEAD_PT[b.level], font=theme.font_heading, bold=True, scale=scale,
        )
    if isinstance(b, Paragraph):
        return _text_block(slide, b.text, theme, numbers, x, y, w, max_h, scale=scale)
    if isinstance(b, BulletList) and b.display == "grid" and 3 <= len(b.items) <= 6:
        return _bullet_grid_block(slide, b, theme, numbers, x, y, w, max_h, scale=scale)
    if isinstance(b, NumberedList) and b.display == "timeline" and 2 <= len(b.items) <= 6:
        return _numbered_timeline_block(slide, b, theme, numbers, x, y, w, max_h, scale=scale)
    if isinstance(b, (BulletList, NumberedList)):
        return _list_block(
            slide, b, theme, numbers, x, y, w, max_h, isinstance(b, NumberedList),
            scale=scale,
        )
    if isinstance(b, Quote):
        return _quote_block(slide, b, theme, numbers, x, y, w, max_h)
    if isinstance(b, Code):
        return _code_block(slide, b, theme, x, y, w, max_h)
    if isinstance(b, Table):
        return _table_block(slide, b, theme, numbers, x, y, w, max_h, solo=solo)
    if isinstance(b, Callout):
        return _callout_block(slide, b, theme, numbers, x, y, w, max_h, scale=scale)
    if isinstance(b, Image):
        return _image_block_or_placeholder(slide, b, theme, x, y, w, max_h)
    if isinstance(b, Chart):
        return _chart_block(slide, b, theme, numbers, x, y, w, max_h, solo=solo)
    if isinstance(b, StatRow):
        return _stats_block(slide, b, theme, x, y, w, max_h, solo=solo)
    if isinstance(b, Diagram):
        # Native, editable PPTX shapes (P2: docs/diagram-plan.md section 4b).
        # All layout/fit/font-floor-degradation/fallback/hash-stamp logic
        # lives in diagram_pptx.py; this hook only solves once at "full"
        # detail and hands off, matching add_diagram(d, solved, theme, ...).
        if not b.nodes:
            return 0.0
        try:
            solved = diagram_svg.solve_ir(
                b, diagram_pptx.theme_dict(theme),
                target_aspect=(w / max_h if max_h > 0 else 2.0),
            )
        except Exception:
            warnings.warn(
                f"pptx: diagram {b.id!r} failed to solve; placeholder shown",
                stacklevel=2,
            )
            return diagram_pptx.placeholder(slide, b, theme, x, y, w, max_h)
        return diagram_pptx.add_diagram(slide, b, solved, theme, x, y, w, max_h)
    if isinstance(b, Artifact):
        # An Artifact is a reference to content the author explicitly put on
        # the slide (a diagram, an infographic); unlike a plain Image slot
        # it never represents a deliberately-empty placeholder, so every
        # failure mode here draws a visible placeholder instead of the
        # block just vanishing with no trace (P5 audit defect 1, confirmed
        # live: an Artifact with no path rendered nothing at all).
        if b.path:
            h = _image_block_or_placeholder(
                slide, Image(path=b.path, alt=b.alt, caption=b.caption),
                theme, x, y, w, max_h,
            )
            if h > 0.0:
                return h
        warnings.warn(
            f"pptx: unresolved artifact (kind={b.kind!r}, no usable path); "
            "placeholder shown",
            stacklevel=2,
        )
        return _placeholder_block(slide, x, y, w, max_h, theme, b.alt, b.caption)
    if isinstance(b, Divider):
        _rect(slide, x, y + 0.06, w, 0.02, _divider_color(theme))
        return 0.14
    raise RenderError(f"unhandled block type {type(b).__name__}")


def _natural_h(b: Block, w: float, scale: float = 1.0,
               theme=None, max_h: float | None = None) -> float:
    """Estimate a block's natural (unclamped) height in inches, so _body can
    tell when a slide is underfull and rebalance the whitespace. `scale`
    mirrors the font-size growth _block applies for the same block types, so
    the grow pass can verify a candidate scale against this same formula.

    `theme`/`max_h` let a Diagram report its TRUE fitted height (a wide, short
    architecture diagram fits to the slide width and renders far shorter than
    the flat DIAGRAM_H_IN reserve): without this, the layout reserves 4.6in for
    a ~1.5in band and leaves the rest as dead space at the slide bottom instead
    of centering the content."""
    if isinstance(b, Heading):
        s = HEAD_PT[b.level] * scale
        return _est_lines(plain(b.text), s, w) * _line_h(s)
    if isinstance(b, Paragraph):
        s = BODY_PT * scale
        return _est_lines(plain(b.text), s, w) * _line_h(s)
    if isinstance(b, BulletList) and b.display == "grid" and 3 <= len(b.items) <= 6:
        # Mirrors _bullet_grid_block's own geometry (item 4b): the plain
        # per-item vertical-list formula below badly understates a 2x2/3x2
        # grid's real footprint (or, before the content-sized-card fix,
        # badly overstated it), throwing off _body's sparse/centering math.
        items = b.items[:6]
        cols, rows, _, card_h, gap, *_ = _grid_card_geometry(items, w, scale)
        return card_h * rows + gap * (rows - 1) if rows else 0.0
    if isinstance(b, NumberedList) and b.display == "timeline" and 2 <= len(b.items) <= 6:
        # Mirrors _numbered_timeline_block's own geometry (item 5): a fixed
        # horizontal band, not a per-item stack -- reporting the generic
        # list formula here (taller than the band actually is) is what let
        # the band "float" with large, uncentered dead space above/below it.
        return min(max_h, 2.0) if max_h else 2.0
    if isinstance(b, (BulletList, NumberedList)):
        s = BODY_PT * scale
        item_gap = LIST_ITEM_GAP_PT / 72 * scale
        return sum(
            _est_lines(plain(it.text), s, w - 0.35 * (it.level + 1))
            * _line_h(s) + item_gap
            for it in b.items
        )
    if isinstance(b, Quote):
        h = _est_lines(plain(b.text), 15, w - 0.4) * _line_h(15)
        return h + (QUOTE_ATTR_H_IN if b.attribution else 0.0)
    if isinstance(b, Code):
        return len(b.code.split("\n")) * _line_h(12) + 0.24
    if isinstance(b, Table):
        _, rows = normalize_table(b.header, b.rows)
        return (len(rows) + 1) * _table_row_h(TABLE_PT) + (CAPTION_H_IN if b.caption else 0.0)
    if isinstance(b, Callout):
        # Mirrors _callout_block's own geometry (item 5): size BODY_PT+1,
        # card width 85% of w, edge_w 0.1, pad 0.24 (2 * pad = 0.48).
        s = (BODY_PT + 1) * scale
        card_w = w * 0.85
        return _est_lines(plain(b.text), s, card_w - 0.1 - 0.48) * _line_h(s) + 0.48
    if isinstance(b, StatRow):
        if not b.items:
            return 0.0
        if len(b.items) == 1:
            # Mirrors _big_number_block's own geometry (item 3): a lone stat
            # is always the oversized hero-numeral treatment now, so its
            # natural height must reflect that -- not the flat compact-card
            # height, which used to make a genuinely ~3in numeral look like a
            # sparse ~1.4in block to the centering pass above.
            st = b.items[0]
            room = max_h if max_h else 3.0
            value_pt = _big_number_pt(st.value, w, room)
            total_h = _line_h(value_pt)
            if st.label:
                total_h += 18 * 1.3 / 72 + 0.08
            if st.delta:
                total_h += 14 * 1.3 / 72 + 0.06
            return total_h
        return LAYOUT["stat_card_h_in"]
    if isinstance(b, Chart):
        return LAYOUT["chart_max_h_in"] + (CAPTION_H_IN if b.caption else 0.0)
    if isinstance(b, Diagram):
        cap = CAPTION_H_IN if b.caption else 0.0
        # Aspect-aware: solve once and return the height the diagram will
        # ACTUALLY occupy after fitting to width `w`, capped at DIAGRAM_H_IN.
        # A wide LR pipeline fits to width and renders as a short band, so this
        # is much less than 4.6in -- letting _body center the content rather
        # than stranding it at the top of a phantom 4.6in reserve.
        if theme is not None and max_h and max_h > 0 and b.nodes:
            try:
                s = diagram_svg.solve_ir(
                    b, diagram_pptx.theme_dict(theme),
                    target_aspect=(w / max_h),
                )
                cw, ch = s.width / 96.0, s.height / 96.0
                if cw > 0 and ch > 0:
                    k = min(w / cw, max(0.1, max_h - cap) / ch)
                    return min(DIAGRAM_H_IN, k * ch) + cap
            except Exception:
                pass  # fall back to the flat reserve on any solve failure
        return DIAGRAM_H_IN + cap
    if isinstance(b, (Image, Artifact)):
        # a resolved image tends to fill much of the content area; estimate high
        # so a lone image centers near the top instead of floating low.
        cap = IMAGE_CAPTION_H_IN if b.caption else 0.0
        if b.path and Path(b.path).is_file():
            return 4.6 + cap
        # An unresolved Artifact still renders a real placeholder box now
        # (P5 audit defect 1), never nothing, so it must reserve real
        # layout room too -- an unresolved Image slot, unlike an Artifact,
        # stays a deliberate, genuinely weightless no-op (see
        # _image_block_or_placeholder), so it alone keeps the 0.0 estimate.
        return (1.6 + cap) if isinstance(b, Artifact) else 0.0
    if isinstance(b, Divider):
        return 0.14
    return _line_h(BODY_PT * scale)


def _grow_scale(blocks: list[Block], w: float, avail: float) -> float:
    """The underfull-slide grow factor for laying `blocks` out in `avail`
    inches: how much to scale up text-block font sizes so sparse content
    fills the space instead of floating small in the top third.

    Rendered height grows ~quadratically with font size (see _est_lines'
    chars-per-line and _line_h), so the candidate is a sqrt closed form, then
    verified with the same _natural_h formula the renderer uses and clamped
    down until every block actually fits: a block that fits at scale 1.0 must
    never be dropped by this feature.

    Capped at the smaller FIXED_NEIGHBOR_GROW_CAP (instead of the full
    GROW_CAP) when `blocks` also carries a fixed-size block (a table, code,
    a chart, stat cards, an image): a full GROW_CAP-scale grow of the prose
    next to a block whose size the content itself dictates is what produced
    23.8pt paragraphs beside 12pt code on the same slide (P5 audit defect
    4) -- but suppressing growth entirely there instead stranded the prose
    at 1.0x beside a large, deliberately-grown fixed block, reading smaller
    than the slide's own hierarchy intends. Modest growth (capped ~1.15)
    keeps the prose readable without rivaling the fixed block or the
    title."""
    if not blocks:
        return 1.0
    has_fixed = any(_is_fixed_size(b) for b in blocks)
    n = len(blocks)
    nat = [max(0.0, _natural_h(b, w)) for b in blocks]
    text_h = sum(h for b, h in zip(blocks, nat) if _is_growable(b))
    fixed_h = sum(nat) - text_h
    total = sum(nat) + GAP * (n - 1)
    if not (total < avail * 0.65 and text_h > 0.2):
        return 1.0
    cap = FIXED_NEIGHBOR_GROW_CAP if has_fixed else GROW_CAP
    # A prose-only slide targets ~90% fill; beside a fixed block the fixed
    # block itself already owns most of the "make this slide feel full"
    # job (see _table_block/_stats_block's own solo-fill growth), so the
    # prose only needs a lighter top-up.
    fill_target = 0.82 if has_fixed else 0.90
    target = avail * fill_target - fixed_h - GAP * (n - 1)
    candidate = (target / text_h) ** 0.5 if target > 0 else 1.0
    scale = max(1.0, min(cap, candidate))
    # Absolute, size-relative safety net independent of GROW_CAP: grown body
    # text must never approach the title's own size (P5 audit defect 4).
    scale = min(scale, MAX_GROWN_PT / BODY_PT)
    for _ in range(20):  # bounded: at most (GROW_CAP - 1) / 0.05 steps
        scaled_text_h = sum(
            _natural_h(b, w, scale) for b in blocks if _is_growable(b)
        )
        if scale <= 1.0 or scaled_text_h + fixed_h + GAP * (n - 1) <= avail:
            return scale
        scale = max(1.0, round(scale - 0.05, 2))
    return scale


def _body_top_offset(blocks: list[Block], theme, w, avail, scale: float) -> float:
    """Peek at the top offset `_body` would use to anchor `blocks` in `avail`
    inches at `scale`, without drawing anything. two_column uses this to
    compute ONE shared offset from its taller (denser) column -- the column
    with less slack -- and pass it to both `_body` calls, so paired headings
    in the two columns land at the same y instead of each column picking its
    own offset from its own (different) slack (item 1)."""
    if not blocks:
        return 0.0
    n = len(blocks)
    solo_chart = n == 1 and isinstance(blocks[0], Chart)
    nat = [max(0.0, _natural_h(b, w, theme=theme, max_h=avail)) for b in blocks]
    if solo_chart:
        cap_h = CAPTION_H_IN if blocks[0].caption else 0.0
        nat[0] = max(nat[0], avail - cap_h)
    text_h = sum(h for b, h in zip(blocks, nat) if _is_growable(b))
    fixed_h = sum(nat) - text_h
    scaled_text_h = text_h if scale <= 1.0 else sum(
        _natural_h(b, w, scale) for b in blocks if _is_growable(b)
    )
    scaled_total = scaled_text_h + fixed_h + GAP * (n - 1)
    if not (0 < scaled_total < avail):
        return 0.0
    slack = avail - scaled_total
    # Mirrors _body's own sparse/nearly-full split (item 1) so two_column's
    # shared offset actually CENTERS the two-column unit as a whole in the
    # available height (item 2) instead of only ever nudging it slightly off
    # the title -- the old flat `min(0.45, slack*0.2)` top-anchored both
    # columns regardless of how much of the slide they actually filled,
    # stranding a dead lower half under a short pair of columns.
    if scaled_total < avail * 0.9:
        return slack * 0.45
    return min(0.45, slack * 0.2)


def _body(slide, blocks: list[Block], theme, numbers, x, y, w,
         scale: float | None = None, top_offset: float | None = None) -> float:
    """Lay body blocks into [y, slide bottom]. When the content is sparse,
    anchor the block group close to its title (a small, capped offset from
    `y`) and spend the rest of any slack INSIDE the content -- larger seam
    gaps between logical groups, plus the grow pass's own font growth --
    instead of optical-centering the whole group and leaving most of the
    slack as dead space below the last block. That old (avail - used) * 0.42
    strand routinely left the bottom ~1/3+ of the column empty regardless of
    how full the slide already was: the #1 "generic PowerPoint" tell the
    audit named (item 1).

    `scale` grows text-block font sizes on an underfull slide; pass an
    explicit value to share one scale across multiple columns of a slide
    (two_column does), otherwise it is computed from `blocks` alone.
    `top_offset`, similarly, lets two_column pass one shared offset (from
    _body_top_offset, computed against its taller column) so both columns'
    titles/first blocks land at the same y; otherwise this slide/column
    computes its own from its own slack."""
    bottom = SLIDE_H - MARGIN
    if not blocks:
        return y
    avail = bottom - y
    n = len(blocks)
    # A lone chart fills the whole body (see _chart_block's solo mode)
    # instead of stopping at the general cap and leaving a dead void below
    # it (P5 audit defect 6); treat its natural height as "fills avail" so
    # the slack pass below doesn't also push it down first. A lone
    # table/stat-row is the other "fixed block owns a sparse slide" case
    # (item 2): _table_block/_stats_block grow themselves toward the
    # available room instead of floating small when told they are solo.
    solo = n == 1
    solo_chart = solo and isinstance(blocks[0], Chart)
    nat = [max(0.0, _natural_h(b, w, theme=theme, max_h=avail)) for b in blocks]
    if solo_chart:
        cap_h = CAPTION_H_IN if blocks[0].caption else 0.0
        nat[0] = max(nat[0], avail - cap_h)
    elif solo and isinstance(blocks[0], StatRow) and 2 <= len(blocks[0].items) <= 4:
        # The upgraded solo stat-card row (item 4) grows to ~min(2.4in, avail);
        # its _natural_h is the flat 1.4in compact height, so without this the
        # solo_visual centering (and its max_h cap) would size the row as a
        # smaller block than it actually draws, mis-centering it (re-audit #1).
        nat[0] = max(nat[0], min(2.4, avail))
    elif solo and isinstance(blocks[0], Table):
        # The solo table path (_table_block) grows font/row-height toward
        # budget*0.72; mirror that grown height here so centering doesn't push
        # the table down from an under-estimated natural height (re-audit #2).
        _, _trows = normalize_table(blocks[0].header, blocks[0].rows)
        _cap = CAPTION_H_IN if blocks[0].caption else 0.0
        nat[0] = max(nat[0], min((avail - _cap) * 0.72, (len(_trows) + 1) * 1.3) + _cap)
    # A solo Diagram/Image/Artifact has no seams to spend slack into (the
    # "anchor near the title, spend slack in bigger seams" convention below
    # needs multiple blocks) -- it is centered in the FULL available space
    # instead, which is the one site of vertical centering for it (paired
    # with the block loop below capping its own max_h at this same nat, so
    # the emitter's internal centering becomes a no-op). Without this
    # special case the small title-anchoring offset left a large, un-spent
    # void below the one block on the slide (item 1).
    # StatRow joins this centering group too: its own solo/dominant height
    # (item 4's big-number/upgraded-card treatment) is hard-clamped at
    # 2.4in, so unlike a table (which grows its own row padding to fill the
    # space -- see _table_block) it cannot grow far enough to make the
    # small-top-anchor-only convention below look intentional either.
    # A solo numbered-timeline is the same shape of problem as StatRow: its
    # own band height is naturally capped (~2in of nodes/dots/labels)
    # regardless of how much room the slide actually has, so it needs the
    # same centering rather than the "anchor near title" convention meant
    # for prose that can still grow into its slack.
    solo_timeline = (
        solo and isinstance(blocks[0], NumberedList) and blocks[0].display == "timeline"
    )
    solo_visual = (
        solo and (isinstance(blocks[0], (Diagram, Image, Artifact, StatRow)) or solo_timeline)
        and nat[0] >= MIN_VISUAL_BLOCK_H_IN
    )
    text_h = sum(h for b, h in zip(blocks, nat) if _is_growable(b))
    fixed_h = sum(nat) - text_h
    if scale is None:
        scale = _grow_scale(blocks, w, avail)
    scaled_text_h = text_h if scale <= 1.0 else sum(
        _natural_h(b, w, scale) for b in blocks if _is_growable(b)
    )
    scaled_total = scaled_text_h + fixed_h + GAP * (n - 1)

    # Seam gaps, one per block boundary, instead of one uniform gap: the seam
    # right after a Heading stays tighter than the rest (keeps a heading
    # visually attached to the list/paragraph it introduces) -- but neither
    # seam is ever INFLATED with slack. An earlier scheme spent ~78% of the
    # slide's slack space inflating these seams, which on a sparse two-group
    # slide (a callout + a bullet list, say) pushed the two groups to
    # opposite ends and stranded a wide dead band in the middle (item 1). The
    # group as a whole is instead kept as one cohesive, tightly-spaced unit
    # and the ENTIRE slack is spent positioning that unit in `avail`.
    seam_gaps = [GAP * 0.5 if isinstance(b, Heading) else GAP for b in blocks[:-1]]
    y0 = y
    if 0 < scaled_total < avail:
        slack = avail - scaled_total
        if solo_visual:
            offset = slack / 2 if top_offset is None else top_offset
        elif top_offset is not None:
            offset = top_offset
        elif scaled_total < avail * 0.9:
            # Sparse: the content unit is materially smaller than the room
            # it has. Optically center it -- a touch above true-center (45%
            # of the slack above, 55% below) reads more balanced under a
            # title than a dead-even split, and keeps the group from ever
            # pooling entirely into one contiguous top/bottom dead band.
            offset = slack * 0.45
        else:
            # Nearly full: stay anchored close to the title instead of
            # spending the (now small) remaining slack on a centering shift
            # that would barely be visible anyway.
            offset = min(0.45, slack * 0.2)
        y0 = y + max(0.0, offset)
    elif top_offset is not None:
        y0 = y + max(0.0, top_offset)

    # Finding A: this loop used to hand each block whatever remained after
    # its predecessors ("greedy" allocation), then drop the current block
    # outright ("break") the moment that remainder fell under 0.3in -- an
    # authored block (a chart's own trailing sibling, a subtitle-shrunk
    # slide's last paragraph) could vanish from the XML with zero trace and
    # no lint signal, because an EARLIER fixed-size block (a chart, capped
    # at its own max) was free to eat the entire budget first.
    #
    # Fix: reserve MIN_BLOCK_RESERVE_IN (+ its seam gap) for every block
    # still to come, so the block being laid out right now can never be
    # given more room than leaves its successors with nothing. A block that
    # only needs less than its reservation still gets exactly what it needs
    # (this only lowers the CEILING offered to the current block, not its
    # actual rendered size), so a comfortably-fitting slide is unaffected --
    # this only bites when the slide is genuinely too full, in which case an
    # earlier block (a chart, a table, a placeholder) shrinks/self-clamps to
    # the smaller ceiling instead of starving a later block to nothing.
    total_nat = sum(nat) + GAP * (n - 1) if n > 1 else sum(nat)
    if total_nat > avail + 0.01:
        warnings.warn(
            f"pptx: slide content (~{total_nat:.2f}in) exceeds the available "
            f"body height (~{avail:.2f}in); blocks are shrunk/squeezed to fit "
            "instead of being dropped -- consider splitting the slide",
            stacklevel=2,
        )
    yy = y0
    for i, b in enumerate(blocks):
        remaining = bottom - yy
        n_later = n - i - 1
        reserve = n_later * (MIN_BLOCK_RESERVE_IN + GAP)
        block_max_h = max(MIN_BLOCK_RESERVE_IN, remaining - reserve)
        if solo and isinstance(b, (Diagram, Image, Artifact)) and nat[i] >= MIN_VISUAL_BLOCK_H_IN:
            # A solo Diagram/Image/Artifact is the ONLY block on this
            # slide/column: the group-level `y0` offset above is already the
            # one site of vertical centering for it. Handing the emitter the
            # full remaining box here let it ALSO center internally
            # (diagram_pptx centers vertically within (w, max_h)), doubling
            # the visual gap into a large void (item 1). Capping max_h at the
            # block's own fitted natural height makes that inner centering a
            # no-op: the box IS the content.
            #
            # Only trusted when nat[i] itself already clears the legibility
            # floor: an unresolved Image (nat 0.0 -- a deliberate weightless
            # no-op, see _natural_h) or a diagram whose simple aspect-fit
            # estimate undershoots what its own font-floor-degradation ladder
            # actually needs must NOT be capped down here -- that starved
            # both of the room they need and caused them to be dropped
            # entirely by the legibility-floor check just below, instead of
            # embedding (or placeholder-warning) at their old, uncapped size.
            block_max_h = min(block_max_h, nat[i])
        if isinstance(b, _VISUAL_BLOCKS) and block_max_h < MIN_VISUAL_BLOCK_H_IN:
            # Squeezed past the point of legibility: drop it (never shrink
            # it into a speck) and say so by name, rather than emit a
            # picture/diagram/chart/stat row no one could ever read.
            ident = f" (id={b.id!r})" if getattr(b, "id", None) else ""
            warnings.warn(
                f"pptx: dropping a {type(b).__name__.lower()} block{ident} "
                f"that would render at only {block_max_h:.2f}in tall (floor "
                f"{MIN_VISUAL_BLOCK_H_IN}in) -- too small to be legible; "
                "split the slide or move it to its own slide",
                stacklevel=2,
            )
            continue  # no shape drawn, no space consumed
        yy += _block(slide, b, theme, numbers, x, yy, w, block_max_h, scale, solo=solo)
        if i < n - 1:
            yy += seam_gaps[i]
    return yy


# ----------------------------------------------------------------- layouts


def _slide_accent(s: Slide, theme: Theme) -> str:
    """Per-slide accent override for rules/edges; invalid hex falls back."""
    c = (s.accent or "").strip().lstrip("#")
    if len(c) == 6 and all(ch in "0123456789abcdefABCDEF" for ch in c):
        return "#" + c.upper()
    return theme.primary


def _usable_image(img: Image | None) -> bool:
    return img is not None and bool(img.path) and Path(img.path).is_file()


def _logo(
    slide, img: Image, theme: Theme, *,
    max_h: float = LOGO_MAX_H, corner: str = "top_right", scrim: bool = False,
) -> None:
    """Place a brand logo in a slide corner, scaled to `max_h` tall (small
    source logos upscale to fill it; the max_w guard still caps a very wide
    logo, keeping aspect ratio either way).

    `corner` is "top_right" (default) or "top_left" (image_right puts its
    image pane at top-right, so the logo moves to the opposite corner there
    instead of landing on the image). `scrim` draws a small contrasting
    plate behind the logo, reusing `_rect`, for full-bleed layouts (section's
    solid fill, hero's cover image) where the corner isn't a plain
    background and the logo would otherwise risk being illegible."""
    try:
        pic = slide.shapes.add_picture(str(img.path), 0, 0)
    except Exception:
        return  # unreadable image (e.g. an unembeddable SVG): skip, don't fail the render
    _set_alt(pic, img.alt)
    scale = min(Inches(max_h) / pic.height, Inches(LOGO_MAX_W) / pic.width)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    w_in, h_in = pic.width / 914400, pic.height / 914400
    left = MARGIN if corner == "top_left" else SLIDE_W - MARGIN - w_in
    top = MARGIN
    if scrim:
        # A stark theme.background (white) plate around a small logo reads
        # as a pasted-on sticker, not brand (P5 audit defect 11), especially
        # on the primary-blue section slides. A soft tint of theme.primary
        # with a larger pad reads as an intentional chip instead.
        pad = 0.14
        plate = _rect(
            slide, left - pad, top - pad, w_in + 2 * pad, h_in + 2 * pad,
            _tint(theme.primary, 0.12),
        )
        pic._element.addprevious(plate._element)  # plate behind the logo picture
    pic.top = Inches(top)
    pic.left = Inches(left)


def _doc_logo(slide, doc: Document, s: Slide, theme: Theme) -> None:
    """Stamp the document's brand logo in a slide corner, small and the same
    ~0.5in target size on every layout.

    section (solid theme.primary fill) and hero (full-bleed cover image, or a
    solid theme.primary fill when the slide has no usable image -- see
    _hero_slide) get a contrasting scrim plate behind the logo since the
    corner there is not a plain background. image_right's image pane sits at
    the top-right
    corner, so its logo moves to top-left instead of landing on the image;
    image_left's pane is already on the left, so the default top-right
    corner is safe as-is. The title layout places its own logo (see
    _title_slide, which also falls back to doc.logo) so it is skipped here
    to avoid stamping it twice."""
    if s.layout == "title":
        return
    if not _usable_image(doc.logo):
        return
    if s.layout == "section":
        _logo(slide, doc.logo, theme, max_h=LOGO_MAX_H, scrim=True)
    elif s.layout == "hero":
        # hero always paints a full-bleed backdrop now (a photo, or a solid
        # theme.primary fill when there is no usable image -- see
        # _hero_slide), never a plain page background, so the logo needs the
        # same contrast scrim in both cases.
        _logo(slide, doc.logo, theme, max_h=LOGO_MAX_H, scrim=True)
    elif s.layout == "image_right" and _usable_image(s.image):
        _logo(slide, doc.logo, theme, max_h=LOGO_MAX_H, corner="top_left")
    else:
        _logo(slide, doc.logo, theme, max_h=LOGO_MAX_H)


def _logo_reserve(doc: Document) -> float:
    """Extra right-edge width `_title_band` should leave blank so a top-right
    logo never overlaps title text. Based on the logo's REAL scaled width
    (a logo is height-capped to LOGO_MAX_H, so a normal ~square logo is only
    ~0.5in wide), not the full-slide LOGO_MAX_W cap, so a small logo does not
    needlessly carve down a narrow image-side title column."""
    if not _usable_image(doc.logo):
        return 0.0
    try:
        from PIL import Image as _PILImage  # python-pptx already depends on Pillow

        with _PILImage.open(doc.logo.path) as im:
            nw, nh = im.size
        scaled_w = min(LOGO_MAX_W, LOGO_MAX_H * (nw / nh)) if nh else LOGO_MAX_H
    except Exception:
        scaled_w = LOGO_MAX_H  # conservative small default if it can't be measured
    return scaled_w + GAP


def _title_band(
    slide, title: str | None, theme: Theme,
    x: float = MARGIN, w: float | None = None, accent: str | None = None,
    reserve: float = 0.0, top: float = 0.42,
) -> float:
    """`reserve` shrinks the title's width to dodge a corner logo sharing the
    same horizontal band; `top` instead pushes the whole band down to dodge a
    logo sitting directly above it. Callers should use whichever actually
    matches the logo's position: carving x as well as width for a logo that
    is above (not beside) the title is what misaligned a title against its
    own body text one column over (P5 audit defect 3)."""
    if not title:
        return MARGIN
    w = SLIDE_W - x - MARGIN if w is None else w
    w = max(1.5, w - reserve)  # room for a corner logo without crowding the text past legibility
    min_h, size = 0.62, LAYOUT["title_pt"]
    box_h = max(min_h, _est_lines(title, size, w) * _line_h(size))
    tf = _box(slide, x, top, w, box_h)
    _runs(
        tf.paragraphs[0], title, theme, {},
        size, theme.text, theme.font_heading, bold=True,
    )
    rule_y = top + box_h + 0.10
    # A ~1.2in kicker, not a full-width rule: a rule spanning the whole
    # title band read as a page-wide divider disconnecting the title from
    # its own body below it, rather than a mark that belongs to the title
    # (item 3). The tightened post-rule gap (0.252 -> 0.15) closes the same
    # gap the other way: the body now visibly belongs to its title instead
    # of floating a generic distance below it.
    kicker_w = min(1.2, w)
    _rect(slide, x, rule_y, kicker_w, 0.028, accent or theme.primary)
    return rule_y + 0.028 + 0.15


def _subtitle_line(slide, subtitle: str | None, theme: Theme, x: float, y: float, w: float) -> float:
    """Render a slide's s.subtitle as a muted caption line beneath its title
    band, in the same treatment _hero_slide/_section_slide/_title_slide
    already give it. Returns the height consumed (0.0, drawing nothing, if
    there is no subtitle) so callers can just do `y += _subtitle_line(...)`.

    _content_slide, _image_side_slide, and two_column used to never call
    this at all -- the only layouts, besides title/section/hero/quote, that
    can carry a subtitle, and s.subtitle was silently dropped on every one
    of them despite being present in the IR (same defect class as finding 8's
    quote-slide bug)."""
    if not subtitle:
        return 0.0
    size = 15
    h = _est_lines(subtitle, size, w) * _line_h(size) + SUBTITLE_PAD_IN
    tf = _box(slide, x, y, w, h)
    _runs(tf.paragraphs[0], subtitle, theme, {}, size, theme.muted, theme.font_body)
    return h


def _title_slide(slide, s: Slide, doc: Document, theme: Theme,
                 numbers: dict[str, int], accent: str) -> None:
    # brand logo (or any title-slide image) → corner; fall back to doc.logo
    # so a logo bound after generation still shows up here at export time
    logo_img = s.image if _usable_image(s.image) else doc.logo
    if _usable_image(logo_img):
        _logo(slide, logo_img, theme)
        if logo_img.caption:
            # This slot is deliberately treated as a small corner logo, not
            # a captioned figure -- there is no legible place to put a
            # caption next to a 0.5in mark. Warn by name instead of
            # silently dropping it (matching the "drawn, or warned about"
            # invariant), rather than squeeze illegible caption text next
            # to the logo just to say something got drawn.
            warnings.warn(
                f"pptx: title slide image caption {logo_img.caption!r} is "
                "not rendered (the image is shown as a small corner logo, "
                "not a captioned figure); move it into a body block if it "
                "must appear on the slide",
                stacklevel=2,
            )
    tx = 1.1
    tw = SLIDE_W - tx - MARGIN
    title_text = s.title or doc.title
    subtitle = s.subtitle or doc.subtitle
    byline = "  \u00b7  ".join(p for p in (", ".join(doc.authors), doc.date) if p)

    # Title: 54-64pt, a real display-scale cover title instead of the old
    # flat 40pt (the same size a content slide's own heading could read at,
    # a "generic PowerPoint" tell) -- shrinking only once the title actually
    # wraps past two lines (item 3).
    title_pt = 64
    while title_pt > 54 and _est_lines(title_text, title_pt, tw) > 2:
        title_pt -= 2
    title_h = max(1.4, _est_lines(title_text, title_pt, tw) * _line_h(title_pt))
    sub_h = (_est_lines(subtitle, 20, tw) * _line_h(20) + 0.18) if subtitle else 0.0
    by_h = (0.5 + 0.35) if byline else 0.0  # 0.5in gap under whatever precedes it, then its own line

    # Title + subtitle + byline as ONE vertically-centered group (not the
    # title alone pinned at a fixed y with the byline separately floor-
    # pinned to y>=6.3, which routinely orphaned the two from each other on
    # a short title) -- an optical-center bias matching the house
    # convention _body itself uses (item 3).
    group_h = title_h + sub_h + by_h
    ty = max(0.7, (SLIDE_H - group_h) / 2 - 0.3)

    # A real composition element -- a ~4in accent panel flanking the title
    # group -- replaces the old 0.3in full-slide-height sliver, which read
    # as a stray ruled line rather than an intentional shape (item 3).
    panel_h = min(4.2, group_h + 1.0)
    _rect(slide, 0, max(0.0, ty - 0.5), 0.14, panel_h, accent)

    tf = _box(slide, tx, ty, tw, title_h)
    _runs(
        tf.paragraphs[0], title_text, theme, {},
        title_pt, theme.text, theme.font_heading, bold=True,
    )
    y = ty + title_h
    if subtitle:
        tf = _box(slide, tx, y, tw, sub_h)
        _runs(tf.paragraphs[0], subtitle, theme, {}, 20, theme.muted, theme.font_body)
        y += sub_h
    if byline:
        y += 0.5
        tf = _box(slide, tx, y, tw, 0.35)
        _runs(tf.paragraphs[0], byline, theme, {}, 14, theme.muted, theme.font_body)
        y += 0.35
    # Finding 2: _title_slide never read s.blocks/s.right at all -- unlike
    # every other layout (including section and quote, both fixed in
    # earlier audit passes), a title slide's entire body block list vanished
    # from the deck with zero trace. There is no established "body zone" on
    # a cover slide, so give it the remaining room below whatever was
    # actually drawn (subtitle and/or byline), down to the slide's own
    # bottom margin, using the same _body layout every content-bearing
    # layout already uses.
    blocks = s.blocks + s.right
    if blocks:
        body_y = y + 0.25
        if body_y < SLIDE_H - MARGIN - MIN_BLOCK_RESERVE_IN:
            _body(slide, blocks, theme, numbers, tx, body_y, tw)
        else:
            warnings.warn(
                f"pptx: title slide has no room left for its {len(blocks)} "
                "body block(s) below the title/subtitle/byline; move them "
                "to a content slide",
                stacklevel=2,
            )


def _section_slide(slide, s: Slide, theme: Theme, numbers: dict[str, int],
                   index: int = 1) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(theme.primary)
    band_theme = _band_theme(theme, theme.primary)

    # An oversized, heavily-muted section numeral behind the title group so
    # a deck with several section dividers stops reading as identical
    # copies of the same slide (item 3). A subtle tint of theme.primary
    # itself (not a light/dark swap toward the background) keeps the
    # numeral visible-but-muted against the primary field.
    num_w = 4.4
    num_tf = _box(slide, SLIDE_W - MARGIN - num_w, 0.3, num_w, SLIDE_H - 0.6)
    num_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    num_tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    _runs(
        num_tf.paragraphs[0], f"{max(1, index):02d}", band_theme, {},
        150, _tint(theme.primary, 0.13), theme.font_heading, bold=True,
    )

    tw = SLIDE_W - 2 * MARGIN
    title_h = max(0.9, _est_lines(s.title or "", 36, tw) * _line_h(36))
    sub_h = (_est_lines(s.subtitle, 18, tw) * _line_h(18) + 0.15) if s.subtitle else 0.0
    rule_gap = 0.05 + 0.28
    group_h = rule_gap + title_h + sub_h

    # Vertically center the title group instead of pinning it at a fixed
    # y=3.0/y=4.25 -- a short one-line title used to leave the same amount
    # of dead space below it regardless of whether the slide carried any
    # body blocks at all (item 3).
    y = max(0.6, (SLIDE_H - group_h) / 2)
    _rect(slide, MARGIN, y, 1.2, 0.05, theme.accent)
    y += rule_gap
    tf = _box(slide, MARGIN, y, tw, title_h)
    _runs(
        tf.paragraphs[0], s.title or "", theme, {},
        36, theme.background, theme.font_heading, bold=True,
    )
    y += title_h
    if s.subtitle:
        tf = _box(slide, MARGIN, y, tw, sub_h)
        _runs(tf.paragraphs[0], s.subtitle, theme, {}, 18, theme.background, theme.font_body)
        y += sub_h
    # A section slide used to render only title/subtitle: any block the
    # author put in s.blocks/s.right (unlike every other layout, which at
    # least falls back through _content_slide) rendered nothing at all (P5
    # audit defect 2, a genuine data-loss bug). Draw them in the same
    # light-on-dark swap _hero_slide uses, since the background here is a
    # solid theme.primary fill, not the page background. `primary` also
    # swaps to `accent`: a bullet marker (drawn in theme.primary with
    # nothing behind it but the slide fill) would otherwise be theme.primary
    # on a theme.primary background -- genuinely invisible, not just low
    # contrast, caught by actually rendering and looking at this fix.
    blocks = s.blocks + s.right
    if blocks:
        _body(slide, blocks, band_theme, numbers, MARGIN, y + 0.3, tw)


def _quote_slide(slide, s: Slide, theme: Theme, numbers: dict[str, int]) -> None:
    """Two different authoring shapes both need to land on the same designed
    pull-quote, not just one of them: a real `Quote` block placed in
    s.blocks/s.right (the attribution living on the block itself), or the
    more obvious `Slide(layout="quote", title=..., subtitle=...)` with no
    Quote block at all. The latter used to fall straight through this
    function untouched: s.title rendered as a stray 16pt caption top-left,
    s.subtitle was never read anywhere in this function, and with q None
    nothing else drew at all -- a ~95% blank white slide with the
    attribution gone from the XML entirely (P5 audit finding 8). Normalize
    both shapes into one (quote_text, attribution, label) triple up front so
    there is exactly one pull-quote treatment: a full-height accent bar,
    real display-scale italic type, and the attribution set as a proper rule
    beneath it, matching the house style used by _quote_block elsewhere in
    this file.
    """
    body = s.blocks + s.right  # right column would otherwise be silently dropped
    q = next((b for b in body if isinstance(b, Quote)), None)
    rest = [b for b in body if b is not q]
    if q is not None:
        # A block-authored quote can still carry a slide title as a small
        # eyebrow label above it (e.g. title="Customer voice"); s.subtitle
        # only ever fills in as the attribution when the block itself didn't
        # set one, so it is never silently dropped either way -- except when
        # BOTH are authored: the block's own attribution wins and s.subtitle
        # is genuinely unused here (report renderers still keep it via
        # flatten_slides, so silence would be a cross-format drop). Warn by
        # name rather than let it vanish with no trace.
        if q.attribution and s.subtitle:
            warnings.warn(
                f"pptx: quote slide subtitle {s.subtitle!r} is not rendered "
                "(the Quote block already has its own attribution, which "
                "takes precedence); drop one of the two",
                stacklevel=2,
            )
        quote_text, attribution, label = q.text, q.attribution or s.subtitle, s.title
    else:
        # The obvious `title=`/`subtitle=` authoring with no Quote block:
        # the title IS the quote, the subtitle IS the attribution. No small
        # label is drawn in this shape -- the title already fills that role
        # as the quote itself, at full pull-quote scale.
        quote_text, attribution, label = s.title, s.subtitle, None
    if label:
        tf = _box(slide, MARGIN, 0.5, SLIDE_W - 2 * MARGIN, 0.4)
        _runs(tf.paragraphs[0], label, theme, {}, 16, theme.muted, theme.font_heading)
    y = 3.0
    if quote_text:
        qx, qw = 2.3, SLIDE_W - 4.6
        attr_h = 0.6 if attribution else 0.0
        avail = SLIDE_H - MARGIN - 1.4 - attr_h
        size = next(
            (pt for pt in (30, 24, 20, 16)
             if _est_lines(plain(quote_text), pt, qw) * _line_h(pt) <= avail),
            14,  # floor the ladder at 14pt; _fit_text_frames measures and bakes any further shrink
        )
        qh = min(avail, _est_lines(plain(quote_text), size, qw) * _line_h(size))
        y = max(1.4, (SLIDE_H - qh - 0.3 - attr_h) / 2)
        _rect(slide, qx - 0.35, y, 0.07, qh, theme.accent)
        tf = _box(slide, qx, y, qw, qh)
        _runs(tf.paragraphs[0], quote_text, theme, numbers, size, theme.text, theme.font_body,
              italic=True)
        y += qh + 0.25
        if attribution:
            tf = _box(slide, qx, y, qw, 0.35)
            _runs(
                tf.paragraphs[0], "\u2014 " + attribution, theme, {},
                16, theme.muted, theme.font_body,
            )
            y += 0.55
    if rest:
        _body(slide, rest, theme, numbers, MARGIN, y, SLIDE_W - 2 * MARGIN)


def _content_slide(slide, s: Slide, doc: Document, theme: Theme,
                   numbers: dict[str, int], accent: str) -> None:
    y = _title_band(slide, s.title, theme, accent=accent, reserve=_logo_reserve(doc))
    y += _subtitle_line(slide, s.subtitle, theme, MARGIN, y, SLIDE_W - 2 * MARGIN)
    _body(slide, s.blocks + s.right, theme, numbers, MARGIN, y, SLIDE_W - 2 * MARGIN)


def _cover_fit(pic, box_x: float, box_y: float, box_w: float, box_h: float,
               max_scale: float | None = None) -> None:
    """Scale `pic` to cover a box, cropping whichever dimension overflows
    instead of distorting the image to force-fit both dimensions. `max_scale`
    caps upscaling (pass 1.0 to never enlarge a low-res image beyond its
    native size; the axis that still doesn't reach the box is centered)."""
    img_w, img_h = pic.width, pic.height
    bw, bh = Inches(box_w), Inches(box_h)
    scale = max(bw / img_w, bh / img_h)
    if max_scale is not None:
        scale = min(scale, max_scale)
    disp_w, disp_h = img_w * scale, img_h * scale
    crop_x = 1 - bw / disp_w if disp_w > bw else 0.0
    crop_y = 1 - bh / disp_h if disp_h > bh else 0.0
    pic.crop_left = pic.crop_right = crop_x / 2
    pic.crop_top = pic.crop_bottom = crop_y / 2
    pic.width = int(min(disp_w, bw))
    pic.height = int(min(disp_h, bh))
    pic.left = Inches(box_x) + (bw - pic.width) // 2
    pic.top = Inches(box_y) + (bh - pic.height) // 2


def _contain_fit(pic, box_x: float, box_y: float, box_w: float, box_h: float,
                 max_scale: float | None = None) -> None:
    """Scale `pic` to fit entirely inside a box (no crop, no distortion) and
    center it. Use for content images (diagrams, charts, screenshots) where
    cropping an edge would lose information; the surrounding matte should be
    filled by the caller so the residual band reads as a frame, not an error."""
    img_w, img_h = pic.width, pic.height
    bw, bh = Inches(box_w), Inches(box_h)
    scale = min(bw / img_w, bh / img_h)
    if max_scale is not None:
        scale = min(scale, max_scale)
    pic.width = int(img_w * scale)
    pic.height = int(img_h * scale)
    pic.left = Inches(box_x) + (bw - pic.width) // 2
    pic.top = Inches(box_y) + (bh - pic.height) // 2


def _hero_slide(slide, s: Slide, doc: Document, theme: Theme, numbers: dict[str, int],
                accent: str) -> None:
    """A hero's whole design is a bold full-bleed backdrop behind a title
    band. That backdrop is s.image when there is one; when there is not (no
    image at all, or the file failed to embed), it used to fall all the way
    through to _content_slide -- a hero WITHOUT an image never reached this
    function at all (the dispatcher gated entry on _usable_image(s.image)),
    landing on a layout with no full-bleed treatment and (until fixed there
    too) that silently dropped s.subtitle from the XML entirely, leaving a
    slide like Slide(layout="hero", title=..., subtitle=...) ~90% blank.
    Fall back to a solid theme.primary fill instead -- the same full-bleed
    language _section_slide already uses -- so an imageless hero still gets
    the real hero treatment (the dark title band, subtitle, and blocks)
    rather than degrading to a generic bullet layout."""
    pic = None
    if _usable_image(s.image):
        try:
            pic = slide.shapes.add_picture(str(s.image.path), 0, 0)
        except Exception:
            pic = None  # unreadable file: fall back to the solid-fill hero below
    if pic is not None:
        _set_alt(pic, s.image.alt)
        _cover_fit(pic, 0.0, 0.0, SLIDE_W, SLIDE_H)
    else:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _rgb(theme.primary)
    blocks = s.blocks + s.right
    # A caption-only hero (image + caption, no title/subtitle/blocks) still has
    # content to draw: s.image.caption is authored content every other renderer
    # keeps, so the caption band below must be reached instead of returning here.
    if not (s.title or s.subtitle or blocks or (pic is not None and s.image and s.image.caption)):
        return
    # scrim: true fill transparency isn't exposed by python-pptx, so the title
    # sits in a solid theme.text band flush with the bottom edge instead
    tw = SLIDE_W - 2 * MARGIN
    pad = 0.3
    title_pt = LAYOUT["hero_title_pt"]
    title_h = (
        _est_lines(s.title, title_pt, tw) * _line_h(title_pt) if s.title else 0.0
    )
    sub_h = _est_lines(s.subtitle, 20, tw) * _line_h(20) + 0.08 if s.subtitle else 0.0
    # Finding 1 (strongest case of the silent-drop invariant): s.image.caption
    # is genuine authored content -- DOCX/HTML/MD all keep it (flatten_slides
    # turns a deck-only document's s.image into a real Image block for the
    # report renderers) -- but this function never drew it at all. Only
    # meaningful when there IS a picture (pic is not None); reserve its strip
    # up front like every other caption in this file.
    cap_h = 0.0
    if pic is not None and s.image.caption:
        cap_h = _est_lines(s.image.caption, 11, tw) * _line_h(11) + 0.08
    if pic is not None or not blocks:
        # A photo-backed hero (or an imageless one with no blocks at all)
        # keeps the original short caption-strip sizing: capped at 60% of
        # the slide so most of the photo stays visible, and a flat 1.7in
        # guess for blocks since they are meant to be a brief takeaway
        # layered over the photo, not a full body.
        band_h = min(SLIDE_H * 0.6, 2 * pad + title_h + sub_h + cap_h + (1.7 if blocks else 0.0))
    else:
        # Finding C: an imageless hero WITH blocks used to get the same
        # flat 1.7in guess meant for a short caption over a photo. There is
        # no photo to protect here (the whole slide is already a solid
        # theme.primary fill), so give blocks their real estimated height
        # -- the same formula every other layout uses -- instead of
        # crushing a diagram into a box so small its label font fell below
        # the legibility floor and silently degraded to an unreadable
        # raster image. Capped only by the slide's own bottom margin.
        blocks_h = sum(max(0.0, _natural_h(b, tw)) for b in blocks)
        blocks_h += GAP * max(0, len(blocks) - 1)
        band_h = min(SLIDE_H - MARGIN, 2 * pad + title_h + sub_h + cap_h + blocks_h)
    band_y = SLIDE_H - band_h
    scrim = _rect(slide, 0, band_y, SLIDE_W, band_h, theme.text)
    if pic is not None:
        # A real alpha-transparency band (item 6): an opaque band used to
        # chop the photo off entirely behind the title; a semi-transparent
        # one lets the image show through while the title/subtitle still
        # read clearly against it. An imageless hero has nothing to show
        # through (the whole slide is already a solid theme.primary fill),
        # so it keeps the fully-opaque band.
        _set_fill_alpha(scrim, 78)
    y = band_y + pad
    if s.title:
        th = min(title_h, SLIDE_H - pad - y)
        tf = _box(slide, MARGIN, y, tw, th)
        _runs(
            tf.paragraphs[0], s.title, theme, {},
            title_pt, theme.background, theme.font_heading, bold=True,
        )
        y += th + 0.08
    if s.subtitle and y + 0.3 <= SLIDE_H - pad:
        tf = _box(slide, MARGIN, y, tw, min(sub_h, SLIDE_H - pad - y))
        _runs(tf.paragraphs[0], s.subtitle, theme, {}, 20, theme.surface, theme.font_body)
        y += sub_h
    if cap_h > 0:
        # band_h already reserved exactly this much room (see cap_h above),
        # so clamp-and-draw like the subtitle above rather than gate on a
        # fixed epsilon: a hardcoded threshold here previously skipped the
        # caption on the exact boundary case band_h itself computed as
        # fitting (off by float slack, not a real space shortage).
        box_h = max(0.0, min(cap_h - 0.08, SLIDE_H - pad - y))
        if box_h > 0.02:
            tf = _box(slide, MARGIN, y, tw, box_h)
            _runs(
                tf.paragraphs[0], s.image.caption, theme, {}, 11, theme.surface,
                theme.font_body, italic=True,
            )
        y += cap_h
    if blocks:
        # blocks resolve their own fills/foregrounds against the band's
        # actual paint (theme.text), not the document's light background --
        # see _band_theme's docstring for why the old text/muted-only swap
        # shipped a contrast regression.
        _body(slide, blocks, _band_theme(theme, theme.text), numbers, MARGIN, y, tw)


def _image_side_slide(slide, s: Slide, doc: Document, theme: Theme,
                      numbers: dict[str, int], accent: str) -> None:
    pane_w = SLIDE_W * LAYOUT["image_pane_ratio"]
    left_side = s.layout == "image_left"
    px = 0.0 if left_side else SLIDE_W - pane_w
    pad = 0.35
    cap_h = IMAGE_CAPTION_H_IN if s.image.caption else 0.0
    pic = _add_picture(slide, Path(s.image.path), theme, px, 0.0)
    if pic is None:
        # Unembeddable file (an SVG with no rasterizer extra installed, or a
        # corrupt raster): previously this fell back to _content_slide, which
        # silently dropped the image, its alt text, and its caption, and drew
        # a full-width title that the image_right corner logo then overlapped
        # (a P5-class regression, since _content_slide never reads s.image).
        # Keep the side-pane layout and draw a placeholder in the pane instead,
        # so the logo/title geometry above still accounts for this being an
        # image layout, and nothing authored is silently lost.
        warnings.warn(
            f"pptx: image could not be embedded ({s.image.path!r}); "
            "placeholder shown", stacklevel=2,
        )
        box_h = min(1.6, SLIDE_H - 2 * pad - cap_h)
        ph_y = (SLIDE_H - (box_h + cap_h)) / 2
        _placeholder_block(
            slide, px + pad, ph_y, pane_w - 2 * pad, box_h + cap_h, theme,
            s.image.alt, s.image.caption,
        )
    else:
        _set_alt(pic, s.image.alt)
        # Contain-fit, not cover-fit: a side image is often a diagram or chart, and
        # cropping its edge would silently drop content. Never upscale past native.
        # Reserve the top/bottom pad and the caption strip BEFORE fitting: fitting
        # into the full SLIDE_H and only afterward re-anchoring the picture inside
        # the pad let a tall/portrait image's fitted height alone reach 7.5in, so
        # the pad push and the caption strip both landed off the bottom of the
        # slide (P5-class regression, measured live).
        avail_h = SLIDE_H - 2 * pad - cap_h
        _contain_fit(pic, px, pad, pane_w, avail_h, max_scale=1.0)
        # Matte a band that hugs the fitted image's height, not the full 7.5in
        # pane: contain-fitting a wide/short image (a banner) into this pane
        # otherwise leaves 80% of the pane as dead gray around a sliver of
        # image (P5 audit defect 7, measured live). _contain_fit already
        # centered the picture within the reduced-height box, so a matte band
        # the same height as the fitted image, centered the same way, encloses
        # it with a deliberate pad instead of a near-empty pane.
        #
        # Finding B: this layout never drew s.image.caption at all -- DOCX/HTML/
        # MD all keep it (flatten_slides turns a deck-only document's s.image
        # into a real Image block for the report renderers), only PPTX silently
        # dropped it. Reserve a caption strip below the image (same cap_h
        # literal _place_picture uses) inside the matte, instead of only
        # bracketing the image itself.
        fitted_h = pic.height / 914400
        matte_h = min(SLIDE_H, fitted_h + 2 * pad + cap_h)
        matte_y = (SLIDE_H - matte_h) / 2
        matte = _rect(slide, px, matte_y, pane_w, matte_h, theme.surface)
        pic._element.addprevious(matte._element)  # matte behind the picture
        # Re-anchor the picture to the top of its own pad within the (now
        # possibly taller) matte -- _contain_fit centered it across the
        # reduced-height box, which no longer matches once cap_h grows the
        # matte asymmetrically to make room for the caption strip below.
        pic.top = Inches(matte_y + pad)
        if s.image.caption:
            cap_tf = _box(slide, px, matte_y + pad + fitted_h + 0.05, pane_w, 0.22)
            cap_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            _runs(
                cap_tf.paragraphs[0], s.image.caption, theme, {}, 11, theme.muted,
                theme.font_body, italic=True,
            )
    tx = pane_w + MARGIN if left_side else MARGIN
    tw = SLIDE_W - pane_w - 2 * MARGIN
    logo_present = _usable_image(doc.logo)
    if left_side:
        # image_left keeps the logo top-right, over this pane's own outer
        # edge, so a width-only reserve (never carving x) keeps the title
        # band's left edge aligned with the body text below it.
        reserve = _logo_reserve(doc)
        y = _title_band(slide, s.title, theme, x=tx, w=tw, accent=accent, reserve=reserve)
    else:
        # image_right flips the logo to top-left, i.e. inside this same text
        # column: carving the title's x (the old behavior) staggered it
        # 0.66in right of the body text it introduces and cost it a wrapped
        # line for no reason (P5 audit defect 3, measured live). Push the
        # whole band down below the logo instead, at the column's real x/w.
        top = MARGIN + LOGO_MAX_H + GAP if (logo_present and s.title) else 0.42
        y = _title_band(slide, s.title, theme, x=tx, w=tw, accent=accent, top=top)
    # a titleless image slide starts its body at the very top, where the corner
    # logo also sits (top-left for image_right, top-right for image_left), so
    # push the body below the logo band to avoid overlapping it
    if not s.title and logo_present:
        y = max(y, MARGIN + LOGO_MAX_H + GAP)
    y += _subtitle_line(slide, s.subtitle, theme, tx, y, tw)
    _body(slide, s.blocks + s.right, theme, numbers, tx, y, tw)


def _render_slide(prs, blank, s: Slide, doc: Document, theme: Theme,
                  numbers: dict[str, int], section_index: int = 1) -> None:
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(theme.background)
    accent = _slide_accent(s, theme)
    if s.layout == "title":
        _title_slide(slide, s, doc, theme, numbers, accent)
    elif s.layout == "section":
        _section_slide(slide, s, theme, numbers, section_index)
    elif s.layout == "quote":
        _quote_slide(slide, s, theme, numbers)
    elif s.layout == "hero":
        _hero_slide(slide, s, doc, theme, numbers, accent)
    elif s.layout in ("image_left", "image_right") and _usable_image(s.image):
        _image_side_slide(slide, s, doc, theme, numbers, accent)
    elif s.layout == "two_column":
        y = _title_band(slide, s.title, theme, accent=accent, reserve=_logo_reserve(doc))
        y += _subtitle_line(slide, s.subtitle, theme, MARGIN, y, SLIDE_W - 2 * MARGIN)
        col_w = (SLIDE_W - 2 * MARGIN - 0.5) / 2
        avail = (SLIDE_H - MARGIN) - y
        col_scales = [
            _grow_scale(cols, col_w, avail) for cols in (s.blocks, s.right) if cols
        ]
        col_scale = min(col_scales) if col_scales else 1.0
        # ONE shared vertical offset from the TALLER (denser -- less slack,
        # so a smaller offset) column, instead of each _body call picking
        # its own offset from its own column's slack: the two calls used to
        # disagree by ~0.1-0.2in whenever the columns' content differed in
        # length, staggering the two columns' paired headings against each
        # other (item 1).
        col_offsets = [
            _body_top_offset(cols, theme, col_w, avail, col_scale)
            for cols in (s.blocks, s.right) if cols
        ]
        shared_offset = min(col_offsets) if col_offsets else 0.0
        _body(slide, s.blocks, theme, numbers, MARGIN, y, col_w, col_scale, shared_offset)
        _body(slide, s.right, theme, numbers, MARGIN + col_w + 0.5, y, col_w, col_scale,
              shared_offset)
    else:  # content, or an image layout without a usable image path
        _content_slide(slide, s, doc, theme, numbers, accent)
    _doc_logo(slide, doc, s, theme)
    if s.notes:
        slide.notes_slide.notes_text_frame.text = s.notes


def _sources_slide(prs, blank, doc: Document, theme: Theme) -> None:
    seen: dict[str, int] = {}
    lines: list[str] = []
    for src in doc.sources:
        if src.id in seen:
            continue  # dedupe so numbering matches the citation superscripts
        seen[src.id] = len(seen) + 1
        line = f"{seen[src.id]}. {src.title}"
        if src.publisher:
            line += f", {src.publisher}"
        if src.date:
            line += f" ({src.date})"
        if src.url:
            line += f", {src.url}"
        lines.append(line)

    # P5 audit defect 10: raise the size (12 -> 13 read as a footnote dump
    # under 5in of empty space) and hang-indent so a wrapped URL lines up
    # under the title text, not back under the number.
    size, item_gap, w = 13, 6 / 72, SLIDE_W - 2 * MARGIN
    hang = 0.34  # room for "12. " at this size before the hanging-indent wrap point
    title, tf, y, bottom, first = "Sources", None, 0.0, 0.0, True
    for line in lines:
        h = _est_lines(line, size, w - hang) * _line_h(size) + item_gap
        if tf is None or y + h > bottom:
            slide = prs.slides.add_slide(blank)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = _rgb(theme.background)
            top = _title_band(slide, title, theme)
            bottom = SLIDE_H - MARGIN
            tf = _box(slide, MARGIN, top, w, bottom - top)
            title, y, first = "Sources (cont.)", top, True
            # Every other slide in the deck carries the brand logo; the
            # sources slide was the sole exception (P5 audit defect 10).
            if _usable_image(doc.logo):
                _logo(slide, doc.logo, theme, max_h=LOGO_MAX_H)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(6)
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(Inches(hang)))
        pPr.set("indent", str(-Inches(hang)))
        _runs(p, line, theme, {}, size, theme.muted, theme.font_body)
        y += h


def _frame_paras(tf) -> list[textfit.ParaSpec] | None:
    """Extract measurable ParaSpecs from a text frame. Returns None when any
    run lacks an explicit size or font name -- _runs always sets both, so a
    frame that doesn't match that shape is skipped rather than mis-measured."""
    paras = []
    for p in tf.paragraphs:
        runs = []
        for r in p.runs:
            if r.font.size is None or not r.font.name:
                return None
            runs.append(textfit.RunSpec(
                r.text or "", r.font.name, r.font.size.pt,
                bool(r.font.bold), bool(r.font.italic),
            ))
        pPr = p._p.pPr  # may be None; marL/indent are EMU string attributes
        marL = int(pPr.get("marL", "0")) if pPr is not None else 0
        indent = int(pPr.get("indent", "0")) if pPr is not None else 0
        sa = p.space_after
        paras.append(textfit.ParaSpec(
            runs=tuple(runs),
            space_after_pt=sa.pt if sa is not None else 0.0,
            first_indent_in=max(0, marL + indent) / 914400,
            cont_indent_in=max(0, marL) / 914400,
        ))
    return paras


def _fit_one_frame(shape) -> None:
    tf = shape.text_frame
    bodyPr = tf._txBody.bodyPr  # CT_TextBody.bodyPr is a required child (verified)
    if bodyPr.find(qn("a:normAutofit")) is None:
        return  # only _box frames opt in; diagram/table/chart text never carries it
    inner_w = (shape.width - tf.margin_left - tf.margin_right) / 914400
    inner_h = (shape.height - tf.margin_top - tf.margin_bottom) / 914400
    if inner_w <= 0.05 or inner_h <= 0.05:
        return
    paras = _frame_paras(tf)
    if not paras:
        return
    need = textfit.required_height_pt(paras, inner_w, 1.0, 0.0)
    if need <= inner_h * 72.0 + FIT_SLACK_PT:
        return  # fits (within slack): leave today's XML byte-identical
    res = textfit.fit_scale(paras, inner_w, inner_h, MIN_FIT_PT)
    if res.scale >= 1.0 and res.lnspc_reduction == 0.0:
        return
    for p in tf.paragraphs:
        if res.lnspc_reduction:
            p.line_spacing = 1.0 - res.lnspc_reduction  # <a:lnSpc><a:spcPct val="90000"/>
        for r in p.runs:
            # floor to half-point, never round up past the measured fit; the
            # MIN_FIT_PT clamp is defense-in-depth against half-point
            # truncation ever landing a run under the legibility floor that
            # fit_scale's own (protected-run-derived) scale already honours
            # -- but ONLY for runs that started at/above the floor. A run
            # authored below MIN_FIT_PT to begin with (a citation
            # superscript) is deliberate small typography, not a legibility
            # bug: it must scale down proportionally with the rest of the
            # frame, not get bumped back up above its own authored size.
            orig_pt = r.font.size.pt
            scaled_pt = int(orig_pt * res.scale * 2) / 2
            if orig_pt >= MIN_FIT_PT - 1e-9:
                scaled_pt = max(MIN_FIT_PT, scaled_pt)
            r.font.size = Pt(scaled_pt)
    # After baking, fontScale MUST be 100000: PowerPoint multiplies run sz by
    # fontScale at render time, so echoing the applied scale here would shrink
    # the text a second time in PowerPoint while every other renderer draws
    # the baked size once. 100000 == "already fits; keep autofit for edits".
    for tag in ("a:normAutofit", "a:spAutoFit", "a:noAutofit"):
        for el in bodyPr.findall(qn(tag)):
            bodyPr.remove(el)
    bodyPr.insert(0, bodyPr.makeelement(
        qn("a:normAutofit"), {"fontScale": "100000", "lnSpcReduction": "0"},
    ))
    if not res.fits:
        warnings.warn(
            f"pptx: text still exceeds its box at the {MIN_FIT_PT:g}pt "
            "legibility floor; it is drawn at the floor size and may overflow "
            "-- split the slide or shorten the text",
            stacklevel=2,
        )


def _fit_text_frames(prs) -> None:
    """Measured autofit pass over every _box text frame (see textfit.py's
    module docstring for why the bare normAutofit python-pptx writes is a
    lie everywhere except interactive desktop PowerPoint). Best-effort per
    frame: a corrupt font file or unexpected XML shape must never lose a
    render, so each frame failure degrades to today's behavior."""
    if not TEXTFIT_ENABLED:
        return
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            try:
                _fit_one_frame(shape)
            except Exception:
                continue


def render(doc: Document, theme: Theme, out_path: Path) -> Path:
    if not doc.slides:
        raise RenderError("document has no slides; add slides[] to render pptx")
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]
    numbers = source_numbers(doc)
    section_n = 0
    for s in doc.slides:
        if s.layout == "section":
            section_n += 1
        _render_slide(prs, blank, s, doc, theme, numbers, section_n or 1)
    if doc.sources and cited_ids(doc):
        _sources_slide(prs, blank, doc, theme)
    _fit_text_frames(prs)
    prs.save(str(out_path))
    return Path(out_path)
