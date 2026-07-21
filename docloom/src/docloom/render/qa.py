"""Reference-free geometric QA pass over a built PPTX deck (AutoPresent/
SlidesBench-style metrics): insurance against shipping a broken demo deck.

Every function here is a PURE function over an already-built python-pptx
`Presentation` -- no Document IR, no Theme object, nothing but the shapes
actually drawn. This module is callable standalone (see the tests) and is
never imported by render/pptx.py: pptx.py owns rendering, this module owns
auditing whatever pptx.py produced. It is NOT guaranteed correct over decks
authored by an arbitrary other PPTX writer: group-shape geometry in
particular assumes the python-pptx convention of chOff/chExt == off/ext (see
_iter_flat), which a different writer is free to violate under OOXML. Findings
reuse lint.py's own Finding/Severity shape (rule/severity/where/message) so
callers can fold these into the same finding list an LLM already
self-corrects against.

Checks:
  * qa/off-slide       -- a shape's bbox extends past the slide bounds
  * qa/shape-overlap    -- two non-decorative shapes overlap by more than a
                           budget, once containment (a card's own text, a
                           node inside its group container, an edge's
                           accent bar) is filtered out
  * qa/font-family-sprawl / qa/palette-sprawl -- deck-wide palette discipline
  * qa/low-contrast     -- a text run's color against its resolved
                           background fill, WCAG 2 (not APCA -- APCA is
                           non-normative and was dropped from the WCAG 3
                           draft; the algorithm WCAG 3 will actually use is
                           still undetermined, so it is not a fit gate here)

CRITICAL: every finding below ships at severity="warning", never "error".
lint.py:349-365 documents a standing prohibition on this class of rule ever
hard-blocking export: cli.py refuses the whole deck (exit 2, no output) on
any error-severity finding, and artifacts.py returns HTTP 422. A reference-
free geometric heuristic -- exactly the kind of "this looks off" signal that
can misfire -- must never carry that authority. That is also why the
overlap rule below goes out of its way to filter containment rather than
just narrowing thresholds: a false-positive machine that cries wolf on every
card-with-text-on-it deck is worse than no check at all.
"""

from __future__ import annotations

from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..lint import Finding
from ..theme import contrast_ratio

EMU_PER_IN = 914400

# ------------------------------------------------------------- geometry

# A shape thinner than this in EITHER dimension carries no real visual area
# to collide with anything: a divider rule, a title underline, a diagram
# node's kind-accent bar, a zero-height connector bbox. Measured on real
# decks: the thinnest deliberate CONTENT shape (a legend color swatch, 0.12in
# square) sits comfortably above this, while every rule/underline/accent bar
# found sits at 0.01-0.04in.
DECORATIVE_MIN_DIM_IN = 0.05
DECORATIVE_MIN_DIM_EMU = int(DECORATIVE_MIN_DIM_IN * EMU_PER_IN)

# Two shapes whose overlap covers at least this fraction of the SMALLER
# shape's own area are treated as one being deliberately placed ON or INSIDE
# the other (a stat card's value text, a node inside its labeled group
# container, an edge's accent bar) -- legitimate composition, not a defect.
# Below this, but above MAX_OVERLAP_FRAC, is the genuinely ambiguous "these
# two things were not meant to share this much space" zone this rule targets.
CONTAINMENT_FRAC = 0.92
MAX_OVERLAP_FRAC = 0.15

# Rounding slack for the off-slide check: python-pptx/EMU math on a shape
# sized to land exactly on a slide edge (a full-bleed band, an edge-to-edge
# divider) can be off by a handful of EMU from float-to-int rounding; treat
# anything within this margin as flush with the edge, not bleeding past it.
BLEED_TOLERANCE_IN = 0.02
BLEED_TOLERANCE_EMU = int(BLEED_TOLERANCE_IN * EMU_PER_IN)

Bbox = tuple[int, int, int, int]  # (left, top, width, height), all EMU


def _in(emu: float) -> float:
    return emu / EMU_PER_IN


def _bbox(shape) -> Bbox | None:
    """`shape`'s absolute (left, top, width, height) in EMU, or None if any
    dimension is unset (an inherited-from-layout placeholder with no
    explicit position -- docloom's own renderers never leave one of these
    behind, but a QA pass over an arbitrary PPTX must not crash on one)."""
    try:
        l, t, w, h = shape.left, shape.top, shape.width, shape.height
    except Exception:
        return None
    if l is None or t is None or w is None or h is None:
        return None
    return int(l), int(t), int(w), int(h)


def _area(bbox: Bbox) -> int:
    _, _, w, h = bbox
    return max(0, w) * max(0, h)


def _intersection_area(a: Bbox, b: Bbox) -> int:
    al, at, aw, ah = a
    bl, bt, bw, bh = b
    ix = max(0, min(al + aw, bl + bw) - max(al, bl))
    iy = max(0, min(at + ah, bt + bh) - max(at, bt))
    return ix * iy


def _has_solid_fill(shape) -> bool:
    try:
        return shape.fill.type == MSO_FILL_TYPE.SOLID
    except Exception:
        return False


def _fill_hex(shape) -> str | None:
    """`shape`'s own solid fill color as "#RRGGBB", or None (no fill, a
    picture/line/chart/table frame with no simple fill, or a non-RGB fill
    e.g. a theme-color reference this reference-free pass cannot resolve)."""
    try:
        if shape.fill.type != MSO_FILL_TYPE.SOLID:
            return None
        color = shape.fill.fore_color
        if color.type != MSO_COLOR_TYPE.RGB:
            return None
        return "#" + str(color.rgb)
    except Exception:
        return None


def _is_group(shape) -> bool:
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.GROUP
    except Exception:
        return False


def _is_line_like(shape) -> bool:
    try:
        st = shape.shape_type
    except Exception:
        return False
    return st in (MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.FREEFORM)


def _is_decorative(shape, bbox: Bbox) -> bool:
    """A shape with no meaningful visual area to genuinely collide with
    anything: a divider rule, a title underline, a kind-accent bar, or a
    connector/freeform edge line (whose bounding box -- an elbow or diagonal
    route -- can span a large rectangle even though the actual drawn line is
    a thin stroke, so bbox-thinness alone would miss it; those are only ever
    unfilled routes in docloom's renderers, so "line/freeform with no solid
    fill" is treated the same as a thin rule)."""
    _, _, w, h = bbox
    if min(w, h) <= DECORATIVE_MIN_DIM_EMU:
        return True
    return _is_line_like(shape) and not _has_solid_fill(shape)


def _iter_flat(shapes, prefix: str = ""):
    """Yield (shape, where, bbox) for every LEAF shape under `shapes`, in
    draw (z/back-to-front) order, recursing into group shapes. A group
    shape's own bbox is always exactly the union of its children's (python-
    pptx's add_group_shape sets chOff/chExt equal to off/ext, so a child's
    raw left/top/width/height IS already its absolute position -- no
    transform math needed, checked empirically), so the group wrapper itself
    is never yielded: checking it would be entirely redundant with checking
    its children, and would falsely accuse the group of "overlapping"
    whichever sibling shape it happens to sit near.

    ASSUMES chOff/chExt == off/ext on every group, which holds for every
    group this codebase's own renderers construct (checked empirically) but
    is NOT an OOXML invariant: a group with a genuine child-offset transform
    (chOff/chExt != off/ext, which a different PPTX writer is free to emit)
    will have its children's raw left/top/width/height mis-measured as
    absolute slide coordinates by this function. This module does not detect
    or correct for that case -- see the module docstring."""
    for i, shape in enumerate(shapes):
        where = f"{prefix}shapes[{i}]"
        if _is_group(shape):
            yield from _iter_flat(shape.shapes, where + ".")
            continue
        bbox = _bbox(shape)
        if bbox is not None:
            yield shape, where, bbox


def _slide_base_bg(slide) -> str | None:
    """The slide's own background fill, resolved the same way _fill_hex
    resolves a shape's: solid RGB only. Every layout in render/pptx.py sets
    this explicitly (slide.background.fill.solid()), so it is almost always
    available; a hero with a photo backdrop or any slide left on the
    inherited master background returns None, and callers must treat that as
    "cannot verify" rather than guessing white."""
    try:
        fill = slide.background.fill
        if fill.type != MSO_FILL_TYPE.SOLID:
            return None
        color = fill.fore_color
        if color.type != MSO_COLOR_TYPE.RGB:
            return None
        return "#" + str(color.rgb)
    except Exception:
        return None


def check_bleed(prs, *, whitelist=None) -> list[Finding]:
    """Flag any shape extending past [0, 0, slide_width, slide_height].
    `whitelist(shape, slide_index) -> bool`, if given, lets a caller exempt
    a specific shape it KNOWS is an intentional full-bleed element."""
    out: list[Finding] = []
    sw, sh = prs.slide_width, prs.slide_height
    for si, slide in enumerate(prs.slides):
        for shape, where, bbox in _iter_flat(slide.shapes):
            if whitelist is not None and whitelist(shape, si):
                continue
            l, t, w, h = bbox
            r, b = l + w, t + h
            over = []
            if l < -BLEED_TOLERANCE_EMU:
                over.append(f"left by {_in(-l):.2f}in")
            if t < -BLEED_TOLERANCE_EMU:
                over.append(f"top by {_in(-t):.2f}in")
            if r > sw + BLEED_TOLERANCE_EMU:
                over.append(f"right by {_in(r - sw):.2f}in")
            if b > sh + BLEED_TOLERANCE_EMU:
                over.append(f"bottom by {_in(b - sh):.2f}in")
            if over:
                out.append(Finding(
                    rule="qa/off-slide", severity="warning",
                    where=f"slides[{si}].{where}",
                    message=f"shape extends past the slide edge: {', '.join(over)}; "
                            "either move/resize it, or whitelist it if the bleed "
                            "is intentional",
                ))
    return out


def check_overlap(
    prs, *, max_overlap_frac: float = MAX_OVERLAP_FRAC,
    containment_frac: float = CONTAINMENT_FRAC,
) -> list[Finding]:
    """Flag pairs of NON-DECORATIVE shapes whose bboxes intersect by more
    than `max_overlap_frac` of the smaller one's area -- UNLESS the overlap
    is at or above `containment_frac`, which means the smaller shape is
    (near enough) entirely inside the larger one: text sitting on a card,
    a node inside its labeled group container, a callout's edge bar inside
    its own fill. That is ordinary composition, not a layout defect, and is
    the single most important false-positive class this rule must not flag
    (a whole class of legitimate decks would otherwise fail every time)."""
    out: list[Finding] = []
    for si, slide in enumerate(prs.slides):
        leaves = [
            (shape, where, bbox) for shape, where, bbox in _iter_flat(slide.shapes)
            if _area(bbox) > 0 and not _is_decorative(shape, bbox)
        ]
        for i in range(len(leaves)):
            _s1, w1, b1 = leaves[i]
            a1 = _area(b1)
            for j in range(i + 1, len(leaves)):
                _s2, w2, b2 = leaves[j]
                a2 = _area(b2)
                inter = _intersection_area(b1, b2)
                if inter <= 0:
                    continue
                smaller = min(a1, a2)
                frac = inter / smaller
                if frac >= containment_frac or frac <= max_overlap_frac:
                    continue
                out.append(Finding(
                    rule="qa/shape-overlap", severity="warning",
                    where=f"slides[{si}].{w1} x {w2}",
                    message=f"shapes overlap across {frac * 100:.0f}% of the "
                            f"smaller shape's area (budget {max_overlap_frac * 100:.0f}%, "
                            "unless one fully contains the other); check the "
                            "slide for unintentionally colliding content",
                ))
    return out


# --------------------------------------------------------------- palette

# Grayscale (including pure black/white): channels within this many levels
# of each other. Neutrals do not count against the non-neutral fill budget
# -- chrome like table banding (background/surface) and body-text colors
# are not "the palette" a viewer perceives as brand color.
NEUTRAL_CHANNEL_SPREAD = 10
MAX_FONT_FAMILIES = 3
MAX_NON_NEUTRAL_FILLS = 6


def _is_neutral(hex_color: str) -> bool:
    h = hex_color.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return max(r, g, b) - min(r, g, b) <= NEUTRAL_CHANNEL_SPREAD


def _text_font_names(shape) -> set[str]:
    names: set[str] = set()
    try:
        if not shape.has_text_frame:
            return names
    except Exception:
        return names
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.name:
                names.add(run.font.name)
    return names


def _table_cells(shape):
    """Yield every (row, col) grid position of a table GraphicFrame that is
    a genuinely distinct cell, skipping the placeholder grid positions a
    merged range covers.

    Two things had to be fixed here, both about identity:

    1. python-pptx's Table.cell() CONSTRUCTS A NEW _Cell WRAPPER OBJECT ON
       EVERY CALL -- the wrapper is not cached, so id(cell) is the address
       of a just-allocated, immediately-garbage-collectable Python object.
       CPython routinely reuses a freed address for the very next
       allocation, so an `id(cell) in seen` dedup fires on totally
       unrelated cells (silently skipping cells that were never actually
       visited before) -- a previous version of this function had exactly
       that bug.

    2. The fix is NOT "dedup on cell._tc" -- checked empirically, every
       grid position in a merged range (origin AND the positions it spans)
       has its OWN distinct <a:tc> XML element; PowerPoint's merge model
       marks the spanned positions' own tc with hMerge/vMerge="1" rather
       than removing or aliasing them, so tc identity never collides
       either. The actual signal for "this grid position is a merge
       shadow, not a real cell" is `cell.is_spanned` (True for every
       position a merge-origin cell covers, False for the origin itself
       and for any ordinary unmerged cell) -- that is what this function
       filters on, not identity of any kind."""
    if not getattr(shape, "has_table", False):
        return
    table = shape.table
    for r in range(len(table.rows)):
        for c in range(len(table.columns)):
            cell = table.cell(r, c)
            if cell.is_spanned:
                continue
            yield r, c, cell


def check_palette(
    prs, *, max_font_families: int = MAX_FONT_FAMILIES,
    max_non_neutral_fills: int = MAX_NON_NEUTRAL_FILLS,
) -> list[Finding]:
    """Deck-wide palette discipline: too many font families or too many
    non-neutral fill colors reads as an unbranded, thrown-together deck."""
    fonts: set[str] = set()
    fills: set[str] = set()
    for slide in prs.slides:
        bg = _slide_base_bg(slide)
        if bg:
            fills.add(bg.upper())
        for shape, _where, _bb in _iter_flat(slide.shapes):
            fonts |= _text_font_names(shape)
            fill = _fill_hex(shape)
            if fill:
                fills.add(fill.upper())
            for _r, _c, cell in _table_cells(shape):
                fonts |= _text_font_names(cell)
                cfill = _fill_hex(cell)
                if cfill:
                    fills.add(cfill.upper())
    non_neutral = {c for c in fills if not _is_neutral(c)}
    out: list[Finding] = []
    if len(fonts) > max_font_families:
        out.append(Finding(
            rule="qa/font-family-sprawl", severity="warning", where="deck",
            message=f"{len(fonts)} distinct font families used ({sorted(fonts)}); "
                    f"keep to <= {max_font_families} for a consistent deck",
        ))
    if len(non_neutral) > max_non_neutral_fills:
        out.append(Finding(
            rule="qa/palette-sprawl", severity="warning", where="deck",
            message=f"{len(non_neutral)} distinct non-neutral fill colors used "
                    f"({sorted(non_neutral)}); keep to <= {max_non_neutral_fills} "
                    "for palette discipline",
        ))
    return out


# -------------------------------------------------------------- contrast

MIN_CONTRAST = 4.5  # WCAG 2 AA, normal text


def _run_fg_hex(run) -> str | None:
    try:
        color = run.font.color
        if color.type != MSO_COLOR_TYPE.RGB:
            return None
        return "#" + str(color.rgb)
    except Exception:
        return None


def _resolve_bg(ordered: list[tuple], idx: int, bbox: Bbox, base_bg: str | None) -> str | None:
    """The background a shape at `ordered[idx]` actually renders against:
    the nearest PRECEDING shape (i.e. drawn earlier, so it sits behind) that
    has its own solid fill and whose bbox contains this shape's bbox, walked
    nearest-first so a shape stacked directly on top of another (a callout's
    text on its fill rect, drawn right after it) resolves to that one and
    not something further back. Falls back to the slide's own background."""
    for k in range(idx - 1, -1, -1):
        other_shape, _where, other_bbox = ordered[k]
        area = _area(bbox)
        if area <= 0:
            continue
        if _intersection_area(bbox, other_bbox) / area < CONTAINMENT_FRAC:
            continue
        hexcolor = _fill_hex(other_shape)
        if hexcolor:
            return hexcolor
    return base_bg


def _check_runs_against_bg(
    shape, where: str, bg: str, min_ratio: float, si: int, out: list[Finding],
) -> None:
    for p_ix, para in enumerate(shape.text_frame.paragraphs):
        for r_ix, run in enumerate(para.runs):
            text = run.text
            if not text or not text.strip():
                continue
            fg = _run_fg_hex(run)
            if fg is None:
                continue
            ratio = contrast_ratio(fg, bg)
            if ratio < min_ratio:
                out.append(Finding(
                    rule="qa/low-contrast", severity="warning",
                    where=f"slides[{si}].{where}.paragraphs[{p_ix}].runs[{r_ix}]",
                    message=f"text {fg} on background {bg} is {ratio:.1f}:1 "
                            f"(WCAG 2 AA needs {min_ratio}:1): "
                            f'"{text.strip()[:40]}"',
                ))


def check_contrast(prs, *, min_ratio: float = MIN_CONTRAST) -> list[Finding]:
    """WCAG 2 (not APCA -- see module docstring) contrast between every text
    run's own color and its resolved background fill. Advisory: a shape
    sitting on an unresolvable background (a photo, a gradient, a theme-
    color reference this reference-free pass cannot look up) is skipped
    rather than guessed at, so this only ever reports what it can actually
    verify."""
    out: list[Finding] = []
    for si, slide in enumerate(prs.slides):
        base_bg = _slide_base_bg(slide)
        ordered = list(_iter_flat(slide.shapes))
        for idx, (shape, where, bbox) in enumerate(ordered):
            if getattr(shape, "has_table", False):
                for r, c, cell in _table_cells(shape):
                    cell_bg = _fill_hex(cell)
                    if cell_bg is None:
                        continue
                    _check_runs_against_bg(
                        cell, f"{where}.table.cell[{r}][{c}]", cell_bg,
                        min_ratio, si, out,
                    )
                continue
            try:
                has_text = shape.has_text_frame
            except Exception:
                has_text = False
            if not has_text:
                continue
            bg = _fill_hex(shape) or _resolve_bg(ordered, idx, bbox, base_bg)
            if bg is None:
                continue
            _check_runs_against_bg(shape, where, bg, min_ratio, si, out)
    return out


def audit(
    prs, *, bleed_whitelist=None,
    max_overlap_frac: float = MAX_OVERLAP_FRAC,
    containment_frac: float = CONTAINMENT_FRAC,
    max_font_families: int = MAX_FONT_FAMILIES,
    max_non_neutral_fills: int = MAX_NON_NEUTRAL_FILLS,
    min_contrast: float = MIN_CONTRAST,
) -> list[Finding]:
    """Run every geometric QA check over a built Presentation and return
    the combined finding list (all severity="warning" -- see module
    docstring). This is the one entry point most callers want; the
    individual check_* functions above stay public for callers/tests that
    want just one axis."""
    out: list[Finding] = []
    out += check_bleed(prs, whitelist=bleed_whitelist)
    out += check_overlap(
        prs, max_overlap_frac=max_overlap_frac, containment_frac=containment_frac,
    )
    out += check_palette(
        prs, max_font_families=max_font_families,
        max_non_neutral_fills=max_non_neutral_fills,
    )
    out += check_contrast(prs, min_ratio=min_contrast)
    return out
