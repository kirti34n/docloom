"""Tests for docloom.render.diagram_pptx (docs/diagram-plan.md section 4b,
P2: native PPTX shape emitter) and its single dispatch hook in pptx.py.

Structural verification (per docs/diagram-plan.md section 7(a)): unzip the
saved .pptx, parse slide XML, and assert real stCxn/endCxn connector glue
exists (proven here, not just asserted against a mock) plus a:tailEnd
arrowheads, the docloom:diagram: hash stamp on the group/shape name, all
shape coordinates within slide EMU bounds, and a fitted node-label font size
>= 8pt for an in-budget diagram. Also covers the font-floor degradation
ladder, the raster-picture fallback, the placeholder-of-last-resort, and
that nothing here ever raises.
"""

from __future__ import annotations

import re
import warnings
import zipfile
from io import BytesIO

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

from docloom import (
    Diagram, DiagramEdge, DiagramGroup, DiagramNode, Document, Slide, Theme,
    diagram_hash,
)
from docloom.render import diagram_pptx, diagram_svg, pptx as pptx_mod

SLIDE_W_IN, SLIDE_H_IN = 13.333, 7.5


# ---------------------------------------------------------------------------
# fixtures: small enough to clear the native 8pt font floor at a typical
# content-slide box, deliberately mixing all three non-ROUNDED_RECTANGLE
# node kinds (store/queue/external) so the shape-mapping and connection-site
# convention get exercised on every preset this module uses.
# ---------------------------------------------------------------------------


def _small_diagram(**overrides) -> Diagram:
    kwargs = dict(
        id="small",
        title="Edge to vault",
        direction="LR",
        nodes=[
            DiagramNode(id="a", label="Merchant", type="client"),
            DiagramNode(id="b", label="Gateway", type="service", group="g1"),
            DiagramNode(id="c", label="Queue", type="queue", group="g1"),
            DiagramNode(id="d", label="Vault", type="store"),
            DiagramNode(id="e", label="Card Network", type="external"),
        ],
        edges=[
            DiagramEdge(source="a", target="b", label="req", style="solid"),
            DiagramEdge(source="b", target="c", label="enqueue", style="dashed"),
            DiagramEdge(source="b", target="d", label="tokenize", style="secure"),
            DiagramEdge(source="d", target="e", label="authorize", style="emphasis"),
        ],
        groups=[DiagramGroup(id="g1", label="us-east-1", kind="region")],
        caption="Figure: edge to vault",
        alt="edge to vault diagram",
    )
    kwargs.update(overrides)
    return Diagram(**kwargs)


def _dense_diagram(n_nodes: int = 14) -> Diagram:
    """A diagram dense enough that no detail level on the native ladder can
    reach the 8pt node-label floor inside a normal content-slide box (the
    same "spec3-sized diagrams get sub-8pt labels" property documented in
    docs/diagram-plan.md section 5, P5.2) -- exercises the raster/placeholder
    fallback path."""
    nodes = [DiagramNode(id=f"n{i}", label=f"Service {i} with a longer name",
                          sublabel="detail line here", type="service")
             for i in range(n_nodes)]
    edges = [DiagramEdge(source=f"n{i}", target=f"n{i + 1}")
             for i in range(n_nodes - 1)]
    return Diagram(id="dense", title="Dense chain", direction="LR",
                   nodes=nodes, edges=edges, caption="dense fallback fixture")


def _new_slide():
    prs = Presentation()
    prs.slide_width = Emu(int(SLIDE_W_IN * 914400))
    prs.slide_height = Emu(int(SLIDE_H_IN * 914400))
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def _save_and_read_slide_xml(prs) -> str:
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    with zipfile.ZipFile(buf) as z:
        return z.read("ppt/slides/slide1.xml").decode("utf-8")


def _walk_shapes(shapes):
    for sh in shapes:
        if sh.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            yield from _walk_shapes(sh.shapes)
        else:
            yield sh


# --------------------------------------------------------------- theme_dict


def test_theme_dict_adapts_the_six_keys():
    t = Theme(primary="#112233", accent="#445566", surface="#778899",
              text="#000000", muted="#ABCDEF", background="#FFFFFF")
    td = diagram_pptx.theme_dict(t)
    assert td == {
        "primary": "#112233", "accent": "#445566", "surface": "#778899",
        "text": "#000000", "muted": "#ABCDEF", "background": "#FFFFFF",
    }


# ------------------------------------------------------- native connector glue


def test_native_diagram_has_real_connector_glue_matching_edge_count():
    d = _small_diagram()
    theme = Theme()
    solved = diagram_svg.solve(d, diagram_pptx.theme_dict(theme),
                               target_aspect=SLIDE_W_IN / 5.6)
    prs, slide = _new_slide()
    h = diagram_pptx.add_diagram(slide, d, solved, theme, 0.6, 1.0, 12.133, 5.6)
    assert h > 0.0
    xml = _save_and_read_slide_xml(prs)
    n_edges = len(d.edges)
    assert xml.count("<p:cxnSp>") == n_edges
    assert len(re.findall(r"<a:stCxn ", xml)) == n_edges
    assert len(re.findall(r"<a:endCxn ", xml)) == n_edges
    assert xml.count("a:tailEnd") == n_edges


def test_native_diagram_group_name_carries_tier1_tier2_hash_stamp():
    d = _small_diagram()
    theme = Theme()
    solved = diagram_svg.solve(d, diagram_pptx.theme_dict(theme))
    prs, slide = _new_slide()
    diagram_pptx.add_diagram(slide, d, solved, theme, 0.6, 1.0, 12.133, 5.6)
    names = [sh.name for sh in slide.shapes]
    expected = f"docloom:diagram:{d.id}:{diagram_hash(d)}"
    assert expected in names
    grp = next(sh for sh in slide.shapes if sh.name == expected)
    assert grp.shape_type == 6  # GROUP


def test_native_diagram_shapes_stay_within_slide_emu_bounds():
    d = _small_diagram()
    theme = Theme()
    solved = diagram_svg.solve(d, diagram_pptx.theme_dict(theme))
    prs, slide = _new_slide()
    diagram_pptx.add_diagram(slide, d, solved, theme, 0.6, 1.0, 12.133, 5.6)
    sw, sh_ = prs.slide_width, prs.slide_height
    for shp in _walk_shapes(slide.shapes):
        if shp.left is None:
            continue
        assert -1000 <= shp.left
        assert -1000 <= shp.top
        assert shp.left + shp.width <= sw + 1000
        assert shp.top + shp.height <= sh_ + 1000


def test_native_diagram_fitted_node_label_meets_8pt_floor():
    d = _small_diagram()
    theme = Theme()
    solved = diagram_svg.solve(d, diagram_pptx.theme_dict(theme))
    prs, slide = _new_slide()
    diagram_pptx.add_diagram(slide, d, solved, theme, 0.6, 1.0, 12.133, 5.6)
    seen = 0
    for shp in _walk_shapes(slide.shapes):
        if shp.name.startswith("docloom:node:") and shp.has_text_frame:
            run = shp.text_frame.paragraphs[0].runs[0]
            assert run.font.size.pt >= 8.0
            seen += 1
    assert seen == len(d.nodes)


def test_native_diagram_shape_kind_mapping():
    from pptx.enum.shapes import MSO_SHAPE

    d = _small_diagram()
    theme = Theme()
    solved = diagram_svg.solve(d, diagram_pptx.theme_dict(theme))
    prs, slide = _new_slide()
    diagram_pptx.add_diagram(slide, d, solved, theme, 0.6, 1.0, 12.133, 5.6)
    by_id = {}
    for shp in _walk_shapes(slide.shapes):
        if shp.name.startswith("docloom:node:"):
            by_id[shp.name.split(":", 2)[2]] = shp
    assert by_id["d"].auto_shape_type == MSO_SHAPE.FLOWCHART_MAGNETIC_DISK  # store
    assert by_id["c"].auto_shape_type == MSO_SHAPE.FLOWCHART_MULTIDOCUMENT  # queue
    assert by_id["e"].auto_shape_type == MSO_SHAPE.FLOWCHART_TERMINATOR    # external
    assert by_id["a"].auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE       # client
    assert by_id["b"].auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE       # service


def test_native_diagram_edge_labels_have_opaque_halo_fill():
    d = _small_diagram()
    theme = Theme()
    solved = diagram_svg.solve(d, diagram_pptx.theme_dict(theme))
    prs, slide = _new_slide()
    diagram_pptx.add_diagram(slide, d, solved, theme, 0.6, 1.0, 12.133, 5.6)
    texts = {
        shp.text_frame.text for shp in _walk_shapes(slide.shapes)
        if shp.has_text_frame and shp.fill.type is not None
        and shp.name.startswith("TextBox")
    }
    # every edge label text should appear on some opaque-filled textbox
    for e in d.edges:
        assert e.label in texts


def test_native_diagram_caption_drawn_when_room_available():
    d = _small_diagram()
    theme = Theme()
    # a wide target_aspect keeps the fitted diagram well short of the full
    # box height, leaving room for the caption below it.
    solved = diagram_svg.solve(d, diagram_pptx.theme_dict(theme),
                               target_aspect=12.133 / 3.0)
    prs, slide = _new_slide()
    diagram_pptx.add_diagram(slide, d, solved, theme, 0.6, 1.0, 12.133, 6.5)
    texts = [
        sh.text_frame.text for sh in slide.shapes
        if sh.has_text_frame and sh.shape_type != 6
    ]
    assert any(d.caption in t for t in texts)


def test_tb_direction_uses_top_bottom_connection_sites():
    d = Diagram(
        id="tb", direction="TB",
        nodes=[
            DiagramNode(id="a", label="Top"),
            DiagramNode(id="b", label="Middle"),
            DiagramNode(id="c", label="Bottom"),
        ],
        edges=[
            DiagramEdge(source="a", target="b"),
            DiagramEdge(source="b", target="c"),
        ],
    )
    theme = Theme()
    solved = diagram_svg.solve(d, diagram_pptx.theme_dict(theme),
                               target_aspect=4.0 / 6.0)
    assert solved.direction == "TB"
    prs, slide = _new_slide()
    h = diagram_pptx.add_diagram(slide, d, solved, theme, 0.6, 1.0, 8.0, 6.5)
    assert h > 0.0
    xml = _save_and_read_slide_xml(prs)
    assert xml.count("<p:cxnSp>") == len(d.edges)
    assert len(re.findall(r"<a:stCxn ", xml)) == len(d.edges)
    assert len(re.findall(r"<a:endCxn ", xml)) == len(d.edges)


def test_freeform_mode_builds_shapes_without_connector_glue():
    d = _small_diagram()
    theme = Theme()
    solved = diagram_svg.solve(d, diagram_pptx.theme_dict(theme))
    prs, slide = _new_slide()
    h = diagram_pptx.add_diagram(
        slide, d, solved, theme, 0.6, 1.0, 12.133, 5.6, mode="freeform"
    )
    assert h > 0.0
    xml = _save_and_read_slide_xml(prs)
    assert xml.count("<p:cxnSp>") == 0  # no connectors: freeform is not glued
    assert xml.count("a:tailEnd") == len(d.edges)  # arrowheads still drawn


# --------------------------------------------------------------- degradation


def test_font_floor_degradation_ladder_climbs_through_every_detail_level(
    monkeypatch,
):
    """A diagram whose full-detail fit misses the 8pt floor makes add_diagram
    climb the "full" -> "label+sub" -> "label" ladder (docs/diagram-plan.md
    section 4b), re-solving at each rung, before falling back. Asserted by
    recording every detail level solve() is actually called with, rather
    than depending on exact font metrics landing at a specific rung."""
    d = Diagram(
        id="ladder", title="t", direction="LR",
        nodes=[
            DiagramNode(id="a", label="A", sublabel="a fairly long sublabel that widens the box"),
            DiagramNode(id="b", label="B", sublabel="another long descriptive sublabel"),
        ],
        edges=[DiagramEdge(source="a", target="b")],
    )
    theme = Theme()
    td = diagram_pptx.theme_dict(theme)
    solved = diagram_svg.solve(d, td, target_aspect=3.0 / 5.6)
    real_solve = diagram_svg.solve
    seen_details = []

    def tracking_solve(*args, **kwargs):
        seen_details.append(kwargs.get("detail", "full"))
        return real_solve(*args, **kwargs)

    monkeypatch.setattr(diagram_pptx, "solve", tracking_solve)
    prs, slide = _new_slide()
    # a narrow box (3in) keeps the fitted label size below the 8pt floor at
    # every detail level for this fixture, so the ladder climbs all the way.
    h = diagram_pptx.add_diagram(slide, d, solved, theme, 0.6, 1.0, 3.0, 5.6)
    assert h > 0.0  # something was drawn (native or a fallback), never nothing
    # the ladder climbed through both non-"full" rungs, in order; a trailing
    # "full" call is the raster fallback's own full-fidelity re-solve.
    assert seen_details[:2] == ["label+sub", "label"]


# --------------------------------------------------------------- fallback


def test_dense_diagram_falls_back_to_raster_picture_when_available():
    if not __import__("docloom.render.raster", fromlist=["available"]).available():
        pytest.skip("resvg not installed in this environment")
    d = _dense_diagram()
    theme = Theme()
    solved = diagram_svg.solve(d, diagram_pptx.theme_dict(theme),
                               target_aspect=12.133 / 4.6)
    prs, slide = _new_slide()
    h = diagram_pptx.add_diagram(slide, d, solved, theme, 0.6, 1.0, 12.133, 4.6)
    assert h > 0.0
    xml = _save_and_read_slide_xml(prs)
    assert xml.count("<p:cxnSp>") == 0  # not native: no connectors at all
    assert xml.count("<p:pic>") == 1
    names = [sh.name for sh in slide.shapes]
    assert f"docloom:diagram:{d.id}:{diagram_hash(d)}" in names


def test_dense_diagram_falls_back_to_placeholder_without_raster(monkeypatch):
    monkeypatch.setattr(diagram_pptx.raster, "svg_to_png", lambda *a, **k: None)
    d = _dense_diagram()
    theme = Theme()
    solved = diagram_svg.solve(d, diagram_pptx.theme_dict(theme),
                               target_aspect=12.133 / 4.6)
    prs, slide = _new_slide()
    h = diagram_pptx.add_diagram(slide, d, solved, theme, 0.6, 1.0, 12.133, 4.6)
    assert h > 0.0
    xml = _save_and_read_slide_xml(prs)
    assert xml.count("<p:cxnSp>") == 0
    assert xml.count("<p:pic>") == 0
    names = [sh.name for sh in slide.shapes]
    assert f"docloom:diagram:{d.id}:{diagram_hash(d)}" in names


def test_placeholder_never_raises_and_shows_alt_text():
    d = _small_diagram(alt="a helpful alt description")
    theme = Theme()
    prs, slide = _new_slide()
    h = diagram_pptx.placeholder(slide, d, theme, 0.6, 1.0, 12.133, 5.6)
    assert h > 0.0
    texts = [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame]
    assert any("a helpful alt description" in t for t in texts)
    names = [sh.name for sh in slide.shapes]
    assert f"docloom:diagram:{d.id}:{diagram_hash(d)}" in names


def test_add_diagram_returns_zero_for_no_nodes():
    d = Diagram(id="empty", nodes=[], edges=[])
    theme = Theme()
    solved = diagram_svg.solve(
        Diagram(id="x", nodes=[DiagramNode(id="a", label="A")], edges=[]),
        diagram_pptx.theme_dict(theme),
    )
    prs, slide = _new_slide()
    h = diagram_pptx.add_diagram(slide, d, solved, theme, 0.6, 1.0, 12.133, 5.6)
    assert h == 0.0
    assert len(slide.shapes) == 0


def test_add_diagram_never_raises_when_resolve_ladder_itself_raises(monkeypatch):
    """If re-solving at a lower detail level raises (defensive: should not
    happen for a diagram lint already passed, but must never propagate),
    add_diagram falls through to the raster/placeholder path instead of
    crashing the whole render."""
    d = Diagram(
        id="tiny", direction="LR",
        nodes=[DiagramNode(id="a", label="A" * 80, sublabel="s" * 80)],
        edges=[],
    )
    theme = Theme()
    td = diagram_pptx.theme_dict(theme)
    # solve() at "full" succeeds and is passed in, but the internal ladder's
    # re-solve calls must never propagate even if diagram_svg.solve starts
    # raising after that first call.
    solved = diagram_svg.solve(d, td)
    real_solve = diagram_svg.solve

    calls = {"n": 0}

    def flaky_solve(*args, **kwargs):
        calls["n"] += 1
        if calls["n"] > 1:
            raise RuntimeError("boom")
        return real_solve(*args, **kwargs)

    monkeypatch.setattr(diagram_pptx, "solve", flaky_solve)
    prs, slide = _new_slide()
    # a tiny box forces the floor check to fail and the ladder to engage
    h = diagram_pptx.add_diagram(slide, d, solved, theme, 0.6, 1.0, 0.5, 0.5)
    assert h >= 0.0  # never raised


# --------------------------------------------------------------- pptx.py hook


def test_pptx_block_dispatch_renders_native_diagram(tmp_path):
    d = _small_diagram()
    doc = Document(title="T", slides=[Slide(layout="content", blocks=[d])])
    out = pptx_mod.render(doc, Theme(), tmp_path / "out.pptx")
    assert out.is_file()
    with zipfile.ZipFile(out) as z:
        xml = z.read("ppt/slides/slide1.xml").decode("utf-8")
    assert xml.count("<p:cxnSp>") == len(d.edges)
    assert f"docloom:diagram:{d.id}:{diagram_hash(d)}" in xml


def test_pptx_block_dispatch_falls_back_for_dense_diagram(tmp_path):
    d = _dense_diagram()
    doc = Document(title="T", slides=[Slide(layout="content", blocks=[d])])
    out = pptx_mod.render(doc, Theme(), tmp_path / "out.pptx")
    with zipfile.ZipFile(out) as z:
        xml = z.read("ppt/slides/slide1.xml").decode("utf-8")
    assert xml.count("<p:cxnSp>") == 0
    assert f"docloom:diagram:{d.id}:{diagram_hash(d)}" in xml


def test_pptx_block_dispatch_empty_diagram_is_silent(tmp_path):
    d = Diagram(id="empty", nodes=[], edges=[])
    doc = Document(title="T", slides=[Slide(layout="content", blocks=[d])])
    out = pptx_mod.render(doc, Theme(), tmp_path / "out.pptx")
    assert out.is_file()  # must not raise


def test_pptx_block_dispatch_warns_and_places_placeholder_on_solve_failure(
    tmp_path, monkeypatch
):
    d = _small_diagram()
    doc = Document(title="T", slides=[Slide(layout="content", blocks=[d])])

    def boom(*a, **k):
        raise RuntimeError("solve exploded")

    monkeypatch.setattr(pptx_mod.diagram_svg, "solve", boom)
    with warnings.catch_warnings(record=True) as caught:
        warnings.simplefilter("always")
        out = pptx_mod.render(doc, Theme(), tmp_path / "out.pptx")
    assert out.is_file()
    assert any("failed to solve" in str(w.message) for w in caught)
    with zipfile.ZipFile(out) as z:
        xml = z.read("ppt/slides/slide1.xml").decode("utf-8")
    assert f"docloom:diagram:{d.id}:{diagram_hash(d)}" in xml


def test_natural_h_mirrors_diagram_h_in_constant():
    d = _small_diagram()
    assert pptx_mod._natural_h(d, 10.0) == pytest.approx(
        pptx_mod.DIAGRAM_H_IN + 0.26
    )
    d_no_caption = _small_diagram(caption=None)
    assert pptx_mod._natural_h(d_no_caption, 10.0) == pytest.approx(
        pptx_mod.DIAGRAM_H_IN
    )


# --------------------------------------------------- docs/diagram-status.md
# regression tests for findings 3, 4, 5, 10, 11 (+ the legend part of 16)


def test_caption_renders_under_default_height_bound_fit_and_never_overlaps():
    """Finding 5: diagram_pptx.py:468 guarded the caption with
    "if d.caption and h + 0.26 <= max_h_in", but h = min(max_h_in,
    canvas_h_in * k) and k is solved so that whenever height binds,
    canvas_h_in * k == max_h_in exactly -- the guard was structurally
    always false. solve()'s own default target_aspect (no override, unlike
    most fixtures in this file which deliberately dodge the bug with a wide
    target_aspect) reproduces the height-bound case directly.

    Also guards the companion bug caught only by rendering through
    LibreOffice and looking at the pixels: the diagram is vertically
    centered inside its box whenever width binds, so the caption must sit
    below the diagram's TRUE rendered bottom edge (the group shape's own
    top + height), not a naive y_in + canvas_h_in * k that ignores the
    centering offset -- otherwise the caption text lands across the
    diagram's own lower half instead of below it."""
    d = _small_diagram()
    theme = Theme()
    solved = diagram_svg.solve(d, diagram_pptx.theme_dict(theme))
    prs, slide = _new_slide()
    diagram_pptx.add_diagram(slide, d, solved, theme, 0.6, 1.0, 12.133, 5.6)
    grp = next(sh for sh in slide.shapes if sh.shape_type == 6)  # GROUP
    caption = next(
        sh for sh in slide.shapes
        if sh.has_text_frame and d.caption in (sh.text_frame.text or "")
    )
    assert caption.top >= grp.top + grp.height - Emu(1000)  # EMU rounding slack


def test_raster_fallback_sets_descr_from_alt_text():
    """Finding 4: the fallback picture carried 0 text runs and
    descr="image.png" (python-pptx's own filename-derived default), so a
    screen reader announced the filename instead of the multi-sentence
    Diagram.alt. python-pptx has no high-level alt-text setter; this
    asserts the descr attribute directly."""
    if not __import__("docloom.render.raster", fromlist=["available"]).available():
        pytest.skip("resvg not installed in this environment")
    d = _dense_diagram().model_copy(
        update={"alt": "a screen-reader description of the dense chain"}
    )
    theme = Theme()
    solved = diagram_svg.solve(d, diagram_pptx.theme_dict(theme),
                               target_aspect=12.133 / 4.6)
    prs, slide = _new_slide()
    diagram_pptx.add_diagram(slide, d, solved, theme, 0.6, 1.0, 12.133, 4.6)
    pic = next(sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE)
    assert pic._element.nvPicPr.cNvPr.get("descr") == (
        "a screen-reader description of the dense chain"
    )
    assert pic._element.nvPicPr.cNvPr.get("descr") != "image.png"


def test_raster_fallback_warns_with_the_fitted_font_size():
    """Finding 3: the 8pt font floor correctly rejected the diagram from
    the native path, then handed it to a raster fallback with no floor of
    its own and no warning -- render(doc, 'pptx') emitted zero warnings
    even though labels came out well under 8pt. Verified here the same way
    the finding was verified: with warnings.simplefilter("always")."""
    if not __import__("docloom.render.raster", fromlist=["available"]).available():
        pytest.skip("resvg not installed in this environment")
    d = _dense_diagram()
    theme = Theme()
    solved = diagram_svg.solve(d, diagram_pptx.theme_dict(theme),
                               target_aspect=12.133 / 4.6)
    prs, slide = _new_slide()
    with warnings.catch_warnings(record=True) as caught:
        warnings.simplefilter("always")
        diagram_pptx.add_diagram(slide, d, solved, theme, 0.6, 1.0, 12.133, 4.6)
    msgs = [str(w.message) for w in caught]
    assert any("does not clear the 8pt node-label floor" in m for m in msgs)
    assert any(d.id in m for m in msgs)


def test_native_diagram_renders_legend_kind_swatches_and_edge_style_key():
    """Finding 11: "legend" appeared 4 times in diagram_svg.py and 0 times
    in diagram_pptx.py -- native slides had no legend and no kind bars, so
    the same Diagram block looked like two different products depending on
    whether it went native or raster."""
    d = _small_diagram()
    theme = Theme()
    solved = diagram_svg.solve(d, diagram_pptx.theme_dict(theme))
    prs, slide = _new_slide()
    diagram_pptx.add_diagram(slide, d, solved, theme, 0.6, 1.0, 12.133, 5.6)
    texts = {
        sh.text_frame.text for sh in _walk_shapes(slide.shapes) if sh.has_text_frame
    }
    for kind in sorted({n.type for n in d.nodes}):
        assert kind in texts
    for name in ("flow", "async / return", "primary path", "secure"):
        assert name in texts


def test_native_diagram_kind_bar_drawn_for_every_node_except_external():
    """Finding 11: mirrors diagram_svg.node_shape()'s barpath() accent,
    which every kind except "external" gets."""
    d = _small_diagram()
    theme = Theme()
    solved = diagram_svg.solve(d, diagram_pptx.theme_dict(theme))
    prs, slide = _new_slide()
    diagram_pptx.add_diagram(slide, d, solved, theme, 0.6, 1.0, 12.133, 5.6)
    n_rects = sum(
        1 for sh in _walk_shapes(slide.shapes)
        if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and sh.name.startswith("Rectangle")  # python-pptx's default autoshape name
    )
    n_bar_worthy = sum(1 for n in d.nodes if n.type != "external")
    # every bar-worthy node contributes exactly one accent-bar rectangle;
    # default-named rectangles also include the title rule, the legend
    # separator and each legend chip's own bar, so assert a lower bound
    # rather than an exact count.
    assert n_rects >= n_bar_worthy


def test_secure_edge_honors_dash_dot_natively():
    """Finding 11 (second half): diagram_pptx.py:413 unpacked EDGE_STYLE's
    dash spec into a throwaway `_dash` and discarded it, so a "secure" edge
    rendered as a solid line natively while the SVG path drew it
    dash-dot. Verified via the raw XML since the connector shape itself
    isn't individually addressable by python-pptx's high-level API."""
    d = _small_diagram()  # has exactly one style="secure" edge (b -> d)
    theme = Theme()
    solved = diagram_svg.solve(d, diagram_pptx.theme_dict(theme))
    prs, slide = _new_slide()
    diagram_pptx.add_diagram(slide, d, solved, theme, 0.6, 1.0, 12.133, 5.6)
    xml = _save_and_read_slide_xml(prs)
    # one for the secure edge itself, one for the legend's "secure" key
    assert xml.count('<a:prstDash val="dashDot"/>') >= 2


def test_store_node_uses_left_mid_connection_site_as_lr_target():
    """Finding 10: at 300 DPI the "write txn" connector doglegged up and
    terminated on top of the Ledger cylinder's cap pointing sideways,
    never entering the body. Verified empirically (rendered every
    FLOWCHART_MAGNETIC_DISK connection-site index 0-3 through LibreOffice):
    index 2 is the clean left-mid entry point on the cylinder body; index 1
    (the rectangle convention this module used to apply uniformly) lands
    near the top of the cap instead."""
    d = _small_diagram()  # edge c(service) -> d(store), direction LR
    theme = Theme()
    solved = diagram_svg.solve(d, diagram_pptx.theme_dict(theme))
    prs, slide = _new_slide()
    diagram_pptx.add_diagram(slide, d, solved, theme, 0.6, 1.0, 12.133, 5.6)
    xml = _save_and_read_slide_xml(prs)
    store_shape_id = next(
        sh.shape_id for sh in _walk_shapes(slide.shapes) if sh.name == "docloom:node:d"
    )
    m = re.search(
        rf'<a:endCxn id="{store_shape_id}" idx="(\d+)"/>', xml
    )
    assert m is not None
    assert m.group(1) == "2"


def test_store_node_uses_bottom_connection_site_as_lr_source():
    """Finding 10, continued: FLOWCHART_MAGNETIC_DISK has no right-mid
    connection site at all (verified empirically across every in-range
    index); a store used as an LR source has to leave from the bottom
    (index 3) instead, the only index that produces a single clean elbow
    rather than a connector that loops backward around the shape."""
    d = _small_diagram()  # edge d(store) -> e(external), direction LR
    theme = Theme()
    solved = diagram_svg.solve(d, diagram_pptx.theme_dict(theme))
    prs, slide = _new_slide()
    diagram_pptx.add_diagram(slide, d, solved, theme, 0.6, 1.0, 12.133, 5.6)
    xml = _save_and_read_slide_xml(prs)
    store_shape_id = next(
        sh.shape_id for sh in _walk_shapes(slide.shapes) if sh.name == "docloom:node:d"
    )
    m = re.search(
        rf'<a:stCxn id="{store_shape_id}" idx="(\d+)"/>', xml
    )
    assert m is not None
    assert m.group(1) == "3"


def test_store_node_uses_top_bottom_connection_sites_in_tb():
    """Finding 10, TB direction: a store's top (index 0) and bottom (index
    3) connection sites happen to coincide with the rectangle convention
    already used for TB, but this locks that in explicitly rather than by
    coincidence, since LR does not coincide (see the two tests above)."""
    d = Diagram(
        id="tbstore", direction="TB",
        nodes=[
            DiagramNode(id="a", label="Ingest"),
            DiagramNode(id="b", label="Store", type="store"),
            DiagramNode(id="c", label="Reporting"),
        ],
        edges=[
            DiagramEdge(source="a", target="b"),
            DiagramEdge(source="b", target="c"),
        ],
    )
    theme = Theme()
    solved = diagram_svg.solve(d, diagram_pptx.theme_dict(theme), target_aspect=0.6)
    prs, slide = _new_slide()
    # a generous box, scaled to the solved canvas, keeps this native rather
    # than falling back to raster (the connection-site XML this test checks
    # for only exists on the native path).
    w_in, h_in = solved.width / 96 + 1, solved.height / 96 + 1
    diagram_pptx.add_diagram(slide, d, solved, theme, 0.3, 0.3, w_in, h_in)
    xml = _save_and_read_slide_xml(prs)
    store_shape_id = next(
        sh.shape_id for sh in _walk_shapes(slide.shapes) if sh.name == "docloom:node:b"
    )
    end_m = re.search(rf'<a:endCxn id="{store_shape_id}" idx="(\d+)"/>', xml)
    beg_m = re.search(rf'<a:stCxn id="{store_shape_id}" idx="(\d+)"/>', xml)
    assert end_m is not None and end_m.group(1) == "0"
    assert beg_m is not None and beg_m.group(1) == "3"


# --------------------------------------------------------- legend_h seam (D)


def test_native_diagram_legend_false_draws_no_legend_band():
    """The legend_h seam (2026-07-16 re-audit, finding D): _emit_native used
    to gate on `s.legend` (the kind list, ALWAYS populated regardless of
    legend_h) and position the band from the module constant LEGEND_H, so a
    caller that solve()d with legend=False (legend_h == 0.0, no band
    reserved in the canvas height) would still get a legend drawn at a y
    computed from the nonzero LEGEND_H -- landing on top of the diagram's
    own lowest node row instead of not existing at all. Fixed by gating on
    `s.legend_h > 0` and positioning from `s.legend_h` itself, so "is a band
    reserved" and "where is it" can never disagree again."""
    d = _small_diagram()
    theme = Theme()
    solved = diagram_svg.solve(d, diagram_pptx.theme_dict(theme), legend=False)
    assert solved.legend_h == 0.0
    prs, slide = _new_slide()
    diagram_pptx.add_diagram(slide, d, solved, theme, 0.6, 1.0, 12.133, 5.6)
    texts = {
        sh.text_frame.text for sh in _walk_shapes(slide.shapes) if sh.has_text_frame
    }
    for name in ("flow", "async / return", "primary path", "secure"):
        assert name not in texts
    for kind in ("client", "service", "queue", "store", "external"):
        assert kind not in texts


def test_native_diagram_legend_true_still_draws_legend_band():
    """Companion to the test above: legend=True (solve()'s own default) must
    keep drawing the band exactly as before -- the fix must not have turned
    `s.legend_h > 0` into a gate that is never true."""
    d = _small_diagram()
    theme = Theme()
    solved = diagram_svg.solve(d, diagram_pptx.theme_dict(theme), legend=True)
    assert solved.legend_h > 0.0
    prs, slide = _new_slide()
    diagram_pptx.add_diagram(slide, d, solved, theme, 0.6, 1.0, 12.133, 5.6)
    texts = {
        sh.text_frame.text for sh in _walk_shapes(slide.shapes) if sh.has_text_frame
    }
    assert "flow" in texts and "secure" in texts


def test_ladder_and_raster_reslve_preserve_the_callers_legend_choice(monkeypatch):
    """Finding D, continued: add_diagram's own degradation-ladder re-solve
    and _raster_fallback's re-solve both used to call solve() without a
    `legend` kwarg, silently reverting to solve()'s legend=True default no
    matter what the caller's original `solved` asked for. Forces a real
    ladder climb (same fixture as the degradation-ladder test above) with
    legend=False and asserts every re-solve along the way -- ladder rungs
    and the raster fallback's own re-solve -- was called with legend=False,
    by recording the actual kwargs solve() was invoked with rather than
    inferring it from output geometry."""
    d = Diagram(
        id="ladder", title="t", direction="LR",
        nodes=[
            DiagramNode(id="a", label="A", sublabel="a fairly long sublabel that widens the box"),
            DiagramNode(id="b", label="B", sublabel="another long descriptive sublabel"),
        ],
        edges=[DiagramEdge(source="a", target="b")],
    )
    theme = Theme()
    td = diagram_pptx.theme_dict(theme)
    solved = diagram_svg.solve(d, td, target_aspect=3.0 / 5.6, legend=False)
    assert solved.legend_h == 0.0
    real_solve = diagram_svg.solve
    seen_legend = []

    def tracking_solve(*args, **kwargs):
        seen_legend.append(kwargs.get("legend", True))
        return real_solve(*args, **kwargs)

    monkeypatch.setattr(diagram_pptx, "solve", tracking_solve)
    prs, slide = _new_slide()
    h = diagram_pptx.add_diagram(slide, d, solved, theme, 0.6, 1.0, 3.0, 5.6)
    assert h > 0.0
    assert seen_legend  # the ladder climbed and re-solved at least once
    assert all(v is False for v in seen_legend)
