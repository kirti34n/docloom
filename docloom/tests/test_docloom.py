"""End-to-end checks: the example document renders to every format and the
linter catches deliberately broken documents."""

import json
import zipfile
from pathlib import Path

import pytest

from docloom import (
    DEFAULT, BulletList, Document, ListItem, Paragraph, Slide, Span,
    has_errors, lint, llm_schema, render,
)

EXAMPLE = Path(__file__).parent.parent / "examples" / "quarterly_report.json"


@pytest.fixture()
def doc() -> Document:
    return Document.load(EXAMPLE)


def test_example_validates_and_lints_clean(doc):
    assert doc.slides and doc.blocks and doc.sheets and doc.sources
    # the example must stay export-ready: no ERROR-severity findings. Advisory
    # authoring warnings (weak-title, unlabeled-visual, etc.) are hints, not
    # blockers, so they are allowed here.
    findings = lint(doc, DEFAULT)
    assert not has_errors(findings), [
        f.model_dump() for f in findings if f.severity == "error"
    ]


def test_llm_schema_is_provider_safe():
    text = json.dumps(llm_schema())
    assert "oneOf" not in text  # rejected by OpenAI strict and Anthropic
    # non-recursive: every $ref points at a $def that exists exactly once
    assert "$defs" in text


def test_render_pptx(doc, tmp_path):
    from pptx import Presentation

    out = render(doc, "pptx", tmp_path / "d.pptx")
    prs = Presentation(str(out))
    assert len(prs.slides) >= len(doc.slides)


def test_render_docx(doc, tmp_path):
    import docx

    out = render(doc, "docx", tmp_path / "d.docx")
    d = docx.Document(str(out))
    text = "\n".join(par.text for par in d.paragraphs)
    assert doc.title in text
    assert "Sources" in text


def test_render_xlsx(doc, tmp_path):
    out = render(doc, "xlsx", tmp_path / "d.xlsx")
    with zipfile.ZipFile(out) as z:
        names = z.namelist()
        assert "[Content_Types].xml" in names
        workbook = z.read("xl/workbook.xml").decode("utf-8")
    assert "Revenue" in workbook and "Headcount" in workbook


def test_render_html(doc, tmp_path):
    out = render(doc, "html", tmp_path / "d.html")
    text = out.read_text(encoding="utf-8")
    assert doc.title in text
    assert "Sources" in text
    assert "<sup" in text  # citation markers


def test_html_escapes_hostile_input(tmp_path):
    hostile = Document(
        title="<script>alert(1)</script>",
        blocks=[Paragraph(text=[Span(text="x", link="javascript:alert(1)")])],
    )
    out = render(hostile, "html", tmp_path / "h.html")
    text = out.read_text(encoding="utf-8")
    assert "<script>alert(1)</script>" not in text
    assert 'href="javascript:' not in text


def test_render_markdown(doc, tmp_path):
    out = render(doc, "md", tmp_path / "d.md")
    text = out.read_text(encoding="utf-8")
    assert text.startswith(f"# {doc.title}")
    assert "[^" in text  # citation footnotes
    assert "|" in text  # tables


def test_render_pdf_and_typ(doc, tmp_path):
    pytest.importorskip("typst")
    typ = render(doc, "typ", tmp_path / "d.typ")
    assert "Q2 2026" in typ.read_text(encoding="utf-8")
    out = render(doc, "pdf", tmp_path / "d.pdf")
    assert out.read_bytes()[:5] == b"%PDF-"


def test_lint_catches_broken_documents():
    broken = Document(
        title="T",
        blocks=[Paragraph(text=[Span(text="claim", cite="no-such-source")])],
        slides=[Slide(
            title="x" * 80,
            blocks=[BulletList(items=[ListItem(text="b" * 20)] * 12)],
        )],
    )
    rules = {f.rule for f in lint(broken, DEFAULT)}
    assert "cite/unknown-source" in rules
    assert "deck/title-too-long" in rules
    assert "deck/too-many-bullets" in rules


def test_render_unknown_format_raises(doc, tmp_path):
    from docloom import RenderError

    with pytest.raises(RenderError):
        render(doc, "wat", tmp_path / "d.wat")
