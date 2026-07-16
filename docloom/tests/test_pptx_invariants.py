"""Spanning invariants for the PPTX renderer (the silent-content-loss CLASS
audit, 2026-07-16 wave). Four consecutive prior waves each fixed several
individually-listed defects and shipped exactly one regression, because
every test asserted STRUCTURE (connector counts, fill distinctness) for one
instance at a time; nothing spanned every layout x block-type combination,
and nothing checked that drawn text was actually legible. These two tests
are that span.

Invariant A: everything authored is drawn, or warned about. Built
table-driven over LAYOUTS x BLOCK_CLASSES -- both read straight off the IR's
own Literal/Union definitions (see docloom.ir.Slide's `layout` field and
docloom.ir.Block), so a newly added layout or block type is automatically
exercised here, or fails loudly until pptx.py handles it. For every
combination, a slide carrying that block (with every optional field filled)
is rendered to PPTX, DOCX, HTML, and MD; any authored marker string that
reaches ANY of DOCX/HTML/MD must also reach the raw PPTX package (shape
text, table cells, group-nested diagram shapes, and picture/group `descr`
alt-text attributes are all in scope, via a raw XML-parts text search) or a
warning naming the drop must have been raised. No silent losses.

Invariant B: every text run PPTX actually draws meets the WCAG AA body-text
contrast floor (4.5:1) against the fill actually behind it -- the shape's
own fill if it paints one (a callout's wash, a diagram edge-label pill, a
diagram node), else the nearest enclosing filled shape (a hero/section
band's own rectangle), else the slide background. This is the check whose
absence shipped the hero contrast regression: two individually-correct
fixes (the imageless-hero full treatment, and the callout fill wash) broke
each other on landing, and every existing test still passed because none of
them looked at contrast, only shape counts.
"""

from __future__ import annotations

import typing
import warnings
import zipfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE

from docloom import (
    Artifact, BulletList, Callout, Chart, Code, Diagram, DiagramEdge,
    DiagramGroup, DiagramNode, Divider, Document, Heading, Image as ImageBlock,
    ListItem, NumberedList, Paragraph, Quote, Series, Slide, Stat, StatRow,
    Table as TableBlock, Theme, render,
)
from docloom.ir import Block
from docloom.render.pptx import _band_theme, _callout_fill_color
from docloom.theme import contrast_ratio

LAYOUTS = typing.get_args(Slide.model_fields["layout"].annotation)
BLOCK_CLASSES = typing.get_args(Block)
# Slide.image is only ever drawn by these four layouts (see ir.py's own
# docstring on Slide.image); authoring it elsewhere is not a supported shape,
# so it is exercised in its own dedicated slides, separate from the
# layout x block-type matrix below.
IMAGE_SLOT_LAYOUTS = ("title", "hero", "image_left", "image_right")


# --------------------------------------------------------------- factories


def _png(path, w=40, h=40, color=(30, 30, 30)):
    from PIL import Image as PILImage

    PILImage.new("RGB", (w, h), color).save(path)
    return str(path)


def _make_block_and_markers(cls, tag: str, tmp_path: Path):
    """One instance of `cls` with every optional field filled (a unique
    marker string per field), and the list of marker strings it authored.
    Marker prefixes are chosen so none is a substring of another (e.g.
    QUOTEMARK vs QUOTEATTRMARK): an accidental containment would make a
    genuinely-dropped field's substring search pass by finding it hiding
    inside an unrelated, correctly-drawn sibling string."""
    if cls is Heading:
        return Heading(level=2, text=f"HEADMARK_{tag}"), [f"HEADMARK_{tag}"]
    if cls is Paragraph:
        return Paragraph(text=f"PARAMARK_{tag}"), [f"PARAMARK_{tag}"]
    if cls is BulletList:
        return (
            BulletList(items=[ListItem(text=f"BULLETMARK_{tag}")]),
            [f"BULLETMARK_{tag}"],
        )
    if cls is NumberedList:
        return (
            NumberedList(items=[ListItem(text=f"NUMBEREDMARK_{tag}")]),
            [f"NUMBEREDMARK_{tag}"],
        )
    if cls is Quote:
        return (
            Quote(text=f"QUOTEMARK_{tag}", attribution=f"QUOTEATTRIBUTIONMARK_{tag}"),
            [f"QUOTEMARK_{tag}", f"QUOTEATTRIBUTIONMARK_{tag}"],
        )
    if cls is Code:
        return Code(code=f"CODEMARK_{tag}"), [f"CODEMARK_{tag}"]
    if cls is TableBlock:
        return (
            TableBlock(header=[f"THEADMARK_{tag}"], rows=[[f"TCELLMARK_{tag}"]],
                       caption=f"TCAPTIONMARK_{tag}"),
            [f"THEADMARK_{tag}", f"TCELLMARK_{tag}", f"TCAPTIONMARK_{tag}"],
        )
    if cls is ImageBlock:
        return (
            ImageBlock(path=_png(tmp_path / f"{tag}_img.png"),
                       alt=f"IMGALTMARK_{tag}", caption=f"IMGCAPTIONMARK_{tag}"),
            [f"IMGALTMARK_{tag}", f"IMGCAPTIONMARK_{tag}"],
        )
    if cls is Callout:
        return (
            Callout(style="warning", text=f"CALLOUTMARK_{tag}"),
            [f"CALLOUTMARK_{tag}"],
        )
    if cls is Divider:
        return Divider(), []  # carries no authored text
    if cls is Chart:
        # A single-series chart draws no legend natively (nor in the SVG
        # painter used by html/docx/md), so a series name would never
        # "reach" any report format at all -- title and caption are the
        # only fields every renderer reliably surfaces regardless of series
        # count, so those are what this invariant can actually govern.
        return (
            Chart(chart="bar", title=f"CHARTTITLEMARK_{tag}", labels=["A"],
                  series=[Series(name="s", values=[1.0])],
                  caption=f"CHARTCAPTIONMARK_{tag}"),
            [f"CHARTTITLEMARK_{tag}", f"CHARTCAPTIONMARK_{tag}"],
        )
    if cls is StatRow:
        return (
            StatRow(items=[Stat(label=f"STATLABELMARK_{tag}", value=f"STATVALUEMARK_{tag}",
                                delta=f"STATDELTAMARK_{tag}")]),
            [f"STATLABELMARK_{tag}", f"STATVALUEMARK_{tag}", f"STATDELTAMARK_{tag}"],
        )
    if cls is Artifact:
        return (
            Artifact(kind="diagram", path=_png(tmp_path / f"{tag}_art.png"),
                      alt=f"ARTALTMARK_{tag}", caption=f"ARTCAPTIONMARK_{tag}"),
            [f"ARTALTMARK_{tag}", f"ARTCAPTIONMARK_{tag}"],
        )
    if cls is Diagram:
        d = Diagram(
            id=tag, title=f"DIAGTITLEMARK_{tag}", direction="LR",
            nodes=[
                DiagramNode(id="a", label=f"NODEONEMARK_{tag}",
                            sublabel=f"SUBLABELMARK_{tag}", group="g1"),
                DiagramNode(id="b", label=f"NODETWOMARK_{tag}"),
            ],
            edges=[DiagramEdge(source="a", target="b", label=f"EDGELABELMARK_{tag}")],
            groups=[DiagramGroup(id="g1", label=f"GROUPLABELMARK_{tag}")],
            caption=f"DIAGCAPTIONMARK_{tag}", alt=f"DIAGALTMARK_{tag}",
        )
        return d, [
            f"DIAGTITLEMARK_{tag}", f"NODEONEMARK_{tag}", f"SUBLABELMARK_{tag}",
            f"NODETWOMARK_{tag}", f"EDGELABELMARK_{tag}", f"GROUPLABELMARK_{tag}",
            f"DIAGCAPTIONMARK_{tag}", f"DIAGALTMARK_{tag}",
        ]
    raise AssertionError(f"no factory registered for block class {cls}")


def _build_block(cls, tag: str, tmp_path: Path):
    """_make_block_and_markers, plus a stable `id` on the block so a
    block-level drop warning (which names the block by id, not by every
    individual field's marker text) can still be matched back to this
    exact (layout, block type) row -- see the "warned" check below."""
    block, markers = _make_block_and_markers(cls, tag, tmp_path)
    block.id = tag
    return block, markers


# ------------------------------------------------------------- extraction


def _zip_xml_blob(path) -> str:
    """Concatenated text of every XML part inside an OOXML package (PPTX or
    DOCX are both zips of XML parts). A raw-text search over this catches
    shape text runs, table cells, shapes nested inside a group (a diagram),
    AND non-text metadata like a picture's docPr `descr` alt-text attribute
    -- everything a marker could possibly reach, in one pass, without having
    to separately walk python-pptx/python-docx object trees for each case."""
    with zipfile.ZipFile(path) as z:
        parts = [n for n in z.namelist() if n.endswith(".xml")]
        return "\n".join(z.read(n).decode("utf-8", "replace") for n in parts)


# --------------------------------------------------------- Invariant A test


def test_invariant_a_everything_authored_is_drawn_or_warned_about(tmp_path):
    slides: list[Slide] = []
    # (marker, tag, field-description) for every string this document authors.
    authored: list[tuple[str, str, str]] = []

    for layout in LAYOUTS:
        for cls in BLOCK_CLASSES:
            tag = f"{layout}__{cls.__name__}"
            block, markers = _build_block(cls, tag, tmp_path)
            kwargs: dict = dict(
                layout=layout,
                title=f"SLIDETITLEMARK_{tag}",
                subtitle=f"SLIDESUBMARK_{tag}",
                blocks=[block],
            )
            authored.append((f"SLIDETITLEMARK_{tag}", tag, "slide.title"))
            authored.append((f"SLIDESUBMARK_{tag}", tag, "slide.subtitle"))
            if layout == "two_column":
                # exercise the right column too, alongside the block under test
                kwargs["right"] = [Paragraph(text=f"RIGHTCOLMARK_{tag}")]
                authored.append((f"RIGHTCOLMARK_{tag}", tag, "slide.right"))
            if layout in ("image_left", "image_right"):
                # these two layouts only reach their dedicated image-pane
                # renderer when a usable image is present; without one they
                # fall back to plain content and this row would not exercise
                # the code path its own name promises.
                kwargs["image"] = ImageBlock(path=_png(tmp_path / f"{tag}_slotimg.png"))
            slides.append(Slide(**kwargs))
            for m in markers:
                authored.append((m, tag, f"{cls.__name__} block"))

    # slide.image (path/alt/caption) is only meaningful on these four layouts
    # (ir.py's own docstring); exercised once each, separate from the matrix
    # above so the title.image.caption / hero.image.caption findings are
    # each isolated to their own slide.
    for layout in IMAGE_SLOT_LAYOUTS:
        tag = f"{layout}__imageslot"
        img = ImageBlock(path=_png(tmp_path / f"{tag}.png"),
                          alt=f"SLOTALTMARK_{tag}", caption=f"SLOTCAPTIONMARK_{tag}")
        slides.append(Slide(layout=layout, title=f"SLIDETITLEMARK_{tag}",
                            subtitle=f"SLIDESUBMARK_{tag}", image=img))
        authored.append((f"SLIDETITLEMARK_{tag}", tag, "slide.title"))
        authored.append((f"SLIDESUBMARK_{tag}", tag, "slide.subtitle"))
        authored.append((f"SLOTALTMARK_{tag}", tag, "slide.image.alt"))
        authored.append((f"SLOTCAPTIONMARK_{tag}", tag, "slide.image.caption"))

    doc = Document(title="Invariant A exhaustive matrix", slides=slides)

    with warnings.catch_warnings(record=True) as caught:
        warnings.simplefilter("always")
        pptx_out = render(doc, "pptx", tmp_path / "invariant_a.pptx")
    warning_text = "\n".join(str(w.message) for w in caught)

    docx_out = render(doc, "docx", tmp_path / "invariant_a.docx")
    html_out = render(doc, "html", tmp_path / "invariant_a.html")
    md_out = render(doc, "md", tmp_path / "invariant_a.md")

    pptx_blob = _zip_xml_blob(pptx_out)
    docx_blob = _zip_xml_blob(docx_out)
    html_blob = html_out.read_text(encoding="utf-8")
    md_blob = md_out.read_text(encoding="utf-8")

    silent_losses = []
    for marker, tag, field in authored:
        reaches_report = marker in docx_blob or marker in html_blob or marker in md_blob
        if not reaches_report:
            continue  # not authored content this invariant governs
        if marker in pptx_blob:
            continue  # drawn
        # Warned about: either the exact field was named (a diagram's own
        # "does not clear the Npt floor" message, the title-slide-logo-
        # caption warning, ...), or the whole block was warned about by its
        # own id/tag (a block dropped wholesale for being squeezed below
        # the visual-legibility floor names the block, not each of its
        # individual field markers).
        if marker in warning_text or tag in warning_text:
            continue
        silent_losses.append((tag, field, marker))

    assert not silent_losses, (
        f"{len(silent_losses)} authored field(s) reached DOCX/HTML/MD but "
        "were silently dropped from PPTX (neither drawn nor warned about):\n"
        + "\n".join(f"  {tag}: {field} ({marker!r})" for tag, field, marker in silent_losses)
    )


# --------------------------------------------------------- Invariant B test


def _walk_shapes(shapes):
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(sh.shapes)
        else:
            yield sh


def _own_fill_hex(shape) -> str | None:
    try:
        if shape.fill.type == MSO_FILL_TYPE.SOLID:
            return "#" + str(shape.fill.fore_color.rgb).upper()
    except Exception:
        return None
    return None


def _rect_of(shape):
    try:
        return (shape.left, shape.top, shape.width, shape.height)
    except Exception:
        return None


def _contains(outer, inner) -> bool:
    if outer is None or inner is None:
        return False
    if any(v is None for v in (*outer, *inner)):
        return False
    ox, oy, ow, oh = outer
    ix, iy, iw, ih = inner
    return ox <= ix and oy <= iy and (ox + ow) >= (ix + iw) and (oy + oh) >= (iy + ih)


def _fill_behind(shape, filled_shapes, slide_bg_hex: str) -> str:
    """The effective fill color behind `shape`: its own solid fill if it
    paints one (a callout's wash rect, a diagram edge-label pill textbox),
    else the SMALLEST earlier-drawn solid-filled shape whose rect fully
    contains `shape`'s rect (the nearest enclosing band/card), else the
    slide background."""
    own = _own_fill_hex(shape)
    if own is not None:
        return own
    rect = _rect_of(shape)
    best, best_area = None, None
    for other, fill in filled_shapes:
        if other is shape:
            continue
        orect = _rect_of(other)
        if _contains(orect, rect):
            area = orect[2] * orect[3]
            if best_area is None or area < best_area:
                best_area, best = area, fill
    return best if best is not None else slide_bg_hex


def _text_runs_with_bg(slide):
    """Yield (text, run_color_hex, fill_behind_hex) for every non-blank text
    run drawn on `slide` -- table cells excluded (this renderer's contrast
    bug lives entirely in textbox/autoshape text; tables are not part of the
    documented repro and use a different fill mechanism)."""
    bg = "#" + str(slide.background.fill.fore_color.rgb).upper()
    all_shapes = list(_walk_shapes(slide.shapes))
    filled = [(sh, f) for sh in all_shapes if (f := _own_fill_hex(sh)) is not None]
    for sh in all_shapes:
        if not getattr(sh, "has_text_frame", False):
            continue
        behind = _fill_behind(sh, filled, bg)
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                if not run.text.strip():
                    continue
                try:
                    color = run.font.color
                    hexcolor = "#" + str(color.rgb).upper() if color.type is not None else None
                except Exception:
                    hexcolor = None
                if hexcolor is None:
                    continue
                yield run.text, hexcolor, behind


def test_invariant_b_hero_band_contrast_floor(tmp_path):
    # The exact combination the regression shipped in: callouts (all four
    # styles) and a diagram (whose sublabels are the "PostgreSQL 16"/
    # "Go 1.23" repro) sharing an imageless hero's inverted band. Measured
    # before the fix: info/success/warning/danger callouts at 1.17-1.28:1,
    # diagram edge-label pill at 1.10:1.
    theme = Theme()
    doc = Document(title="T", slides=[
        Slide(layout="hero", title="Contrast floor check", subtitle="on an inverted band",
              blocks=[
                  Callout(style="info", text="info callout body copy"),
                  Callout(style="success", text="success callout body copy"),
                  Callout(style="warning", text="warning callout body copy"),
                  Callout(style="danger", text="danger callout body copy"),
                  Diagram(
                      id="svc", direction="LR",
                      nodes=[
                          DiagramNode(id="a", label="API", sublabel="Go 1.23"),
                          DiagramNode(id="b", label="DB", sublabel="PostgreSQL 16"),
                      ],
                      edges=[DiagramEdge(source="a", target="b", label="reads")],
                      caption="edge label pill contrast check",
                  ),
              ]),
    ])
    out = render(doc, "pptx", tmp_path / "hero_contrast.pptx", theme=theme)
    slide = Presentation(str(out)).slides[0]

    checked_prefixes = (
        "info callout", "success callout", "warning callout", "danger callout",
        "on an inverted band", "Contrast floor check", "API", "DB",
        "Go 1.23", "PostgreSQL 16", "reads",
    )
    failures = []
    seen = 0
    for text, fg, bg in _text_runs_with_bg(slide):
        if not any(text.strip().startswith(p) or p in text for p in checked_prefixes):
            continue
        seen += 1
        ratio = contrast_ratio(fg, bg)
        if ratio < 4.5:
            failures.append((text, fg, bg, round(ratio, 2)))

    assert seen >= len(checked_prefixes), (
        f"only matched {seen}/{len(checked_prefixes)} expected text runs; "
        "the checker's shape-walking may not be finding them"
    )
    assert not failures, f"text runs below the 4.5:1 contrast floor: {failures}"


def test_invariant_b_section_band_contrast_floor(tmp_path):
    theme = Theme()
    doc = Document(title="T", slides=[
        Slide(layout="section", title="Section band", subtitle="callout on section",
              blocks=[Callout(style="danger", text="danger callout on a section band")]),
    ])
    out = render(doc, "pptx", tmp_path / "section_contrast.pptx", theme=theme)
    slide = Presentation(str(out)).slides[0]
    failures = [
        (text, fg, bg, round(contrast_ratio(fg, bg), 2))
        for text, fg, bg in _text_runs_with_bg(slide)
        if "callout on a section band" in text
        and contrast_ratio(fg, bg) < 4.5
    ]
    assert not failures, failures


def test_invariant_b_callout_wash_contrasts_for_every_style_and_band():
    # White-box companion to the render-and-inspect tests above: pins the
    # exact color math for every (band fill, callout style) pair, including
    # the plain (non-band) document background, so a future change to
    # _callout_fill_color/_band_theme regresses here even before it would
    # show up as invisible pixels.
    theme = Theme()
    bands = {"document background": theme.background, "hero (theme.text)": theme.text,
             "section (theme.primary)": theme.primary}
    for band_name, fill in bands.items():
        bt = _band_theme(theme, fill) if fill != theme.background else theme
        for style in ("info", "success", "warning", "danger"):
            callout_fill = _callout_fill_color(style, bt)
            ratio = contrast_ratio(bt.text, callout_fill)
            assert ratio >= 4.5, (
                f"{style} callout on {band_name}: text {bt.text} on fill "
                f"{callout_fill} is {ratio:.2f}:1"
            )


def test_invariant_b_diagram_node_label_contrast_independent_of_band():
    # White-box companion for diagram_pptx._readable_fg: node/sublabel text
    # must contrast against the node's OWN fill (kind_palette's near-white
    # tints, ~0.955 lightness, independent of theme.background/surface) no
    # matter what "theme.text" currently means to the caller.
    from docloom.render import diagram_pptx, diagram_svg

    theme = Theme()
    band_theme = _band_theme(theme, theme.text)
    palette = diagram_svg.kind_palette(diagram_pptx.theme_dict(band_theme))
    for kind, colors in palette.items():
        fg = diagram_pptx._readable_fg(band_theme, colors["fill"])
        ratio = contrast_ratio(fg, colors["fill"])
        assert ratio >= 4.5, f"node kind {kind!r} on a hero band: {fg} on {colors['fill']} is {ratio:.2f}:1"
