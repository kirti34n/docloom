"""Regression tests for the PPTX renderer layout engine.

Each test maps to one confirmed layout bug: silent block-dropping in the
underfull-slide grow pass, a multi-line title overrunning its accent rule,
non-uniform image distortion on the hero cover-fit layout, tables
overflowing the slide instead of truncating visibly, and two_column slides
landing its two columns on mismatched font sizes."""

import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from docloom import Document, Image, Paragraph, Slide, Table, render
from docloom.render.pptx import LAYOUT

SLIDE_W = LAYOUT["slide_w_in"]
SLIDE_H = LAYOUT["slide_h_in"]
MARGIN = LAYOUT["margin_in"]
CONTENT_BOTTOM = Inches(SLIDE_H - MARGIN)
TOL = Inches(0.05)  # rounding slack for float-inch -> EMU conversions


def _xml(path: Path, member: str = "slides/slide") -> str:
    with zipfile.ZipFile(path) as z:
        return "".join(
            z.read(n).decode("utf-8", "replace")
            for n in z.namelist() if member in n
        )


# --------------------------------------------------------- bug 1: grow pass


def test_grow_pass_never_drops_a_block_that_fits_at_scale_one(tmp_path):
    # Reproduces the confirmed repro: 3 paragraphs of ~325 chars each on a
    # content slide (title band leaves ~5.48in available). The old linear
    # scale formula grew text to ~23.8pt, whose *actual* (quadratic) height
    # blew past the slide and silently dropped the third paragraph.
    markers = ["ALPHAMARK", "BRAVOMARK", "CHARLIEMARK"]
    filler = " word" * 63  # ~315 chars; marker + filler ~= 325 chars
    doc = Document(title="T", slides=[
        Slide(layout="content", title="Grow pass", blocks=[
            Paragraph(text=m + filler) for m in markers
        ]),
    ])
    out = render(doc, "pptx", tmp_path / "grow.pptx")
    text = _xml(out)
    for m in markers:
        assert m in text, f"{m} was dropped by the underfull-slide grow pass"


# ------------------------------------------------------ bug 2: title band


def test_multiline_title_does_not_overlap_its_accent_rule(tmp_path):
    long_title = (
        "This Is A Deliberately Long Slide Title Written To Wrap Across "
        "At Least Two Full Lines At The Title Font Size"
    )
    doc = Document(title="T", slides=[
        Slide(layout="content", title=long_title, blocks=[Paragraph(text="body")]),
    ])
    out = render(doc, "pptx", tmp_path / "title.pptx")
    slide = Presentation(str(out)).slides[0]

    title_box = next(
        s for s in slide.shapes
        if s.has_text_frame and long_title.split()[0] in s.text_frame.text
    )
    # the box itself must grow for a 2-line title: the old code always used a
    # fixed 0.62in box (sized for one line) regardless of wrap, so its
    # bounding box never reflected the true (overflowing) text extent
    assert title_box.height > Inches(0.7), (
        "title box did not grow for a multi-line title (still the 1-line height)"
    )
    # the accent rule is a thin (0.028in tall) rectangle drawn under the title
    rule = next(
        s for s in slide.shapes
        if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and abs(s.height - Inches(0.028)) < Inches(0.01)
    )
    assert rule.top >= title_box.top + title_box.height, (
        "accent rule overlaps a multi-line title instead of sitting below it"
    )


# ------------------------------------------------------- bug 4: hero cover-fit


def test_hero_image_cover_fit_preserves_aspect_ratio(tmp_path):
    from PIL import Image as PILImage

    img_path = tmp_path / "portrait.png"
    src_w, src_h = 200, 400  # tall and narrow: far from the 16:9 slide
    PILImage.new("RGB", (src_w, src_h), "blue").save(img_path)

    doc = Document(title="T", slides=[
        Slide(layout="hero", title="Hero", image=Image(path=str(img_path))),
    ])
    out = render(doc, "pptx", tmp_path / "hero.pptx")
    slide = Presentation(str(out)).slides[0]
    pic = next(s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)

    # forced non-uniform scaling (the bug) never crops; a correct cover-fit
    # crops exactly the axis that overflows (here: vertical) and leaves the
    # other uncropped
    assert pic.crop_left == 0 and pic.crop_right == 0
    assert pic.crop_top > 0 and pic.crop_bottom > 0

    # the visible (uncropped) slice of the source must carry the same aspect
    # ratio as the box it's displayed in -- otherwise the image is stretched
    visible_h_fraction = 1 - pic.crop_top - pic.crop_bottom
    src_aspect = src_w / (src_h * visible_h_fraction)
    disp_aspect = pic.width / pic.height
    assert abs(src_aspect - disp_aspect) < 0.02

    # covers the full slide with no letterboxing
    assert abs(pic.width - Inches(SLIDE_W)) <= TOL
    assert abs(pic.height - Inches(SLIDE_H)) <= TOL


# -------------------------------------------------------------- bug 5: tables


def test_long_table_truncates_visibly_instead_of_overflowing(tmp_path):
    header = ["Row", "Value"]
    rows = [[f"ROWMARK{i}", str(i)] for i in range(40)]
    doc = Document(title="T", slides=[
        Slide(layout="content", title="Big table", blocks=[
            Table(header=header, rows=rows),
        ]),
    ])
    out = render(doc, "pptx", tmp_path / "table.pptx")
    slide = Presentation(str(out)).slides[0]

    table_shape = next(s for s in slide.shapes if s.has_table)
    assert table_shape.top + table_shape.height <= CONTENT_BOTTOM + TOL, (
        "table overflows the slide's content margin"
    )
    text = _xml(out)
    assert "more row" in text, "rows were dropped with no visible notice"
    assert "ROWMARK39" not in text  # confirms rows past the fold were cut


# ------------------------------------------------------ bug 6: two_column


def test_two_column_slides_share_one_font_scale(tmp_path):
    # left column: one short paragraph that would want to grow a lot if
    # scaled on its own; right column: enough text that it wants to grow
    # much less. Both must render at the same (shared, smaller) scale.
    left_text = "SHORTMARK tiny"
    right_text = "LONGMARK " + " filler word" * 60
    doc = Document(title="T", slides=[
        Slide(layout="two_column", title="Split",
              blocks=[Paragraph(text=left_text)],
              right=[Paragraph(text=right_text)]),
    ])
    out = render(doc, "pptx", tmp_path / "two_col.pptx")
    slide = Presentation(str(out)).slides[0]

    def _size_of(marker):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    if marker in run.text:
                        return run.font.size
        raise AssertionError(f"{marker!r} not found on the rendered slide")

    left_size = _size_of("SHORTMARK")
    right_size = _size_of("LONGMARK")
    assert left_size == right_size, (
        f"two_column columns rendered at different font sizes: "
        f"{left_size} vs {right_size}"
    )
