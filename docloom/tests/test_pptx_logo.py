"""Regression: a title-slide image (brand logo) must actually render into the
PPTX. Prior to this, generation put the logo in slides[0].image but the pptx
renderer's _title_slide ignored it, so the logo was silently dropped."""

import base64
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from docloom import render
from docloom.ir import Document, Image, Slide

# a minimal valid 1x1 PNG python-pptx can read dimensions from
_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAAC0lEQVR42mP8"
    "z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
)


def _logo(tmp_path: Path) -> str:
    p = tmp_path / "logo.png"
    p.write_bytes(_PNG)
    return str(p)


def _pics(slide):
    return [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]


def test_title_slide_renders_logo(tmp_path):
    doc = Document(title="Branded", slides=[
        Slide(layout="title", title="Q3", subtitle="Co", image=Image(path=_logo(tmp_path))),
        Slide(layout="content", title="Body", blocks=[{"type": "paragraph", "text": "hi"}]),
    ])
    out = render(doc, "pptx", tmp_path / "out.pptx")
    prs = Presentation(str(out))
    pics = _pics(prs.slides[0])
    assert len(pics) == 1, "brand logo was dropped from the title slide"
    # top-right corner, above the title band
    assert pics[0].top / 914400 < 2.0
    assert pics[0].left / 914400 > 6.0


def test_title_slide_without_image_has_no_picture(tmp_path):
    doc = Document(title="Plain", slides=[Slide(layout="title", title="X")])
    out = render(doc, "pptx", tmp_path / "plain.pptx")
    prs = Presentation(str(out))
    assert _pics(prs.slides[0]) == []


def test_image_layout_slide_renders_bound_picture(tmp_path):
    """A slide with an image bound to it (the deck editor's asset picker) must
    place that picture on the slide."""
    doc = Document(title="D", slides=[
        Slide(layout="image_left", title="T", image=Image(path=_logo(tmp_path)),
              blocks=[{"type": "paragraph", "text": "body"}]),
    ])
    out = render(doc, "pptx", tmp_path / "img.pptx")
    prs = Presentation(str(out))
    assert len(_pics(prs.slides[0])) == 1


def test_doc_logo_stamps_every_slide(tmp_path):
    """A document-level logo appears top-right on every slide layout: the
    title, content slides, and full-bleed section slides (over a scrim)."""
    logo = Image(path=_logo(tmp_path))
    doc = Document(title="Branded", logo=logo, slides=[
        Slide(layout="title", title="Q3"),
        Slide(layout="content", title="A", blocks=[{"type": "paragraph", "text": "x"}]),
        Slide(layout="content", title="B", blocks=[{"type": "paragraph", "text": "y"}]),
        Slide(layout="section", title="Break"),
    ])
    out = render(doc, "pptx", tmp_path / "branded.pptx")
    prs = Presentation(str(out))
    for i in range(4):
        pics = _pics(prs.slides[i])
        assert len(pics) == 1, f"slide {i} should carry the brand logo"
        # top-right corner, modest size
        assert pics[0].left / 914400 > 6.0
        assert pics[0].height / 914400 <= 0.9
