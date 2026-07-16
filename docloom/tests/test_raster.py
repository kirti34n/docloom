"""The optional SVG rasterizer seam (docloom.render.raster, extra: diagrams).

Two halves, and both matter:
  * with resvg installed, a Chart in DOCX and an SVG Image in PPTX become real
    pictures instead of a text placeholder / a silent drop;
  * with it missing (simulated by poisoning the import), every format still
    renders, with the old fallbacks and no exception. That second half is what
    keeps the core install (no optional extras) safe.
"""

import sys
import zipfile

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from docloom import Chart, Document, Image, Series, render
from docloom.render import raster
from docloom.theme import Theme

SVG = (
    '<svg xmlns="http://www.w3.org/2000/svg" width="80" height="40">'
    '<rect width="80" height="40" fill="#1D4ED8"/>'
    '<text x="8" y="26" font-family="Arial" font-size="14" fill="#FFFFFF">hi</text>'
    "</svg>"
)


def _chart() -> Chart:
    return Chart(
        title="Revenue",
        labels=["FY23", "FY24"],
        series=[Series(name="usd", values=[1234567.0, 1000000.0])],
        caption="source: internal",
    )


def _chart_doc() -> Document:
    return Document(title="T", blocks=[_chart()])


def _slide_doc(block) -> Document:
    return Document(
        title="T",
        slides=[{"layout": "content", "title": "S", "blocks": [block]}],
    )


def _svg_file(tmp_path):
    p = tmp_path / "diagram.svg"
    p.write_text(SVG, encoding="utf-8")
    return p


def _pictures(pptx_path):
    prs = Presentation(str(pptx_path))
    return [
        sh
        for slide in prs.slides
        for sh in slide.shapes
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]


def _media(docx_path) -> list[str]:
    with zipfile.ZipFile(docx_path) as z:
        return [n for n in z.namelist() if n.startswith("word/media/")]


def _no_resvg(monkeypatch):
    """Make `import resvg_py` fail the way it does on a core install: None in
    sys.modules raises ImportError at the import statement itself."""
    monkeypatch.setitem(sys.modules, "resvg_py", None)


# --------------------------------------------------- the seam itself


def test_svg_to_png_returns_png_bytes():
    pytest.importorskip("resvg_py")
    png = raster.svg_to_png(SVG, width=160)
    assert png is not None
    assert png.startswith(raster.PNG_MAGIC)


def test_svg_to_png_is_none_without_the_extra(monkeypatch):
    _no_resvg(monkeypatch)
    assert raster.svg_to_png(SVG) is None
    assert raster.available() is False


def test_svg_to_png_never_raises_on_garbage():
    assert raster.svg_to_png("") is None
    assert raster.svg_to_png("not an svg at all <<<") is None


def test_svg_file_to_png_missing_file_is_none(tmp_path):
    assert raster.svg_file_to_png(tmp_path / "nope.svg") is None


def test_theme_font_files_skips_woff2_and_missing(tmp_path):
    ttf = tmp_path / "brand.ttf"
    ttf.write_bytes(b"\x00")
    theme = Theme(font_heading_src=str(ttf), font_body_src=str(tmp_path / "brand.woff2"))
    assert raster.theme_font_files(theme) == [str(ttf)]


# ------------------------------------------- with the rasterizer present


def test_docx_chart_embeds_a_real_picture(tmp_path):
    pytest.importorskip("resvg_py")
    import docx as docx_lib

    out = render(_chart_doc(), "docx", tmp_path / "c.docx")
    media = _media(out)
    assert media, "chart did not reach the docx as an image part"
    assert any(n.endswith(".png") for n in media)

    d = docx_lib.Document(str(out))
    body = d.element.body.xml
    assert "<w:drawing" in body and "a:blip" in body  # an inline picture
    assert not d.tables  # not the data-table fallback any more
    paragraphs = [p.text for p in d.paragraphs]
    assert "Revenue" in paragraphs  # title kept
    assert any("source: internal" in p for p in paragraphs)  # caption kept


def test_pptx_svg_image_block_is_a_picture_not_a_silent_drop(tmp_path):
    pytest.importorskip("resvg_py")
    doc = _slide_doc(Image(path=str(_svg_file(tmp_path)), alt="a blue box", caption="Fig 1"))
    pics = _pictures(render(doc, "pptx", tmp_path / "i.pptx"))
    assert len(pics) == 1
    assert pics[0].image.ext == "png"  # rasterized on the way in


def test_pptx_chart_falls_back_to_a_rasterized_chart_not_a_table(tmp_path):
    pytest.importorskip("resvg_py")
    # a multi-series pie cannot be a native pptx chart, so this exercises the
    # fallback chain: it used to end in a data table, it now paints the chart
    doc = _slide_doc(
        Chart(
            chart="pie",
            title="Split",
            labels=["a", "b"],
            series=[
                Series(name="s1", values=[1.0, 2.0]),
                Series(name="s2", values=[3.0, 4.0]),
            ],
        )
    )
    out = render(doc, "pptx", tmp_path / "p.pptx")
    assert len(_pictures(out)) == 1


# ------------------------------------------- with the rasterizer missing


def test_docx_chart_falls_back_to_the_table_without_the_extra(tmp_path, monkeypatch):
    import docx as docx_lib

    _no_resvg(monkeypatch)
    out = render(_chart_doc(), "docx", tmp_path / "c.docx")
    d = docx_lib.Document(str(out))
    text = "\n".join(c.text for t in d.tables for row in t.rows for c in row.cells)
    assert "1,234,567" in text  # the data still reaches the reader
    assert not _media(out)  # and nothing was embedded


def test_pptx_svg_image_falls_back_without_the_extra(tmp_path, monkeypatch):
    _no_resvg(monkeypatch)
    doc = _slide_doc(Image(path=str(_svg_file(tmp_path)), alt="a blue box"))
    out = render(doc, "pptx", tmp_path / "i.pptx")  # must not raise
    assert _pictures(out) == []  # same as today: nothing embeddable


def test_every_format_still_renders_without_the_extra(tmp_path, monkeypatch):
    _no_resvg(monkeypatch)
    doc = Document(
        title="T",
        blocks=[_chart(), Image(path=str(_svg_file(tmp_path)), alt="a blue box")],
        slides=[{
            "layout": "content", "title": "S",
            "blocks": [_chart(), Image(path=str(_svg_file(tmp_path)), alt="x")],
        }],
    )
    for fmt in ("docx", "pptx", "html", "md", "xlsx"):
        assert render(doc, fmt, tmp_path / f"o.{fmt}").exists()


def test_pptx_png_image_path_is_untouched_by_the_rasterizer(tmp_path, monkeypatch):
    # a non-SVG image must embed exactly as before, extra installed or not
    png = tmp_path / "dot.png"
    png.write_bytes(raster.svg_to_png(SVG) or b"")
    if not png.read_bytes():
        pytest.skip("rasterizer unavailable; no PNG fixture to embed")
    _no_resvg(monkeypatch)
    doc = _slide_doc(Image(path=str(png), alt="a dot"))
    assert len(_pictures(render(doc, "pptx", tmp_path / "p.pptx"))) == 1
