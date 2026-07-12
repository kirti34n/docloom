"""Re-audit regression for the PPTX image_left/image_right layout.

The pane matte rectangle was drawn before add_picture. When python-pptx could
not embed the image (an SVG, a corrupt/unsupported raster) add_picture raised
and the except branch fell back to a clean content layout, but the already
committed full-height, ~45%-width surface matte was never removed, painting a
stray colored block over the fallback slide. The fix deletes the matte on the
failure path, mirroring how _hero_slide adds its picture first."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from docloom import render
from docloom.ir import Document, Image, Slide
from docloom.render.pptx import LAYOUT

SLIDE_H = LAYOUT["slide_h_in"]

# a real on-disk file that python-pptx cannot embed: add_picture raises on it,
# which is exactly the failure path the fix has to clean up after
_SVG = (
    '<?xml version="1.0" encoding="UTF-8"?>'
    '<svg xmlns="http://www.w3.org/2000/svg" width="120" height="80">'
    '<rect width="120" height="80" fill="#cccccc"/></svg>'
)


def _full_height_autoshapes(slide):
    # the pane matte is a full-slide-height rectangle autoshape; nothing else on
    # a content-layout fallback comes close (the title rule is ~0.028in tall)
    return [
        sh for sh in slide.shapes
        if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and sh.height >= Inches(SLIDE_H * 0.9)
    ]


def test_image_side_unembeddable_image_leaves_no_stray_matte(tmp_path: Path):
    svg = tmp_path / "diagram.svg"
    svg.write_text(_SVG, encoding="utf-8")
    doc = Document(title="D", slides=[
        Slide(layout="image_left", title="FALLBACKMARK",
              image=Image(path=str(svg)),
              blocks=[{"type": "paragraph", "text": "body"}]),
    ])
    out = render(doc, "pptx", tmp_path / "side.pptx")
    slide = Presentation(str(out)).slides[0]

    assert not _full_height_autoshapes(slide), (
        "stray full-height pane matte left over the content-layout fallback"
    )
    # the fallback content layout actually ran: its title band carries the title
    text = " ".join(
        sh.text_frame.text for sh in slide.shapes if sh.has_text_frame
    )
    assert "FALLBACKMARK" in text, "image_side fallback did not draw a content slide"
