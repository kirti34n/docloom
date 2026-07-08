"""Regression tests for review findings (core layer)."""

import json

from docloom import (
    DEFAULT, BulletList, Callout, Code, Document, ListItem, NumberedList,
    Paragraph, Quote, Slide, Source, Span, Table, lint, llm_schema,
)
from docloom.render import slug


def test_save_load_roundtrip_preserves_block_types(tmp_path):
    # save() once stripped the Literal type tags (exclude_defaults), turning
    # paragraphs/quotes/callouts into headings and numbered lists into bullets
    doc = Document(
        title="T",
        blocks=[
            Paragraph(text="para"),
            Quote(text="q"),
            Callout(text="c"),
            NumberedList(items=[ListItem(text="one")]),
        ],
        slides=[Slide(layout="two_column", blocks=[Paragraph(text="left")])],
    )
    path = tmp_path / "d.json"
    doc.save(path)
    loaded = Document.load(path)
    assert [type(b).__name__ for b in loaded.blocks] == [
        "Paragraph", "Quote", "Callout", "NumberedList",
    ]
    assert type(loaded.slides[0].blocks[0]).__name__ == "Paragraph"


def test_load_tolerates_utf8_bom(tmp_path):
    path = tmp_path / "bom.json"
    payload = json.dumps({"title": "T", "blocks": [{"type": "paragraph", "text": "x"}]})
    path.write_bytes(b"\xef\xbb\xbf" + payload.encode("utf-8"))
    assert Document.load(path).title == "T"


def test_lint_duplicate_source_ids():
    doc = Document(
        title="T",
        blocks=[Paragraph(text=[Span(text="x", cite="a")])],
        sources=[Source(id="a", title="A"), Source(id="a", title="B")],
    )
    assert "cite/duplicate-source" in {f.rule for f in lint(doc)}


def test_lint_ragged_table():
    doc = Document(title="T", blocks=[Table(header=["a", "b"], rows=[["1", "2", "3"]])])
    assert "table/ragged" in {f.rule for f in lint(doc)}


def test_lint_ignored_blocks_on_title_slide():
    doc = Document(
        title="T",
        slides=[Slide(layout="title", title="T", blocks=[Paragraph(text="lost")])],
    )
    assert "deck/ignored-blocks" in {f.rule for f in lint(doc)}


def test_lint_two_column_gets_half_budget():
    # 660 chars in one column: under the 800 full-width budget, over the
    # 400-per-column budget for two_column slides
    items = [ListItem(text="b" * 110) for _ in range(6)]
    doc = Document(
        title="T",
        slides=[Slide(layout="two_column", title="t",
                      blocks=[BulletList(items=items)],
                      right=[Paragraph(text="short")])],
    )
    assert "deck/overflow" in {f.rule for f in lint(doc)}


def test_lint_counts_code_blocks():
    doc = Document(
        title="T",
        slides=[Slide(title="t", blocks=[Code(code="x = 1\n" * 200)])],
    )
    assert "deck/overflow" in {f.rule for f in lint(doc)}


def test_llm_schema_objects_are_closed():
    schema = llm_schema()
    objects = [schema] + [
        d for d in schema["$defs"].values() if d.get("type") == "object"
    ]
    assert all(node.get("additionalProperties") is False for node in objects)


def test_slug_keeps_unicode():
    assert slug("Отчёт за квартал") != "document"
    assert slug("四半期レポート") != "document"


def test_default_theme_lints_clean():
    assert not [f for f in lint(Document(title="T"), DEFAULT)
                if f.rule == "theme/low-contrast"]


# ---------------------------------------------------------------- renderers


RAGGED = Document(
    title="Ragged",
    blocks=[Table(header=["a"], rows=[["1", "EXTRA-CELL"], ["2"]]),
            Table(header=[], rows=[["ORPHAN-ROW", "x"]])],
    sheets=[{
        "name": "S",
        "columns": [{"header": "only"}],
        "rows": [["v1", "WIDE-CELL"]],
    }],
)


def test_ragged_tables_lose_no_cells_in_any_format(tmp_path):
    from docloom import render

    md = render(RAGGED, "md", tmp_path / "r.md").read_text(encoding="utf-8")
    assert "EXTRA-CELL" in md and "ORPHAN-ROW" in md

    html = render(RAGGED, "html", tmp_path / "r.html").read_text(encoding="utf-8")
    assert "EXTRA-CELL" in html and "ORPHAN-ROW" in html

    import docx as docx_lib

    d = docx_lib.Document(str(render(RAGGED, "docx", tmp_path / "r.docx")))
    docx_text = "\n".join(c.text for t in d.tables for row in t.rows for c in row.cells)
    assert "EXTRA-CELL" in docx_text and "ORPHAN-ROW" in docx_text

    with __import__("zipfile").ZipFile(render(RAGGED, "xlsx", tmp_path / "r.xlsx")) as z:
        shared = z.read("xl/sharedStrings.xml").decode("utf-8")
    assert "WIDE-CELL" in shared


def test_typst_compiles_ragged_and_enum_lookalike(tmp_path):
    import pytest as _pytest

    _pytest.importorskip("typst")
    from docloom import render

    doc = Document(
        title="Enum",
        blocks=[Paragraph(text="1. Introduction covers scope."),
                Table(header=[], rows=[["a", "b"]])],
    )
    out = render(doc, "pdf", tmp_path / "e.pdf")
    assert out.read_bytes()[:5] == b"%PDF-"


def test_docx_numbered_lists_restart(tmp_path):
    import docx as docx_lib

    from docloom import render

    doc = Document(
        title="Nums",
        blocks=[
            NumberedList(items=[ListItem(text="first-a"), ListItem(text="first-b")]),
            Paragraph(text="between"),
            NumberedList(items=[ListItem(text="second-a")]),
        ],
    )
    d = docx_lib.Document(str(render(doc, "docx", tmp_path / "n.docx")))
    starts = [p.text for p in d.paragraphs if p.text.startswith("1.")]
    assert len(starts) == 2  # both lists restart at 1


def test_pptx_long_quote_stays_on_canvas(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    from docloom import Quote as QuoteBlock
    from docloom import render

    doc = Document(
        title="Q",
        slides=[Slide(layout="quote",
                      blocks=[QuoteBlock(text="wise words " * 200,
                                         attribution="Someone")])],
    )
    prs = Presentation(str(render(doc, "pptx", tmp_path / "q.pptx")))
    for slide in prs.slides:
        for shape in slide.shapes:
            assert shape.left >= 0 and shape.top >= 0
            assert shape.left + shape.width <= Inches(13.333) + 1
            assert shape.top + shape.height <= Inches(7.5) + 1


def test_xlsx_apostrophe_sheet_name(tmp_path):
    from docloom import render

    doc = Document(
        title="A",
        sheets=[{"name": "'Costs'", "columns": [{"header": "h"}], "rows": [["v"]]}],
    )
    assert render(doc, "xlsx", tmp_path / "a.xlsx").exists()


def test_markdown_escapes_hostile_text(tmp_path):
    from docloom import render
    from docloom import Table as TableBlock

    doc = Document(
        title="Hostile",
        blocks=[
            Paragraph(text="# not a heading"),
            Paragraph(text="- not a list item"),
            TableBlock(header=["h"], rows=[["cell | with pipe"]]),
        ],
    )
    md = render(doc, "md", tmp_path / "h.md").read_text(encoding="utf-8")
    lines = md.splitlines()
    assert not any(line.startswith("# not") for line in lines)
    assert not any(line.startswith("- not") for line in lines)
    # the pipe inside the cell must be escaped so the table keeps one column
    table_lines = [ln for ln in lines if "with pipe" in ln]
    assert table_lines and "\\|" in table_lines[0]


# ------------------------------------------------- test-campaign regressions


def test_control_chars_stripped_at_ir_boundary():
    from docloom import Column, Formula, Sheet

    doc = Document(
        title="a\x07b",
        authors=["\x1b[31m"],
        blocks=[Paragraph(text=[Span(text="x\x00y", link="http://a\x01b.com")])],
        sheets=[Sheet(name="s\x00heet", columns=[Column(header="h", format="0\x1b0")],
                      rows=[[Formula(formula="=X\x00Y")]])],
    )
    assert doc.title == "ab"
    assert doc.blocks[0].text[0].text == "xy"
    assert doc.blocks[0].text[0].link == "http://ab.com"
    assert doc.sheets[0].name == "sheet"
    assert doc.sheets[0].columns[0].format == "00"
    assert doc.sheets[0].rows[0][0].formula == "=XY"
    assert Document(title="tab\tand\nnewline ok").title == "tab\tand\nnewline ok"


def test_llm_schema_provider_contract():
    schema = llm_schema()
    text = json.dumps(schema)
    for keyword in ('"minimum"', '"maximum"', '"exclusiveMinimum"',
                    '"exclusiveMaximum"', '"multipleOf"', '"minLength"',
                    '"maxLength"', '"pattern"', '"oneOf"'):
        assert keyword not in text, keyword
    # type tags must be required so untagged blocks can't reach the union
    for name, node in schema["$defs"].items():
        if "const" in node.get("properties", {}).get("type", {}):
            assert "type" in node.get("required", []), name


def test_lint_empty_formula():
    from docloom import Column, Formula, Sheet

    doc = Document(title="t", sheets=[
        Sheet(name="s", columns=[Column(header="h")], rows=[[Formula(formula=" ")]]),
    ])
    assert "sheet/empty-formula" in {f.rule for f in lint(doc)}


def test_md_caption_cannot_inject_structure(tmp_path):
    from docloom import Table as TableBlock
    from docloom import render

    doc = Document(
        title="Cap",
        blocks=[
            TableBlock(header=["A"], rows=[["1"]],
                       caption="line1\n# fake heading\n```\nstolen"),
            Paragraph(text="AFTER"),
        ],
    )
    md = render(doc, "md", tmp_path / "c.md").read_text(encoding="utf-8")
    lines = md.splitlines()
    assert not any(line.startswith("# fake") for line in lines)
    fences = sum(1 for line in lines if line.startswith("```"))
    assert fences % 2 == 0  # no unclosed fence swallowing content
    assert "AFTER" in md


def test_md_code_language_sanitized(tmp_path):
    from docloom import render

    doc = Document(
        title="Lang",
        blocks=[Code(code="x = 1", language="py`thon\nevil = 1"),
                Paragraph(text="AFTER TEXT")],
    )
    md = render(doc, "md", tmp_path / "l.md").read_text(encoding="utf-8")
    assert "x = 1" in md
    assert "evil" not in md  # injected via language must be dropped entirely
    # the paragraph after the code block must survive as text, not code
    assert "AFTER TEXT" in md
    fences = [ln for ln in md.splitlines() if ln.startswith("```")]
    assert len(fences) % 2 == 0
    assert all("`" not in ln[3:] for ln in fences)  # info string has no backtick


def test_typst_skips_undecodable_image(tmp_path):
    import pytest as _pytest

    _pytest.importorskip("typst")
    from docloom import Image as ImageBlock
    from docloom import render

    bad = tmp_path / "bad.png"
    bad.write_bytes(b"not a png")
    doc = Document(title="Img", blocks=[
        Paragraph(text="before"),
        ImageBlock(path=str(bad), alt="broken"),
        Paragraph(text="over"),
    ])
    out = render(doc, "pdf", tmp_path / "i.pdf")
    assert out.read_bytes()[:5] == b"%PDF-"


def test_pdf_links_in_wide_tables_render(tmp_path):
    import pytest as _pytest

    _pytest.importorskip("typst")
    from docloom import Table as TableBlock
    from docloom import render

    # typst 0.15 panics on wrapped emphasized links inside wide tables;
    # docloom de-links table cells as a workaround
    cell = [Span(text="a very long linked label that will wrap " * 3,
                 italic=True, link="https://example.com/x")]
    doc = Document(title="T", blocks=[
        TableBlock(header=["a", "b", "c", "d", "e", "f"],
                   rows=[[cell, "2", "3", "4", "5", "6"]]),
    ])
    out = render(doc, "pdf", tmp_path / "t.pdf")
    assert out.read_bytes()[:5] == b"%PDF-"


def test_office_renderers_drop_unsafe_link_schemes(tmp_path):
    import zipfile

    from docloom import render

    doc = Document(
        title="Links",
        blocks=[Paragraph(text=[Span(text="bad", link="javascript:alert(1)"),
                                Span(text=" and "),
                                Span(text="good", link="https://example.com")])],
        slides=[Slide(title="s", blocks=[
            Paragraph(text=[Span(text="bad", link="javascript:alert(1)")]),
        ])],
    )
    for fmt in ("docx", "pptx"):
        path = render(doc, fmt, tmp_path / f"x.{fmt}")
        with zipfile.ZipFile(path) as z:
            blob = b"".join(z.read(n) for n in z.namelist() if n.endswith(".rels"))
        assert b"javascript:" not in blob


def test_parse_llm_output_tolerates_model_wrapping():
    # failure modes observed live from local Ollama models whose structured
    # output was silently unenforced
    from docloom import parse_llm_output

    bare = '{"title": "T", "blocks": [{"type": "paragraph", "text": "x"}]}'
    fenced = f"```json\n{bare}\n```"
    prose = f"Here is your document:\n\n{bare}\n\nLet me know if you need changes!"
    wrapped = f'{{"document": {bare}}}'
    for text in (bare, fenced, prose, wrapped):
        doc = parse_llm_output(text)
        assert doc.title == "T"
        assert type(doc.blocks[0]).__name__ == "Paragraph"

    import pytest as _pytest

    with _pytest.raises(Exception):
        parse_llm_output("```json\n{\"title\": ")  # truncated output stays an error


def test_parse_llm_output_normalizes_type_aliases():
    # observed live: qwen3.5:9b tags bullet lists "bulletlist"
    from docloom import parse_llm_output

    text = ('{"title": "T", "blocks": ['
            '{"type": "bulletlist", "items": [{"text": "a"}]},'
            '{"type": "code_block", "code": "x"},'
            '{"type": "blockquote", "text": "q"}]}')
    doc = parse_llm_output(text)
    assert [type(b).__name__ for b in doc.blocks] == [
        "BulletList", "Code", "Quote",
    ]


def test_parse_llm_output_unknown_type_is_one_clear_error():
    import pytest as _pytest

    from docloom import parse_llm_output

    with _pytest.raises(ValueError) as exc:
        parse_llm_output('{"title": "T", "blocks": [{"type": "wibble"}]}')
    message = str(exc.value)
    assert 'unknown block type "wibble"' in message
    assert "paragraph" in message  # lists the valid types for self-correction


def test_parse_llm_output_filters_union_errors_to_tagged_member():
    # live finding: a valid tag with an invalid body made Pydantic report all
    # ten union branches, leading with "Input should be 'heading'" — which
    # sent the self-correcting model in the wrong direction every round
    import pytest as _pytest

    from docloom import parse_llm_output

    with _pytest.raises(ValueError) as exc:
        parse_llm_output(
            '{"title": "T", "blocks": ['
            '{"type": "table", "header": ["a"], "rows": "oops"}]}'
        )
    message = str(exc.value)
    assert "Table" in message and "rows" in message
    assert "heading" not in message.lower()  # no wrong-member noise


def test_parse_llm_output_filters_cell_union_errors():
    # live finding: a dict cell missing "formula" reported bool/int/float
    # branches too — only the Formula branch is relevant for a dict cell
    import pytest as _pytest

    from docloom import parse_llm_output

    with _pytest.raises(ValueError) as exc:
        parse_llm_output(
            '{"title": "T", "sheets": [{"name": "s", '
            '"columns": [{"header": "h"}], "rows": [[{"wrong": 1}]]}]}'
        )
    message = str(exc.value)
    assert "Formula" in message and "formula" in message
    assert "boolean" not in message and "integer" not in message


def test_xlsx_survives_nan_inf(tmp_path):
    import zipfile

    from docloom import Column, Sheet
    from docloom import render

    doc = Document(title="n", sheets=[
        Sheet(name="s", columns=[Column(header="h")],
              rows=[[float("nan")], [float("inf")]]),
    ])
    path = render(doc, "xlsx", tmp_path / "n.xlsx")
    with zipfile.ZipFile(path) as z:
        assert "[Content_Types].xml" in z.namelist()
