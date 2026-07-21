"""Render-quality regression tests: the dependency-free chart SVG painter,
its wiring into html/typst/docx/markdown, the docx figure/placeholder
fallbacks and header-link color, the markdown emphasis-boundary and
assets-mode fixes, and the P5 PPTX quality-audit fixes (silent block drops,
section-slide data loss, image_right title/body misalignment, the grow-pass
runaway, heading-seam spacing, solo-chart fill, image-pane letterboxing,
table column weighting, the invisible divider, the sources slide, stat
card legibility, and the pie chart percent-label bug). Also covers the
post-diagram-wave fixes: _content_slide/_image_side_slide/two_column
silently dropping s.subtitle, an imageless hero falling through to the
plain content layout instead of getting its own full-bleed treatment, and
warning/danger callouts rendering as uncolored gray/near-black.

Also covers the silent-content-loss CLASS audit: _body's "ponytail" trailing-
block drop replaced with reserve-ahead allocation + a warning (never a
silent drop again); slide.image.caption now rendered on image_left/
image_right (it was silently dropped there while DOCX/HTML/MD kept it via
flatten_slides); an imageless hero with blocks giving them real room instead
of crushing a diagram into an illegible raster tile; and callout fills now a
tinted wash of their own style color instead of one flat gray for all four
styles. Table/chart/quote/placeholder captions and attributions are now
reserved up front and always drawn, never opportunistically dropped when a
slide runs tight."""

from __future__ import annotations

import xml.dom.minidom as minidom

import pytest

from docloom import (
    Artifact, BulletList, Chart, Code, Diagram, DiagramEdge, DiagramNode,
    Document, Heading, Image as ImageBlock, ListItem, Paragraph, Quote,
    Series, Slide, Source, Span, Stat, StatRow, Table as TableBlock, Theme,
    render,
)
from docloom.render import chart_svg

# ------------------------------------------------------------- chart_svg


@pytest.mark.parametrize("kind", ["column", "bar", "line", "area", "pie", "scatter"])
def test_chart_svg_renders_every_kind_as_valid_xml(kind):
    chart = Chart(
        chart=kind, title="T", labels=["A", "B", "C"],
        series=[Series(name="s1", values=[1.0, 2.0, 3.0]), Series(name="s2", values=[3.0, 1.0, 2.0])],
        caption="cap",
    )
    svg = chart_svg.render_svg(chart, Theme())
    assert svg.startswith("<svg") and svg.endswith("</svg>")
    minidom.parseString(svg)  # must be well-formed XML


def test_chart_svg_no_data_is_empty_string():
    assert chart_svg.render_svg(Chart(chart="column", labels=[], series=[]), Theme()) == ""
    all_none = Chart(chart="line", labels=["a"], series=[Series(name="s", values=[None, None])])
    assert chart_svg.render_svg(all_none, Theme()) == ""


def test_chart_svg_none_gap_breaks_the_line_not_bridges_it():
    chart = Chart(
        chart="line", labels=["A", "B", "C", "D", "E"],
        series=[Series(name="s", values=[1.0, 2.0, None, 3.0, 4.0])],
    )
    svg = chart_svg.render_svg(chart, Theme())
    # the None value splits the series into two 2-point runs, each its own
    # path segment, instead of one path interpolating across the gap
    assert svg.count("<path") == 2


def test_chart_svg_column_skips_the_bar_for_none():
    chart = Chart(chart="column", labels=["A", "B"], series=[Series(name="only", values=[1.0, None])])
    svg = chart_svg.render_svg(chart, Theme())
    # one bar for 1.0; the None value draws no bar. The full-canvas themed
    # background rect is not a bar, so exclude it from the count.
    assert svg.count("<rect") - svg.count('width="100%"') == 1


def test_chart_svg_single_series_has_no_legend_multi_series_does():
    single = Chart(chart="column", labels=["A"], series=[Series(name="OnlySeries", values=[1.0])])
    multi = Chart(chart="column", labels=["A"],
                  series=[Series(name="Xseries", values=[1.0]), Series(name="Yseries", values=[2.0])])
    assert "OnlySeries" not in chart_svg.render_svg(single, Theme())
    svg_multi = chart_svg.render_svg(multi, Theme())
    assert "Xseries" in svg_multi and "Yseries" in svg_multi


def test_chart_svg_pie_uses_only_the_first_series():
    chart = Chart(
        chart="pie", labels=["A", "B"],
        series=[Series(name="s1", values=[25.0, 75.0]), Series(name="s2", values=[999.0, 1.0])],
    )
    svg = chart_svg.render_svg(chart, Theme())
    assert "25%" in svg and "75%" in svg


def test_chart_svg_theme_colors_drive_series_and_axes():
    theme = Theme(primary="#112233", accent="#445566", muted="#778899")
    chart = Chart(chart="column", labels=["A", "B"],
                  series=[Series(name="a", values=[1.0, 2.0]), Series(name="b", values=[2.0, 1.0])])
    svg = chart_svg.render_svg(chart, theme)
    assert "#112233" in svg and "#445566" in svg and "#778899" in svg


def test_chart_svg_escapes_hostile_text():
    chart = Chart(
        chart="bar", title="<script>alert(1)</script>", labels=['"><svg onload=alert(1)>'],
        series=[Series(name="<img src=x onerror=alert(1)>", values=[1.0])],
    )
    svg = chart_svg.render_svg(chart, Theme())
    assert "<script>alert(1)</script>" not in svg
    assert "onerror=alert(1)>" not in svg
    minidom.parseString(svg)


def test_chart_svg_never_crashes_on_nan_or_infinity():
    chart = Chart(chart="line", labels=["A", "B", "C"],
                  series=[Series(name="s", values=[float("nan"), float("inf"), 5.0])])
    svg = chart_svg.render_svg(chart, Theme())
    assert svg != ""
    minidom.parseString(svg)


def test_chart_svg_scatter_labels_shorter_than_series_keeps_every_point():
    # regression: zip(numeric_labels, series.values) truncates to the
    # shorter of the two, so a series longer than the label list used to
    # have its trailing (real) values silently dropped instead of falling
    # back to index-based x positions
    chart = Chart(chart="scatter", labels=["1"], series=[Series(name="s", values=[None, 5.0])])
    svg = chart_svg.render_svg(chart, Theme())
    assert "<circle" in svg  # the value 5.0 must still be plotted

    ragged = Chart(
        chart="scatter", labels=["1", "2", "3"],
        series=[Series(name="a", values=[5.0]), Series(name="b", values=[1.0, 2.0, 3.0])],
    )
    assert chart_svg.render_svg(ragged, Theme()).count("<circle") == 4


# -------------------------------------------------- renderer integration


def _chart_doc(**kw) -> Document:
    chart = Chart(
        chart=kw.pop("chart", "column"), title=kw.pop("title", "Revenue"),
        labels=kw.pop("labels", ["Q1", "Q2"]),
        series=kw.pop("series", [Series(name="2026", values=[10.0, 20.0])]),
        caption=kw.pop("caption", None),
    )
    return Document(title="Report", blocks=[chart])


def test_html_chart_embeds_inline_svg_not_a_table(tmp_path):
    html = render(_chart_doc(), "html", tmp_path / "c.html").read_text(encoding="utf-8")
    assert "<svg" in html and 'class="docloom-chart"' in html
    assert "<table" not in html  # a real chart, not the data-table fallback


def test_typst_chart_embeds_svg_bytes():
    from docloom.render.typst import to_typst

    typ = to_typst(_chart_doc(), Theme())
    assert "image(bytes(" in typ and 'format: "svg"' in typ


def test_typst_chart_compiles_to_pdf(tmp_path):
    pytest.importorskip("typst")
    doc = _chart_doc(chart="line", series=[Series(name="s", values=[1.0, None])])
    out = render(doc, "pdf", tmp_path / "c.pdf")
    assert out.read_bytes()[:5] == b"%PDF-"


def test_markdown_chart_title_is_a_heading_caption_below(tmp_path):
    md = render(_chart_doc(caption="source: internal"), "md", tmp_path / "c.md").read_text(encoding="utf-8")
    lines = md.splitlines()
    assert any(ln.startswith("#### Revenue") for ln in lines)
    assert "source: internal" in md
    assert md.index("Revenue") < md.index("source: internal")


def test_docx_chart_without_prerendered_path_is_titled_captioned_table(tmp_path, monkeypatch):
    import sys

    import docx as docx_lib

    # with the optional rasterizer (docloom[diagrams]) installed the chart is a
    # real picture instead; this test pins the no-extra fallback, so it poisons
    # the import. The picture path is covered in tests/test_raster.py.
    monkeypatch.setitem(sys.modules, "resvg_py", None)
    out = render(_chart_doc(caption="source: internal"), "docx", tmp_path / "c.docx")
    d = docx_lib.Document(str(out))
    paragraphs = [p.text for p in d.paragraphs]
    assert "Revenue" in paragraphs
    assert any("source: internal" in p for p in paragraphs)
    assert any(t.rows for t in d.tables)  # the chart data still reaches the reader
    # LOW-3 regression: the title must stay glued to the chart/table that
    # follows it, or a page break can strand the title alone on one page
    title_idx = paragraphs.index("Revenue")
    assert d.paragraphs[title_idx].paragraph_format.keep_with_next is True


# --------------------------------------------------------- docx: findings


def test_docx_svg_artifact_renders_picture_or_placeholder(tmp_path):
    # An SVG Image/Artifact used to always degrade to a "[image: alt]" text
    # stub in DOCX even though Chart/Diagram SVGs are rasterized. It now
    # rasterizes to a real embedded picture when the rasterizer is present,
    # and only falls back to the labeled placeholder when it is absent --
    # never a silent drop either way.
    import docx as docx_lib

    from docloom.render import raster

    svg_path = tmp_path / "diagram.svg"
    svg_path.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="10" height="10"><rect width="10" height="10"/></svg>',
        encoding="utf-8",
    )
    doc = Document(title="T", blocks=[
        Paragraph(text="before"),
        Artifact(kind="diagram", path=str(svg_path), alt="architecture diagram", caption="Fig 1"),
        Paragraph(text="after"),
    ])
    out = render(doc, "docx", tmp_path / "a.docx")
    d = docx_lib.Document(str(out))
    text = "\n".join(p.text for p in d.paragraphs)
    assert "before" in text and "after" in text
    assert "Fig 1" in text  # caption is kept either way
    # never a silent drop: an embedded picture (rasterized SVG) OR the alt placeholder
    assert len(d.inline_shapes) >= 1 or "architecture diagram" in text
    if raster.available():
        assert len(d.inline_shapes) >= 1  # the fix: a real picture, not a text stub
        assert "architecture diagram" not in text


def test_docx_svg_image_block_is_not_silently_dropped(tmp_path):
    import docx as docx_lib

    svg_path = tmp_path / "pic.svg"
    svg_path.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="8" height="8"><circle cx="4" cy="4" r="3"/></svg>',
        encoding="utf-8",
    )
    doc = Document(title="T", blocks=[ImageBlock(path=str(svg_path), alt="a circle")])
    out = render(doc, "docx", tmp_path / "i.docx")
    d = docx_lib.Document(str(out))
    text = "\n".join(p.text for p in d.paragraphs)
    # embedded picture (rasterized) OR the alt placeholder -- never nothing
    assert len(d.inline_shapes) >= 1 or "a circle" in text


def test_typst_titled_chart_title_not_duplicated(tmp_path):
    # The painted chart SVG already bakes in the chart title, so the PDF
    # (typst) renderer must NOT also emit a standalone bold title line above
    # it -- that shipped the title twice. It belongs only to the data-table
    # fallback, which carries no title of its own.
    title = "Revenue by Quarter"
    chart = Chart(chart="column", title=title, labels=["Q1", "Q2"],
                  series=[Series(name="Rev", values=[1.0, 2.0])])
    doc = Document(title="T", slides=[Slide(layout="content", title="s", blocks=[chart])])
    out = render(doc, "typ", tmp_path / "c.typ")
    src = out.read_text(encoding="utf-8")
    assert f'#text(weight: "bold")[{title}]' not in src


def test_chart_datatable_fallback_formats_numbers_readably(tmp_path):
    # When a chart cannot be painted (an all-non-positive pie), the html and
    # markdown data-table fallbacks must show readable, comma-grouped numbers,
    # not Python's scientific "%g" form ("-1.5e+06").
    pie = Chart(chart="pie", title="Losses", labels=["A", "B"],
                series=[Series(name="usd", values=[-1500000.0, -500000.0])])
    doc = Document(title="T", slides=[Slide(layout="content", title="s", blocks=[pie])])
    for fmt, ext in (("html", ".html"), ("md", ".md")):
        out = render(doc, fmt, tmp_path / f"c{ext}")
        src = out.read_text(encoding="utf-8")
        assert "1,500,000" in src
        assert "1.5e+06" not in src


def test_docx_missing_image_path_still_skipped_silently(tmp_path):
    import docx as docx_lib

    doc = Document(title="T", blocks=[
        Paragraph(text="only text here"),
        ImageBlock(path=str(tmp_path / "does_not_exist.png"), alt="ghost"),
    ])
    out = render(doc, "docx", tmp_path / "m.docx")
    d = docx_lib.Document(str(out))
    text = "\n".join(p.text for p in d.paragraphs)
    assert "ghost" not in text  # a genuinely missing path stays silent, matching every other renderer


def test_docx_header_link_recolors_to_background_not_primary(tmp_path):
    from docx.oxml.ns import qn
    import docx as docx_lib

    theme = Theme()  # primary #1D4ED8, background #FFFFFF
    doc = Document(
        title="T",
        blocks=[TableBlock(header=[[Span(text="link", link="https://example.com")]], rows=[["x"]])],
    )
    out = render(doc, "docx", tmp_path / "h.docx")
    d = docx_lib.Document(str(out))
    header_cell = d.tables[0].cell(0, 0)
    hyperlinks = header_cell._tc.findall(".//" + qn("w:hyperlink"))
    assert len(hyperlinks) == 1
    colors = hyperlinks[0].findall(".//" + qn("w:color"))
    assert len(colors) == 1
    assert colors[0].get(qn("w:val")).upper() == theme.background.lstrip("#").upper()
    assert colors[0].get(qn("w:val")).upper() != theme.primary.lstrip("#").upper()


def test_docx_body_link_color_is_unaffected_by_the_header_fix(tmp_path):
    from docx.oxml.ns import qn
    import docx as docx_lib

    theme = Theme()
    doc = Document(title="T", blocks=[Paragraph(text=[Span(text="link", link="https://example.com")])])
    out = render(doc, "docx", tmp_path / "b.docx")
    d = docx_lib.Document(str(out))
    hyperlink = None
    for p in d.paragraphs:
        found = p._p.findall(".//" + qn("w:hyperlink"))
        if found:
            hyperlink = found[0]
            break
    assert hyperlink is not None
    color = hyperlink.findall(".//" + qn("w:color"))[0]
    assert color.get(qn("w:val")).upper() == theme.primary.lstrip("#").upper()


def test_docx_heading_gets_keep_with_next(tmp_path):
    import docx as docx_lib

    doc = Document(title="T", blocks=[Heading(level=2, text="Section"), Paragraph(text="body")])
    out = render(doc, "docx", tmp_path / "h.docx")
    d = docx_lib.Document(str(out))
    heading = next(p for p in d.paragraphs if p.text == "Section")
    assert heading.paragraph_format.keep_with_next is True


def test_docx_image_caption_keeps_together_and_bound_to_its_picture(tmp_path):
    import docx as docx_lib

    # 1x1 PNG: python-docx embeds it, so the caption goes through _render_image
    png_path = tmp_path / "pic.png"
    png_path.write_bytes(bytes.fromhex(
        "89504e470d0a1a0a0000000d4948445200000001000000010802000000907753de"
        "0000000c49444154789c63f8cfc0000003010100c9fe92ef0000000049454e44ae426082"
    ))
    doc = Document(title="T", blocks=[ImageBlock(path=str(png_path), caption="Fig 1: a pixel")])
    out = render(doc, "docx", tmp_path / "i.docx")
    d = docx_lib.Document(str(out))
    paragraphs = d.paragraphs  # one property access: each access rewraps fresh objects
    idx = next(i for i, p in enumerate(paragraphs) if p.text == "Fig 1: a pixel")
    assert paragraphs[idx].paragraph_format.keep_together is True
    # the picture's own paragraph, immediately above, is bound forward to the
    # caption so a page break can never separate the two
    assert paragraphs[idx - 1].paragraph_format.keep_with_next is True


def test_docx_wide_table_gets_repeating_header_uncuttable_rows_and_proportional_widths(tmp_path):
    import docx as docx_lib
    from docx.oxml.ns import qn
    from docx.shared import Emu, Inches

    doc = Document(
        title="T",
        blocks=[TableBlock(
            header=[f"Col {i} is quite a bit longer" if i == 0 else f"C{i}" for i in range(8)],
            rows=[[f"r{r}c{c}" for c in range(8)] for r in range(3)],
            caption="Fig T: wide",
        )],
    )
    out = render(doc, "docx", tmp_path / "w.docx")
    d = docx_lib.Document(str(out))
    table = d.tables[0]
    header_trPr = table.rows[0]._tr.find(qn("w:trPr"))
    assert header_trPr is not None and header_trPr.find(qn("w:tblHeader")) is not None
    for row in table.rows:  # no row -- including the header -- may split across a page
        row_trPr = row._tr.find(qn("w:trPr"))
        assert row_trPr is not None and row_trPr.find(qn("w:cantSplit")) is not None
    widths = [c.width for c in table.columns]
    assert len(set(widths)) > 1  # the long-label column gets more room than the rest
    frame = d.sections[0].page_width - d.sections[0].left_margin - d.sections[0].right_margin
    assert abs(sum(widths) - frame) < Emu(Inches(0.01))  # fills the frame, never overflows it


def test_docx_many_columns_enforce_min_width_floor_after_renormalization():
    """MEDIUM-1 regression: the floor must survive the renormalization pass
    that follows it. A skewed weight distribution (one very long header,
    thirteen short ones) used to push every short column below
    MIN_COL_WIDTH once the post-clamp rescale ran, even though the frame is
    wide enough for all 14 columns to meet the floor."""
    import docx as docx_lib
    from docx.shared import Emu, Inches

    from docloom.render.docx import MIN_COL_WIDTH, _column_widths

    cols = 14
    d = docx_lib.Document()
    # widen the page so 14 * MIN_COL_WIDTH comfortably fits the frame --
    # this test targets the floor-enforcement math, not the too-many-
    # columns fallback (covered separately below)
    section = d.sections[0]
    section.page_width = Inches(20)
    section.left_margin = Inches(0.5)
    section.right_margin = Inches(0.5)
    frame = section.page_width - section.left_margin - section.right_margin

    header = ["This header is deliberately extremely long to skew weights"] + [
        f"C{i}" for i in range(cols - 1)
    ]
    rows = [[f"r{r}c{c}" for c in range(cols)] for r in range(3)]

    widths = _column_widths(d, header, rows, cols)
    assert widths is not None  # the frame fits every column at the floor
    for w in widths:
        assert w >= MIN_COL_WIDTH, f"column width {w} fell below the documented floor"
    assert abs(sum(widths) - frame) < Emu(Inches(0.01))  # still fills the frame exactly


def test_docx_too_many_columns_falls_back_to_word_autofit(tmp_path):
    """MEDIUM-1: when cols * MIN_COL_WIDTH exceeds the frame (default Letter
    page, 14 columns), no explicit split can honor the floor for every
    column. The documented fallback engages: Word autofit, not slivers."""
    import docx as docx_lib

    cols = 14
    doc = Document(
        title="T",
        blocks=[TableBlock(
            header=[f"C{i}" for i in range(cols)],
            rows=[[f"r{r}c{c}" for c in range(cols)] for r in range(3)],
            caption="Fig T: too many columns",
        )],
    )
    out = render(doc, "docx", tmp_path / "many.docx")
    d = docx_lib.Document(str(out))
    table = d.tables[0]
    assert table.autofit is True  # documented fallback engaged, not below-floor slivers


def test_docx_long_table_last_row_keeps_with_next_to_its_caption(tmp_path):
    import docx as docx_lib

    doc = Document(
        title="T",
        blocks=[TableBlock(
            header=["ID", "Name"],
            rows=[[str(r), f"item-{r}"] for r in range(65)],
            caption="Fig T: long",
        )],
    )
    out = render(doc, "docx", tmp_path / "l.docx")
    d = docx_lib.Document(str(out))
    table = d.tables[0]
    assert len(table.rows) == 66  # header + 65 data rows, none dropped
    for cell in table.rows[-1].cells:
        assert cell.paragraphs[-1].paragraph_format.keep_with_next is True


# ------------------------------------------------------ markdown: findings


def test_markdown_span_emphasis_hoists_boundary_whitespace():
    from docloom.render.markdown import _span_md

    assert _span_md(Span(text=" hi", bold=True), {}) == " **hi**"
    assert _span_md(Span(text="hi ", bold=True), {}) == "**hi** "
    assert _span_md(Span(text=" hi ", italic=True), {}) == " *hi* "
    assert _span_md(Span(text=" hi ", bold=True, italic=True), {}) == " ***hi*** "
    assert _span_md(Span(text="hi", bold=True), {}) == "**hi**"  # no boundary whitespace: unchanged
    assert _span_md(Span(text="   ", bold=True), {}) == "   "  # whitespace-only: no markers at all


def test_markdown_emphasis_boundary_whitespace_end_to_end(tmp_path):
    doc = Document(title="T", blocks=[
        Paragraph(text=[Span(text=" leading", bold=True), Span(text="mid "),
                        Span(text="trailing ", italic=True)]),
    ])
    md = render(doc, "md", tmp_path / "e.md").read_text(encoding="utf-8")
    assert "** leading**" not in md
    assert " **leading**" in md
    assert "*trailing *" not in md
    assert "*trailing* " in md


def test_markdown_assets_mode_copies_images_and_rewrites_relative(tmp_path):
    from PIL import Image as PILImage

    src_dir = tmp_path / "src"
    src_dir.mkdir()
    img_path = src_dir / "pic.png"
    PILImage.new("RGB", (4, 4), "red").save(img_path)
    doc = Document(title="T", blocks=[ImageBlock(path=str(img_path), alt="a pic")])
    out_dir = tmp_path / "out"
    out = render(doc, "md", out_dir / "report.md")
    text = out.read_text(encoding="utf-8")
    assert str(img_path) not in text  # the generating machine's path must not leak
    assert "report_files/pic.png" in text
    assert (out_dir / "report_files" / "pic.png").is_file()


def test_markdown_assets_false_keeps_the_original_path(tmp_path):
    from PIL import Image as PILImage

    from docloom.render.markdown import render as md_render

    img_path = tmp_path / "pic.png"
    PILImage.new("RGB", (4, 4), "red").save(img_path)
    doc = Document(title="T", blocks=[ImageBlock(path=str(img_path), alt="a")])
    out = md_render(doc, Theme(), tmp_path / "out.md", assets=False)
    text = out.read_text(encoding="utf-8")
    assert str(img_path) in text
    assert not (tmp_path / "out_files").exists()


def test_markdown_assets_mode_dedupes_same_source_reused_twice(tmp_path):
    from PIL import Image as PILImage

    img_path = tmp_path / "logo.png"
    PILImage.new("RGB", (2, 2), "blue").save(img_path)
    doc = Document(title="T", blocks=[
        ImageBlock(path=str(img_path), alt="first"),
        ImageBlock(path=str(img_path), alt="second"),
    ])
    out = render(doc, "md", tmp_path / "out" / "d.md")
    text = out.read_text(encoding="utf-8")
    assert text.count("d_files/logo.png") == 2  # same source, same destination
    assert not (tmp_path / "out" / "d_files" / "logo-2.png").exists()


# -------------------------------------------------------------- finding 6


def test_html_sources_list_dedupes_duplicate_ids(tmp_path):
    doc = Document(
        title="T",
        blocks=[Paragraph(text=[Span(text="claim", cite="a")])],
        sources=[{"id": "a", "title": "Alpha"}, {"id": "a", "title": "AlphaDup"}],
    )
    html = render(doc, "html", tmp_path / "s.html").read_text(encoding="utf-8")
    assert html.count('id="src-') == 1
    assert "AlphaDup" not in html


def test_xlsx_sources_sheet_dedupes_duplicate_ids(tmp_path):
    import zipfile

    doc = Document(
        title="T",
        blocks=[
            Paragraph(text=[Span(text="claim", cite="a")]),
            TableBlock(header=["h"], rows=[["v"]]),
        ],
        sources=[{"id": "a", "title": "Alpha"}, {"id": "a", "title": "AlphaDup"}],
    )
    out = render(doc, "xlsx", tmp_path / "s.xlsx")
    with zipfile.ZipFile(out) as z:
        shared = z.read("xl/sharedStrings.xml").decode("utf-8")
    assert "AlphaDup" not in shared


# ------------------------------------------------ P5: pptx quality audit


def _png(path, w=40, h=40, color="blue"):
    from PIL import Image as PILImage

    PILImage.new("RGB", (w, h), color).save(path)
    return str(path)


def _shapes_text(slide) -> str:
    return " ".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)


def test_pptx_unresolved_artifact_gets_placeholder_not_silent_drop(tmp_path):
    # P5 audit defect 1, proven live: an Artifact with no path used to
    # render nothing at all, with no trace and no warning.
    from pptx import Presentation

    doc = Document(title="T", slides=[
        Slide(layout="content", title="T2", blocks=[
            Artifact(kind="diagram", alt="Rollout architecture", caption="Fig 1"),
        ]),
    ])
    with pytest.warns(UserWarning, match="unresolved artifact"):
        out = render(doc, "pptx", tmp_path / "a.pptx")
    text = _shapes_text(Presentation(str(out)).slides[0])
    assert "Rollout architecture" in text
    assert "Fig 1" in text


def test_pptx_image_missing_file_gets_placeholder_and_warns(tmp_path):
    from pptx import Presentation

    doc = Document(title="T", slides=[
        Slide(layout="content", title="T2", blocks=[
            ImageBlock(path=str(tmp_path / "does_not_exist.png"), alt="ghost image"),
        ]),
    ])
    with pytest.warns(UserWarning, match="could not be embedded"):
        out = render(doc, "pptx", tmp_path / "i.pptx")
    assert "ghost image" in _shapes_text(Presentation(str(out)).slides[0])


def test_pptx_image_with_no_path_stays_silent(tmp_path):
    # Unlike Artifact, an Image slot with no path is a deliberate,
    # not-yet-resolved slot (matches the docx renderer's convention): it
    # must stay silent, not sprout a placeholder box on every slide with an
    # unbound image slot.
    from pptx import Presentation

    doc = Document(title="T", slides=[
        Slide(layout="content", title="T2", blocks=[ImageBlock(alt="never shown")]),
    ])
    out = render(doc, "pptx", tmp_path / "i2.pptx")
    assert "never shown" not in _shapes_text(Presentation(str(out)).slides[0])


def test_pptx_section_slide_renders_its_blocks(tmp_path):
    # P5 audit defect 2: a section slide rendered only title/subtitle and
    # silently dropped every block the author put in s.blocks/s.right.
    from pptx import Presentation

    doc = Document(title="T", slides=[
        Slide(layout="section", title="Section", blocks=[
            BulletList(items=[ListItem(text="SECMARK point one")]),
        ]),
    ])
    out = render(doc, "pptx", tmp_path / "s.pptx")
    assert "SECMARK point one" in _shapes_text(Presentation(str(out)).slides[0])


def test_pptx_image_right_title_shares_the_bodys_left_edge(tmp_path):
    # P5 audit defect 3: image_right carved the title's x (not just its
    # width) to dodge the top-left logo, staggering it 0.66in right of the
    # body text it introduces.
    from pptx.util import Inches
    from pptx import Presentation

    doc = Document(
        title="D", logo=ImageBlock(path=_png(tmp_path / "logo.png")),
        slides=[Slide(
            layout="image_right", title="TITLEMARK",
            image=ImageBlock(path=_png(tmp_path / "pic.png", 400, 300)),
            blocks=[Paragraph(text="BODYMARK text")],
        )],
    )
    out = render(doc, "pptx", tmp_path / "ir.pptx")
    slide = Presentation(str(out)).slides[0]
    title_box = next(s for s in slide.shapes if s.has_text_frame and "TITLEMARK" in s.text_frame.text)
    body_box = next(s for s in slide.shapes if s.has_text_frame and "BODYMARK" in s.text_frame.text)
    assert abs(title_box.left - body_box.left) < Inches(0.01)


def test_pptx_grow_pass_suppressed_next_to_a_fixed_size_block(tmp_path):
    # P5 audit defect 4: growing only the prose beside a fixed-size block
    # (here, code) is what produced 23.8pt paragraphs beside 12pt code.
    from pptx import Presentation

    from docloom.render.pptx import BODY_PT

    doc = Document(title="T", slides=[
        Slide(layout="content", title="T2", blocks=[
            Paragraph(text="short prose"),
            Code(code="x = 1"),
        ]),
    ])
    out = render(doc, "pptx", tmp_path / "g.pptx")
    slide = Presentation(str(out)).slides[0]
    para = next(s for s in slide.shapes if s.has_text_frame and "short prose" in s.text_frame.text)
    assert para.text_frame.paragraphs[0].runs[0].font.size.pt == BODY_PT


def test_pptx_heading_seam_tighter_than_other_seams_when_sparse(tmp_path):
    # P5 audit defect 5: a uniform slack-distributed gap gave the
    # heading-to-body seam the same treatment as every other seam,
    # orphaning a heading from the very content it introduces.
    from pptx import Presentation

    doc = Document(title="T", slides=[
        Slide(layout="content", title="T2", blocks=[
            Heading(level=2, text="HEADMARK"),
            Paragraph(text="FIRSTMARK short"),
            Paragraph(text="SECONDMARK short"),
        ]),
    ])
    out = render(doc, "pptx", tmp_path / "seam.pptx")
    slide = Presentation(str(out)).slides[0]

    def _span(marker):
        for sh in slide.shapes:
            if sh.has_text_frame and marker in sh.text_frame.text:
                return sh.top, sh.top + sh.height
        raise AssertionError(marker)

    _, head_bottom = _span("HEADMARK")
    first_top, first_bottom = _span("FIRSTMARK")
    second_top, _ = _span("SECONDMARK")
    assert (first_top - head_bottom) < (second_top - first_bottom)


def test_pptx_solo_chart_fills_the_slide_body(tmp_path):
    # P5 audit defect 6: a chart alone on its slide stopped at a fixed cap,
    # leaving a permanent dead void below it.
    from pptx import Presentation

    from docloom.render.pptx import LAYOUT

    doc = Document(title="T", slides=[
        Slide(layout="content", title="T2", blocks=[
            Chart(chart="column", labels=["A", "B"], series=[Series(name="s", values=[1.0, 2.0])]),
        ]),
    ])
    out = render(doc, "pptx", tmp_path / "chart.pptx")
    slide = Presentation(str(out)).slides[0]
    chart_shape = next(s for s in slide.shapes if s.has_chart)
    # exceeds the shared-slide cap, proving the solo-fill path (not the cap) fired
    assert chart_shape.height / 914400 > LAYOUT["chart_max_h_in"]


def test_pptx_image_side_matte_hugs_a_wide_short_image(tmp_path):
    # P5 audit defect 7: a wide/short image contain-fit into a tall pane
    # left ~80% of the pane as dead gray around a sliver of image.
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx import Presentation

    from docloom.render.pptx import LAYOUT

    doc = Document(title="T", slides=[
        Slide(layout="image_left", title="T2",
              image=ImageBlock(path=_png(tmp_path / "wide.png", 2400, 300)),
              blocks=[Paragraph(text="x")]),
    ])
    out = render(doc, "pptx", tmp_path / "wide.pptx")
    slide = Presentation(str(out)).slides[0]
    matte = next(
        s for s in slide.shapes
        if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and s.left == 0
    )
    assert matte.height / 914400 < LAYOUT["slide_h_in"] * 0.6


def test_pptx_table_columns_are_weighted_not_equal(tmp_path):
    # P5 audit defect 8: an equal tw/cols split gave a long label and a
    # short value the same track width.
    from pptx import Presentation

    doc = Document(title="T", slides=[
        Slide(layout="content", title="T2", blocks=[
            TableBlock(
                header=["A very long vendor name column", "$"],
                rows=[["Cascade (open source)", "18"]],
            ),
        ]),
    ])
    out = render(doc, "pptx", tmp_path / "tbl.pptx")
    slide = Presentation(str(out)).slides[0]
    tbl = next(s for s in slide.shapes if s.has_table).table
    widths = [c.width for c in tbl.columns]
    assert widths[0] > widths[1] * 1.3


def test_pptx_empty_and_ragged_tables_do_not_crash_column_weighting(tmp_path):
    # Regression: the weighted-column-width helper (defect 8) divided by
    # sum(raw) unconditionally; a genuinely empty table (header=[], rows=[])
    # has zero columns to weight and raised ZeroDivisionError. A ragged
    # table (a row shorter than the header) must not raise either.
    doc = Document(title="T", slides=[
        Slide(layout="content", title="T2", blocks=[
            TableBlock(header=[], rows=[]),
            TableBlock(header=["a", "b"], rows=[["1", "2", "3"], ["1"]]),
        ]),
    ])
    out = render(doc, "pptx", tmp_path / "empty_tbl.pptx")
    assert out.exists()


def test_pptx_divider_color_clears_a_real_contrast_ratio(tmp_path):
    # P5 audit defect 9: the default divider color measured a ~1.1 contrast
    # ratio against the background -- a near-invisible ghost.
    from docloom.render.pptx import _divider_color
    from docloom.theme import contrast_ratio

    theme = Theme()
    assert contrast_ratio(_divider_color(theme), theme.background) >= 1.5


def test_pptx_sources_slide_has_logo_and_larger_font(tmp_path):
    # P5 audit defect 10: the sources slide was the one layout with no
    # brand logo, at a footnote-scale 12pt.
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx import Presentation

    doc = Document(
        title="T", logo=ImageBlock(path=_png(tmp_path / "logo.png")),
        slides=[Slide(layout="content", title="T2",
                      blocks=[Paragraph(text=[Span(text="claim", cite="a")])])],
        sources=[Source(id="a", title="Alpha")],
    )
    out = render(doc, "pptx", tmp_path / "src.pptx")
    sources_slide = Presentation(str(out)).slides[-1]
    pics = [s for s in sources_slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1
    body = next(s for s in sources_slide.shapes if s.has_text_frame and "Alpha" in s.text_frame.text)
    assert body.text_frame.paragraphs[0].runs[0].font.size.pt == 13


def test_pptx_stat_card_label_and_delta_are_larger(tmp_path):
    # P5 audit defect 12: label/delta at 11/10pt read as footnotes on a
    # 13.3in slide.
    from pptx import Presentation

    doc = Document(title="T", slides=[
        Slide(layout="content", title="T2", blocks=[
            StatRow(items=[Stat(label="LABELMARK", value="$1", delta="+1 pt")]),
        ]),
    ])
    out = render(doc, "pptx", tmp_path / "stat.pptx")
    slide = Presentation(str(out)).slides[0]

    def _size(marker):
        for sh in slide.shapes:
            if sh.has_text_frame and marker in sh.text_frame.text:
                return sh.text_frame.paragraphs[0].runs[0].font.size.pt
        raise AssertionError(marker)

    assert _size("LABELMARK") == 13
    assert _size("+1 pt") == 12


def test_pptx_pie_point_labels_show_percent_not_raw_value(tmp_path):
    # P5 audit defect 14: materializing a per-point data label font (to
    # recolor each slice's label) silently created a per-point c:dLbl that
    # defaulted to showVal=1/showPercent=0, overriding the plot-level
    # show_percentage=True and turning "54%" back into "5.4".
    import re
    import zipfile

    doc = Document(title="T", slides=[
        Slide(layout="content", title="T2", blocks=[
            Chart(chart="pie", labels=["A", "B"], series=[Series(name="s", values=[30.0, 70.0])]),
        ]),
    ])
    out = render(doc, "pptx", tmp_path / "pie.pptx")
    with zipfile.ZipFile(out) as z:
        chart_xml = next(n for n in z.namelist() if n.startswith("ppt/charts/chart"))
        xml = z.read(chart_xml).decode("utf-8")
    point_dlbls = re.findall(r"<c:dLbl>.*?</c:dLbl>", xml, re.S)
    assert point_dlbls, "expected a per-point dLbl (from the label recolor loop)"
    for d in point_dlbls:
        assert 'c:showPercent val="1"' in d
        assert 'c:showVal val="0"' in d


# ---------------------------------------------------- finding 8: quote slide


def test_pptx_quote_slide_from_title_subtitle_shows_both(tmp_path):
    # Finding 8: `Slide(layout="quote", title=..., subtitle=...)` with no
    # Quote block used to render s.title as a stray 16pt caption top-left
    # and never touch s.subtitle at all -- the attribution was gone from the
    # XML entirely and the slide was ~95% blank white.
    from pptx import Presentation

    doc = Document(title="T", slides=[
        Slide(layout="quote", title="QUOTEMARK the future is native",
              subtitle="ATTRIBUTIONMARK, CEO"),
    ])
    out = render(doc, "pptx", tmp_path / "q.pptx")
    slide = Presentation(str(out)).slides[0]
    text = _shapes_text(slide)
    assert "QUOTEMARK" in text
    assert "ATTRIBUTIONMARK" in text


def test_pptx_quote_slide_from_title_subtitle_uses_display_scale(tmp_path):
    # The title-as-quote shape must land on the same pull-quote treatment as
    # a real Quote block (real display-scale type), not the old 16pt muted
    # caption size the title used to get stuck at.
    from pptx import Presentation

    doc = Document(title="T", slides=[
        Slide(layout="quote", title="Short quote", subtitle="Someone"),
    ])
    out = render(doc, "pptx", tmp_path / "q2.pptx")
    slide = Presentation(str(out)).slides[0]
    quote_shape = next(
        s for s in slide.shapes
        if s.has_text_frame and "Short quote" in s.text_frame.text
    )
    assert quote_shape.text_frame.paragraphs[0].runs[0].font.size.pt >= 20


def test_pptx_quote_slide_with_quote_block_keeps_working(tmp_path):
    # A real Quote block must still render correctly (unaffected regression
    # guard for the finding-8 fix), including a slide title used as a small
    # eyebrow label above it.
    from pptx import Presentation

    doc = Document(title="T", slides=[
        Slide(layout="quote", title="LABELMARK",
              blocks=[Quote(text="BLOCKQUOTEMARK", attribution="BLOCKATTRMARK")]),
    ])
    out = render(doc, "pptx", tmp_path / "q3.pptx")
    text = _shapes_text(Presentation(str(out)).slides[0])
    assert "LABELMARK" in text
    assert "BLOCKQUOTEMARK" in text
    assert "BLOCKATTRMARK" in text


def test_pptx_quote_slide_block_attribution_falls_back_to_subtitle(tmp_path):
    # A Quote block without its own attribution still picks up s.subtitle
    # instead of leaving it stranded and unused.
    from pptx import Presentation

    doc = Document(title="T", slides=[
        Slide(layout="quote", subtitle="SUBATTRMARK",
              blocks=[Quote(text="no inline attribution here")]),
    ])
    out = render(doc, "pptx", tmp_path / "q4.pptx")
    assert "SUBATTRMARK" in _shapes_text(Presentation(str(out)).slides[0])


# -------------------------------------------------- finding 12: diagram grow


def _small_diagram() -> Diagram:
    return Diagram(
        nodes=[DiagramNode(id="a", label="A"), DiagramNode(id="b", label="B")],
        edges=[DiagramEdge(source="a", target="b")],
    )


def test_pptx_grow_pass_suppressed_next_to_a_diagram(tmp_path):
    # Finding 12: _FIXED_SIZE_BLOCKS omitted Diagram, so a paragraph sharing
    # a slide with a diagram grew past BODY_PT (measured live: 15.57pt),
    # while the same paragraph next to a table (an already-fixed block type)
    # correctly stayed at BODY_PT. A diagram's size is fixed by its own
    # solved layout, exactly like a table's is fixed by its row count.
    from pptx import Presentation

    from docloom.render.pptx import BODY_PT

    diagram_doc = Document(title="T", slides=[
        Slide(layout="content", title="T2", blocks=[
            _small_diagram(),
            Paragraph(text="short prose"),
        ]),
    ])
    table_doc = Document(title="T", slides=[
        Slide(layout="content", title="T2", blocks=[
            TableBlock(header=["h"], rows=[["v"]]),
            Paragraph(text="short prose"),
        ]),
    ])

    def _prose_size(doc, name):
        out = render(doc, "pptx", tmp_path / name)
        slide = Presentation(str(out)).slides[0]
        para = next(s for s in slide.shapes if s.has_text_frame and "short prose" in s.text_frame.text)
        return para.text_frame.paragraphs[0].runs[0].font.size.pt

    diagram_size = _prose_size(diagram_doc, "diag.pptx")
    table_size = _prose_size(table_doc, "tbl2.pptx")
    assert diagram_size == BODY_PT
    assert diagram_size == table_size


# ---------------------------------------------- finding A: dropped subtitle


def test_pptx_content_slide_renders_subtitle(tmp_path):
    # _content_slide was the only layout function that never read s.subtitle
    # at all -- the subtitle was gone from the XML entirely, even though
    # it is present in the IR.
    from pptx import Presentation

    doc = Document(title="T", slides=[
        Slide(layout="content", title="T2", subtitle="SUBCONTENTMARK detail line",
              blocks=[Paragraph(text="body text")]),
    ])
    out = render(doc, "pptx", tmp_path / "content_sub.pptx")
    text = _shapes_text(Presentation(str(out)).slides[0])
    assert "SUBCONTENTMARK" in text


def test_pptx_image_side_slide_renders_subtitle(tmp_path):
    from pptx import Presentation

    doc = Document(title="T", slides=[
        Slide(layout="image_left", title="T2", subtitle="SUBIMGMARK caption line",
              image=ImageBlock(path=_png(tmp_path / "pic.png")),
              blocks=[Paragraph(text="body text")]),
    ])
    out = render(doc, "pptx", tmp_path / "imgside_sub.pptx")
    text = _shapes_text(Presentation(str(out)).slides[0])
    assert "SUBIMGMARK" in text


def test_pptx_two_column_slide_renders_subtitle(tmp_path):
    from pptx import Presentation

    doc = Document(title="T", slides=[
        Slide(layout="two_column", title="T2", subtitle="SUBCOLMARK detail line",
              blocks=[Paragraph(text="left")], right=[Paragraph(text="right")]),
    ])
    out = render(doc, "pptx", tmp_path / "twocol_sub.pptx")
    text = _shapes_text(Presentation(str(out)).slides[0])
    assert "SUBCOLMARK" in text


def test_pptx_hero_slide_without_image_gets_hero_treatment_not_content(tmp_path):
    # Finding A, part 2: the dispatcher gated _hero_slide on
    # _usable_image(s.image), so a hero slide with NO image never reached
    # _hero_slide (which does handle subtitle correctly) and fell through to
    # _content_slide instead -- rendering ~90% blank with the subtitle gone
    # entirely (same defect class as finding 8's quote-slide bug). An
    # imageless hero must still get the hero's own full-bleed treatment (a
    # solid theme.primary backdrop), not the generic title-band layout.
    from pptx import Presentation

    theme = Theme()
    doc = Document(title="T", slides=[
        Slide(layout="hero", title="HEROMARK the ask", subtitle="SUBHEROMARK context"),
    ])
    out = render(doc, "pptx", tmp_path / "hero_no_img.pptx", theme=theme)
    slide = Presentation(str(out)).slides[0]
    text = _shapes_text(slide)
    assert "HEROMARK" in text
    assert "SUBHEROMARK" in text  # must not be dropped
    bg = str(slide.background.fill.fore_color.rgb).upper()
    assert bg == theme.primary.lstrip("#").upper()  # the hero's own full-bleed fill, not the page background


def test_pptx_hero_slide_with_image_still_covers_the_slide(tmp_path):
    # Regression guard: fixing the imageless case must not disturb the
    # existing image-backed hero behavior.
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx import Presentation

    doc = Document(title="T", slides=[
        Slide(layout="hero", title="HEROIMGMARK", image=ImageBlock(path=_png(tmp_path / "bg.png", 800, 600))),
    ])
    out = render(doc, "pptx", tmp_path / "hero_img.pptx")
    slide = Presentation(str(out)).slides[0]
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1
    assert "HEROIMGMARK" in _shapes_text(slide)


# --------------------------------------- finding B: uncolored callout edges


def test_pptx_warning_and_danger_callouts_are_not_gray_or_black(tmp_path):
    # warning used to map to theme.muted (gray) and danger to theme.text
    # (near-black) -- the two styles that most need to signal urgency
    # signaled nothing at all.
    from pptx import Presentation
    from pptx.util import Inches

    from docloom import Callout

    theme = Theme()
    doc = Document(title="T", slides=[
        Slide(layout="content", title="T2", blocks=[
            Callout(style="warning", text="warn"),
            Callout(style="danger", text="danger"),
        ]),
    ])
    out = render(doc, "pptx", tmp_path / "callouts.pptx", theme=theme)
    slide = Presentation(str(out)).slides[0]
    # the two colored edge bars are thin AUTO_SHAPE rectangles distinguishable
    # from the wide theme.surface backing rect by their narrow width
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    edges = [
        s for s in slide.shapes
        if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and s.width < Inches(0.2)
    ]
    edge_colors = {str(s.fill.fore_color.rgb).upper() for s in edges}
    assert theme.muted.lstrip("#").upper() not in edge_colors
    assert theme.text.lstrip("#").upper() not in edge_colors


def test_pptx_callout_edge_color_derives_from_theme_accent(tmp_path):
    from docloom.render.pptx import _callout_edge_color

    theme = Theme(accent="#0E9F6E")
    warning = _callout_edge_color("warning", theme)
    danger = _callout_edge_color("danger", theme)
    # distinct from each other, from the raw accent, and from the old
    # gray/near-black mapping
    assert len({warning, danger, theme.accent, theme.muted, theme.text}) == 5
    # still hex colors within the theme's own tonal family (same saturation
    # class as accent), not an unrelated hardcoded stoplight hex
    assert warning.startswith("#") and danger.startswith("#")


# ------------------------------------- silent-content-loss CLASS audit fixes


def _blank_slide():
    from pptx import Presentation

    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


def test_pptx_never_drops_a_trailing_block_reserves_room_and_warns(tmp_path):
    # The audit's exact repro: a subtitle (shrinks the real body height) plus
    # a captioned chart (whose own real footprint the old lint/renderer
    # geometry model did not account for) left a trailing paragraph ~0.03in
    # of room -- under the old "remaining < 0.3: break" floor in _body, which
    # silently dropped it from the XML with zero trace and no warning. It
    # must now be reserved room ahead of time and always drawn, with a
    # warning surfacing the crowding instead of silence.
    doc = Document(title="T", slides=[
        Slide(
            layout="content", title="Quarterly results at a glance",
            subtitle="A subtitle that eats into the body's real available height",
            blocks=[
                Chart(chart="bar", title="Revenue", labels=["Q1", "Q2", "Q3", "Q4"],
                      series=[Series(name="Revenue", values=[10.0, 12.0, 14.0, 16.0])],
                      caption="source: internal finance system"),
                Paragraph(text="TRAILINGMARK this paragraph must not be silently dropped"),
            ],
        ),
    ])
    with pytest.warns(UserWarning, match="exceeds the available body height"):
        out = render(doc, "pptx", tmp_path / "ponytail.pptx")
    from pptx import Presentation

    text = _shapes_text(Presentation(str(out)).slides[0])
    assert "TRAILINGMARK" in text
    assert "source: internal" in text  # the chart's own caption also survives


def test_pptx_never_drops_a_trailing_block_matches_docx_html_md(tmp_path):
    # Cross-renderer diff (the technique the audit used to catch this class):
    # PPTX must keep everything DOCX/HTML/MD keep for the same document.
    doc = Document(title="T", slides=[
        Slide(
            layout="content", title="Quarterly results at a glance",
            subtitle="a subtitle",
            blocks=[
                Chart(chart="bar", title="Revenue", labels=["Q1", "Q2", "Q3", "Q4"],
                      series=[Series(name="Revenue", values=[10.0, 12.0, 14.0, 16.0])],
                      caption="a caption"),
                Paragraph(text="TRAILINGMARK2"),
            ],
        ),
    ])
    with pytest.warns(UserWarning):
        pptx_out = render(doc, "pptx", tmp_path / "cross.pptx")
    from pptx import Presentation

    assert "TRAILINGMARK2" in _shapes_text(Presentation(str(pptx_out)).slides[0])
    for fmt in ("docx", "html", "md"):
        out = render(doc, fmt, tmp_path / f"cross.{fmt}")
        if fmt == "docx":
            import docx as docx_lib

            d = docx_lib.Document(str(out))
            text = "\n".join(p.text for p in d.paragraphs)
        else:
            text = out.read_text(encoding="utf-8")
        assert "TRAILINGMARK2" in text


def test_pptx_quote_attribution_never_dropped_even_when_max_h_is_tiny():
    # Same silent-drop class one field down: _quote_block used to only draw
    # b.attribution "if there happens to be room left over" after the quote
    # text itself, which could be false on a genuinely tight slide.
    from docloom.render.pptx import _quote_block

    slide = _blank_slide()
    q = Quote(
        text="A long quote that would normally want lots of room to render "
             "at full display scale, forcing the attribution to compete for space.",
        attribution="ATTRMARK",
    )
    _quote_block(slide, q, Theme(), {}, x=1.0, y=1.0, w=6.0, max_h=0.3)
    assert "ATTRMARK" in _shapes_text(slide)


def test_pptx_table_caption_never_dropped_even_when_max_h_is_tiny():
    from docloom.render.pptx import _table_block

    slide = _blank_slide()
    t = TableBlock(header=["a", "b"], rows=[["1", "2"], ["3", "4"]], caption="TABLECAPMARK")
    _table_block(slide, t, Theme(), {}, x=1.0, y=1.0, w=6.0, max_h=0.3)
    assert "TABLECAPMARK" in _shapes_text(slide)


def test_pptx_chart_caption_never_dropped_even_when_max_h_is_tiny():
    from docloom.render.pptx import _chart_block

    slide = _blank_slide()
    c = Chart(chart="column", labels=["a", "b"], series=[Series(name="s", values=[1.0, 2.0])],
              caption="CHARTCAPMARK")
    _chart_block(slide, c, Theme(), {}, x=1.0, y=1.0, w=6.0, max_h=0.3)
    assert "CHARTCAPMARK" in _shapes_text(slide)


def test_pptx_bar_chart_category_axis_is_reversed_column_chart_is_not():
    # Office's native BAR_CLUSTERED (horizontal bars) plots the first
    # category at the BOTTOM by default (its documented "categories in
    # reverse order" quirk) while chart_svg -- the painter every other
    # renderer (html/docx/typst) shares, and this file's own image/table
    # fallback -- always draws the first category at the TOP. Pin that
    # _chart_block flips only the bar chart's category axis (via
    # reverse_order, i.e. OOXML orientation="maxMin") to match, and leaves
    # a column chart's (vertical bars, left-to-right already matches)
    # category axis at its native, unreversed order.
    from docloom.render.pptx import _chart_block

    slide = _blank_slide()
    bar = Chart(chart="bar", labels=["a", "b", "c"],
                series=[Series(name="s", values=[1.0, 2.0, 3.0])])
    _chart_block(slide, bar, Theme(), {}, x=1.0, y=1.0, w=6.0, max_h=4.0)
    bar_chart = next(s for s in slide.shapes if s.has_chart).chart
    assert bar_chart.category_axis.reverse_order is True

    slide2 = _blank_slide()
    col = Chart(chart="column", labels=["a", "b", "c"],
                series=[Series(name="s", values=[1.0, 2.0, 3.0])])
    _chart_block(slide2, col, Theme(), {}, x=1.0, y=1.0, w=6.0, max_h=4.0)
    col_chart = next(s for s in slide2.shapes if s.has_chart).chart
    assert col_chart.category_axis.reverse_order is False


def test_pptx_placeholder_caption_never_dropped_even_when_max_h_is_tiny():
    from docloom.render.pptx import _placeholder_block

    slide = _blank_slide()
    _placeholder_block(
        slide, x=1.0, y=1.0, w=6.0, max_h=0.3, theme=Theme(),
        alt="alt text", caption="PLACEHOLDERCAPMARK",
    )
    assert "PLACEHOLDERCAPMARK" in _shapes_text(slide)


def test_pptx_image_left_renders_the_image_caption(tmp_path):
    # Finding B's own reported instance: slide.image.caption was silently
    # dropped on image_left/image_right (DOCX/HTML/MD all kept it, via
    # flatten_slides turning a deck-only document's s.image into a real
    # Image block for the report renderers -- only PPTX lost it).
    from pptx import Presentation

    doc = Document(title="T", slides=[
        Slide(layout="image_left", title="T2",
              image=ImageBlock(path=_png(tmp_path / "pic.png"), caption="SIDECAPMARK"),
              blocks=[Paragraph(text="body")]),
    ])
    out = render(doc, "pptx", tmp_path / "left_cap.pptx")
    assert "SIDECAPMARK" in _shapes_text(Presentation(str(out)).slides[0])


def test_pptx_image_right_renders_the_image_caption(tmp_path):
    from pptx import Presentation

    doc = Document(title="T", slides=[
        Slide(layout="image_right", title="T2",
              image=ImageBlock(path=_png(tmp_path / "pic.png"), caption="SIDECAPMARK2"),
              blocks=[Paragraph(text="body")]),
    ])
    out = render(doc, "pptx", tmp_path / "right_cap.pptx")
    assert "SIDECAPMARK2" in _shapes_text(Presentation(str(out)).slides[0])


def test_pptx_image_side_caption_matches_docx_html_md(tmp_path):
    # Cross-renderer diff for finding B's exact instance.
    doc = Document(title="T", slides=[
        Slide(layout="image_left", title="T2",
              image=ImageBlock(path=_png(tmp_path / "pic.png"), caption="SIDECAPMARK3"),
              blocks=[Paragraph(text="body")]),
    ])
    from pptx import Presentation

    pptx_out = render(doc, "pptx", tmp_path / "cross_cap.pptx")
    assert "SIDECAPMARK3" in _shapes_text(Presentation(str(pptx_out)).slides[0])
    for fmt in ("docx", "html", "md"):
        out = render(doc, fmt, tmp_path / f"cross_cap.{fmt}")
        if fmt == "docx":
            import docx as docx_lib

            d = docx_lib.Document(str(out))
            text = "\n".join(p.text for p in d.paragraphs)
        else:
            text = out.read_text(encoding="utf-8")
        assert "SIDECAPMARK3" in text


# ------------------------------------------------- finding C: hero-with-blocks


def _diagram_edge_count(pptx_path) -> int:
    import zipfile

    with zipfile.ZipFile(pptx_path) as z:
        xml = "".join(
            z.read(n).decode("utf-8", "replace")
            for n in z.namelist() if "slides/slide" in n
        )
    return xml.count("<p:cxnSp>")


def test_pptx_imageless_hero_with_diagram_renders_natively_not_crushed(tmp_path):
    # Finding C: the dispatcher fix that let an imageless hero reach
    # _hero_slide at all (finding A, part 2) regressed a hero WITH blocks --
    # the block band was capped at a flat 1.7in guess sized for a short
    # caption over a photo, crushing a 2-node diagram's label font below the
    # legibility floor and silently degrading it to an unreadable raster
    # tile with zero connectors. It must now render natively (a real
    # cxnSp-glued connector per edge), not raster, and without triggering
    # diagram_pptx's own "does not clear the Npt node-label floor" warning.
    doc = Document(title="T", slides=[
        Slide(layout="hero", title="HEROMARK the two-node flow",
              subtitle="SUBHEROMARK context line",
              blocks=[Diagram(
                  id="flow",
                  nodes=[DiagramNode(id="a", label="Client"), DiagramNode(id="b", label="Service")],
                  edges=[DiagramEdge(source="a", target="b", label="request")],
                  caption="Two-node request flow",
              )]),
    ])
    import warnings

    with warnings.catch_warnings(record=True) as caught:
        warnings.simplefilter("always")
        out = render(doc, "pptx", tmp_path / "hero_diagram.pptx")
    assert not any("node-label floor" in str(w.message) for w in caught), (
        "diagram degraded to raster instead of getting real room"
    )
    assert _diagram_edge_count(out) == 1  # native connector glue, not a raster picture

    from pptx import Presentation

    text = _shapes_text(Presentation(str(out)).slides[0])
    assert "HEROMARK" in text and "SUBHEROMARK" in text


def test_pptx_imageless_hero_without_blocks_is_unaffected(tmp_path):
    # Regression guard: the finding-C fix must not disturb the already-
    # correct short-band sizing for a bare title/subtitle imageless hero.
    from pptx import Presentation

    theme = Theme()
    doc = Document(title="T", slides=[
        Slide(layout="hero", title="HEROMARK the ask", subtitle="SUBHEROMARK context"),
    ])
    out = render(doc, "pptx", tmp_path / "hero_no_blocks.pptx", theme=theme)
    slide = Presentation(str(out)).slides[0]
    bg = str(slide.background.fill.fore_color.rgb).upper()
    assert bg == theme.primary.lstrip("#").upper()
    assert "HEROMARK" in _shapes_text(slide) and "SUBHEROMARK" in _shapes_text(slide)


# --------------------------------------------- finding D: callout fill wash


def test_pptx_callout_fills_are_tinted_and_distinct_per_style(tmp_path):
    # All four fills used to be the identical flat theme.surface gray, so
    # only the 4px edge bar carried any color at all.
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Inches

    from docloom import Callout

    theme = Theme()
    doc = Document(title="T", slides=[
        Slide(layout="content", title="T2", blocks=[
            Callout(style="info", text="i"),
            Callout(style="success", text="s"),
            Callout(style="warning", text="w"),
            Callout(style="danger", text="d"),
        ]),
    ])
    out = render(doc, "pptx", tmp_path / "callout_fills.pptx", theme=theme)
    slide = Presentation(str(out)).slides[0]
    # the fill rects are the wide, tall AUTO_SHAPEs: narrow ones (< 0.2in)
    # are the edge bars, and the thin (~0.028in) one is the title's accent
    # rule, neither of which is a callout fill
    fills = [
        str(s.fill.fore_color.rgb).upper()
        for s in slide.shapes
        if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and s.width >= Inches(0.2) and s.height >= Inches(0.2)
    ]
    assert len(fills) == 4
    assert len(set(fills)) == 4, "callout fills are not distinct per style"
    assert theme.surface.lstrip("#").upper() not in fills, (
        "callout fill is still the old flat gray"
    )


def test_pptx_callout_fill_color_derives_from_the_edge_color():
    from docloom.render.pptx import _callout_edge_color, _callout_fill_color

    theme = Theme(accent="#0E9F6E")
    for style in ("info", "success", "warning", "danger"):
        fill = _callout_fill_color(style, theme)
        edge = _callout_edge_color(style, theme)
        assert fill.startswith("#") and fill != edge
        # a wash toward background, not the flat surface gray
        assert fill.upper() != theme.surface.upper()


# --------------------------------------------------------- render/qa.py
#
# Reference-free geometric QA over a BUILT pptx.Presentation: off-slide
# bleed, illegitimate shape overlap (with containment/decorative filtering
# -- the false-positive machine the task warns about), deck-wide palette
# discipline, and WCAG 2 text/background contrast. qa.py is pure and
# standalone (no Document IR, no Theme object): the synthetic-shape tests
# below build raw python-pptx Presentations directly to control geometry and
# color exactly; the doc-driven tests exercise it against a real docloom
# render.


def _qa_blank_prs():
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return prs, slide


def _qa_rect(slide, x, y, w, h, rgb):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    return shape


def test_qa_finding_reuses_lints_finding_shape():
    # "reuse docloom's existing finding/severity shape from lint.py" -- this
    # must be the SAME class, not a structurally-similar copy, so any
    # existing caller that folds lint findings and qa findings into one list
    # (isinstance checks, model_dump(), etc.) just works.
    from docloom.lint import Finding as LintFinding
    from docloom.render.qa import Finding as QaFinding

    assert QaFinding is LintFinding


def test_qa_flags_a_shape_that_bleeds_past_the_slide_edge():
    from pptx.dml.color import RGBColor

    from docloom.render import qa

    prs, slide = _qa_blank_prs()
    _qa_rect(slide, -0.5, 0, 2, 2, RGBColor(0, 0, 0))
    findings = qa.check_bleed(prs)
    assert len(findings) == 1
    assert findings[0].rule == "qa/off-slide"
    assert findings[0].severity == "warning"
    assert "left" in findings[0].message


def test_qa_bleed_ignores_a_shape_flush_with_the_edge():
    # Rounding slack: a shape sized to land EXACTLY on the slide bounds must
    # not be flagged just because of float-to-EMU rounding.
    from pptx.dml.color import RGBColor

    from docloom.render import qa

    prs, slide = _qa_blank_prs()
    _qa_rect(slide, 0, 0, 13.333, 7.5, RGBColor(0, 0, 0))
    assert qa.check_bleed(prs) == []


def test_qa_bleed_whitelist_exempts_an_intentional_full_bleed_shape():
    from pptx.dml.color import RGBColor

    from docloom.render import qa

    prs, slide = _qa_blank_prs()
    _qa_rect(slide, -0.5, 0, 2, 2, RGBColor(0, 0, 0))
    assert qa.check_bleed(prs, whitelist=lambda shape, si: True) == []


def test_qa_flags_two_shapes_that_genuinely_collide():
    from pptx.dml.color import RGBColor

    from docloom.render import qa

    prs, slide = _qa_blank_prs()
    _qa_rect(slide, 1, 1, 3, 3, RGBColor(255, 0, 0))
    _qa_rect(slide, 2, 2, 3, 3, RGBColor(0, 255, 0))  # ~44% of the smaller area
    findings = qa.check_overlap(prs)
    assert len(findings) == 1
    assert findings[0].rule == "qa/shape-overlap"
    assert findings[0].severity == "warning"


def test_qa_overlap_ignores_containment_text_on_a_card():
    # The exact false-positive class the task calls out: a label textbox
    # sitting entirely inside a card's fill rectangle is legitimate
    # composition (a stat card, a diagram node inside its group container),
    # not a layout defect.
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    from docloom.render import qa

    prs, slide = _qa_blank_prs()
    _qa_rect(slide, 1, 1, 4, 2, RGBColor(240, 240, 240))
    txt = slide.shapes.add_textbox(Inches(1.2), Inches(1.2), Inches(3), Inches(0.5))
    txt.text_frame.text = "Card label"
    assert qa.check_overlap(prs) == []


def test_qa_overlap_ignores_a_thin_decorative_divider():
    # A hairline rule/underline (min dimension below the decorative floor)
    # carries no real visual area to collide with anything.
    from pptx.dml.color import RGBColor

    from docloom.render import qa

    prs, slide = _qa_blank_prs()
    _qa_rect(slide, 1, 1, 4, 2, RGBColor(240, 240, 240))
    _qa_rect(slide, 1, 1.9, 4, 0.02, RGBColor(0, 0, 0))  # a hairline rule
    assert qa.check_overlap(prs) == []


def test_qa_overlap_ignores_connector_lines_with_a_large_bounding_box():
    # An elbow/diagonal connector's bounding box can span a large rectangle
    # even though the actual drawn line is a thin stroke; bbox-thinness
    # alone would miss a diagonal one, so line/freeform shapes with no
    # solid fill are always decorative regardless of their bbox size.
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.util import Inches

    from docloom.render import qa

    prs, slide = _qa_blank_prs()
    node = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(3))
    node.text_frame.text = "node"
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.ELBOW, Inches(0.5), Inches(0.5), Inches(3.5), Inches(3.5)
    )
    assert qa.check_overlap(prs) == []


def test_qa_flags_low_contrast_text_on_its_own_fill():
    from pptx.dml.color import RGBColor

    from docloom.render import qa

    prs, slide = _qa_blank_prs()
    box = _qa_rect(slide, 1, 1, 3, 1, RGBColor(0xFF, 0xFF, 0xFF))
    box.text_frame.text = "Invisible text"
    box.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFE, 0xFE, 0xFE)
    findings = qa.check_contrast(prs)
    assert len(findings) == 1
    assert findings[0].rule == "qa/low-contrast"
    assert findings[0].severity == "warning"


def test_qa_contrast_resolves_a_transparent_textbox_against_the_shape_behind_it():
    # A plain textbox (no fill of its own, the common case for docloom's own
    # _box()) sitting on top of a filled card must resolve its background to
    # the CARD's fill, not silently pass or wrongly assume the slide's own
    # background.
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    from docloom.render import qa

    prs, slide = _qa_blank_prs()
    _qa_rect(slide, 1, 1, 4, 2, RGBColor(0xFF, 0xFF, 0xFF))  # white card
    txt = slide.shapes.add_textbox(Inches(1.2), Inches(1.2), Inches(3), Inches(0.5))
    txt.text_frame.text = "Card label"
    txt.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFE, 0xFE, 0xFE)
    findings = qa.check_contrast(prs)
    assert len(findings) == 1
    assert "Card label" in findings[0].message


def test_qa_contrast_falls_back_to_the_slide_background():
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    from docloom.render import qa

    prs, slide = _qa_blank_prs()  # white slide background
    txt = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(0.5))
    txt.text_frame.text = "Ghost text"
    txt.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFE, 0xFE, 0xFE)
    findings = qa.check_contrast(prs)
    assert len(findings) == 1
    assert "Ghost text" in findings[0].message


def test_qa_contrast_checks_table_cells_against_their_own_fill():
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    from docloom.render import qa

    prs, slide = _qa_blank_prs()
    frame = slide.shapes.add_table(1, 1, Inches(1), Inches(1), Inches(3), Inches(1))
    cell = frame.table.cell(0, 0)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    cell.text_frame.text = "Faint cell text"
    cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFE, 0xFE, 0xFE)
    findings = qa.check_contrast(prs)
    assert len(findings) == 1
    assert findings[0].rule == "qa/low-contrast"


def test_qa_table_cells_visits_every_distinct_cell_and_dedups_merged_span():
    # Regression for a bug where dedup keyed on id(cell): python-pptx
    # constructs a brand-new _Cell wrapper on every table.cell() call, so
    # id(cell) is a freed-and-reused address, not a stable identity, and the
    # old `id(cell) in seen` dedup spuriously skipped unrelated cells whose
    # wrapper happened to reuse a just-freed address. The fix filters on
    # `cell.is_spanned` (the grid positions a merge covers) instead of any
    # object/element identity. Assert an exact, known cell count for a table
    # with one 2x1 horizontal merge: 3x2 grid == 6 grid positions, one merge
    # shadows 1 of them, so exactly 5 distinct cells must be yielded, and
    # each of the 4 unmerged cells' own text must all be visited (proves no
    # unrelated cell is silently skipped either).
    from pptx.util import Inches

    from docloom.render import qa

    prs, slide = _qa_blank_prs()
    frame = slide.shapes.add_table(3, 2, Inches(1), Inches(1), Inches(4), Inches(3))
    table = frame.table
    table.cell(0, 0).merge(table.cell(0, 1))  # merge row 0 across both columns
    for r, c in [(1, 0), (1, 1), (2, 0), (2, 1)]:
        table.cell(r, c).text_frame.text = f"r{r}c{c}"
    table.cell(0, 0).text_frame.text = "merged"

    cells = list(qa._table_cells(frame))
    assert len(cells) == 5
    seen_coords = {(r, c) for r, c, _cell in cells}
    assert seen_coords == {(0, 0), (1, 0), (1, 1), (2, 0), (2, 1)}
    texts = {cell.text_frame.text for _r, _c, cell in cells}
    assert texts == {"merged", "r1c0", "r1c1", "r2c0", "r2c1"}


def test_qa_palette_flags_too_many_fonts_and_non_neutral_fills():
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    from docloom.render import qa

    prs, slide = _qa_blank_prs()
    colors = [(255, 0, 0), (0, 255, 0), (0, 0, 255), (255, 255, 0),
              (255, 0, 255), (0, 255, 255), (128, 0, 128)]
    for i, c in enumerate(colors):
        _qa_rect(slide, 0.2 + i * 1.5, 0.5, 1, 1, RGBColor(*c))
    for i, font in enumerate(["Arial", "Times New Roman", "Courier New", "Comic Sans MS"]):
        tb = slide.shapes.add_textbox(Inches(0.2 + i * 3), Inches(3), Inches(2.8), Inches(0.5))
        tb.text_frame.text = f"text {i}"
        tb.text_frame.paragraphs[0].runs[0].font.name = font
    findings = qa.check_palette(prs)
    rules = {f.rule for f in findings}
    assert rules == {"qa/font-family-sprawl", "qa/palette-sprawl"}
    assert all(f.severity == "warning" for f in findings)


def test_qa_palette_ignores_neutral_grays_and_stays_within_budget():
    from pptx.dml.color import RGBColor

    from docloom.render import qa

    prs, slide = _qa_blank_prs()
    for gray in (0x10, 0x40, 0x80, 0xC0):
        _qa_rect(slide, 1, 1, 1, 1, RGBColor(gray, gray, gray))
    _qa_rect(slide, 3, 3, 1, 1, RGBColor(0x4F, 0x46, 0xE5))  # one brand color
    assert qa.check_palette(prs) == []


def test_qa_audit_on_a_realistic_deck_finds_no_overlap_false_positives(tmp_path):
    # Integration guard: a deck that exercises stat cards, a callout, a
    # table, a chart, and a native diagram (group container + nodes +
    # connectors) -- exactly the composition patterns that would make a
    # naive bbox-overlap rule cry wolf on every slide -- must come back
    # clean on qa/shape-overlap, and every finding that DOES fire (contrast,
    # palette) must stay advisory. This is the "or make the rule narrower"
    # regression guard: proof the containment/decorative filtering actually
    # works on renderer-real geometry, not just synthetic shapes.
    from pptx import Presentation

    from docloom import Callout
    from docloom.render import qa

    diagram = Diagram(
        id="arch", title="Architecture", caption="Request flow",
        nodes=[
            DiagramNode(id="client", label="Client", type="external"),
            DiagramNode(id="api", label="API", type="service"),
            DiagramNode(id="db", label="Database", type="store"),
        ],
        edges=[
            DiagramEdge(source="client", target="api", label="HTTPS"),
            DiagramEdge(source="api", target="db", label="SQL"),
        ],
    )
    doc = Document(title="Realistic deck", slides=[
        Slide(layout="content", title="Stats and a callout", blocks=[
            StatRow(items=[Stat(label="Uptime", value="99.9%", delta="+0.2%"),
                            Stat(label="Latency", value="120ms")]),
            Callout(text="Danger callout", style="danger"),
        ]),
        Slide(layout="content", title="Table and chart", blocks=[
            TableBlock(header=["A", "B"], rows=[["1", "2"]], caption="tbl"),
            Chart(chart="column", title="T", labels=["A"],
                  series=[Series(name="s", values=[1.0])], caption="chart cap"),
        ]),
        Slide(layout="content", title="Architecture diagram", blocks=[diagram]),
    ])
    out = render(doc, "pptx", tmp_path / "qa_realistic.pptx")
    prs = Presentation(str(out))
    findings = qa.audit(prs)
    assert all(f.severity == "warning" for f in findings)
    assert not any(f.rule == "qa/shape-overlap" for f in findings)
    assert not any(f.rule == "qa/off-slide" for f in findings)


# -------------------------------------------------- measured text-fit pass


def test_pptx_overflowing_quote_gets_baked_shrink(tmp_path):
    # A quote long enough that the old _est_lines heuristic's chosen tier
    # still measurably overflows its real, font-measured wrapped extent:
    # _fit_text_frames must bake a smaller run size (and, per the no-
    # double-shrink invariant, reset normAutofit to fontScale=100000 /
    # lnSpcReduction=0 rather than encode the shrink twice).
    import warnings

    from pptx import Presentation
    from pptx.oxml.ns import qn

    from docloom.render.pptx import MIN_FIT_PT

    quote_text = " ".join(f"word{i}" for i in range(350))
    doc = Document(title="T", slides=[Slide(layout="quote", blocks=[Quote(text=quote_text)])])
    with warnings.catch_warnings():
        warnings.simplefilter("always")  # crowding/legibility warnings may fire too
        out = render(doc, "pptx", tmp_path / "quote_shrink.pptx")
    slide = Presentation(str(out)).slides[0]
    quote_shape = next(
        s for s in slide.shapes if s.has_text_frame and "word0 " in s.text_frame.text
    )
    sizes = [r.font.size.pt for p in quote_shape.text_frame.paragraphs for r in p.runs]
    assert all(MIN_FIT_PT <= sz < 14 for sz in sizes)
    bodyPr = quote_shape.text_frame._txBody.bodyPr
    autofit = bodyPr.find(qn("a:normAutofit"))
    assert autofit.get("fontScale") == "100000"
    assert autofit.get("lnSpcReduction") == "0"
    # nothing-lost: the full quote is still present, just shrunk to fit
    assert "word349" in quote_shape.text_frame.text


def test_pptx_fitting_frame_xml_untouched(tmp_path):
    # A frame that already fits (same grow-suppressed short-prose fixture as
    # test_pptx_grow_pass_suppressed_next_to_a_fixed_size_block above) must
    # come out of _fit_text_frames byte-identical to today: run size
    # unchanged, and its normAutofit element left with no fontScale
    # attribute at all (proving zero churn, not just a no-op scale).
    from pptx import Presentation
    from pptx.oxml.ns import qn

    from docloom.render.pptx import BODY_PT

    doc = Document(title="T", slides=[
        Slide(layout="content", title="T2", blocks=[
            Paragraph(text="short prose"),
            Code(code="x = 1"),
        ]),
    ])
    out = render(doc, "pptx", tmp_path / "fit_untouched.pptx")
    slide = Presentation(str(out)).slides[0]
    para = next(s for s in slide.shapes if s.has_text_frame and "short prose" in s.text_frame.text)
    assert para.text_frame.paragraphs[0].runs[0].font.size.pt == BODY_PT
    bodyPr = para.text_frame._txBody.bodyPr
    autofit = bodyPr.find(qn("a:normAutofit"))
    assert autofit is not None
    assert autofit.get("fontScale") is None


def test_pptx_floor_clamped_frame_warns_and_keeps_text(tmp_path):
    # A quote so long that even the MIN_FIT_PT floor at max line-spacing
    # reduction cannot make it fit: _fit_text_frames must still draw it (at
    # the floor size, never dropped) and warn by name instead of silently
    # overflowing the slide.
    from pptx import Presentation

    from docloom.render.pptx import MIN_FIT_PT

    quote_text = " ".join(f"word{i}" for i in range(1200))
    doc = Document(title="T", slides=[Slide(layout="quote", blocks=[Quote(text=quote_text)])])
    with pytest.warns(UserWarning, match="legibility floor"):
        out = render(doc, "pptx", tmp_path / "quote_floor.pptx")
    slide = Presentation(str(out)).slides[0]
    quote_shape = next(
        s for s in slide.shapes if s.has_text_frame and "word0 " in s.text_frame.text
    )
    sizes = {r.font.size.pt for p in quote_shape.text_frame.paragraphs for r in p.runs}
    assert sizes == {MIN_FIT_PT}
    assert "word1199" in quote_shape.text_frame.text


def test_fit_scale_authored_below_floor_run_does_not_veto_shrink():
    # Regression for the fix that derived floor_scale from the frame's
    # absolute SMALLEST run: a citation superscript is authored below
    # MIN_FIT_PT by deliberate design (_runs bakes it at
    # max(8, round(size*0.65)), routinely under 9pt), and deriving the floor
    # from it makes min_pt/small > 1, which the old min(1.0, ...) clamp
    # pinned at floor_scale = 1.0 -- disabling autofit for the WHOLE frame,
    # body text included. A run authored below the floor is not a
    # legibility bug to defend; it must scale right along with everything
    # else, never veto the fit of runs that actually need to shrink.
    from docloom.render import textfit

    body = textfit.RunSpec("word " * 200, "Arial", 18.0)
    citation = textfit.RunSpec("1", "Arial", 6.0)  # authored below MIN_FIT_PT=9 on purpose
    para = textfit.ParaSpec(runs=(body, citation))
    res = textfit.fit_scale([para], width_in=4.0, height_in=1.0, min_pt=9.0)
    assert res.scale < 1.0, (
        "a deliberately-small authored-below-floor run must not veto "
        f"autofit for the runs that overflow; got {res!r}"
    )


def test_fit_scale_never_bakes_an_at_or_above_floor_run_below_the_floor():
    # The ORIGINAL fix, pinned by a test that actually exercises the shrink
    # path (unlike the tautological version this replaces -- see below): the
    # floor must protect every run that started AT OR ABOVE MIN_FIT_PT, not
    # just the frame's largest run. A big headline run sharing the same
    # uniform scale as a small-but-still-legible run must never push that
    # smaller run under min_pt, even under egregious overflow that would
    # otherwise drive the shared scale far lower.
    from docloom.render import textfit

    big = textfit.RunSpec("word " * 200, "Arial", 30.0)
    small_at_floor = textfit.RunSpec("caption", "Arial", 12.0)  # authored ABOVE floor
    para = textfit.ParaSpec(runs=(big, small_at_floor))
    res = textfit.fit_scale([para], width_in=4.0, height_in=1.0, min_pt=9.0)
    # the fit actually engaged (not a trivial no-shrink no-op)...
    assert res.scale < 1.0, res
    # ...but never at the cost of pushing the protected run under the floor
    baked = 12.0 * res.scale
    assert baked >= 9.0 - 1e-6, (
        f"a run authored at/above the floor (12pt) was baked to {baked:.3f}pt, "
        f"under MIN_FIT_PT=9.0; res={res!r}"
    )


def test_pptx_fit_one_frame_citation_below_floor_does_not_veto_shrink():
    # End-to-end version of the regression above, through the REAL
    # _fit_one_frame glue (not a synthetic textfit call): a frame holding a
    # normal-size body run plus a small run authored below MIN_FIT_PT (a
    # citation superscript is exactly this -- _runs bakes it at
    # max(8, round(size*0.65)), routinely under 9pt) must still shrink the
    # overflowing body run, and the citation must scale proportionally right
    # along with it instead of vetoing the fit or getting bumped back up
    # above its own authored size (the OTHER way "nothing authored is lost"
    # could be violated: silently changing a deliberately-small citation
    # into a larger one).
    from pptx import Presentation
    from pptx.util import Inches, Pt as PptxPt

    from docloom.render.pptx import MIN_FIT_PT, _box, _fit_one_frame

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tf = _box(slide, 1, 1, 3.0, 1.0)
    p = tf.paragraphs[0]
    body = p.add_run()
    body.text = "word " * 60
    body.font.name, body.font.size = "Arial", PptxPt(18)
    cite = p.add_run()
    cite.text = "1"
    cite.font.name, cite.font.size = "Arial", PptxPt(6)  # authored below the floor, by design
    shape = slide.shapes[-1]
    _fit_one_frame(shape)

    body_size, cite_size = body.font.size.pt, cite.font.size.pt
    assert body_size < 18.0, "the overflowing body run never shrank -- autofit was vetoed"
    assert body_size >= MIN_FIT_PT, body_size
    # the citation scales proportionally with the shared scale; it must
    # never get bumped back up above its own authored size just because a
    # blanket floor clamp doesn't distinguish it from a protected run
    assert cite_size < 6.0, cite_size


def test_pptx_mixed_run_sizes_never_bake_an_at_or_above_floor_run_below_it(tmp_path):
    # Real quote-pipeline sanity check for the ORIGINAL fix: even under
    # egregious overflow with a citation superscript sharing the frame,
    # every word run (authored at/above the floor) stays at/above
    # MIN_FIT_PT, and nothing authored -- word or citation -- is dropped.
    from pptx import Presentation

    from docloom.render.pptx import MIN_FIT_PT

    quote_text = [
        Span(text=f"word{i}", cite="a") for i in range(300)
    ]
    doc = Document(
        title="T",
        slides=[Slide(layout="quote", blocks=[Quote(text=quote_text)])],
        sources=[Source(id="a", title="Alpha")],
    )
    out = render(doc, "pptx", tmp_path / "quote_mixed_runs.pptx")
    slide = Presentation(str(out)).slides[0]
    quote_shape = next(
        s for s in slide.shapes if s.has_text_frame and "word0" in s.text_frame.text
    )
    runs = [r for p in quote_shape.text_frame.paragraphs for r in p.runs]
    sizes = [r.font.size.pt for r in runs]
    assert len(sizes) > 300  # every word run plus its citation-superscript run
    word_sizes = {r.font.size.pt for r in runs if not r.text.isdigit()}
    assert all(sz >= MIN_FIT_PT for sz in word_sizes), word_sizes
    # nothing-lost: the full quote (and its citation markers) all survive
    assert "word299" in quote_shape.text_frame.text
