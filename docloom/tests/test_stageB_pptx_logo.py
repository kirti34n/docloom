"""Stage B (CONTRACT C6): doc.logo must appear top-right, at a consistent
~0.5in size, on EVERY pptx slide layout.

Before this fix: `_doc_logo` skipped title/section/hero/image_left/
image_right entirely (G2); the title slide only ever showed a logo bound at
*generation* time via slides[0].image, never a doc.logo bound later at
export time (G3); `_logo`'s `min(..., 1.0)` scale cap meant a small source
logo could never upscale to a legible size (G7); and section/hero (full
bleed) and image_right (image pane at the top-right corner) had no safe
corner for a logo at all (G9)."""

import base64
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from docloom import Document, Image, Paragraph, Quote, Slide, render
from docloom.render.pptx import LOGO_MAX_H, SLIDE_W

# a minimal valid 1x1 PNG (same fixture used by test_pptx_logo.py): tiny
# enough that only the upscale fix makes it render at a legible size
_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAAC0lEQVR42mP8"
    "z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
)


def _logo_path(tmp_path: Path) -> str:
    p = tmp_path / "logo.png"
    p.write_bytes(_PNG)
    return str(p)


def _content_image_path(tmp_path: Path, name: str) -> str:
    from PIL import Image as PILImage

    p = tmp_path / name
    PILImage.new("RGB", (400, 300), "blue").save(p)
    return str(p)


def _pics(slide):
    return [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]


def _autoshapes(slide):
    return [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]


def _encloses(outer, inner) -> bool:
    return (
        outer.left <= inner.left and outer.top <= inner.top
        and outer.left + outer.width >= inner.left + inner.width
        and outer.top + outer.height >= inner.top + inner.height
    )


LAYOUT_INDEX = {
    "title": 0, "content": 1, "two_column": 2, "quote": 3,
    "section": 4, "hero": 5, "image_left": 6, "image_right": 7,
}


def _build_doc(tmp_path: Path) -> Document:
    """One deck, one doc.logo (no per-slide image on the title slide), with
    every slide layout the pptx renderer knows about."""
    logo = Image(path=_logo_path(tmp_path), alt="Acme")
    return Document(title="Branded", logo=logo, slides=[
        Slide(layout="title", title="Q3 results"),
        Slide(layout="content", title="A content slide",
              blocks=[Paragraph(text="body copy")]),
        Slide(layout="two_column", title="Split",
              blocks=[Paragraph(text="left")], right=[Paragraph(text="right")]),
        Slide(layout="quote", blocks=[Quote(text="A short quote", attribution="Someone")]),
        Slide(layout="section", title="Part two"),
        Slide(layout="hero", title="Hero moment",
              image=Image(path=_content_image_path(tmp_path, "hero.png"))),
        Slide(layout="image_left", title="Left image",
              image=Image(path=_content_image_path(tmp_path, "left.png")),
              blocks=[Paragraph(text="body")]),
        Slide(layout="image_right", title="Right image",
              image=Image(path=_content_image_path(tmp_path, "right.png")),
              blocks=[Paragraph(text="body")]),
    ])


def test_render_succeeds_with_a_logo_on_every_layout(tmp_path):
    doc = _build_doc(tmp_path)
    out = render(doc, "pptx", tmp_path / "branded.pptx")
    assert out.is_file()
    Presentation(str(out))  # reopens cleanly: no corrupt XML from any layout


def test_logo_appears_on_every_non_title_layout(tmp_path):
    """G2: section/hero/image_left/image_right used to be skipped entirely."""
    doc = _build_doc(tmp_path)
    out = render(doc, "pptx", tmp_path / "branded.pptx")
    prs = Presentation(str(out))
    for layout in ("content", "two_column", "quote", "section"):
        pics = _pics(prs.slides[LAYOUT_INDEX[layout]])
        assert len(pics) == 1, f"{layout}: expected exactly one logo picture"
    for layout in ("hero", "image_left", "image_right"):
        pics = _pics(prs.slides[LAYOUT_INDEX[layout]])
        assert len(pics) == 2, f"{layout}: expected the content image plus the logo"


def test_title_layout_falls_back_to_doc_logo(tmp_path):
    """G3: a logo bound only on the document (not slides[0].image, which is
    unset here) must still show up on the title slide."""
    doc = _build_doc(tmp_path)
    out = render(doc, "pptx", tmp_path / "branded.pptx")
    prs = Presentation(str(out))
    pics = _pics(prs.slides[LAYOUT_INDEX["title"]])
    assert len(pics) == 1


def test_logo_upscales_to_the_shared_target_height(tmp_path):
    """G7: the 1x1 PNG fixture is a tiny fraction of an inch natively; the
    old min(..., 1.0) cap left it that size forever instead of scaling up."""
    doc = _build_doc(tmp_path)
    out = render(doc, "pptx", tmp_path / "branded.pptx")
    prs = Presentation(str(out))
    pic = _pics(prs.slides[LAYOUT_INDEX["content"]])[0]
    assert abs(pic.height / 914400 - LOGO_MAX_H) < 0.05


def test_section_and_hero_get_a_scrim_plate_behind_the_logo(tmp_path):
    """G9: full-bleed layouts (section's solid theme.primary fill, hero's
    cover-fit image) get a contrasting plate drawn behind the logo, reusing
    _rect, so it stays legible."""
    doc = _build_doc(tmp_path)
    out = render(doc, "pptx", tmp_path / "branded.pptx")
    prs = Presentation(str(out))
    for layout in ("section", "hero"):
        slide = prs.slides[LAYOUT_INDEX[layout]]
        logo = min(_pics(slide), key=lambda p: p.height)  # smallest picture = the logo
        shapes = list(slide.shapes)
        plates = [
            sh for sh in _autoshapes(slide)
            if _encloses(sh, logo) and shapes.index(sh) < shapes.index(logo)
        ]
        assert plates, f"{layout}: no scrim plate found behind the logo"


def test_image_right_moves_the_logo_off_the_image_pane(tmp_path):
    """G9: image_right's image pane sits at the top-right corner; the logo
    must move to the opposite corner instead of landing on the image."""
    doc = _build_doc(tmp_path)
    out = render(doc, "pptx", tmp_path / "branded.pptx")
    prs = Presentation(str(out))
    slide = prs.slides[LAYOUT_INDEX["image_right"]]
    logo = min(_pics(slide), key=lambda p: p.height)
    assert logo.left / 914400 < 1.0, "image_right logo should sit at top-left"


def test_image_left_keeps_the_logo_top_right(tmp_path):
    """image_left's image pane is already on the left, so the default
    top-right corner never collides with it."""
    doc = _build_doc(tmp_path)
    out = render(doc, "pptx", tmp_path / "branded.pptx")
    prs = Presentation(str(out))
    slide = prs.slides[LAYOUT_INDEX["image_left"]]
    logo = min(_pics(slide), key=lambda p: p.height)
    assert logo.left / 914400 > SLIDE_W / 2, "image_left logo should stay top-right"


def test_svg_logo_is_skipped_without_crashing(tmp_path):
    """G8: python-pptx cannot embed SVG; the render must still succeed with
    the logo silently skipped, not crash."""
    svg_path = tmp_path / "logo.svg"
    svg_path.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="10" height="10">'
        '<rect width="10" height="10"/></svg>'
    )
    doc = Document(title="Branded", logo=Image(path=str(svg_path)), slides=[
        Slide(layout="content", title="A", blocks=[Paragraph(text="x")]),
    ])
    out = render(doc, "pptx", tmp_path / "svg.pptx")
    prs = Presentation(str(out))
    assert _pics(prs.slides[0]) == []
