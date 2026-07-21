"""Regression coverage for the diagram-emit fix wave (2026-07-16 ultracode
pass): the native PPTX diagram node TAG run was drawn in the near-white
p["line"] tint (an unresolved kind-accent color) instead of a fill-resolved
readable foreground, so tags on the default near-white node fill sat at
~1.5-2.1:1 contrast -- unreadable, and disagreeing with diagram_svg.py and
drawio.py, both of which paint the tag inside a colored chip/box.

These tests render through the full `render(doc, "pptx", ...)` pipeline and
inspect the actual XML, the same style as tests/test_pptx_invariants.py's
Invariant B, but self-contained here per this pass's file-ownership rule
(this module may not edit test_pptx_invariants.py).
"""

from __future__ import annotations

from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE

from docloom import (
    Diagram, DiagramEdge, DiagramNode, Document, Slide, Theme, render,
)
from docloom.theme import contrast_ratio


def _walk_shapes(shapes):
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(sh.shapes)
        else:
            yield sh


def _own_fill_hex(shape) -> str | None:
    try:
        if shape.fill.type == MSO_FILL_TYPE.SOLID:
            return "#" + str(shape.fill.fore_color.rgb).upper()
    except Exception:
        return None
    return None


def _rect_of(shape):
    try:
        return (shape.left, shape.top, shape.width, shape.height)
    except Exception:
        return None


def _contains(outer, inner) -> bool:
    if outer is None or inner is None:
        return False
    if any(v is None for v in (*outer, *inner)):
        return False
    ox, oy, ow, oh = outer
    ix, iy, iw, ih = inner
    return ox <= ix and oy <= iy and (ox + ow) >= (ix + iw) and (oy + oh) >= (iy + ih)


def _fill_behind(shape, filled_shapes, slide_bg_hex: str) -> str:
    """The effective fill color behind `shape`: its own solid fill if it
    paints one (a diagram node's own tint), else the smallest earlier-drawn
    solid-filled shape whose rect fully contains `shape`'s rect, else the
    slide background."""
    own = _own_fill_hex(shape)
    if own is not None:
        return own
    rect = _rect_of(shape)
    best, best_area = None, None
    for other, fill in filled_shapes:
        if other is shape:
            continue
        orect = _rect_of(other)
        if _contains(orect, rect):
            area = orect[2] * orect[3]
            if best_area is None or area < best_area:
                best_area, best = area, fill
    return best if best is not None else slide_bg_hex


def _text_runs_with_bg(slide):
    """Yield (text, run_color_hex, fill_behind_hex) for every non-blank text
    run drawn on `slide`."""
    bg = "#" + str(slide.background.fill.fore_color.rgb).upper()
    all_shapes = list(_walk_shapes(slide.shapes))
    filled = [(sh, f) for sh in all_shapes if (f := _own_fill_hex(sh)) is not None]
    for sh in all_shapes:
        if not getattr(sh, "has_text_frame", False):
            continue
        behind = _fill_behind(sh, filled, bg)
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                if not run.text.strip():
                    continue
                try:
                    color = run.font.color
                    hexcolor = "#" + str(color.rgb).upper() if color.type is not None else None
                except Exception:
                    hexcolor = None
                if hexcolor is None:
                    continue
                yield run.text, hexcolor, behind


def test_invariant_b_hero_band_tag_contrast_floor(tmp_path):
    # `store` is the deliberate discriminator: p["bar"] on the store kind's
    # fill measures 3.19:1 (below the 4.5:1 floor), so a fix that swapped
    # p["line"] for p["bar"] instead of the fill-resolved label_fg would
    # still fail this test on the DB/PostgreSQL tag while passing on the
    # API/Go tag.
    theme = Theme()
    doc = Document(title="T", slides=[
        Slide(layout="hero", title="Tag contrast check", subtitle="on an inverted band",
              blocks=[
                  Diagram(
                      id="svc", direction="LR",
                      nodes=[
                          DiagramNode(id="a", type="service", label="API", tag="Go 1.23"),
                          DiagramNode(id="b", type="store", label="DB", tag="PostgreSQL 16"),
                      ],
                      edges=[DiagramEdge(source="a", target="b")],
                      caption="tag contrast check",
                  ),
              ]),
    ])
    out = render(doc, "pptx", tmp_path / "hero_tag_contrast.pptx", theme=theme)
    slide = Presentation(str(out)).slides[0]

    seen = 0
    failures = []
    for text, fg, bg in _text_runs_with_bg(slide):
        stripped = text.strip()
        if stripped not in ("Go 1.23", "PostgreSQL 16"):
            continue
        seen += 1
        ratio = contrast_ratio(fg, bg)
        if ratio < 4.5:
            failures.append((stripped, fg, bg, round(ratio, 2)))

    assert seen == 2, (
        f"expected exactly 2 tag runs ('Go 1.23', 'PostgreSQL 16'); found {seen} "
        "-- the native path may not have been taken, or a tag was dropped/rasterized"
    )
    assert not failures, f"tag runs below the 4.5:1 contrast floor: {failures}"


def test_invariant_b_diagram_tag_contrast_independent_of_kind():
    # White-box companion: every node kind's tag color must clear 4.5:1
    # against that kind's OWN fill on a hero (inverted) band, mirroring
    # test_pptx_invariants.py::test_invariant_b_diagram_node_label_contrast_independent_of_band
    # for the label, but for the tag run specifically.
    from docloom.render import diagram_pptx, diagram_svg
    from docloom.render.pptx import _band_theme

    theme = Theme()
    band_theme = _band_theme(theme, theme.text)
    palette = diagram_svg.kind_palette(diagram_pptx.theme_dict(band_theme))
    for kind, colors in palette.items():
        tag_fg = diagram_pptx._readable_fg(band_theme, colors["fill"])
        ratio = contrast_ratio(tag_fg, colors["fill"])
        assert ratio >= 4.5, (
            f"node kind {kind!r} tag on a hero band: {tag_fg} on {colors['fill']} "
            f"is {ratio:.2f}:1"
        )
