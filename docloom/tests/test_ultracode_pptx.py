"""Regression tests for the image_left/image_right side-pane layout and the
native PPTX bar chart's category order.

Finding 1: a tall/portrait side image was contain-fit into the full 7.5in
slide height and only afterward re-anchored inside a 0.35in top/bottom pad,
so its bottom (and its caption, drawn below it) landed off the slide.

Finding 2: when a side image could not be embedded (any SVG without the
[diagrams] raster extra, or a corrupt raster), _image_side_slide fell back to
_content_slide, which never reads s.image -- silently dropping the image,
its alt text, and its caption, and leaving image_right's top-left corner
logo to overlap the resulting full-width title.

Finding 3: native PPTX horizontal bar charts render categories first-at-
bottom (Office's default), the reverse of every other docloom chart painter
(chart_svg, used by HTML/DOCX/Typst and this file's own SVG fallback), so
the same Chart IR reads top-to-bottom in every format except native PPTX
bar charts.
"""

import warnings
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from docloom import Chart, Document, Image, Series, Slide, render
from docloom.render.pptx import LAYOUT

SLIDE_H = LAYOUT["slide_h_in"]
EMU_PER_IN = 914400


def _portrait_png(path: Path) -> str:
    from PIL import Image as PILImage

    PILImage.new("RGB", (900, 1400), "green").save(path)
    return str(path)


def _solid_png(path: Path, w: int, h: int, color: str) -> str:
    from PIL import Image as PILImage

    PILImage.new("RGB", (w, h), color).save(path)
    return str(path)


def _corrupt_png(path: Path) -> str:
    path.write_bytes(b"\x89PNG\r\n\x1a\n" + b"garbage, not a real PNG payload")
    return str(path)


def _pics(slide):
    return [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]


def _text_shape(slide, marker: str):
    return next(
        sh for sh in slide.shapes
        if sh.has_text_frame and marker in sh.text_frame.text
    )


def _overlaps(a, b) -> bool:
    ax0, ay0 = a.left, a.top
    ax1, ay1 = a.left + a.width, a.top + a.height
    bx0, by0 = b.left, b.top
    bx1, by1 = b.left + b.width, b.top + b.height
    return ax0 < bx1 and bx0 < ax1 and ay0 < by1 and by0 < ay1


# --------------------------------------------------------- finding 1


def test_portrait_side_image_and_caption_stay_on_the_slide(tmp_path):
    img_path = _portrait_png(tmp_path / "portrait.png")

    for layout in ("image_left", "image_right"):
        doc = Document(title="D", slides=[
            Slide(layout=layout, title="T",
                  image=Image(path=img_path, caption="CAPMARK"),
                  blocks=[{"type": "paragraph", "text": "body"}]),
        ])
        out = render(doc, "pptx", tmp_path / f"{layout}_portrait.pptx")
        slide = Presentation(str(out)).slides[0]

        content_pics = [p for p in _pics(slide) if p.height / EMU_PER_IN > 2]
        assert content_pics, f"{layout}: content picture not found"
        pic = content_pics[0]
        pic_bottom = (pic.top + pic.height) / EMU_PER_IN
        assert pic_bottom <= SLIDE_H + 1e-3, (
            f"{layout}: image bottom {pic_bottom}in is past the {SLIDE_H}in slide edge"
        )

        cap = _text_shape(slide, "CAPMARK")
        cap_bottom = (cap.top + cap.height) / EMU_PER_IN
        assert cap_bottom <= SLIDE_H + 1e-3, (
            f"{layout}: caption bottom {cap_bottom}in is past the {SLIDE_H}in slide edge"
        )


def test_portrait_side_image_without_caption_stays_on_the_slide(tmp_path):
    img_path = _portrait_png(tmp_path / "portrait_nocap.png")

    for layout in ("image_left", "image_right"):
        doc = Document(title="D", slides=[
            Slide(layout=layout, title="T", image=Image(path=img_path),
                  blocks=[{"type": "paragraph", "text": "body"}]),
        ])
        out = render(doc, "pptx", tmp_path / f"{layout}_portrait_nocap.pptx")
        slide = Presentation(str(out)).slides[0]

        content_pics = [p for p in _pics(slide) if p.height / EMU_PER_IN > 2]
        assert content_pics, f"{layout}: content picture not found"
        pic = content_pics[0]
        pic_bottom = (pic.top + pic.height) / EMU_PER_IN
        assert pic_bottom <= SLIDE_H + 1e-3, (
            f"{layout}: image bottom {pic_bottom}in is past the {SLIDE_H}in slide edge"
        )


# --------------------------------------------------------- finding 2


def test_unembeddable_side_image_shows_placeholder_and_does_not_overlap_logo(tmp_path):
    logo_path = _solid_png(tmp_path / "logo.png", 200, 200, "red")
    corrupt_path = _corrupt_png(tmp_path / "corrupt.png")

    for layout in ("image_left", "image_right"):
        doc = Document(title="D", logo=Image(path=logo_path), slides=[
            Slide(layout=layout, title="TITLEMARK",
                  image=Image(path=corrupt_path, alt="ALTMARK", caption="CAPMARK"),
                  blocks=[{"type": "paragraph", "text": "body"}]),
        ])
        with warnings.catch_warnings(record=True) as caught:
            warnings.simplefilter("always")
            out = render(doc, "pptx", tmp_path / f"{layout}_corrupt.pptx")

        assert any(
            "could not be embedded" in str(w.message) for w in caught
        ), f"{layout}: no warning raised for the unembeddable image"

        with zipfile.ZipFile(out) as z:
            blob = "".join(
                z.read(n).decode("utf-8", "replace")
                for n in z.namelist() if n.endswith(".xml")
            )
        assert "ALTMARK" in blob, f"{layout}: alt text silently dropped"
        assert "CAPMARK" in blob, f"{layout}: caption silently dropped"

        slide = Presentation(str(out)).slides[0]
        logo_pic = _pics(slide)[0]  # the only embeddable picture left is the logo
        title_box = _text_shape(slide, "TITLEMARK")
        assert not _overlaps(logo_pic, title_box), (
            f"{layout}: doc.logo overlaps the title textbox"
        )


# --------------------------------------------------------- finding 3


def test_native_bar_chart_categories_run_first_at_top(tmp_path):
    doc = Document(title="T", slides=[
        Slide(layout="content", title="T", blocks=[
            Chart(chart="bar", labels=["ROW_A_first", "ROW_B", "ROW_C_last"],
                  series=[Series(name="S", values=[3, 5, 2])]),
        ]),
    ])
    out = render(doc, "pptx", tmp_path / "bar.pptx")
    with zipfile.ZipFile(out) as z:
        chart_xml = next(n for n in z.namelist() if n.startswith("ppt/charts/chart"))
        xml = z.read(chart_xml).decode("utf-8")

    assert '<c:barDir val="bar"' in xml
    assert '<c:orientation val="maxMin"/>' in xml, (
        "native bar chart still uses Office's default first-at-bottom "
        "category order, diverging from every other docloom chart painter"
    )


def test_native_column_chart_category_order_is_unchanged(tmp_path):
    doc = Document(title="T", slides=[
        Slide(layout="content", title="T", blocks=[
            Chart(chart="column", labels=["A", "B", "C"],
                  series=[Series(name="S", values=[3, 5, 2])]),
        ]),
    ])
    out = render(doc, "pptx", tmp_path / "column.pptx")
    with zipfile.ZipFile(out) as z:
        chart_xml = next(n for n in z.namelist() if n.startswith("ppt/charts/chart"))
        xml = z.read(chart_xml).decode("utf-8")

    assert '<c:orientation val="minMax"/>' in xml
